"""Strict current-only cache boundary for one editable Word-page package."""

from __future__ import annotations

import hashlib
import copy
import json
import os
import posixpath
import re
import stat
import uuid
import zipfile
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Mapping
from xml.etree import ElementTree as ET

from pptx import Presentation


from cache_key import canonical_sha256
from cache_store import CacheStore
from page_pipeline import cache_record


PAGE_PACKAGE_SCHEMA = "editable-page-package-v1"
MAX_ARCHIVE_ENTRIES = 4096
MAX_ARCHIVE_UNCOMPRESSED_BYTES = 512 * 1024 * 1024
MAX_ARCHIVE_ENTRY_BYTES = 128 * 1024 * 1024
MAX_XML_PART_BYTES = 16 * 1024 * 1024
MAX_COMPRESSION_RATIO = 200.0
CURRENT_CACHE_IDENTITY_FIELDS = frozenset({
    "page_source_sha256",
    "style_execution_sha256",
    "generation_parameters",
    "repair_feedback",
    "reconstruction_version",
})
PAGE_PACKAGE_FIELDS = frozenset({
    "schema_version",
    "cache_key",
    "pptx",
    "pptx_sha256",
    "editable_object_count",
    "slide_fingerprint",
})
_REQUIRED_OPC_PARTS = frozenset({
    "[Content_Types].xml",
    "_rels/.rels",
    "ppt/presentation.xml",
    "ppt/_rels/presentation.xml.rels",
})
_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
_PKG_REL = "http://schemas.openxmlformats.org/package/2006/relationships"
_CT = "http://schemas.openxmlformats.org/package/2006/content-types"


class PackageValidationError(ValueError):
    """The current page package or its cache authority is invalid."""


@dataclass(frozen=True)
class PptxInspection:
    slide_count: int
    editable_object_counts: tuple[int, ...]
    slide_fingerprints: tuple[str, ...]
    slide_page_ids: tuple[int | None, ...] = ()


@dataclass(frozen=True)
class PagePackage:
    page_number: int
    cache_key: str
    cache_manifest: Path
    cache_manifest_sha256: str
    cache_artifact: Path
    cache_artifact_sha256: str
    pptx_path: Path
    pptx_sha256: str
    editable_object_count: int
    slide_fingerprint: str
    package_valid: bool = True


def _sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def _valid_sha256(value: object) -> bool:
    return (
        isinstance(value, str)
        and len(value) == 64
        and all(character in "0123456789abcdef" for character in value)
    )


def _is_link_or_reparse(path: Path) -> bool:
    try:
        info = path.lstat()
    except OSError:
        return False
    reparse = bool(
        getattr(info, "st_file_attributes", 0)
        & getattr(stat, "FILE_ATTRIBUTE_REPARSE_POINT", 0x400)
    )
    return stat.S_ISLNK(info.st_mode) or reparse


def _literal_absolute(value: str | Path) -> Path:
    return Path(os.path.abspath(os.fspath(value)))


def require_plain_project(value: str | Path) -> Path:
    """Reject a linked project path before any resolving filesystem access."""
    project = _literal_absolute(value)
    current = Path(project.anchor)
    for part in project.parts[1:]:
        current /= part
        if os.path.lexists(current) and _is_link_or_reparse(current):
            raise PackageValidationError(
                f"workflow project ancestor cannot be a link or reparse point: {current}"
            )
    try:
        info = project.lstat()
    except OSError as exc:
        raise PackageValidationError("workflow project does not exist") from exc
    if not stat.S_ISDIR(info.st_mode):
        raise PackageValidationError("workflow project must be a directory")
    resolved = project.resolve(strict=True)
    if resolved != project:
        raise PackageValidationError("workflow project path is redirected")
    return resolved


def _literal_project_path(project: Path, value: str | Path) -> Path:
    project = require_plain_project(project)
    raw = Path(value)
    path = _literal_absolute(raw if raw.is_absolute() else project / raw)
    try:
        relative = path.relative_to(project)
    except ValueError as exc:
        raise PackageValidationError("page package artifact must be project-local") from exc
    current = project
    for part in relative.parts:
        current = current / part
        if os.path.lexists(current) and _is_link_or_reparse(current):
            raise PackageValidationError(
                f"project path cannot contain a link or reparse point: {current}"
            )
    if path == project:
        raise PackageValidationError("page package artifact must be a regular project-local file")
    return path


def project_output_path(project: Path, value: str | Path) -> Path:
    """Return a lexical project child after rejecting linked existing parents."""
    return _literal_project_path(project, value)


def require_project_file(project: Path, value: str | Path) -> Path:
    return _project_file(project, value)


def _project_file(project: Path, value: str | Path, *, must_exist: bool = True) -> Path:
    path = _literal_project_path(project, value)
    if must_exist and not path.is_file():
        raise PackageValidationError(f"page package artifact is missing: {path}")
    if must_exist:
        info = path.lstat()
        if not stat.S_ISREG(info.st_mode) or info.st_nlink != 1:
            raise PackageValidationError("page package artifact must be a regular unlinked file")
    resolved = path.resolve(strict=must_exist)
    if resolved != path:
        raise PackageValidationError("page package artifact path is redirected")
    return path


def _read_object(path: Path, *, exact_fields: frozenset[str] | None = None) -> dict[str, Any]:
    try:
        value = json.loads(path.read_text(encoding="utf-8"))
    except (OSError, UnicodeDecodeError, json.JSONDecodeError) as exc:
        raise PackageValidationError(f"page package JSON is invalid: {path}") from exc
    if not isinstance(value, dict):
        raise PackageValidationError("page package JSON must contain an object")
    if exact_fields is not None and set(value) != exact_fields:
        raise PackageValidationError("page package descriptor has unexpected fields")
    return value


def _relationship_source(rels_name: str) -> str:
    if rels_name == "_rels/.rels":
        return ""
    directory = posixpath.dirname(rels_name)
    if directory.endswith("/_rels"):
        directory = posixpath.dirname(directory)
    return posixpath.normpath(posixpath.join(directory, posixpath.basename(rels_name)[:-5]))


def _validate_relationships(archive: zipfile.ZipFile, names: set[str]) -> None:
    for rels_name in sorted(name for name in names if name.endswith(".rels")):
        try:
            root = ET.fromstring(_zip_read_bounded(archive, rels_name))
        except (KeyError, ET.ParseError) as exc:
            raise PackageValidationError(f"invalid relationship part: {rels_name}") from exc
        records = root.findall(f"{{{_PKG_REL}}}Relationship")
        ids = [record.get("Id") for record in records]
        if any(not value for value in ids) or len(ids) != len(set(ids)):
            raise PackageValidationError(f"duplicate or missing relationship id: {rels_name}")
        source = _relationship_source(rels_name)
        for record in records:
            if record.get("TargetMode") == "External":
                raise PackageValidationError(f"external relationship is not allowed: {rels_name}")
            target = record.get("Target")
            if not target:
                raise PackageValidationError(f"missing relationship target: {rels_name}")
            resolved = posixpath.normpath(posixpath.join(posixpath.dirname(source), target))
            if resolved.startswith("../") or resolved not in names:
                raise PackageValidationError(
                    f"relationship target does not exist: {rels_name} -> {target}"
                )


def _validate_archive_limits(archive: zipfile.ZipFile) -> None:
    entries = archive.infolist()
    if len(entries) > MAX_ARCHIVE_ENTRIES:
        raise PackageValidationError("PPTX archive exceeds the maximum entry count")
    total = 0
    for item in entries:
        parts = Path(item.filename.replace("\\", "/")).parts
        if item.filename.startswith(("/", "\\")) or ".." in parts:
            raise PackageValidationError(f"PPTX archive contains an unsafe entry: {item.filename}")
        if item.file_size > MAX_ARCHIVE_ENTRY_BYTES:
            raise PackageValidationError(f"PPTX archive entry is oversized: {item.filename}")
        if item.filename.lower().endswith((".xml", ".rels")) and item.file_size > MAX_XML_PART_BYTES:
            raise PackageValidationError(f"PPTX XML part is oversized: {item.filename}")
        total += item.file_size
        if total > MAX_ARCHIVE_UNCOMPRESSED_BYTES:
            raise PackageValidationError("PPTX archive exceeds the aggregate uncompressed limit")
        if item.file_size >= 1024 * 1024:
            if item.compress_size <= 0 or item.file_size / item.compress_size > MAX_COMPRESSION_RATIO:
                raise PackageValidationError(
                    f"PPTX archive entry has an unsafe compression ratio: {item.filename}"
                )
    try:
        content_types = ET.fromstring(_zip_read_bounded(archive, "[Content_Types].xml"))
    except ET.ParseError as exc:
        raise PackageValidationError("PPTX content-types XML is invalid") from exc
    defaults = {
        node.get("Extension", "").lower(): node.get("ContentType", "").lower()
        for node in content_types.findall(f"{{{_CT}}}Default")
    }
    overrides = {
        node.get("PartName", "").lstrip("/"): node.get("ContentType", "").lower()
        for node in content_types.findall(f"{{{_CT}}}Override")
    }

    def declared_xml(item: zipfile.ZipInfo) -> bool:
        extension = posixpath.basename(item.filename).rpartition(".")[2].lower()
        content_type = overrides.get(item.filename, defaults.get(extension, ""))
        return (
            content_type in {"application/xml", "text/xml"}
            or content_type.endswith("+xml")
            or content_type.endswith("/xml")
        )

    for item in entries:
        if declared_xml(item) and item.file_size > MAX_XML_PART_BYTES:
            raise PackageValidationError(f"PPTX XML part is oversized: {item.filename}")


def _zip_read_bounded(
    archive: zipfile.ZipFile,
    name: str,
    *,
    limit: int | None = None,
) -> bytes:
    if limit is None:
        limit = MAX_XML_PART_BYTES
    try:
        info = archive.getinfo(name)
    except KeyError as exc:
        raise PackageValidationError(f"PPTX package part is missing: {name}") from exc
    if info.file_size > limit:
        raise PackageValidationError(f"PPTX XML part is oversized: {name}")
    with archive.open(info) as handle:
        payload = handle.read(limit + 1)
    if len(payload) > limit:
        raise PackageValidationError(f"PPTX XML part exceeds its read bound: {name}")
    return payload


def _zip_member_sha256(archive: zipfile.ZipFile, name: str) -> str:
    digest = hashlib.sha256()
    with archive.open(name) as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def _canonical_xml_bytes(element: ET.Element) -> bytes:
    """Canonicalize OOXML while retaining text-bearing whitespace."""
    normalized = copy.deepcopy(element)
    text_tag = "{http://schemas.openxmlformats.org/drawingml/2006/main}t"
    for node in normalized.iter():
        if node.tag != text_tag and node.text is not None and not node.text.strip():
            node.text = None
        if node.tail is not None and not node.tail.strip():
            node.tail = None
    payload = ET.tostring(normalized, encoding="unicode")
    return ET.canonicalize(payload, rewrite_prefixes=True).encode("utf-8")


def _zip_member_semantic_sha256(archive: zipfile.ZipFile, name: str) -> str:
    payload = _zip_read_bounded(archive, name)
    if name.lower().endswith((".xml", ".rels")):
        try:
            payload = _canonical_xml_bytes(ET.fromstring(payload))
        except ET.ParseError as exc:
            raise PackageValidationError(f"PPTX XML relationship target is invalid: {name}") from exc
    return hashlib.sha256(payload).hexdigest()


def _slide_fingerprint(
    archive: zipfile.ZipFile,
    names: set[str],
    slide_name: str,
    slide: ET.Element,
) -> str:
    """Hash slide content after binding relationship IDs to target content."""
    rels_name = posixpath.join(
        posixpath.dirname(slide_name),
        "_rels",
        f"{posixpath.basename(slide_name)}.rels",
    )
    relationships: dict[str, tuple[str, str]] = {}
    if rels_name in names:
        rels = ET.fromstring(_zip_read_bounded(archive, rels_name))
        source = _relationship_source(rels_name)
        for record in rels.findall(f"{{{_PKG_REL}}}Relationship"):
            relationship_id = record.get("Id")
            target = record.get("Target")
            reltype = record.get("Type", "")
            if not relationship_id or not target:
                continue
            if record.get("TargetMode") == "External":
                raise PackageValidationError(f"external relationship is not allowed: {rels_name}")
            else:
                resolved = posixpath.normpath(posixpath.join(posixpath.dirname(source), target))
                bound = _zip_member_semantic_sha256(archive, resolved)
            relationships[relationship_id] = (reltype, bound)
    normalized = copy.deepcopy(slide)
    content = normalized.find(f"{{{_P}}}cSld")
    if content is not None:
        content.attrib.pop("name", None)
    for node in normalized.iter():
        for attribute, relationship_id in tuple(node.attrib.items()):
            if attribute not in {
                f"{{{_R}}}embed",
                f"{{{_R}}}id",
                f"{{{_R}}}link",
            }:
                continue
            binding = relationships.get(relationship_id)
            if binding is None:
                raise PackageValidationError(
                    f"slide relationship is unresolved: {slide_name} -> {relationship_id}"
                )
            node.set(attribute, f"{binding[0]}:{binding[1]}")
    return hashlib.sha256(_canonical_xml_bytes(normalized)).hexdigest()


def inspect_editable_pptx(path: Path) -> PptxInspection:
    """Mechanically open and validate an editable PPTX without visual scoring."""
    path = _literal_absolute(path)
    if _is_link_or_reparse(path):
        raise PackageValidationError("PPTX package cannot be a link or reparse point")
    if not path.is_file() or path.resolve(strict=True) != path:
        raise PackageValidationError("PPTX package path is redirected or missing")
    try:
        with zipfile.ZipFile(path) as archive:
            _validate_archive_limits(archive)
            listed = archive.namelist()
            names = set(listed)
            if len(listed) != len(names) or archive.testzip() is not None:
                raise PackageValidationError("PPTX package contains duplicate or corrupt parts")
            if not _REQUIRED_OPC_PARTS.issubset(names):
                raise PackageValidationError("PPTX package is missing required OPC parts")
            _validate_relationships(archive, names)
            try:
                presentation = ET.fromstring(_zip_read_bounded(archive, "ppt/presentation.xml"))
                relationships = ET.fromstring(
                    _zip_read_bounded(archive, "ppt/_rels/presentation.xml.rels")
                )
            except ET.ParseError as exc:
                raise PackageValidationError("PPTX presentation XML is invalid") from exc
            rel_targets = {
                record.get("Id"): record.get("Target")
                for record in relationships.findall(f"{{{_PKG_REL}}}Relationship")
            }
            slide_ids = presentation.findall(f".//{{{_P}}}sldId")
            ordered = [rel_targets.get(item.get(f"{{{_R}}}id")) for item in slide_ids]
            if any(not isinstance(target, str) for target in ordered):
                raise PackageValidationError("PPTX slide order contains an unresolved relationship")
            slide_names = [posixpath.normpath(f"ppt/{target}") for target in ordered]
            if len(slide_names) != len(set(slide_names)) or any(name not in names for name in slide_names):
                raise PackageValidationError("PPTX slide order is invalid")
            counts: list[int] = []
            fingerprints: list[str] = []
            page_ids: list[int | None] = []
            editable_tags = {
                f"{{{_P}}}sp",
                f"{{{_P}}}pic",
                f"{{{_P}}}graphicFrame",
                f"{{{_P}}}cxnSp",
                f"{{{_P}}}grpSp",
            }
            for slide_name in slide_names:
                try:
                    slide = ET.fromstring(_zip_read_bounded(archive, slide_name))
                except ET.ParseError as exc:
                    raise PackageValidationError(f"PPTX slide XML is invalid: {slide_name}") from exc
                sp_tree = slide.find(f".//{{{_P}}}spTree")
                children = list(sp_tree) if sp_tree is not None else []
                counts.append(sum(child.tag in editable_tags for child in children))
                fingerprints.append(_slide_fingerprint(archive, names, slide_name, slide))
                content = slide.find(f"{{{_P}}}cSld")
                name = content.get("name", "") if content is not None else ""
                match = re.fullmatch(r"editable-ppt-word-page:(\d+)", name)
                page_ids.append(int(match.group(1)) if match else None)
    except PackageValidationError:
        raise
    except (OSError, KeyError, zipfile.BadZipFile) as exc:
        raise PackageValidationError("PPTX package is not openable") from exc

    try:
        opened = Presentation(path)
    except Exception as exc:
        raise PackageValidationError("PowerPoint package cannot be opened by the PPTX runtime") from exc
    if len(opened.slides) != len(counts):
        raise PackageValidationError("PPTX open result disagrees with its slide relationships")
    return PptxInspection(len(counts), tuple(counts), tuple(fingerprints), tuple(page_ids))


def validate_pptx_canvas(path: Path, profile: Mapping[str, Any]) -> None:
    """Require every reconstructed page to use the confirmed no-crop canvas."""
    if not isinstance(profile, Mapping):
        raise PackageValidationError("confirmed canvas profile is missing")
    if profile.get("fit") != "contain" or profile.get("allow_crop") is not False:
        raise PackageValidationError("confirmed canvas profile must use contain without crop")
    width = profile.get("slide_width_inches")
    height = profile.get("slide_height_inches")
    if not isinstance(width, (int, float)) or not isinstance(height, (int, float)):
        raise PackageValidationError("confirmed canvas slide size is invalid")
    opened = Presentation(path)
    actual_width = float(opened.slide_width) / 914400.0
    actual_height = float(opened.slide_height) / 914400.0
    if abs(actual_width - float(width)) > 0.002 or abs(actual_height - float(height)) > 0.002:
        raise PackageValidationError(
            "reconstructed page slide size does not match the confirmed canvas ratio"
        )


def create_page_package(
    project: Path,
    *,
    page_number: int,
    cache_key: str,
    pptx: Path,
    output: Path,
) -> Path:
    """Create the exact descriptor Task 6 can seal for one reconstructed page."""
    project = require_plain_project(project)
    if type(page_number) is not int or page_number < 1:
        raise PackageValidationError("page number must be a positive integer")
    if not _valid_sha256(cache_key):
        raise PackageValidationError("page cache key must be a lowercase SHA-256 digest")
    pptx_path = _project_file(project, pptx)
    output_path = _project_file(project, output, must_exist=False)
    inspection = inspect_editable_pptx(pptx_path)
    if inspection.slide_count != 1:
        raise PackageValidationError("a reconstructed page package must contain exactly one slide")
    if inspection.editable_object_counts[0] < 1:
        raise PackageValidationError("a reconstructed page package must contain editable objects")
    payload = {
        "schema_version": PAGE_PACKAGE_SCHEMA,
        "cache_key": cache_key,
        "pptx": pptx_path.relative_to(project).as_posix(),
        "pptx_sha256": _sha256(pptx_path),
        "editable_object_count": inspection.editable_object_counts[0],
        "slide_fingerprint": inspection.slide_fingerprints[0],
    }
    output_path.parent.mkdir(parents=True, exist_ok=True)
    temporary = output_path.with_name(f".{output_path.name}.{uuid.uuid4().hex}.tmp")
    try:
        temporary.write_text(
            json.dumps(payload, ensure_ascii=False, sort_keys=True, separators=(",", ":")) + "\n",
            encoding="utf-8",
        )
        os.replace(temporary, output_path)
    finally:
        if temporary.exists():
            temporary.unlink()
    return output_path


def _current_cache(job: Mapping[str, Any]) -> tuple[str, Mapping[str, Any]]:
    cache = job.get("cache")
    if not isinstance(cache, Mapping):
        raise PackageValidationError("completed page has no current cache identity")
    key = cache.get("key")
    identity = cache.get("identity")
    if (
        not _valid_sha256(key)
        or not isinstance(identity, Mapping)
        or set(identity) != CURRENT_CACHE_IDENTITY_FIELDS
        or canonical_sha256(identity) != key
    ):
        raise PackageValidationError("completed page cache identity is stale or invalid")
    return key, identity


def _load_current_run(project: Path) -> dict[str, Any]:
    try:
        run = json.loads((project / "workflow_run.json").read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError) as exc:
        raise PackageValidationError("current workflow state is unavailable") from exc
    if not isinstance(run, dict) or run.get("workflow_contract_version") != "word-only-v1":
        raise PackageValidationError("page package requires the current Word-only workflow")
    return run


def _authoritative_job(
    run: Mapping[str, Any],
    supplied: Mapping[str, Any],
) -> Mapping[str, Any]:
    page_number = supplied.get("page_number")
    jobs = run.get("jobs")
    if not isinstance(jobs, list):
        raise PackageValidationError("current workflow page jobs are invalid")
    current = next(
        (
            item
            for item in jobs
            if isinstance(item, Mapping) and item.get("page_number") == page_number
        ),
        None,
    )
    if current is None:
        raise PackageValidationError("completed page is absent from current workflow state")
    supplied_cache = supplied.get("cache")
    current_cache = current.get("cache")
    if (
        supplied.get("status") != current.get("status")
        or not isinstance(supplied_cache, Mapping)
        or not isinstance(current_cache, Mapping)
        or supplied_cache.get("key") != current_cache.get("key")
    ):
        raise PackageValidationError("supplied page package state is not current")
    return current


def load_completed_page_package(
    project: Path,
    job: Mapping[str, Any],
    *,
    _run: Mapping[str, Any] | None = None,
) -> PagePackage:
    """Reload one completed package through Task 6's immutable current cache."""
    project = require_plain_project(project)
    run = _run if _run is not None else _load_current_run(project)
    job = _authoritative_job(run, job)
    page_number = job.get("page_number")
    if type(page_number) is not int or page_number < 1:
        raise PackageValidationError("completed page number is invalid")
    if job.get("status") != "complete":
        raise PackageValidationError(f"Word page {page_number} is not complete")
    qa = job.get("qa_result")
    if not isinstance(qa, Mapping) or qa.get("status") not in {"pass", "pass_with_advisory"}:
        raise PackageValidationError(f"Word page {page_number} has not passed page QA")
    key, identity = _current_cache(job)
    try:
        expected_cache = cache_record(project, run, job)
    except Exception as exc:
        raise PackageValidationError(
            f"Word page {page_number} current cache identity cannot be recomputed: {exc}"
        ) from exc
    if expected_cache.get("key") != key or expected_cache.get("identity") != dict(identity):
        raise PackageValidationError(f"Word page {page_number} current cache identity changed")
    hit = CacheStore(project).lookup("pages", key)
    if hit is None:
        raise PackageValidationError(f"Word page {page_number} cache package is missing or replaced")
    stored_identity = hit.manifest.get("cache_identity")
    outputs = hit.manifest.get("outputs")
    if (
        not isinstance(stored_identity, Mapping)
        or set(stored_identity) != CURRENT_CACHE_IDENTITY_FIELDS
        or canonical_sha256(stored_identity) != key
        or not isinstance(outputs, Mapping)
        or set(outputs) != {"reconstruction"}
        or not isinstance(outputs.get("reconstruction"), str)
    ):
        raise PackageValidationError(f"Word page {page_number} cache manifest is stale")
    manifest = _project_file(hit.path, "manifest.json")
    descriptor = _project_file(hit.path, outputs["reconstruction"])
    payload = _read_object(descriptor, exact_fields=PAGE_PACKAGE_FIELDS)
    if (
        payload.get("schema_version") != PAGE_PACKAGE_SCHEMA
        or payload.get("cache_key") != key
        or not _valid_sha256(payload.get("pptx_sha256"))
        or type(payload.get("editable_object_count")) is not int
        or payload["editable_object_count"] < 1
        or not _valid_sha256(payload.get("slide_fingerprint"))
    ):
        raise PackageValidationError(f"Word page {page_number} package identity is invalid")
    pptx = _project_file(project, payload.get("pptx", ""))
    actual_sha256 = _sha256(pptx)
    if actual_sha256 != payload["pptx_sha256"]:
        raise PackageValidationError(f"Word page {page_number} PPTX was replaced after reconstruction")
    inspection = inspect_editable_pptx(pptx)
    if inspection.slide_count != 1:
        raise PackageValidationError(f"Word page {page_number} package does not contain one slide")
    if inspection.editable_object_counts != (payload["editable_object_count"],):
        raise PackageValidationError(f"Word page {page_number} editable-object identity changed")
    if inspection.slide_fingerprints != (payload["slide_fingerprint"],):
        raise PackageValidationError(f"Word page {page_number} slide identity changed")
    return PagePackage(
        page_number=page_number,
        cache_key=key,
        cache_manifest=manifest,
        cache_manifest_sha256=_sha256(manifest),
        cache_artifact=descriptor,
        cache_artifact_sha256=_sha256(descriptor),
        pptx_path=pptx,
        pptx_sha256=actual_sha256,
        editable_object_count=payload["editable_object_count"],
        slide_fingerprint=payload["slide_fingerprint"],
    )


def load_locked_page_packages(
    project: Path,
    locked_page_numbers: tuple[int, ...],
    jobs: list[Mapping[str, Any]],
) -> tuple[PagePackage, ...]:
    """Return only completed packages in the authoritative locked Word order."""
    project = require_plain_project(project)
    run = _load_current_run(project)
    current_jobs = run.get("jobs")
    if not isinstance(current_jobs, list):
        raise PackageValidationError("current workflow page jobs are invalid")
    numbers = [job.get("page_number") for job in current_jobs if isinstance(job, Mapping)]
    if numbers != list(locked_page_numbers):
        raise PackageValidationError("locked Word-page order does not match current page jobs")
    if [job.get("page_number") for job in jobs] != numbers:
        raise PackageValidationError("supplied page jobs are not in the current locked order")
    return tuple(
        load_completed_page_package(project, job, _run=run)
        for job in current_jobs
    )
