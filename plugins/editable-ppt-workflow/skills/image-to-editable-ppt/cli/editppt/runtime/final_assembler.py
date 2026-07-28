"""Atomic current-only assembly of completed Word-page reconstruction packages."""

from __future__ import annotations

import copy
import hashlib
import json
import os
import stat
import uuid
from contextlib import contextmanager
from dataclasses import dataclass
from io import BytesIO
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TYPE as RT

try:
    from .editable_page_cache import (
        PackageValidationError,
        PagePackage,
        inspect_editable_pptx,
        load_locked_page_packages,
        project_output_path,
        require_plain_project,
        require_project_file,
    )
except ImportError:  # direct runtime script entrypoint
    from editable_page_cache import (
        PackageValidationError,
        PagePackage,
        inspect_editable_pptx,
        load_locked_page_packages,
        project_output_path,
        require_plain_project,
        require_project_file,
    )


CURRENT_CONTRACT = "word-only-v1"
ASSEMBLY_POLICY_VERSION = "word-pages-v1"
_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
_RELATIONSHIP_ATTRIBUTES = {
    f"{{{_R}}}embed",
    f"{{{_R}}}id",
    f"{{{_R}}}link",
}


@dataclass(frozen=True)
class AssemblyPlan:
    project_root: Path
    locked_page_numbers: tuple[int, ...]
    assembly_policy_version: str = ASSEMBLY_POLICY_VERSION

    def __post_init__(self) -> None:
        project = require_plain_project(self.project_root)
        object.__setattr__(self, "project_root", project)
        if (
            type(self.locked_page_numbers) is not tuple
            or not self.locked_page_numbers
            or any(type(page) is not int or page < 1 for page in self.locked_page_numbers)
            or len(set(self.locked_page_numbers)) != len(self.locked_page_numbers)
        ):
            raise PackageValidationError("locked Word-page order is invalid")
        if self.assembly_policy_version != ASSEMBLY_POLICY_VERSION:
            raise PackageValidationError("unsupported current assembly policy")


@dataclass(frozen=True)
class AssemblyReceipt:
    output_path: Path
    output_sha256: str
    page_count: int
    slide_order: tuple[int, ...]
    writer_count: int
    assembly_policy_version: str
    page_packages: tuple[PagePackage, ...]
    editable_object_counts: tuple[int, ...]


def _sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def _load_run(project: Path) -> dict[str, Any]:
    project = require_plain_project(project)
    path = require_project_file(project, "workflow_run.json")
    try:
        run = json.loads(path.read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError) as exc:
        raise PackageValidationError("current workflow state is unavailable") from exc
    if not isinstance(run, dict) or run.get("workflow_contract_version") != CURRENT_CONTRACT:
        raise PackageValidationError("final assembly accepts only the current Word-page workflow")
    pagination = run.get("pagination")
    jobs = run.get("jobs")
    if not isinstance(pagination, dict) or not isinstance(jobs, list) or not jobs:
        raise PackageValidationError("locked Word-page state is invalid")
    return run


def _load_packages(plan: AssemblyPlan) -> tuple[PagePackage, ...]:
    project = require_plain_project(plan.project_root)
    run = _load_run(project)
    pagination = run["pagination"]
    locked = pagination.get("locked_page_order")
    if (
        locked != list(plan.locked_page_numbers)
        or pagination.get("page_count") != len(plan.locked_page_numbers)
    ):
        raise PackageValidationError("assembly plan no longer matches the locked Word-page contract")
    return load_locked_page_packages(project, plan.locked_page_numbers, run["jobs"])


def _remap_relationships(source_slide, destination_slide) -> dict[str, str]:
    mapping: dict[str, str] = {}
    for relationship in source_slide.part.rels.values():
        old_id = relationship.rId
        if relationship.reltype == RT.SLIDE_LAYOUT:
            continue
        if relationship.reltype == RT.IMAGE:
            _image_part, new_id = destination_slide.part.get_or_add_image_part(
                BytesIO(relationship.target_part.blob)
            )
            mapping[old_id] = new_id
            continue
        if relationship.is_external:
            mapping[old_id] = destination_slide.part.relate_to(
                relationship.target_ref,
                relationship.reltype,
                is_external=True,
            )
            continue
        raise PackageValidationError(
            "reconstructed page uses an unsupported internal slide relationship: "
            f"{relationship.reltype}"
        )
    return mapping


def _rewrite_relationship_ids(element, mapping: dict[str, str]) -> None:
    for node in element.iter():
        for attribute, old_id in tuple(node.attrib.items()):
            if attribute not in _RELATIONSHIP_ATTRIBUTES:
                continue
            if old_id not in mapping:
                raise PackageValidationError(
                    f"reconstructed slide contains an unresolved relationship: {old_id}"
                )
            node.set(attribute, mapping[old_id])


def _copy_page_slide(source_path: Path, destination: Presentation) -> None:
    source = Presentation(source_path)
    if len(source.slides) != 1:
        raise PackageValidationError("reconstructed page package must contain one slide")
    if source.slide_width != destination.slide_width or source.slide_height != destination.slide_height:
        raise PackageValidationError("all reconstructed Word pages must share one slide size")
    source_slide = source.slides[0]
    destination_slide = destination.slides.add_slide(destination.slide_layouts[6])
    mapping = _remap_relationships(source_slide, destination_slide)

    copied_content = copy.deepcopy(source_slide.element.cSld)
    _rewrite_relationship_ids(copied_content, mapping)
    destination_slide.element.replace(destination_slide.element.cSld, copied_content)
    source_color_map = source_slide.element.clrMapOvr
    if source_color_map is not None:
        destination_color_map = destination_slide.element.clrMapOvr
        copied_color_map = copy.deepcopy(source_color_map)
        if destination_color_map is None:
            destination_slide.element.append(copied_color_map)
        else:
            destination_slide.element.replace(destination_color_map, copied_color_map)


@contextmanager
def _single_writer(path: Path):
    flags = os.O_RDWR | os.O_CREAT | getattr(os, "O_BINARY", 0) | getattr(os, "O_NOFOLLOW", 0)
    descriptor = os.open(path, flags, 0o600)
    handle = os.fdopen(descriptor, "r+b", closefd=True)
    locked = False
    try:
        info = path.lstat()
        reparse = bool(
            getattr(info, "st_file_attributes", 0)
            & getattr(stat, "FILE_ATTRIBUTE_REPARSE_POINT", 0x400)
        )
        if (
            path.is_symlink()
            or reparse
            or not stat.S_ISREG(info.st_mode)
            or info.st_nlink != 1
            or path.resolve() != path
        ):
            raise PackageValidationError("final writer lock must be a project-local regular file")
        if handle.seek(0, os.SEEK_END) == 0:
            handle.write(b"\0")
            handle.flush()
            os.fsync(handle.fileno())
        try:
            if os.name == "nt":
                import msvcrt

                handle.seek(0)
                msvcrt.locking(handle.fileno(), msvcrt.LK_NBLCK, 1)
            else:
                import fcntl

                fcntl.flock(handle.fileno(), fcntl.LOCK_EX | fcntl.LOCK_NB)
            locked = True
        except OSError as exc:
            raise FileExistsError(f"another final writer owns {path}") from exc
        yield
    finally:
        if locked:
            if os.name == "nt":
                import msvcrt

                handle.seek(0)
                msvcrt.locking(handle.fileno(), msvcrt.LK_UNLCK, 1)
            else:
                import fcntl

                fcntl.flock(handle.fileno(), fcntl.LOCK_UN)
        handle.close()


def _validate_output(path: Path, packages: tuple[PagePackage, ...]) -> tuple[int, ...]:
    inspection = inspect_editable_pptx(path)
    expected = tuple(package.editable_object_count for package in packages)
    if inspection.slide_count != len(packages):
        raise PackageValidationError("assembled deck does not have the exact Word-page count")
    if inspection.editable_object_counts != expected:
        raise PackageValidationError("assembled deck editable-object inventory changed")
    if inspection.slide_fingerprints != tuple(package.slide_fingerprint for package in packages):
        raise PackageValidationError("assembled deck does not preserve the locked Word-page order")
    return inspection.editable_object_counts


def assemble_deck(plan: AssemblyPlan, output: Path) -> AssemblyReceipt:
    """Assemble exactly the locked Word pages; no synthetic slide is permitted."""
    project = require_plain_project(plan.project_root)
    out = project_output_path(project, output)
    out.parent.mkdir(parents=True, exist_ok=True)
    require_plain_project(out.parent)

    lock_path = project_output_path(project, out.parent / f".{out.name}.writer.lock")
    temporary = project_output_path(project, out.parent / f".{out.name}.{uuid.uuid4().hex}.tmp")
    try:
        with _single_writer(lock_path):
            packages = _load_packages(plan)
            first = Presentation(packages[0].pptx_path)
            destination = Presentation()
            destination.slide_width = first.slide_width
            destination.slide_height = first.slide_height
            for package in packages:
                _copy_page_slide(package.pptx_path, destination)
            destination.save(temporary)
            editable_counts = _validate_output(temporary, packages)

            # Reload every current cache and referenced PPTX after assembly. Any
            # missing/replaced package aborts publication rather than producing a
            # deck from stale authority.
            replayed = _load_packages(plan)
            if replayed != packages:
                raise PackageValidationError("page package authority changed during assembly")
            os.replace(temporary, out)
            try:
                _validate_output(out, replayed)
                final_replay = _load_packages(plan)
                if final_replay != packages:
                    raise PackageValidationError("page package authority changed during publication")
            except BaseException:
                try:
                    out.unlink()
                except FileNotFoundError:
                    pass
                except OSError as cleanup_error:
                    raise PackageValidationError(
                        "rejected final output could not be removed after authority replay"
                    ) from cleanup_error
                raise
            return AssemblyReceipt(
                output_path=out,
                output_sha256=_sha256(out),
                page_count=len(packages),
                slide_order=plan.locked_page_numbers,
                writer_count=1,
                assembly_policy_version=plan.assembly_policy_version,
                page_packages=packages,
                editable_object_counts=editable_counts,
            )
    finally:
        try:
            temporary.unlink()
        except FileNotFoundError:
            pass
