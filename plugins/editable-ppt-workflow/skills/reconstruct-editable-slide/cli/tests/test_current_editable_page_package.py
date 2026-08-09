from __future__ import annotations

import json
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches

from editppt.runtime.editable_page_cache import (
    PackageValidationError,
    create_page_package,
)


def _write_editable_pptx(path, *, slide_count: int) -> None:
    presentation = Presentation()
    presentation.slides.add_slide(presentation.slide_layouts[6])
    for index in range(1, slide_count):
        presentation.slides.add_slide(presentation.slide_layouts[6])
    for index, slide in enumerate(presentation.slides, start=1):
        text_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        text_box.text = f"Editable page {index}"
    presentation.save(path)


def test_create_page_package_records_one_real_editable_slide(tmp_path):
    project = tmp_path / "project"
    project.mkdir()
    pptx = project / "page.pptx"
    descriptor = project / "packages" / "page-001.json"
    _write_editable_pptx(pptx, slide_count=1)

    result = create_page_package(
        project,
        page_number=1,
        cache_key="a" * 64,
        pptx=pptx,
        output=descriptor,
    )

    payload = json.loads(result.read_text(encoding="utf-8"))
    assert payload == {
        "schema_version": "editable-page-package-v1",
        "cache_key": "a" * 64,
        "pptx": "page.pptx",
        "pptx_sha256": payload["pptx_sha256"],
        "editable_object_count": 1,
        "slide_fingerprint": payload["slide_fingerprint"],
    }
    assert len(payload["pptx_sha256"]) == 64
    assert len(payload["slide_fingerprint"]) == 64


def test_create_page_package_rejects_a_multi_slide_reconstruction(tmp_path):
    project = tmp_path / "project"
    project.mkdir()
    pptx = project / "two-slides.pptx"
    _write_editable_pptx(pptx, slide_count=2)

    with pytest.raises(PackageValidationError, match="exactly one slide"):
        create_page_package(
            project,
            page_number=1,
            cache_key="b" * 64,
            pptx=pptx,
            output=project / "packages" / "page-001.json",
        )


def test_installed_runtime_does_not_guess_a_source_tree_sibling():
    """An installed wheel cannot rely on the repository's sibling skill layout."""
    runtime = Path(__file__).resolve().parents[1] / "editppt" / "runtime"

    for name in ("editable_page_cache.py", "finalize_deck_run.py", "record_page_result.py"):
        source = (runtime / name).read_text(encoding="utf-8")
        assert "parents[4]" not in source
        assert '"run-word-to-ppt-workflow" / "scripts"' not in source
