from __future__ import annotations

import json
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Cm

from editppt.runtime.build_pptx_from_manifest import normalize_manifest, write_pptx
from editppt.runtime.fixed_region_runtime import CONTENT_BOX, SLIDE, require_content_box
from fixed_frame import apply_fixed_frame, inspect_fixed_frame
from fixed_region_contract import (
    BODY_BOX_CM,
    CONTRACT_VERSION,
    FOOTER_LINE,
    LOGO_BOX_CM,
    PAGE_NUMBER_BOX_CM,
    TITLE_BOX_CM,
    fixed_frame_execution,
)

SOURCE_WIDTH_PX = 1904
SOURCE_HEIGHT_PX = 896


def _manifest() -> dict:
    return {
        "schema_version": 1,
        "workflow_contract_version": CONTRACT_VERSION,
        "slide": dict(SLIDE),
        "content_box": dict(CONTENT_BOX),
        "source": {"width_px": SOURCE_WIDTH_PX, "height_px": SOURCE_HEIGHT_PX},
        "shapes": [
            {
                "type": "rect",
                "box_px": [0, 0, SOURCE_WIDTH_PX, SOURCE_HEIGHT_PX],
                "fill": "#FFFFFF",
                "stroke": "none",
            }
        ],
        "images": [],
        "text_boxes": [],
    }


def _execution(*, size: float = 28, color: str = "#123456") -> dict:
    return {
        "schema_version": "2.0",
        "hard_constraints": {
            "title_color": color,
            "typography": {
                "heading": {"cjk": "SimHei", "latin": "Arial"},
                "type_scale_pt": {"page_title": size},
            },
        },
        "fixed_frame": {"title_color": color, **fixed_frame_execution()},
    }


def _logo(path: Path) -> Path:
    path.write_text(
        '<svg xmlns="http://www.w3.org/2000/svg" width="120" height="40"><rect width="120" height="40"/></svg>',
        encoding="utf-8",
    )
    return path


def test_17_by_8_manifest_maps_source_pixels_directly_into_the_authoritative_body_box() -> None:
    normalized = normalize_manifest(_manifest())
    shape = normalized["shapes"][0]
    assert shape["left"] == pytest.approx(BODY_BOX_CM["x"] / 2.54)
    assert shape["top"] == pytest.approx(BODY_BOX_CM["y"] / 2.54)
    assert shape["width"] == pytest.approx(BODY_BOX_CM["w"] / 2.54)
    assert shape["height"] == pytest.approx(BODY_BOX_CM["h"] / 2.54)


def test_unexpected_16_by_9_source_is_contained_without_stretching() -> None:
    manifest = _manifest()
    source_size = (1920, 1080)
    manifest["source"] = {"width_px": source_size[0], "height_px": source_size[1]}
    manifest["shapes"][0]["box_px"] = [0, 0, source_size[0], source_size[1]]

    shape = normalize_manifest(manifest)["shapes"][0]

    expected_width_cm = BODY_BOX_CM["h"] * source_size[0] / source_size[1]
    expected_left_cm = BODY_BOX_CM["x"] + (BODY_BOX_CM["w"] - expected_width_cm) / 2
    assert shape["left"] == pytest.approx(expected_left_cm / 2.54)
    assert shape["top"] == pytest.approx(BODY_BOX_CM["y"] / 2.54)
    assert shape["width"] == pytest.approx(expected_width_cm / 2.54)
    assert shape["height"] == pytest.approx(BODY_BOX_CM["h"] / 2.54)
    assert shape["width"] / shape["height"] == pytest.approx(source_size[0] / source_size[1])


def test_content_box_accepts_at_most_point_one_percent_relative_error() -> None:
    within = dict(CONTENT_BOX)
    within["width"] *= 1.001
    require_content_box(within)

    outside = dict(CONTENT_BOX)
    outside["width"] *= 1.00101
    with pytest.raises(ValueError, match="content_box"):
        require_content_box(outside)


@pytest.mark.parametrize(
    "mutation, message",
    [
        (lambda value: value.update(workflow_contract_version="body-frame-v2"), "fixed-canvas-cm-v2"),
        (lambda value: value["content_box"].update(left=0), "content_box"),
        (lambda value: value["slide"].update(width=13.333333), "slide"),
        (lambda value: value["slide"].update(background="#112233"), "background"),
        (lambda value: value["source"].update(width_px=0), "positive integers"),
    ],
)
def test_manifest_rejects_every_non_authoritative_geometry(mutation, message) -> None:
    manifest = _manifest()
    mutation(manifest)
    with pytest.raises(ValueError, match=message):
        normalize_manifest(manifest)


def test_fixed_layer_does_not_move_reconstructed_body_and_is_complete(tmp_path: Path) -> None:
    manifest_path = tmp_path / "manifest.json"
    manifest_path.write_text(json.dumps(_manifest()), encoding="utf-8")
    pptx = tmp_path / "page.pptx"
    write_pptx(_manifest(), pptx, manifest_path)
    before_shape = Presentation(pptx).slides[0].shapes[0]
    before = (before_shape.left, before_shape.top, before_shape.width, before_shape.height)

    apply_fixed_frame(
        pptx,
        page_title="本页结论",
        page_number=7,
        style_execution=_execution(size=72),
        logo_svg=_logo(tmp_path / "logo.svg"),
    )

    deck = Presentation(pptx)
    body = next(shape for shape in deck.slides[0].shapes if not shape.name.startswith("fixed-frame-"))
    assert (body.left, body.top, body.width, body.height) == before
    by_name = {shape.name: shape for shape in deck.slides[0].shapes}
    assert set(by_name) >= {
        "fixed-frame-title",
        "fixed-frame-logo",
        "fixed-frame-footer",
        "fixed-frame-page-number",
    }
    title_run = by_name["fixed-frame-title"].text_frame.paragraphs[0].runs[0]
    assert title_run.font.name == "SimHei"
    assert title_run.font.size.pt == 72
    assert str(title_run.font.color.rgb) == "123456"
    assert str(by_name["fixed-frame-footer"].fill.fore_color.rgb) == "B8C0CC"
    page_run = by_name["fixed-frame-page-number"].text_frame.paragraphs[0].runs[0]
    assert str(page_run.font.color.rgb) == "6B7280"
    assert inspect_fixed_frame(
        pptx,
        expected_title="本页结论",
        expected_page_number=7,
        style_execution=_execution(size=72),
    )["passed"] is True


def test_fixed_layer_rejects_missing_logo_and_out_of_box_body(tmp_path: Path) -> None:
    manifest = _manifest()
    manifest["shapes"][0]["box_px"] = [-1, 0, SOURCE_WIDTH_PX, SOURCE_HEIGHT_PX]
    manifest_path = tmp_path / "manifest.json"
    manifest_path.write_text(json.dumps(manifest), encoding="utf-8")
    pptx = tmp_path / "page.pptx"
    write_pptx(manifest, pptx, manifest_path)
    with pytest.raises(ValueError, match="existing SVG"):
        apply_fixed_frame(
            pptx,
            page_title="标题",
            page_number=1,
            style_execution=_execution(),
            logo_svg=tmp_path / "missing.svg",
        )
    with pytest.raises(ValueError, match="outside the authoritative content_box"):
        apply_fixed_frame(
            pptx,
            page_title="标题",
            page_number=1,
            style_execution=_execution(),
            logo_svg=_logo(tmp_path / "logo.svg"),
        )
