import json
import zipfile
from pathlib import Path

import pytest
from pptx import Presentation

from editppt.runtime.build_pptx_from_manifest import normalize_manifest, write_pptx
from editppt.runtime.fixed_region_runtime import CONTENT_BOX, SLIDE
from editppt.runtime.validate_pptx import ALLOWED_SOURCE_TYPES


def _manifest(width=1700, height=800):
    return {
        "workflow_contract_version": "fixed-canvas-cm-v2",
        "reconstruction_contract_version": "editable-image-v3",
        "slide": dict(SLIDE),
        "content_box": dict(CONTENT_BOX),
        "source": {"width_px": width, "height_px": height},
        "text_boxes": [{"object_id": "word-p1", "name": "body-paragraph-1", "text": "权威正文", "box_px": [80, 60, 700, 90]}],
        "tables": [{
            "object_id": "word-t1", "name": "body-table-1", "box_px": [80, 220, 1100, 300],
            "rows": [["项目", "数值"], ["投资额", "50万元"]],
            "font_size": 12, "font_color": "#000000", "cell_fill": "#FFFFFF", "cell_margin_px": 8,
        }],
        "shapes": [{"object_id": "decor-1", "name": "decorative-panel", "type": "rect", "box_px": [40, 30, 1500, 650], "fill": "#f4f4f4"}],
        "images": [],
    }


def test_authentic_published_source_is_a_first_class_provenance_type():
    assert "authentic-published-source" in ALLOWED_SOURCE_TYPES


def test_v4_rejects_unexpected_16_by_9_instead_of_containing():
    with pytest.raises(ValueError, match="17:8"):
        normalize_manifest(_manifest(1600, 900))


def test_v4_builds_stable_named_text_shape_and_native_table(tmp_path: Path):
    manifest = _manifest()
    manifest_path = tmp_path / "manifest.json"
    manifest_path.write_text(json.dumps(manifest), encoding="utf-8")
    output = tmp_path / "page.pptx"
    write_pptx(manifest, output, manifest_path)

    deck = Presentation(output)
    names = {shape.name for shape in deck.slides[0].shapes}
    assert {"body-paragraph-1", "body-table-1", "decorative-panel"} <= names
    table = next(shape.table for shape in deck.slides[0].shapes if shape.has_table)
    assert table.cell(1, 1).text == "50万元"
    with zipfile.ZipFile(output) as archive:
        xml = archive.read("ppt/slides/slide1.xml").decode("utf-8")
    assert 'descr="object_id:word-p1"' in xml
    assert 'descr="object_id:word-t1"' in xml
    assert '<a:srgbClr val="000000"' in xml
    assert '<a:srgbClr val="FFFFFF"' in xml
    assert 'marL="' in xml and 'marR="' in xml and 'marT="' in xml and 'marB="' in xml


def test_v4_native_table_rejects_unresolved_default_cell_visuals():
    manifest = _manifest()
    for field in ("font_size", "font_color", "cell_fill", "cell_margin_px"):
        manifest["tables"][0].pop(field)
    with pytest.raises(ValueError, match="table cells require explicit"):
        normalize_manifest(manifest)
