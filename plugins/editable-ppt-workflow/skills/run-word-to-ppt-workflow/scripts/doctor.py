"""Report required and optional runtime capabilities without modifying the machine."""

from __future__ import annotations

import argparse
import importlib.util
import json
import os
import platform
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path

PLUGIN_SCRIPTS = Path(__file__).resolve().parents[3] / "scripts"
if str(PLUGIN_SCRIPTS) not in sys.path:
    sys.path.insert(0, str(PLUGIN_SCRIPTS))

from runtime_office import resolve_soffice
from background_text_detector import capability_status


REQUIRED_MODULES = ["flask", "jsonschema", "PIL", "pypdf", "pypdfium2", "docx", "pptx"]
EXPECTED_EDITPPT_CLI_VERSION = "0.3.0"
POWERPOINT_COM_LIFECYCLE_TIMEOUT_SECONDS = 60


def _editable_python_path(editppt_path: str | None) -> Path | None:
    candidates: list[Path] = []
    installed_executable = os.environ.get("EDITPPT_EXE")
    if installed_executable:
        candidates.append(Path(installed_executable).parent / "python.exe")
    if editppt_path:
        candidates.append(Path(editppt_path).parent / "python.exe")
    return next((path for path in candidates if path.is_file()), None)


def background_text_detection_status() -> dict:
    """Report the mandatory local zero-token background text detector."""
    return capability_status()


def resolve_command(command: str) -> str | None:
    if command == "soffice":
        return resolve_soffice()
    home = Path.home()
    runtime = home / ".codex/plugin-runtimes/editable-ppt-workflow-fixed-canvas-cm-v2"
    candidates: dict[str, list[Path]] = {
        "editppt": [
            Path(os.environ["EDITPPT_EXE"]) if os.getenv("EDITPPT_EXE") else Path("__unset__"),
            runtime / "editable-ppt/Scripts/editppt.exe",
            home / ".codex/bin/editppt.CMD",
        ],
        "officecli": [
            Path(os.environ["OFFICECLI_EXE"]) if os.getenv("OFFICECLI_EXE") else Path("__unset__"),
            Path(os.getenv("LOCALAPPDATA", "")) / "OfficeCLI/officecli.exe",
            home / ".codex/bin/officecli.CMD",
        ],
    }
    for candidate in candidates.get(command, []):
        if candidate.is_file():
            return str(candidate.resolve())
    return shutil.which(command)


def command_version(command: str, args: list[str]) -> dict:
    path = resolve_command(command)
    if not path:
        return {"available": False, "path": None, "detail": "not found"}
    try:
        completed = subprocess.run([path, *args], capture_output=True, text=True, timeout=8, check=False)
        detail = (completed.stdout or completed.stderr).strip().splitlines()
        return {"available": completed.returncode == 0, "path": path, "detail": detail[0] if detail else "detected"}
    except Exception as exc:  # diagnostic must keep running
        return {"available": False, "path": path, "detail": str(exc)}


def cjk_font_status() -> dict:
    font_dir = Path(os.getenv("WINDIR", "C:/Windows")) / "Fonts"
    required_families = {
        "Microsoft YaHei": ["msyh.ttc", "msyhbd.ttc"],
        "DengXian": ["Deng.ttf", "Dengb.ttf"],
    }
    families = {name: any((font_dir / filename).is_file() for filename in files) for name, files in required_families.items()}
    return {"font_dir": str(font_dir), "families": families, "ready": any(families.values())}


def _create_officecli_smoke(officecli_path: str, pptx: Path) -> dict:
    commands = [
        [officecli_path, "create", str(pptx)],
        [officecli_path, "add", str(pptx), "/", "--type", "slide", "--prop", "title=运行环境测试"],
        [officecli_path, "add", str(pptx), "/slide[1]", "--type", "shape", "--prop", "text=可编辑对象", "--prop", "x=2cm", "--prop", "y=3cm", "--prop", "w=8cm", "--prop", "h=2cm", "--prop", "font=Microsoft YaHei", "--prop", "size=24"],
        [officecli_path, "validate", str(pptx)],
    ]
    child_env = os.environ.copy()
    child_env["OFFICECLI_NO_AUTO_RESIDENT"] = "1"
    for command in commands:
        try:
            completed = subprocess.run(
                command,
                capture_output=True,
                text=True,
                timeout=30,
                check=False,
                env=child_env,
            )
        except subprocess.TimeoutExpired:
            return {"passed": False, "detail": f"officecli timed out: {command[1]}"}
        if completed.returncode != 0:
            detail = (completed.stderr or completed.stdout).strip()[:500]
            return {"passed": False, "detail": f"officecli returned {completed.returncode}: {detail}"}
    if not pptx.is_file() or pptx.stat().st_size == 0:
        return {"passed": False, "detail": "officecli did not create a non-empty PPTX"}
    return {"passed": True, "detail": "officecli create/edit/validate succeeded"}


def _powerpoint_open_smoke(pptx: Path) -> dict:
    if platform.system() != "Windows":
        return {"available": False, "passed": False, "detail": "PowerPoint COM is Windows only"}
    script = r'''
import json
import sys
import win32com.client

app = None
presentation = None
try:
    app = win32com.client.DispatchEx("PowerPoint.Application")
    presentation = app.Presentations.Open(sys.argv[1], WithWindow=False)
    print(json.dumps({"version": str(app.Version), "slides": int(presentation.Slides.Count)}))
finally:
    if presentation is not None:
        presentation.Close()
    if app is not None:
        app.Quit()
'''
    try:
        completed = subprocess.run(
            [sys.executable, "-c", script, str(pptx)],
            capture_output=True,
            text=True,
            timeout=POWERPOINT_COM_LIFECYCLE_TIMEOUT_SECONDS,
            check=False,
        )
    except subprocess.TimeoutExpired:
        return {
            "available": False,
            "passed": False,
            "detail": (
                "PowerPoint COM lifecycle timed out after "
                f"{POWERPOINT_COM_LIFECYCLE_TIMEOUT_SECONDS} seconds"
            ),
        }
    if completed.returncode != 0:
        detail = (completed.stderr or completed.stdout).strip()[:500]
        return {
            "available": False,
            "passed": False,
            "detail": f"PowerPoint COM returned {completed.returncode}: {detail}",
        }
    try:
        payload = json.loads(completed.stdout.strip().splitlines()[-1])
        version = str(payload["version"])
        slides = int(payload["slides"])
    except (IndexError, KeyError, TypeError, ValueError, json.JSONDecodeError) as exc:
        return {"available": False, "passed": False, "detail": f"PowerPoint COM result was invalid: {exc}"}
    if slides != 1:
        return {
            "available": True,
            "passed": False,
            "detail": f"PowerPoint {version} opened {slides} slides instead of 1",
        }
    return {"available": True, "passed": True, "detail": f"PowerPoint {version} opened 1 slide"}


def _libreoffice_render_smoke(soffice_path: str, pptx: Path, temporary: Path) -> dict:
    output_dir = temporary / "libreoffice-output"
    profile_dir = temporary / "libreoffice-profile"
    output_dir.mkdir()
    profile_dir.mkdir()
    command = [
        soffice_path,
        f"-env:UserInstallation={profile_dir.resolve().as_uri()}",
        "--headless",
        "--nologo",
        "--nodefault",
        "--nofirststartwizard",
        "--nolockcheck",
        "--convert-to",
        "pdf",
        "--outdir",
        str(output_dir),
        str(pptx),
    ]
    try:
        completed = subprocess.run(
            command,
            capture_output=True,
            text=True,
            timeout=60,
            check=False,
        )
    except subprocess.TimeoutExpired:
        return {"available": True, "passed": False, "detail": "LibreOffice timed out after 60 seconds"}
    if completed.returncode != 0:
        detail = (completed.stderr or completed.stdout).strip()[:500]
        return {
            "available": True,
            "passed": False,
            "detail": f"LibreOffice returned {completed.returncode}: {detail}",
        }
    rendered = output_dir / f"{pptx.stem}.pdf"
    if not rendered.is_file() or rendered.stat().st_size == 0:
        return {
            "available": True,
            "passed": False,
            "detail": f"LibreOffice returned success without {rendered.name}",
        }
    try:
        from pypdf import PdfReader
        pages = len(PdfReader(rendered).pages)
    except Exception as exc:
        return {"available": True, "passed": False, "detail": f"Cannot inspect LibreOffice PDF: {exc}"}
    if pages != 1:
        return {"available": True, "passed": False, "detail": f"LibreOffice rendered {pages} pages instead of 1"}
    return {"available": True, "passed": True, "detail": "LibreOffice rendered a non-empty one-page PDF"}


def office_smoke_test(officecli_path: str | None, powerpoint_probe: dict | None = None) -> dict:
    try:
        with tempfile.TemporaryDirectory(prefix="editable-ppt-workflow-") as temporary_name:
            temporary = Path(temporary_name)
            pptx = temporary / "smoke.pptx"
            # OfficeCLI is an optional diagnostic only.  The core smoke must remain
            # deterministic even when a broken preinstalled OfficeCLI is discovered.
            from pptx import Presentation
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            slide.shapes.add_textbox(1000000, 1000000, 5000000, 1000000).text = "Editable PPT runtime smoke"
            presentation.save(pptx)
            created = {"passed": True, "detail": "python-pptx editable smoke created"}
            if not created["passed"]:
                return {
                    "passed": False,
                    "backend": None,
                    "detail": created["detail"],
                    "powerpoint": None,
                }

            if powerpoint_probe is not None and powerpoint_probe.get("available") is False:
                powerpoint = {
                    "available": False,
                    "passed": False,
                    "detail": powerpoint_probe.get("detail", "PowerPoint status probe unavailable"),
                }
            else:
                powerpoint = _powerpoint_open_smoke(pptx)
            if powerpoint["passed"]:
                return {
                    "passed": True,
                    "backend": "powerpoint",
                    "detail": f"{created['detail']}; {powerpoint['detail']}",
                    "powerpoint": powerpoint,
                }
            if powerpoint["available"]:
                return {
                    "passed": False,
                    "backend": "powerpoint",
                    "detail": powerpoint["detail"],
                    "powerpoint": powerpoint,
                }

            soffice_path = resolve_soffice()
            if not soffice_path:
                return {
                    "passed": False,
                    "backend": None,
                    "detail": (
                        f"PowerPoint unavailable ({powerpoint['detail']}); neither PowerPoint COM nor "
                        "LibreOffice is available"
                    ),
                    "powerpoint": powerpoint,
                }
            libreoffice = _libreoffice_render_smoke(soffice_path, pptx, temporary)
            return {
                "passed": bool(libreoffice["passed"]),
                "backend": "libreoffice",
                "detail": f"PowerPoint unavailable ({powerpoint['detail']}); {libreoffice['detail']}",
                "powerpoint": powerpoint,
            }
    except Exception as exc:
        return {"passed": False, "backend": None, "detail": str(exc), "powerpoint": None}


def powerpoint_status() -> dict:
    if platform.system() != "Windows":
        return {"available": False, "detail": "Windows only"}
    script = r'''
import json
import win32com.client

app = None
try:
    app = win32com.client.DispatchEx("PowerPoint.Application")
    print(json.dumps({"version": str(app.Version)}))
finally:
    if app is not None:
        app.Quit()
'''
    try:
        completed = subprocess.run(
            [sys.executable, "-c", script],
            capture_output=True,
            text=True,
            timeout=POWERPOINT_COM_LIFECYCLE_TIMEOUT_SECONDS,
            check=False,
        )
    except subprocess.TimeoutExpired:
        return {
            "available": False,
            "detail": (
                "PowerPoint COM lifecycle timed out after "
                f"{POWERPOINT_COM_LIFECYCLE_TIMEOUT_SECONDS} seconds"
            ),
        }
    if completed.returncode != 0:
        detail = (completed.stderr or completed.stdout).strip()[:500]
        return {
            "available": False,
            "detail": f"PowerPoint status probe returned {completed.returncode}: {detail}",
        }
    try:
        payload = json.loads(completed.stdout.strip().splitlines()[-1])
        version = str(payload["version"])
    except (IndexError, KeyError, TypeError, json.JSONDecodeError) as exc:
        return {"available": False, "detail": f"PowerPoint status result was invalid: {exc}"}
    return {"available": True, "detail": f"PowerPoint {version}"}


def external_skill_status() -> dict:
    candidates = []
    configured = os.getenv("CODEX_GPT_IMAGE_SKILL")
    if configured:
        candidates.append(Path(configured))
    skill_root = Path(__file__).resolve().parents[1]
    candidates.append(skill_root.parent / "generate-slide-body-image")
    editppt_path = resolve_command("editppt")
    if editppt_path:
        candidates.append(Path(editppt_path).resolve().parent.parent / "generate-slide-body-image")
    script = next((path / "scripts/codex_gpt_image.py" for path in candidates if (path / "scripts/codex_gpt_image.py").is_file()), None)
    auth_file = Path(os.getenv("CODEX_AUTH_FILE", Path.home() / ".codex/auth.json"))
    codex_command = resolve_command(os.getenv("EDITABLE_PPT_CODEX_EXECUTABLE", "codex"))
    return {
        "codex_gpt_image_script": str(script) if script else None,
        "codex_gpt_image_available": script is not None,
        "codex_oauth_available": auth_file.is_file(),
        "codex_auth_file": str(auth_file),
        "codex_app_server_available": codex_command is not None,
        "codex_command": codex_command,
        "image_to_editable_ppt_available": editppt_path is not None,
        "officecli_available": resolve_command("officecli") is not None,
    }


def diagnose(check_powerpoint: bool = False, smoke_test: bool = False) -> dict:
    modules = {name: importlib.util.find_spec(name) is not None for name in REQUIRED_MODULES}
    external = external_skill_status()
    editppt_command = command_version("editppt", ["--help"])
    officecli_command = command_version("officecli", ["--version"])
    fonts = cjk_font_status()
    editable_cli_version = None
    editppt_path = editppt_command.get("path")
    if editppt_path:
        editable_python = _editable_python_path(editppt_path)
        if editable_python is not None:
            completed = subprocess.run(
                [str(editable_python), "-c", "import importlib.metadata; print(importlib.metadata.version('image-to-editable-ppt-cli'))"],
                capture_output=True,
                text=True,
                timeout=10,
                check=False,
            )
            if completed.returncode == 0:
                editable_cli_version = completed.stdout.strip()
    defer_powerpoint_to_smoke = bool(check_powerpoint and smoke_test)
    if defer_powerpoint_to_smoke:
        powerpoint = {"available": None, "detail": "checked by the real PPTX open smoke"}
    elif check_powerpoint:
        powerpoint = powerpoint_status()
    else:
        powerpoint = {"available": None, "detail": "use --check-powerpoint"}
    result = {
        "python": {"version": platform.python_version(), "executable": sys.executable, "ok": sys.version_info >= (3, 10)},
        "platform": platform.platform(),
        "modules": modules,
        "commands": {
            "editppt": editppt_command,
            "officecli": officecli_command,
            "soffice": command_version(resolve_soffice() or "soffice", ["--version"]),
        },
        "powerpoint": powerpoint,
        "subscription_runtime": {
            "transport": "codex-app-server",
            "authentication": "managed-by-codex",
            "available": external.get("codex_app_server_available", external["codex_oauth_available"]),
            "api_key_required": False,
        },
        "cjk_fonts": fonts,
        "runtime_versions": {
            "image_to_editable_ppt_cli": editable_cli_version,
            "expected_image_to_editable_ppt_cli": EXPECTED_EDITPPT_CLI_VERSION,
            "officecli": officecli_command.get("detail"),
        },
        "external_skills": external,
        "officecli_optional": {
            "bundled": False,
            "detected": bool(officecli_command.get("path")),
            "status": (
                "healthy" if officecli_command["available"]
                else "broken" if officecli_command.get("path")
                else "absent"
            ),
            "detail": officecli_command.get("detail"),
            "warning": bool(officecli_command.get("path") and not officecli_command["available"]),
        },
        "background_text_detection": background_text_detection_status(),
    }
    result["required_ok"] = bool(
        result["python"]["ok"]
        and all(modules.values())
        and result["background_text_detection"]["available"]
    )
    result["render_backend_available"] = bool(
        result["commands"]["soffice"]["available"] or result["powerpoint"]["available"] is True
    )
    result["workflow_ready"] = bool(
        result["required_ok"]
        and external["codex_gpt_image_available"]
        and external["codex_oauth_available"]
        and external.get("codex_app_server_available", external["codex_oauth_available"])
        and external["image_to_editable_ppt_available"]
        and editppt_command["available"]
        and editable_cli_version == EXPECTED_EDITPPT_CLI_VERSION
    )
    result["high_quality_ready"] = bool(
        result["workflow_ready"] and fonts["ready"] and result["render_backend_available"]
    )
    if smoke_test:
        result["office_smoke_test"] = office_smoke_test(
            officecli_command.get("path"),
            result["powerpoint"] if check_powerpoint and not defer_powerpoint_to_smoke else None,
        )
        if defer_powerpoint_to_smoke and result["office_smoke_test"].get("powerpoint") is not None:
            probe = result["office_smoke_test"]["powerpoint"]
            result["powerpoint"] = {
                "available": bool(probe.get("available")),
                "detail": str(probe.get("detail", "PowerPoint smoke completed")),
            }
        result["render_backend_available"] = bool(result["office_smoke_test"]["passed"])
        result["workflow_ready"] = bool(result["workflow_ready"] and result["render_backend_available"])
        result["high_quality_ready"] = bool(
            result["workflow_ready"] and fonts["ready"] and result["render_backend_available"]
        )
        officecli_path = officecli_command.get("path")
        if officecli_path:
            diagnostic_pptx = Path(tempfile.gettempdir()) / "editable-ppt-officecli-optional-smoke.pptx"
            try:
                diagnostic = _create_officecli_smoke(str(officecli_path), diagnostic_pptx)
            finally:
                diagnostic_pptx.unlink(missing_ok=True)
            result["officecli_optional"]["diagnostic"] = diagnostic
            result["officecli_optional"]["warning"] = not bool(diagnostic.get("passed"))
            if not diagnostic.get("passed"):
                result["officecli_optional"]["status"] = "broken"
        else:
            result["officecli_optional"]["diagnostic"] = {
                "passed": None,
                "detail": "not run; OfficeCLI is not installed",
            }
    return result


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--json", type=Path, help="Also write the report as JSON.")
    parser.add_argument("--check-powerpoint", action="store_true", help="Launch PowerPoint briefly through COM.")
    parser.add_argument("--smoke-test", action="store_true", help="Create, edit, validate, and open a temporary PPTX.")
    parser.add_argument("--require-high-quality", action="store_true", help="Fail unless PaddleOCR, a CJK font, and a render backend are ready.")
    args = parser.parse_args()
    result = diagnose(args.check_powerpoint, args.smoke_test)
    print(json.dumps(result, ensure_ascii=False, indent=2))
    if args.json:
        args.json.parent.mkdir(parents=True, exist_ok=True)
        args.json.write_text(json.dumps(result, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    required = result["high_quality_ready"] if args.require_high_quality else result["workflow_ready"]
    return 0 if required else 2


if __name__ == "__main__":
    raise SystemExit(main())
