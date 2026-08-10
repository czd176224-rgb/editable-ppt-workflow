"""Deterministic fixed layer for the immutable ``fixed-canvas-cm-v2`` canvas."""

from __future__ import annotations

import hashlib
import re
from pathlib import Path
from typing import Any, Mapping
from xml.etree import ElementTree

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.opc.package import Part
from pptx.opc.packuri import PackURI
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls
from pptx.util import Cm, Pt

from fixed_region_contract import (
    BODY_BOX_CM,
    CONTRACT_VERSION,
    FOOTER_LINE,
    LOGO_BOX_CM,
    PAGE_NUMBER_BOX_CM,
    PAGE_NUMBER_STYLE,
    SLIDE_SIZE_CM,
    TITLE_BOX_CM,
)


FRAME_NAMES = (
    "fixed-frame-title",
    "fixed-frame-logo",
    "fixed-frame-footer",
    "fixed-frame-page-number",
)
_EMU_TOLERANCE = 2


def _hex_color(value: Any, fallback: str) -> RGBColor:
    text = value if isinstance(value, str) and re.fullmatch(r"#[0-9A-Fa-f]{6}", value) else fallback
    return RGBColor.from_string(text[1:].upper())


def _assert_close(actual: int, expected: int, label: str) -> None:
    if abs(int(actual) - int(expected)) > _EMU_TOLERANCE:
        raise ValueError(f"{label} does not match {CONTRACT_VERSION}")


def _assert_contract(execution: Mapping[str, Any]) -> Mapping[str, Any]:
    frame = execution.get("fixed_frame")
    if not isinstance(frame, Mapping) or frame.get("geometry_version") != CONTRACT_VERSION:
        raise ValueError(f"fixed frame accepts only geometry_version {CONTRACT_VERSION}")
    expected_boxes = {
        "body_bounds_cm": BODY_BOX_CM,
        "title_bounds_cm": TITLE_BOX_CM,
        "logo_bounds_cm": LOGO_BOX_CM,
        "page_number_bounds_cm": PAGE_NUMBER_BOX_CM,
    }
    for key, expected in expected_boxes.items():
        actual = frame.get(key)
        if not isinstance(actual, Mapping) or any(float(actual.get(axis, -1)) != float(value) for axis, value in expected.items()):
            raise ValueError(f"fixed_frame.{key} does not match {CONTRACT_VERSION}")
    footer = frame.get("footer_line")
    if not isinstance(footer, Mapping) or any(footer.get(key) != value for key, value in FOOTER_LINE.items()):
        raise ValueError(f"fixed_frame.footer_line does not match {CONTRACT_VERSION}")
    return frame


def _assert_slide_and_body(slide, deck: Presentation) -> None:
    _assert_close(deck.slide_width, Cm(SLIDE_SIZE_CM["w"]), "slide width")
    _assert_close(deck.slide_height, Cm(SLIDE_SIZE_CM["h"]), "slide height")
    left = int(Cm(BODY_BOX_CM["x"]))
    top = int(Cm(BODY_BOX_CM["y"]))
    right = left + int(Cm(BODY_BOX_CM["w"]))
    bottom = top + int(Cm(BODY_BOX_CM["h"]))
    body_shapes = [shape for shape in slide.shapes if not shape.name.startswith("fixed-frame-")]
    if not body_shapes:
        raise ValueError("reconstructed page contains no body objects")
    for shape in body_shapes:
        if (
            shape.left < left - _EMU_TOLERANCE
            or shape.top < top - _EMU_TOLERANCE
            or shape.left + shape.width > right + _EMU_TOLERANCE
            or shape.top + shape.height > bottom + _EMU_TOLERANCE
        ):
            raise ValueError(f"body object {shape.name!r} lies outside the authoritative content_box")


def _add_textbox(slide, *, name: str, text: str, box: Mapping[str, float]):
    shape = slide.shapes.add_textbox(Cm(box["x"]), Cm(box["y"]), Cm(box["w"]), Cm(box["h"]))
    shape.name = name
    shape.text_frame.clear()
    shape.text_frame.word_wrap = True
    shape.text_frame.margin_left = shape.text_frame.margin_right = 0
    shape.text_frame.margin_top = shape.text_frame.margin_bottom = 0
    paragraph = shape.text_frame.paragraphs[0]
    paragraph.text = text
    return shape, paragraph


def _svg_part(slide, logo_svg: Path) -> tuple[Part, str]:
    payload = logo_svg.read_bytes()
    digest = hashlib.sha256(payload).hexdigest()[:16]
    partname = PackURI(f"/ppt/media/fixed-logo-{digest}.svg")
    package = slide.part.package
    part = next((item for item in package.iter_parts() if item.partname == partname), None)
    if part is None:
        part = Part(partname, "image/svg+xml", package, payload)
    return part, slide.part.relate_to(part, RT.IMAGE)


def _svg_length(value: str | None) -> float | None:
    if not isinstance(value, str):
        return None
    match = re.fullmatch(r"\s*([0-9]+(?:\.[0-9]+)?)\s*(?:px|pt|pc|mm|cm|in)?\s*", value, re.IGNORECASE)
    if not match:
        return None
    number = float(match.group(1))
    return number if number > 0 else None


def svg_aspect_ratio(logo_svg: Path) -> float:
    """Read intrinsic SVG ratio from viewBox, then width/height."""
    try:
        root = ElementTree.fromstring(Path(logo_svg).read_bytes())
    except (ElementTree.ParseError, OSError) as exc:
        raise ValueError("fixed-frame logo must be a readable SVG") from exc
    view_box = root.get("viewBox")
    if isinstance(view_box, str):
        values = re.split(r"[\s,]+", view_box.strip())
        if len(values) == 4:
            try:
                width, height = float(values[2]), float(values[3])
            except ValueError:
                width = height = 0.0
            if width > 0 and height > 0:
                return width / height
    width = _svg_length(root.get("width"))
    height = _svg_length(root.get("height"))
    if width is not None and height is not None:
        return width / height
    raise ValueError("fixed-frame logo SVG must define a positive viewBox or width/height")


def contained_logo_box(logo_svg: Path) -> dict[str, float]:
    """Contain the logo in its maximum box, right-aligned and vertically centered."""
    ratio = svg_aspect_ratio(logo_svg)
    max_width = float(LOGO_BOX_CM["w"])
    max_height = float(LOGO_BOX_CM["h"])
    width = max_width
    height = width / ratio
    if height > max_height:
        height = max_height
        width = height * ratio
    return {
        "x": float(LOGO_BOX_CM["x"]) + max_width - width,
        "y": float(LOGO_BOX_CM["y"]) + (max_height - height) / 2,
        "w": width,
        "h": height,
    }


def _add_svg_logo(slide, logo_svg: Path) -> None:
    """Add the original SVG as the picture payload without a raster fallback."""
    fitted = contained_logo_box(logo_svg)
    _part, svg_rid = _svg_part(slide, logo_svg)
    shape_id = slide.shapes._next_shape_id
    picture = parse_xml(
        f'<p:pic {nsdecls("a", "p", "r")}>'
        f'<p:nvPicPr><p:cNvPr id="{shape_id}" name="fixed-frame-logo"/>'
        f'<p:cNvPicPr><a:picLocks noChangeAspect="1"/></p:cNvPicPr><p:nvPr/></p:nvPicPr>'
        f'<p:blipFill><a:blip r:embed="{svg_rid}"/><a:stretch><a:fillRect/></a:stretch></p:blipFill>'
        f'<p:spPr><a:xfrm><a:off x="{int(Cm(fitted["x"]))}" y="{int(Cm(fitted["y"]))}"/>'
        f'<a:ext cx="{int(Cm(fitted["w"]))}" cy="{int(Cm(fitted["h"]))}"/></a:xfrm>'
        f'<a:prstGeom prst="rect"><a:avLst/></a:prstGeom></p:spPr>'
        f'</p:pic>'
    )
    slide.shapes._spTree.insert_element_before(picture, "p:extLst")


def _title_style(execution: Mapping[str, Any], frame: Mapping[str, Any]) -> tuple[str, float, str]:
    hard = execution.get("hard_constraints")
    if not isinstance(hard, Mapping):
        raise ValueError("confirmed hard_constraints are missing")
    typography = hard.get("typography")
    if not isinstance(typography, Mapping):
        raise ValueError("confirmed title typography is missing")
    heading = typography.get("heading")
    scale = typography.get("type_scale_pt")
    if not isinstance(heading, Mapping) or not isinstance(heading.get("cjk"), str) or not heading["cjk"].strip():
        raise ValueError("confirmed title font is missing")
    if not isinstance(scale, Mapping) or not isinstance(scale.get("page_title"), (int, float)):
        raise ValueError("confirmed title size is missing")
    size = float(scale["page_title"])
    if size <= 0:
        raise ValueError("confirmed title size must be positive")
    color = frame.get("title_color", hard.get("title_color"))
    if not isinstance(color, str) or not re.fullmatch(r"#[0-9A-Fa-f]{6}", color):
        raise ValueError("confirmed title color is missing")
    return heading["cjk"].strip(), size, color.upper()


def apply_fixed_frame(
    pptx: Path,
    *,
    page_title: str,
    page_number: int,
    style_execution: Mapping[str, Any],
    logo_svg: Path,
) -> Path:
    """Validate direct body placement, then add the four immutable native objects."""
    pptx = Path(pptx).resolve()
    if not isinstance(page_title, str) or not page_title.strip():
        raise ValueError("page title is required")
    if not isinstance(page_number, int) or page_number < 1:
        raise ValueError("page number must be a positive integer")
    logo_svg = Path(logo_svg).resolve()
    if not logo_svg.is_file() or logo_svg.suffix.lower() != ".svg":
        raise ValueError("fixed-frame logo must be an existing SVG file")
    frame = _assert_contract(style_execution)
    deck = Presentation(pptx)
    if len(deck.slides) != 1:
        raise ValueError("fixed frame accepts exactly one reconstructed slide")
    slide = deck.slides[0]
    if any(shape.name.startswith("fixed-frame-") for shape in slide.shapes):
        raise ValueError("fixed frame has already been applied")
    _assert_slide_and_body(slide, deck)
    title_font, title_pt, title_color = _title_style(style_execution, frame)

    title, paragraph = _add_textbox(
        slide,
        name="fixed-frame-title",
        text=page_title.strip(),
        box=TITLE_BOX_CM,
    )
    title.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph.alignment = PP_ALIGN.LEFT
    for run in paragraph.runs:
        run.font.name = title_font
        run.font.size = Pt(title_pt)
        run.font.bold = True
        run.font.color.rgb = _hex_color(title_color, "#0B1727")

    _add_svg_logo(slide, logo_svg)

    footer = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Cm(FOOTER_LINE["x"]),
        Cm(FOOTER_LINE["y"]),
        Cm(FOOTER_LINE["w"]),
        Cm(FOOTER_LINE["h"]),
    )
    footer.name = "fixed-frame-footer"
    footer.line.fill.background()
    footer.fill.solid()
    footer.fill.fore_color.rgb = _hex_color(FOOTER_LINE["color"], "#B8C0CC")

    page_box, page_paragraph = _add_textbox(
        slide,
        name="fixed-frame-page-number",
        text=str(page_number),
        box=PAGE_NUMBER_BOX_CM,
    )
    page_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    page_paragraph.alignment = PP_ALIGN.RIGHT
    for run in page_paragraph.runs:
        run.font.name = str(PAGE_NUMBER_STYLE["font"])
        run.font.size = Pt(float(PAGE_NUMBER_STYLE["size_pt"]))
        run.font.color.rgb = _hex_color(PAGE_NUMBER_STYLE["color"], "#6B7280")

    deck.save(pptx)
    return pptx


def inspect_fixed_frame(
    pptx: Path,
    *,
    expected_title: str,
    expected_page_number: int,
    style_execution: Mapping[str, Any],
    logo_svg: Path | None = None,
) -> dict[str, Any]:
    """Deterministically verify the complete fixed layer and its coordinates."""
    deck = Presentation(Path(pptx))
    if len(deck.slides) != 1:
        return {"passed": False, "issues": ["page package must contain one slide"]}
    slide = deck.slides[0]
    shapes = list(slide.shapes)
    by_name: dict[str, list[Any]] = {}
    for shape in shapes:
        by_name.setdefault(shape.name, []).append(shape)
    issues = [f"{name} must exist exactly once" for name in FRAME_NAMES if len(by_name.get(name, [])) != 1]
    title = by_name.get("fixed-frame-title", [])
    if title and title[0].text != expected_title.strip():
        issues.append("fixed title does not match the locked page title")
    number = by_name.get("fixed-frame-page-number", [])
    if number and number[0].text != str(expected_page_number):
        issues.append("fixed page number does not match the locked page number")
    try:
        frame = _assert_contract(style_execution)
        title_font, title_pt, title_color = _title_style(style_execution, frame)
        _assert_slide_and_body(slide, deck)
    except ValueError as exc:
        issues.append(str(exc))
        title_font, title_pt, title_color = "", 0.0, ""
    if title:
        runs = title[0].text_frame.paragraphs[0].runs
        run = runs[0] if runs else None
        if (
            run is None
            or run.font.name != title_font
            or run.font.size is None
            or abs(run.font.size.pt - title_pt) > 0.05
            or run.font.color.type is None
            or str(run.font.color.rgb) != title_color.lstrip("#")
        ):
            issues.append("fixed title style does not match the confirmed font, size, and color")
    footer = by_name.get("fixed-frame-footer", [])
    if footer and str(footer[0].fill.fore_color.rgb) != FOOTER_LINE["color"].lstrip("#"):
        issues.append("fixed footer does not use the built-in neutral style")
    if number:
        runs = number[0].text_frame.paragraphs[0].runs
        run = runs[0] if runs else None
        if (
            run is None
            or run.font.name != PAGE_NUMBER_STYLE["font"]
            or run.font.size is None
            or abs(run.font.size.pt - float(PAGE_NUMBER_STYLE["size_pt"])) > 0.05
            or run.font.color.type is None
            or str(run.font.color.rgb) != PAGE_NUMBER_STYLE["color"].lstrip("#")
        ):
            issues.append("fixed page number does not use the built-in neutral style")
    logo = by_name.get("fixed-frame-logo", [])
    if logo and logo_svg is not None:
        try:
            blip_fill = logo[0]._element.blipFill
            relationship_id = blip_fill.blip.rEmbed
            part = slide.part.related_part(relationship_id)
            if part.content_type != "image/svg+xml":
                issues.append("fixed logo media is not the original SVG content type")
            if bytes(part.blob) != Path(logo_svg).read_bytes():
                issues.append("fixed logo SVG bytes do not match the locked source")
            if blip_fill.find("{http://schemas.openxmlformats.org/drawingml/2006/main}srcRect") is not None:
                issues.append("fixed logo must not use a crop rectangle")
        except (AttributeError, KeyError, OSError):
            issues.append("fixed logo does not directly embed the locked SVG")
    expected_geometry = {
        "fixed-frame-title": TITLE_BOX_CM,
        "fixed-frame-logo": contained_logo_box(logo_svg) if logo_svg is not None else None,
        "fixed-frame-footer": FOOTER_LINE,
        "fixed-frame-page-number": PAGE_NUMBER_BOX_CM,
    }
    for name, box in expected_geometry.items():
        if box is None:
            continue
        if len(by_name.get(name, [])) != 1:
            continue
        shape = by_name[name][0]
        actual = (shape.left, shape.top, shape.width, shape.height)
        expected = (Cm(box["x"]), Cm(box["y"]), Cm(box["w"]), Cm(box["h"]))
        if any(abs(int(left) - int(right)) > _EMU_TOLERANCE for left, right in zip(actual, expected)):
            issues.append(f"{name} geometry does not match {CONTRACT_VERSION}")
    return {"passed": not issues, "issues": issues}
