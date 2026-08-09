"""Small editable native body builder constrained to the unchanged V2 body rectangle."""

from __future__ import annotations

import json
import math
from pathlib import Path
from typing import Any, Mapping

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR
from pptx.util import Cm, Pt

from fixed_region_contract import BODY_BOX_CM, SLIDE_SIZE_CM
from body_image_profile import mapping_for_source
from page_coverage import CoverageValidationError, validate_render_receipt
from page_qa import validate_qa_issue


class ContentOverflowError(ValueError):
    pass


RECONSTRUCTION_CORRECTION_CODES = frozenset({
    "native_artifact_unchecked",
    "native_body_missing",
    "native_table_missing",
    "attachment_supplement_missing",
    "native_coverage_receipt_mismatch",
    "missing_inline_image",
})


def _hex(value: Any, fallback: str) -> RGBColor:
    text = str(value or "")
    if len(text) == 7 and text.startswith("#"):
        try:
            return RGBColor.from_string(text[1:].upper())
        except ValueError:
            pass
    return RGBColor.from_string(fallback.lstrip("#").upper())


def _visual_style(style_execution: Mapping[str, Any]) -> dict[str, Any]:
    hard = style_execution.get("hard_constraints") if isinstance(style_execution, Mapping) else {}
    hard = hard if isinstance(hard, Mapping) else {}
    palette = hard.get("palette") if isinstance(hard.get("palette"), Mapping) else {}
    typography = hard.get("typography") if isinstance(hard.get("typography"), Mapping) else {}
    body = typography.get("body") if isinstance(typography.get("body"), Mapping) else {}
    scale = typography.get("type_scale_pt") if isinstance(typography.get("type_scale_pt"), Mapping) else {}
    soft = style_execution.get("soft_preferences") if isinstance(style_execution, Mapping) else {}
    density = soft.get("information_density", "balanced") if isinstance(soft, Mapping) else "balanced"
    density_factor = {"airy": 1.12, "balanced": 1.0, "dense": 0.9}.get(str(density), 1.0)
    return {
        "background": _hex(palette.get("background"), "FFFFFF"),
        "panel": _hex(palette.get("secondary_bg"), "F2F5F8"),
        "primary": _hex(palette.get("primary"), "17365D"),
        "body_text": _hex(palette.get("body_text"), "1F2937"),
        "table_header": _hex(palette.get("table_header"), "17365D"),
        "border": _hex(palette.get("border"), "D8E0E8"),
        "font": str(body.get("cjk") or "Microsoft YaHei"),
        "body_pt": max(8.0, float(scale.get("body", 12)) * density_factor),
        "caption_pt": max(7.0, float(scale.get("caption", 9)) * density_factor),
    }


def _rows(markdown: str) -> list[list[str]]:
    rows = [[cell.strip() for cell in line.strip().strip("|").split("|")] for line in markdown.splitlines() if "|" in line]
    return [row for row in rows if row and not all(set(cell) <= {"-", ":"} for cell in row)]


def _match_text(value: Any) -> str:
    return "".join(str(value).split()).replace("％", "%").casefold()


def _receipt_path(output: Path) -> Path:
    return output.with_suffix(".coverage.json")


def load_render_receipt(output: Path) -> list[dict[str, Any]]:
    return json.loads(_receipt_path(Path(output)).read_text(encoding="utf-8"))["receipts"]


def _resolve_image(project_root: Path, value: str | Path) -> Path:
    raw = Path(value)
    path = raw.resolve() if raw.is_absolute() else (project_root / raw).resolve()
    try:
        path.relative_to(project_root)
    except ValueError as exc:
        raise ValueError("native page image must remain inside the project") from exc
    if not path.is_file() or path.is_symlink():
        raise ValueError(f"required native image is unavailable: {path}")
    return path


def _layout_heights(plan: Mapping[str, Any], available: float) -> tuple[float, float, list[float], float]:
    text = str(plan.get("body_text", ""))
    supplement_count = len(plan.get("attachment_supplements", []))
    table_rows = [_rows(value) for value in plan.get("tables", []) if isinstance(value, str) and _rows(value)]
    image_count = len(plan.get("required_images", []))
    blocks = int(bool(text)) + int(bool(supplement_count)) + len(table_rows) + int(bool(image_count))
    gap_total = max(0, blocks - 1) * 0.18
    minimum = (
        (1.4 if text else 0) + (max(0.9, 0.7 * supplement_count) if supplement_count else 0)
        + len(table_rows) * 1.4 + (2.0 if image_count else 0) + gap_total
    )
    if minimum > available:
        raise ContentOverflowError("content_overflow: too many required tables/images for one visible slide")
    weights = ([max(1.4, min(5.0, 1.2 + len(text) / 500))] if text else [])
    if supplement_count:
        weights.append(max(0.9, min(2.4, 0.7 * supplement_count)))
    weights += [max(1.4, min(4.0, 0.55 * len(rows))) for rows in table_rows]
    if image_count:
        weights.append(max(2.0, min(4.5, 2.0 + 0.45 * image_count)))
    usable = available - gap_total
    scale = usable / sum(weights) if weights else 1.0
    heights = [value * scale for value in weights]
    offset = 0
    text_height = heights[offset] if text else 0
    offset += int(bool(text))
    supplement_height = heights[offset] if supplement_count else 0
    offset += int(bool(supplement_count))
    table_heights = heights[offset:offset + len(table_rows)]
    offset += len(table_rows)
    image_height = heights[offset] if image_count else 0
    return text_height, supplement_height, table_heights, image_height


def _coverage_receipt(
    coverage_items: Mapping[str, Mapping[str, Any]], coverage_id: str, object_id: str,
    *, observed_text: str | None = None, observed_asset_id: str | None = None,
) -> dict[str, Any]:
    item = coverage_items.get(coverage_id)
    expected = item.get("expected") if isinstance(item, Mapping) else None
    if not isinstance(expected, Mapping) or not isinstance(expected.get("sha256"), str):
        raise ValueError(f"coverage item has no current expected identity: {coverage_id}")
    receipt: dict[str, Any] = {
        "coverage_id": coverage_id,
        "object_id": object_id,
        "visible": True,
        "expected_sha256": expected["sha256"],
    }
    if observed_text is not None:
        receipt["observed_text"] = observed_text
    if observed_asset_id is not None:
        receipt["observed_asset_id"] = observed_asset_id
    return receipt


def _build_page(
    plan: Mapping[str, Any], output: Path, *, project_root: Path, background: Path | None = None,
    style_execution: Mapping[str, Any] | None = None,
    corrections: list[Mapping[str, Any]] | None = None,
) -> Path:
    correction_codes: list[str] = []
    for correction in corrections or []:
        issue = validate_qa_issue(correction)
        if issue["code"] not in RECONSTRUCTION_CORRECTION_CODES:
            raise ValueError(f"unsupported reconstruction correction: {issue['code']}")
        if issue["code"] not in correction_codes:
            correction_codes.append(issue["code"])
    output = Path(output).resolve()
    output.parent.mkdir(parents=True, exist_ok=True)
    deck = Presentation()
    deck.slide_width, deck.slide_height = Cm(SLIDE_SIZE_CM["w"]), Cm(SLIDE_SIZE_CM["h"])
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    x, y, w, h = (BODY_BOX_CM[key] for key in ("x", "y", "w", "h"))
    visual_style = _visual_style(style_execution or {})
    body_background = slide.background.fill
    body_background.solid()
    body_background.fore_color.rgb = visual_style["background"]
    if background is not None:
        with Image.open(background) as source:
            source_size = source.size
            source.verify()
        mapped = mapping_for_source(int(source_size[0]), int(source_size[1]))["effective_box_cm"]
        visual = slide.shapes.add_picture(
            str(background), Cm(mapped["x"]), Cm(mapped["y"]), Cm(mapped["w"]), Cm(mapped["h"])
        )
        visual.name = "hybrid-visual-background"
    tables = [rows for value in plan.get("tables", []) if isinstance(value, str) and (rows := _rows(value))]
    text_value = str(plan.get("body_text", ""))
    text_height, supplement_height, table_heights, image_height = _layout_heights(plan, h)
    receipts: list[dict[str, Any]] = []
    coverage = plan.get("coverage_contract") if isinstance(plan.get("coverage_contract"), Mapping) else {"required_items": []}
    coverage_items = {
        str(item.get("coverage_id")): item
        for item in coverage.get("required_items", [])
        if isinstance(item, Mapping)
    }
    cursor = y
    if text_value:
        if len(text_value) > int(w * text_height * 22):
            raise ContentOverflowError("content_overflow: Word text cannot remain visible at the minimum type size")
        if background is not None:
            panel = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Cm(x), Cm(cursor), Cm(w), Cm(text_height),
            )
            panel.name = "native-authoritative-text-panel"
            panel.fill.solid()
            panel.fill.fore_color.rgb = visual_style["panel"]
            panel.line.color.rgb = visual_style["border"]
        text = slide.shapes.add_textbox(Cm(x + 0.12), Cm(cursor + 0.08), Cm(w - 0.24), Cm(max(0.2, text_height - 0.16)))
        text.name = "native-body-text"
        frame = text.text_frame
        frame.clear()
        frame.word_wrap = True
        frame.vertical_anchor = MSO_ANCHOR.TOP
        frame.margin_left = frame.margin_right = Cm(0.15)
        frame.margin_top = frame.margin_bottom = Cm(0.1)
        paragraph = frame.paragraphs[0]
        paragraph.text = text_value
        font_size = max(8.0, min(visual_style["body_pt"], 18 - len(text_value) / max(80, w * text_height * 3)))
        for run in paragraph.runs:
            run.font.name = visual_style["font"]
            run.font.size = Pt(font_size)
            run.font.color.rgb = visual_style["body_text"]
        for unit in plan.get("semantic_units", []):
            if (
                isinstance(unit, Mapping)
                and isinstance(unit.get("coverage_id"), str)
                and _match_text(unit.get("text", "")) in _match_text(text_value)
            ):
                receipts.append(_coverage_receipt(
                    coverage_items, unit["coverage_id"], text.name, observed_text=text_value,
                ))
        for index, anchor in enumerate(plan.get("mandatory_anchors", []), start=1):
            if isinstance(anchor, str) and _match_text(anchor) in _match_text(text_value):
                receipts.append(_coverage_receipt(
                    coverage_items, f"anchor:{index:03d}", text.name, observed_text=text_value,
                ))
        cursor += text_height + 0.18
    supplements = [item for item in plan.get("attachment_supplements", []) if isinstance(item, Mapping)]
    if supplements:
        item_height = max(0.32, supplement_height / len(supplements))
        for index, item in enumerate(supplements, start=1):
            supplement_text = str(item.get("text", "")).strip()
            box = slide.shapes.add_textbox(
                Cm(x + 0.12), Cm(cursor + (index - 1) * item_height),
                Cm(w - 0.24), Cm(max(0.25, item_height - 0.06)),
            )
            box.name = f"native-attachment-supplement-{index:03d}"
            box.text_frame.clear()
            box.text_frame.word_wrap = True
            paragraph = box.text_frame.paragraphs[0]
            paragraph.text = supplement_text
            for run in paragraph.runs:
                run.font.name = visual_style["font"]
                run.font.size = Pt(visual_style["caption_pt"])
                run.font.color.rgb = visual_style["body_text"]
            receipts.append(_coverage_receipt(
                coverage_items, str(item.get("coverage_id")), box.name, observed_text=supplement_text,
            ))
        cursor += supplement_height + 0.18
    for table_index, (rows, table_height) in enumerate(zip(tables, table_heights), start=1):
        if len(rows) * 0.38 > table_height:
            raise ContentOverflowError(f"content_overflow: table {table_index} cannot remain visible")
        columns = max(len(row) for row in rows)
        table_shape = slide.shapes.add_table(
            len(rows), columns, Cm(x), Cm(cursor), Cm(w), Cm(table_height)
        )
        table_shape.name = f"native-data-table-{table_index:03d}"
        table = table_shape.table
        for row_index, row in enumerate(rows):
            for column_index in range(columns):
                cell = table.cell(row_index, column_index)
                cell.text = row[column_index] if column_index < len(row) else ""
                cell.fill.solid()
                cell.fill.fore_color.rgb = visual_style["table_header"] if row_index == 0 else visual_style["background"]
                for p in cell.text_frame.paragraphs:
                    for run in p.runs:
                        run.font.name = visual_style["font"]
                        run.font.size = Pt(visual_style["caption_pt"] if row_index else visual_style["body_pt"])
                        run.font.bold = row_index == 0
                        run.font.color.rgb = RGBColor(255, 255, 255) if row_index == 0 else visual_style["body_text"]
        observed_table = "\n".join("|".join(row) for row in rows)
        receipts.append(_coverage_receipt(
            coverage_items, f"table:{table_index:03d}", table_shape.name, observed_text=observed_table,
        ))
        for anchor_index, anchor in enumerate(plan.get("mandatory_anchors", []), start=1):
            if isinstance(anchor, str) and _match_text(anchor) in _match_text("\n".join("|".join(row) for row in rows)) and not any(item["coverage_id"] == f"anchor:{anchor_index:03d}" for item in receipts):
                receipts.append(_coverage_receipt(
                    coverage_items, f"anchor:{anchor_index:03d}", table_shape.name, observed_text=observed_table,
                ))
        cursor += table_height + 0.18
    images = [item for item in plan.get("required_images", []) if isinstance(item, Mapping)]
    if images:
        columns = min(3, max(1, math.ceil(math.sqrt(len(images)))))
        rows_count = math.ceil(len(images) / columns)
        cell_w, cell_h = w / columns, image_height / rows_count
        for index, item in enumerate(images):
            image_path = _resolve_image(project_root, str(item.get("relative_path", "")))
            column, row = index % columns, index // columns
            with Image.open(image_path) as source:
                ratio = source.width / source.height
            target_ratio = cell_w / cell_h
            if ratio >= target_ratio:
                picture_w, picture_h = cell_w, cell_w / ratio
            else:
                picture_h, picture_w = cell_h, cell_h * ratio
            picture_x = x + column * cell_w + (cell_w - picture_w) / 2
            picture_y = cursor + row * cell_h + (cell_h - picture_h) / 2
            picture = slide.shapes.add_picture(str(image_path), Cm(picture_x), Cm(picture_y), Cm(picture_w), Cm(picture_h))
            picture.name = f"native-required-image-{item.get('asset_id')}"
            receipts.append(_coverage_receipt(
                coverage_items, f"image:{item.get('asset_id')}", picture.name,
                observed_asset_id=str(item.get("asset_id")),
            ))
    if coverage.get("required_items"):
        try:
            validate_render_receipt(coverage, receipts)
        except CoverageValidationError as exc:
            raise ContentOverflowError(f"content_overflow: incomplete render receipt {exc.report}") from exc
    deck.save(output)
    _receipt_path(output).write_text(json.dumps({
        "schema_version": "1.0",
        "receipts": receipts,
        "corrections_consumed": correction_codes,
    }, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    return output


def build_native_page(
    plan: Mapping[str, Any], style_execution: Mapping[str, Any], output: Path, *, project_root: Path | None = None,
    corrections: list[Mapping[str, Any]] | None = None,
) -> Path:
    root = Path(project_root).resolve() if project_root is not None else Path(output).resolve().parent
    return _build_page(
        plan, Path(output), project_root=root, style_execution=style_execution,
        corrections=corrections,
    )


def build_overlay_page(
    plan: Mapping[str, Any], visual_assets: Mapping[str, Any], style_execution: Mapping[str, Any], output: Path,
    *, project_root: Path | None = None, corrections: list[Mapping[str, Any]] | None = None,
) -> Path:
    root = Path(project_root).resolve() if project_root is not None else Path(output).resolve().parent
    value = visual_assets.get("background")
    background = _resolve_image(root, value) if value is not None else None
    return _build_page(
        plan, Path(output), project_root=root, background=background,
        style_execution=style_execution, corrections=corrections,
    )
