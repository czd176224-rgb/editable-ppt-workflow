"""Production QA bridge: one deterministic pass with no semantic image review."""

from __future__ import annotations

import hashlib
import json
from pathlib import Path
from typing import Any, Callable, Mapping

from pptx import Presentation

from deterministic_qa import run_deterministic_qa
from background_text_detector import detect_background_text
from page_coverage import (
    CoverageValidationError,
    normalize_coverage_text,
    table_visible_text,
    validate_render_receipt,
)
from qa_orchestrator import decide_qa_path


OCRProvider = Callable[[Path], Mapping[str, Any]]


def _shape_text(shape: Any) -> str:
    if getattr(shape, "has_table", False):
        return "\n".join("|".join(cell.text for cell in row.cells) for row in shape.table.rows)
    if getattr(shape, "has_text_frame", False):
        return str(shape.text)
    return ""


def _native_pptx_observations(
    pptx: Path, route: Mapping[str, Any],
) -> dict[str, Any]:
    plan = route.get("native_plan")
    coverage = route.get("coverage_contract")
    receipt_value = route.get("coverage_receipt")
    if not isinstance(plan, Mapping) or not isinstance(coverage, Mapping) or not isinstance(receipt_value, str):
        raise ValueError("native QA requires the authoritative plan, coverage contract, and receipt")
    deck = Presentation(Path(pptx))
    if len(deck.slides) != 1:
        raise ValueError("native QA requires a one-slide PPTX")
    shapes = {shape.name: shape for shape in deck.slides[0].shapes}
    body_shape = shapes.get("native-body-text")
    expected_body = str(plan.get("body_text", ""))
    body_text = _shape_text(body_shape) if body_shape is not None else ""
    body_present = bool(expected_body.strip()) and (
        normalize_coverage_text(expected_body) in normalize_coverage_text(body_text)
    )
    table_presence: dict[str, bool] = {}
    for index, expected in enumerate(plan.get("tables", []), start=1):
        name = f"native-data-table-{index:03d}"
        shape = shapes.get(name)
        actual = _shape_text(shape) if shape is not None and getattr(shape, "has_table", False) else ""
        table_presence[f"table:{index:03d}"] = bool(actual) and (
            normalize_coverage_text(table_visible_text(str(expected)))
            in normalize_coverage_text(actual)
        )
    image_presence: dict[str, bool] = {}
    for item in plan.get("required_images", []):
        if not isinstance(item, Mapping):
            continue
        asset_id = str(item.get("asset_id", ""))
        shape = shapes.get(f"native-required-image-{asset_id}")
        present = False
        if shape is not None:
            try:
                present = bool(shape.image.blob)
            except (AttributeError, KeyError, ValueError):
                present = False
        image_presence[asset_id] = present
    supplement_presence: dict[str, bool] = {}
    for index, item in enumerate(plan.get("attachment_supplements", []), start=1):
        if not isinstance(item, Mapping):
            continue
        coverage_id = str(item.get("coverage_id", ""))
        shape = shapes.get(f"native-attachment-supplement-{index:03d}")
        actual = _shape_text(shape) if shape is not None else ""
        supplement_presence[coverage_id] = bool(actual) and (
            normalize_coverage_text(item.get("text", "")) in normalize_coverage_text(actual)
        )
    receipt_valid = False
    receipt_error = ""
    try:
        receipt_path = Path(receipt_value).resolve()
        payload = json.loads(receipt_path.read_text(encoding="utf-8"))
        receipts = payload.get("receipts") if isinstance(payload, Mapping) else None
        if not isinstance(receipts, list):
            raise ValueError("coverage receipt entries are missing")
        validate_render_receipt(coverage, receipts)
        for receipt in receipts:
            if not isinstance(receipt, Mapping):
                raise ValueError("coverage receipt entry is invalid")
            object_id = receipt.get("object_id")
            shape = shapes.get(str(object_id))
            if shape is None:
                raise ValueError(f"coverage receipt object is absent from PPTX: {object_id}")
            if "observed_text" in receipt and (
                normalize_coverage_text(receipt.get("observed_text"))
                != normalize_coverage_text(_shape_text(shape))
            ):
                raise ValueError(f"coverage receipt text differs from PPTX object: {object_id}")
            if "observed_asset_id" in receipt:
                try:
                    if not shape.image.blob:
                        raise ValueError(f"coverage receipt image relationship is broken: {object_id}")
                except (AttributeError, KeyError, ValueError) as exc:
                    raise ValueError(f"coverage receipt image relationship is broken: {object_id}") from exc
        receipt_valid = True
    except (OSError, json.JSONDecodeError, ValueError, CoverageValidationError) as exc:
        receipt_error = str(exc)
    return {
        "native_artifact_checked": True,
        "native_body_present": body_present,
        "native_table_presence": table_presence,
        "inline_image_presence": image_presence,
        "native_supplement_presence": supplement_presence,
        "native_coverage_receipt_valid": receipt_valid,
        "native_coverage_receipt_error": receipt_error,
    }


def collect_observations(
    image: Path,
    contract: Mapping[str, Any],
    route: Mapping[str, Any],
    body_image_mapping: Mapping[str, Any] | None = None,
    *,
    ocr_provider: OCRProvider | None = None,
) -> dict[str, Any]:
    """Collect only observations needed by the risk gate; no duplicate review."""
    observations: dict[str, Any] = {
        "ocr_available": False,
        "aspect_mapping": (body_image_mapping or {}).get("mode"),
        "inline_image_presence": {
            str(item.get("asset_id")): True
            for item in contract.get("asset_bindings", [])
            if isinstance(item, Mapping) and item.get("asset_role") == "mandatory_inline_image"
        },
        "native_text_authority": route.get("text_authority") == "native_overlay",
    }
    if route.get("route") in {"native", "image", "hybrid"} and route.get("text_authority") == "native_overlay" and Path(image).suffix.lower() == ".pptx":
        observations.update(_native_pptx_observations(Path(image), route))
    if route.get("route") in {"image", "hybrid"} and Path(image).suffix.lower() != ".pptx":
        detector = ocr_provider or detect_background_text
        try:
            observed = dict(detector(Path(image)))
            observations.update(observed)
            text = observed.get("ocr_text")
            boxes = observed.get("ocr_boxes")
            observations["background_text_detection_available"] = True
            observations["background_text_detected"] = bool(
                observed.get("background_text_detected")
                or (isinstance(text, str) and text.strip())
                or (isinstance(boxes, list) and boxes)
                or observed.get("background_text_regions")
            )
            observations["ocr_available"] = isinstance(text, str)
        except Exception as exc:
            observations["background_text_detection_available"] = False
            observations["background_text_detection_error"] = (
                f"{type(exc).__name__}: {' '.join(str(exc).split())[:240]}"
            )
    return observations


def decide_page_qa(
    image: Path,
    contract: Mapping[str, Any],
    fact_plan: Mapping[str, Any],
    route: Mapping[str, Any],
    body_image_mapping: Mapping[str, Any] | None = None,
    *,
    ocr_provider: OCRProvider | None = None,
) -> dict[str, Any]:
    observations = collect_observations(
        image, contract, route, body_image_mapping, ocr_provider=ocr_provider
    )
    deterministic = run_deterministic_qa(image, contract, fact_plan, route, observations)
    gate_observations = {
        **observations,
        "uncertain": bool(deterministic.get("uncertain")),
        "deterministic_mismatch": deterministic.get("status") == "repair",
    }
    path = decide_qa_path(contract, route, gate_observations)
    report = {
        "schema_version": "1.0",
        "qa_path": path["path"],
        "risk_score": path["risk_score"],
        "semantic_calls": 0,
        "observations_sha256": hashlib.sha256(
            json.dumps(observations, ensure_ascii=False, sort_keys=True, separators=(",", ":")).encode("utf-8")
        ).hexdigest(),
        "observations": observations,
        "result": deterministic,
    }
    return report


def workflow_qa_result(report: Mapping[str, Any]) -> dict[str, Any]:
    result = report.get("result") if isinstance(report.get("result"), Mapping) else {}
    status = str(result.get("status", "pass_with_advisory"))
    issues = list(result.get("issues", []))
    return {
        "status": status,
        "repair_scope": str(result.get("repair_scope", "none")),
        "issues": issues,
        "evidence": list(result.get("evidence", [])),
        "confidence": str(result.get("confidence", "high")),
        "trigger_reason": str(result.get("trigger_reason", "checks_passed")),
        "checked_scope": str(result.get("checked_scope", "full")),
    }
