"""Extract original Word body assets and bind them to locked content pages."""

from __future__ import annotations

import hashlib
import mimetypes
import posixpath
import zipfile
from collections import OrderedDict
from io import BytesIO
from pathlib import Path, PurePosixPath
from typing import Any, Iterable
from xml.etree import ElementTree

from PIL import Image, UnidentifiedImageError

from extract_docx_pages import iter_blocks, relationship_ids
from page_assets import DOCUMENT_MEDIA, SPREADSHEET_MEDIA, classify_page_asset
from docx import Document
from pptx import Presentation


RELATIONSHIPS_NAMESPACE = "http://schemas.openxmlformats.org/package/2006/relationships"
CONTENT_TYPES_NAMESPACE = "http://schemas.openxmlformats.org/package/2006/content-types"
ORIGINAL_ASSET_DIRECTORY = PurePosixPath("00_source/word_assets/original")
DERIVED_ASSET_DIRECTORY = PurePosixPath("00_source/word_assets/derived")
SUPPORTED_GENERATION_MEDIA = frozenset({"image/png", "image/jpeg", "image/webp"})
DERIVABLE_GENERATION_MEDIA = frozenset({"image/bmp"})


def extract_attachment_text(path: Path, media_type: str) -> str:
    """Extract deterministic text from a supported page-local attachment."""
    path = Path(path)
    chunks: list[str] = []
    if media_type == "application/pdf":
        from pypdf import PdfReader
        document = PdfReader(path)
        chunks = [f"[Page {number}]\n{page.extract_text() or ''}" for number, page in enumerate(document.pages, 1)]
    elif media_type in {"application/vnd.openxmlformats-officedocument.wordprocessingml.document", "application/msword"}:
        document = Document(path)
        chunks.extend(
            f"[Paragraph {number}]\n{paragraph.text}"
            for number, paragraph in enumerate(document.paragraphs, 1) if paragraph.text.strip()
        )
        for table_number, table in enumerate(document.tables, 1):
            for row_number, row in enumerate(table.rows, 1):
                text = " | ".join(cell.text for cell in row.cells)
                if text.strip(" |"):
                    chunks.append(f"[Table {table_number} Row {row_number}]\n{text}")
    elif media_type in {"application/vnd.openxmlformats-officedocument.presentationml.presentation", "application/vnd.ms-powerpoint"}:
        deck = Presentation(path)
        for number, slide in enumerate(deck.slides, 1):
            text = [shape.text for shape in slide.shapes if getattr(shape, "has_text_frame", False) and shape.text.strip()]
            chunks.append(f"[Slide {number}]\n" + "\n".join(text))
    elif media_type in {"application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", "application/vnd.ms-excel"}:
        with zipfile.ZipFile(path) as archive:
            names = set(archive.namelist())
            shared: list[str] = []
            if "xl/sharedStrings.xml" in names:
                root = ElementTree.fromstring(archive.read("xl/sharedStrings.xml"))
                shared = ["".join(node.itertext()) for node in root]
            for name in sorted(item for item in names if item.startswith("xl/worksheets/sheet") and item.endswith(".xml")):
                root = ElementTree.fromstring(archive.read(name))
                rows = []
                for row in root.iter():
                    if not row.tag.endswith("}row"):
                        continue
                    values = []
                    for cell in row:
                        if not cell.tag.endswith("}c"):
                            continue
                        value = next((node.text or "" for node in cell if node.tag.endswith("}v")), "")
                        if cell.get("t") == "s" and value.isdigit() and int(value) < len(shared):
                            value = shared[int(value)]
                        values.append(f"{cell.get('r', '?')}={value}")
                    if any(values):
                        rows.append(" | ".join(values))
                chunks.append(f"[{PurePosixPath(name).name}]\n" + "\n".join(rows))
    elif media_type == "text/plain":
        chunks = [path.read_text(encoding="utf-8-sig")]
    text = "\n\n".join(chunk.strip() for chunk in chunks if chunk and chunk.strip())
    return text[:120000]


def _element_text(element: ElementTree.Element) -> str:
    return "".join(node.text or "" for node in element.iter()).strip()


def _chart_record(asset: dict[str, Any], data: bytes) -> dict[str, Any] | None:
    """Extract literal Word chart labels and cached values without producing an image input."""
    try:
        root = ElementTree.fromstring(data)
    except ElementTree.ParseError:
        return None

    def descendants(node: ElementTree.Element, name: str) -> list[ElementTree.Element]:
        return [item for item in node.iter() if item.tag.rsplit("}", 1)[-1] == name]

    title_node = next(iter(descendants(root, "title")), None)
    title = _element_text(title_node) if title_node is not None else ""
    if not title:
        return None
    series: list[dict[str, Any]] = []
    for series_node in descendants(root, "ser"):
        series_title = ""
        tx_node = next(iter(descendants(series_node, "tx")), None)
        if tx_node is not None:
            series_title = _element_text(tx_node)
        values: list[str] = []
        times: list[str] = []
        for child in series_node:
            local = child.tag.rsplit("}", 1)[-1]
            points = descendants(child, "pt")
            if local in {"cat", "xVal"}:
                times = [_element_text(point) for point in points]
            elif local in {"val", "yVal"}:
                values = [_element_text(point) for point in points]
        record: dict[str, Any] = {"series": series_title, "values": values}
        if times:
            record["times"] = times
        series.append(record)
    return {
        "page_numbers": list(asset["page_numbers"]),
        "source_asset_id": asset["asset_id"],
        "title": title,
        "series": series,
    }


def _read_relationships(docx_path: Path) -> dict[str, str | None]:
    """Read document-part relationships; ``None`` denotes an external target."""
    with zipfile.ZipFile(docx_path) as archive:
        root = ElementTree.fromstring(archive.read("word/_rels/document.xml.rels"))
    relationships: dict[str, str | None] = {}
    for relationship in root.findall(f"{{{RELATIONSHIPS_NAMESPACE}}}Relationship"):
        relationship_id = relationship.get("Id")
        target = relationship.get("Target")
        if not relationship_id or target is None:
            continue
        relationships[relationship_id] = None if relationship.get("TargetMode") == "External" else target
    return relationships


def _read_content_types(docx_path: Path) -> tuple[dict[str, str], dict[str, str]]:
    with zipfile.ZipFile(docx_path) as archive:
        root = ElementTree.fromstring(archive.read("[Content_Types].xml"))
    defaults: dict[str, str] = {}
    overrides: dict[str, str] = {}
    for element in root:
        if element.tag == f"{{{CONTENT_TYPES_NAMESPACE}}}Default":
            extension = element.get("Extension")
            content_type = element.get("ContentType")
            if extension and content_type:
                defaults[extension.lower()] = content_type
        elif element.tag == f"{{{CONTENT_TYPES_NAMESPACE}}}Override":
            part_name = element.get("PartName")
            content_type = element.get("ContentType")
            if part_name and content_type:
                overrides[part_name.lstrip("/")] = content_type
    return defaults, overrides


def _safe_source_part(target: str) -> str:
    normalized_target = target.replace("\\", "/")
    candidate = PurePosixPath(normalized_target)
    if candidate.is_absolute() or not normalized_target or any(part in {"", ".", ".."} for part in candidate.parts):
        raise ValueError(f"unsafe DOCX relationship target: {target}")
    source_part = PurePosixPath(posixpath.normpath(str(PurePosixPath("word") / candidate)))
    if not str(source_part).startswith("word/"):
        raise ValueError(f"unsafe DOCX relationship target: {target}")
    return str(source_part)


def _media_type(source_part: str, defaults: dict[str, str], overrides: dict[str, str]) -> str:
    if source_part in overrides:
        return overrides[source_part]
    extension = PurePosixPath(source_part).suffix.removeprefix(".").lower()
    return defaults.get(extension) or mimetypes.types_map.get(f".{extension}") or "application/octet-stream"


def _dimensions(data: bytes) -> tuple[int, int] | None:
    try:
        with Image.open(BytesIO(data)) as image:
            image.load()
            return image.size
    except (UnidentifiedImageError, OSError, ValueError):
        return None


def _page_bindings(pages_payload: dict[str, Any]) -> dict[str, tuple[list[int], list[int]]]:
    """Use only explicit marker evidence, or explicit Microsoft Word block evidence."""
    pagination_mode = pages_payload.get("pagination_mode")
    backend = pages_payload.get("pagination_backend")
    can_bind = pagination_mode != "physical_rendered_pages" or backend == "microsoft-word"
    bindings: dict[str, tuple[list[int], list[int]]] = {}
    if not can_bind:
        return bindings
    if pagination_mode == "physical_rendered_pages" and backend == "microsoft-word":
        for block in pages_payload.get("block_page_evidence", []):
            page_number = block.get("page_number")
            source_block_index = block.get("source_block_index")
            relationship_id_values = block.get("relationship_ids", [])
            if (
                not isinstance(page_number, int) or isinstance(page_number, bool) or page_number < 1
                or not isinstance(source_block_index, int) or isinstance(source_block_index, bool)
                or not isinstance(relationship_id_values, list)
            ):
                continue
            for relationship_id in relationship_id_values:
                if not isinstance(relationship_id, str):
                    continue
                page_numbers, block_indexes = bindings.setdefault(relationship_id, ([], []))
                if page_number not in page_numbers:
                    page_numbers.append(page_number)
                if source_block_index not in block_indexes:
                    block_indexes.append(source_block_index)
        return bindings
    for page in pages_payload.get("pages", []):
        page_number = page.get("page_number")
        if not isinstance(page_number, int) or isinstance(page_number, bool) or page_number < 1:
            continue
        for block in page.get("blocks", []):
            source_block_index = block.get("source_block_index")
            relationship_id_values = block.get("relationship_ids", [])
            if not isinstance(source_block_index, int) or isinstance(source_block_index, bool):
                continue
            if not isinstance(relationship_id_values, list):
                continue
            for relationship_id in relationship_id_values:
                if not isinstance(relationship_id, str):
                    continue
                page_numbers, block_indexes = bindings.setdefault(relationship_id, ([], []))
                if page_number not in page_numbers:
                    page_numbers.append(page_number)
                if source_block_index not in block_indexes:
                    block_indexes.append(source_block_index)
    return bindings


def _unique_filename(original_filename: str, used_names: set[str]) -> str:
    if not original_filename or original_filename in {".", ".."} or "/" in original_filename or "\\" in original_filename:
        raise ValueError(f"unsafe DOCX asset filename: {original_filename}")
    candidate = original_filename
    counter = 2
    while candidate.casefold() in used_names:
        candidate = f"{Path(original_filename).stem}_{counter}{Path(original_filename).suffix}"
        counter += 1
    used_names.add(candidate.casefold())
    return candidate


def _safe_relative_path(relative_path: str) -> str:
    normalized_path = relative_path.replace("\\", "/")
    candidate = PurePosixPath(normalized_path)
    if (
        candidate.is_absolute()
        or any(part in {"", ".", ".."} for part in candidate.parts)
        or not candidate.is_relative_to(ORIGINAL_ASSET_DIRECTORY)
    ):
        raise ValueError(f"unsafe source asset relative path: {relative_path}")
    return str(candidate)


def build_manifest(occurrences: Iterable[dict[str, Any]], source_file: str | None = None) -> dict[str, Any]:
    """Build a stable manifest from source-part occurrences without converting bytes."""
    grouped: OrderedDict[str, dict[str, Any]] = OrderedDict()
    for occurrence in occurrences:
        source_part = occurrence["source_part"]
        if not isinstance(source_part, str) or not source_part:
            raise ValueError("source_part must be a non-empty string")
        data = occurrence.get("data")
        if data is not None and not isinstance(data, bytes):
            raise ValueError("source asset data must be bytes or None for an unresolved external asset")
        record = grouped.setdefault(source_part, {**occurrence, "page_numbers": [], "source_block_indexes": []})
        if record["data"] != data:
            raise ValueError(f"inconsistent bytes for source part: {source_part}")
        for page_number in occurrence.get("page_numbers", []):
            if page_number not in record["page_numbers"]:
                record["page_numbers"].append(page_number)
        for source_block_index in occurrence.get("source_block_indexes", []):
            if source_block_index not in record["source_block_indexes"]:
                record["source_block_indexes"].append(source_block_index)

    assets: list[dict[str, Any]] = []
    used_names: set[str] = set()
    for index, record in enumerate(grouped.values(), start=1):
        data = record["data"]
        original_filename = record.get("original_filename") or PurePosixPath(record["source_part"]).name
        filename = _unique_filename(original_filename, used_names)
        relative_path = _safe_relative_path(
            record.get("relative_path") or str(ORIGINAL_ASSET_DIRECTORY / filename)
        )
        page_numbers = sorted(record["page_numbers"])
        source_block_indexes = sorted(record["source_block_indexes"])
        asset = {
            "asset_id": f"word_asset_{index:03d}",
            "relationship_id": record["relationship_id"],
            "source_part": record["source_part"],
            "original_filename": original_filename,
            "media_type": record["media_type"],
            "sha256": hashlib.sha256(data).hexdigest() if data is not None else None,
            "byte_size": len(data) if data is not None else None,
            "page_numbers": page_numbers,
            "source_block_indexes": source_block_indexes,
            "binding_status": "bound" if data is not None and page_numbers else "unresolved",
            "relative_path": relative_path,
        }
        dimensions = _dimensions(data) if data is not None else None
        if dimensions is not None:
            asset["width"], asset["height"] = dimensions
        asset.update(
            classify_page_asset(
                asset["media_type"],
                binding_status=asset["binding_status"],
                has_generation_input=(
                    data is not None
                    and asset["media_type"] in SUPPORTED_GENERATION_MEDIA | DERIVABLE_GENERATION_MEDIA
                ),
            )
        )
        assets.append(asset)
    manifest: dict[str, Any] = {"schema_version": "1.0", "assets": assets}
    if source_file is not None:
        manifest["source_file"] = source_file
    return manifest


def extract_source_assets(docx_path: Path, pages_payload: dict, output_dir: Path) -> dict:
    """Copy each unique embedded Word body part once and return its page bindings."""
    docx_path = Path(docx_path)
    relationships = _read_relationships(docx_path)
    defaults, overrides = _read_content_types(docx_path)
    page_bindings = _page_bindings(pages_payload)

    referenced_relationship_ids: list[str] = []
    for block in iter_blocks(Document(docx_path)):
        for relationship_id in relationship_ids(block):
            if relationship_id not in referenced_relationship_ids:
                referenced_relationship_ids.append(relationship_id)

    occurrences: list[dict[str, Any]] = []
    with zipfile.ZipFile(docx_path) as archive:
        archive_names = set(archive.namelist())
        for relationship_id in referenced_relationship_ids:
            if relationship_id not in relationships:
                raise ValueError(f"missing DOCX relationship: {relationship_id}")
            target = relationships[relationship_id]
            if target is None:
                page_numbers, source_block_indexes = page_bindings.get(relationship_id, ([], []))
                occurrences.append(
                    {
                        "relationship_id": relationship_id,
                        "source_part": f"external:{relationship_id}",
                        "original_filename": f"external_{relationship_id}.bin",
                        "media_type": "application/octet-stream",
                        "data": None,
                        "page_numbers": page_numbers,
                        "source_block_indexes": source_block_indexes,
                    }
                )
                continue
            source_part = _safe_source_part(target)
            if source_part not in archive_names:
                raise ValueError(f"missing DOCX relationship target: {source_part}")
            page_numbers, source_block_indexes = page_bindings.get(relationship_id, ([], []))
            occurrences.append(
                {
                    "relationship_id": relationship_id,
                    "source_part": source_part,
                    "original_filename": PurePosixPath(source_part).name,
                    "media_type": _media_type(source_part, defaults, overrides),
                    "data": archive.read(source_part),
                    "page_numbers": page_numbers,
                    "source_block_indexes": source_block_indexes,
                }
            )

    manifest = build_manifest(occurrences, source_file=docx_path.name)
    chart_records: list[dict[str, Any]] = []
    for asset in manifest["assets"]:
        if asset["sha256"] is None:
            continue
        destination = output_dir / Path(asset["relative_path"])
        destination.parent.mkdir(parents=True, exist_ok=True)
        source_part = asset["source_part"]
        data = next(item["data"] for item in occurrences if item["source_part"] == source_part)
        if destination.exists() and destination.read_bytes() != data:
            raise ValueError(f"refusing to overwrite different source asset: {destination}")
        destination.write_bytes(data)
        if asset["media_type"] == "application/vnd.openxmlformats-officedocument.drawingml.chart+xml":
            chart = _chart_record(asset, data)
            if chart is not None:
                chart_records.append(chart)
        if asset["media_type"] in SUPPORTED_GENERATION_MEDIA:
            asset["generation_input"] = {
                "relative_path": asset["relative_path"],
                "sha256": asset["sha256"],
                "media_type": asset["media_type"],
                "derivation": "original_supported",
            }
        elif asset["media_type"] in DERIVABLE_GENERATION_MEDIA:
            derivative_relative = str(DERIVED_ASSET_DIRECTORY / f"{asset['asset_id']}.png")
            derivative = output_dir / derivative_relative
            derivative.parent.mkdir(parents=True, exist_ok=True)
            try:
                with Image.open(BytesIO(data)) as image:
                    image.seek(0)
                    converted = image.convert("RGBA" if "A" in image.getbands() else "RGB")
                    converted.save(derivative, format="PNG", optimize=False, compress_level=9)
                asset["generation_input"] = {
                    "relative_path": derivative_relative,
                    "sha256": hashlib.sha256(derivative.read_bytes()).hexdigest(),
                    "media_type": "image/png",
                    "derivation": "deterministic_png",
                    "source_sha256": asset["sha256"],
                }
            except (UnidentifiedImageError, OSError, ValueError):
                asset.update(
                    classify_page_asset(
                        asset["media_type"], binding_status=asset["binding_status"], has_generation_input=False
                    )
                )
        elif asset["media_type"] in DOCUMENT_MEDIA | SPREADSHEET_MEDIA:
            derivative_relative = str(DERIVED_ASSET_DIRECTORY / f"{asset['asset_id']}.txt")
            derivative = output_dir / derivative_relative
            derivative.parent.mkdir(parents=True, exist_ok=True)
            try:
                extracted = extract_attachment_text(destination, asset["media_type"])
            # Embedded packages may be malformed, encrypted, legacy binary files,
            # or unsupported by the corresponding reader. Attachment extraction
            # is deliberately advisory and must never block unrelated pages.
            except Exception:
                extracted = ""
            if extracted:
                derivative.write_text(extracted + "\n", encoding="utf-8")
                asset["generation_input"] = {
                    "relative_path": derivative_relative,
                    "sha256": hashlib.sha256(derivative.read_bytes()).hexdigest(),
                    "media_type": "text/plain",
                    "derivation": "text_extraction",
                    "source_sha256": asset["sha256"],
                }
    if chart_records:
        manifest["chart_records"] = chart_records
    return manifest
