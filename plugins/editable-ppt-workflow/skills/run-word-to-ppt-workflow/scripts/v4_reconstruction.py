"""Sealed object-level reconstruction boundary for Word-to-PPT V4."""

from __future__ import annotations

import hashlib
import html
import hmac
import json
import math
import os
import shutil
import secrets
import stat
import sys
import re
import posixpath
from pathlib import Path, PurePosixPath
from typing import Any, Mapping
from zipfile import ZipFile
from xml.etree import ElementTree as ET

from pptx import Presentation
from PIL import Image

from cache_key import canonical_sha256
from fixed_frame import FRAME_NAMES, apply_fixed_frame, inspect_fixed_frame
from fixed_region_contract import BODY_BOX_CM, CONTRACT_VERSION, SLIDE_SIZE_CM
from page_material_bundle_v4 import verify_page_material_bundle_seal
from page_coverage import verify_coverage_contract
from workflow_v4_contract import validate_v4_artifact
from workflow_contract import PAGE_CACHE_CONTRACT_VERSION


EDITPPT_CLI = Path(__file__).resolve().parents[2] / "reconstruct-editable-slide" / "cli"
if str(EDITPPT_CLI) not in sys.path:
    sys.path.append(str(EDITPPT_CLI))
EDITPPT_RUNTIME = EDITPPT_CLI / "editppt" / "runtime"
if str(EDITPPT_RUNTIME) not in sys.path:
    sys.path.append(str(EDITPPT_RUNTIME))

from editppt.runtime.build_pptx_from_manifest import write_pptx  # noqa: E402
from editppt.runtime.editable_page_cache import inspect_editable_pptx  # noqa: E402


WORK_ITEM_VERSION = "editable-reconstruction-work-item-v1"
MANIFEST_VERSION = "editable-reconstruction-manifest-v1"
SIGNED_BUNDLE_VERSION = "editable-reconstruction-bundle-v1"
RECEIPT_VERSION = "editable-receipt-v1"
_SECRET_FILE = ".private/reconstruction-v1.key"
_BODY = BODY_BOX_CM
_EMU_PER_CM = 360000


_OwnedIdentity = tuple[str, int, int, int, int]


def _filesystem_identity(path: Path, value: os.stat_result) -> _OwnedIdentity | None:
    """Return a conservative identity for a non-reparse filesystem object."""
    try:
        mode = value.st_mode
        attributes = getattr(value, "st_file_attributes", 0)
        device = int(value.st_dev)
        inode = int(value.st_ino)
        birth_ns = getattr(value, "st_birthtime_ns", None)
        if birth_ns is None:
            birth = getattr(value, "st_birthtime", None)
            birth_ns = None if birth is None else int(float(birth) * 1_000_000_000)
        canonical = os.path.normcase(str(path.resolve(strict=True)))
    except Exception:
        return None
    reparse_flag = getattr(stat, "FILE_ATTRIBUTE_REPARSE_POINT", 0x400)
    if stat.S_ISLNK(mode) or attributes & reparse_flag:
        return None
    if device == 0 or inode == 0 or birth_ns is None:
        return None
    return canonical, device, inode, int(birth_ns), stat.S_IFMT(mode)


def _owned_handle_identity(path: Path, handle: Any) -> _OwnedIdentity | None:
    try:
        return _filesystem_identity(path, os.fstat(handle.fileno()))
    except (OSError, ValueError):
        return None


def _owned_path_identity(path: Path) -> _OwnedIdentity | None:
    try:
        return _filesystem_identity(path, os.lstat(path))
    except (FileNotFoundError, OSError):
        return None


def _still_owned(path: Path, identity: _OwnedIdentity) -> bool:
    return _owned_path_identity(path) == identity


def _sha(path: Path) -> str:
    return hashlib.sha256(path.read_bytes()).hexdigest()


def _project_file(project: Path, value: str | Path, *, must_exist: bool = True) -> Path:
    project = Path(project).resolve()
    path = Path(value)
    path = path if path.is_absolute() else project / path
    path = path.resolve(strict=must_exist)
    if project not in path.parents and path != project:
        raise ValueError("reconstruction artifact must remain inside the project")
    if must_exist and not path.is_file():
        raise ValueError("reconstruction artifact is missing")
    return path


def _read(project: Path, value: str | Path) -> tuple[Path, dict[str, Any]]:
    path = _project_file(project, value)
    try:
        payload = json.loads(path.read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError) as exc:
        raise ValueError(f"invalid reconstruction JSON: {path.name}") from exc
    if not isinstance(payload, dict):
        raise ValueError("reconstruction JSON root must be an object")
    return path, payload


def _artifact(project: Path, path: Path) -> dict[str, str]:
    return {"path": path.relative_to(project).as_posix(), "sha256": _sha(path)}


def _normalized(value: Any) -> str:
    return " ".join(str(value).split())


def _paragraphs(body_text: str) -> list[dict[str, str]]:
    values = [_normalized(line) for line in body_text.splitlines() if _normalized(line)]
    if not values and _normalized(body_text):
        values = [_normalized(body_text)]
    return [{"source_id": f"word-p{index}", "text": text} for index, text in enumerate(values, 1)]


def build_reconstruction_work_item(
    project: Path,
    *,
    material_bundle: Path,
    generation_receipt: Path,
    qa_receipt: Path,
    page_contract: Path,
    style_execution: Path,
    logo_svg: Path,
    output: Path,
) -> Path:
    """Bind the reconstructor to the accepted generation and all locked authorities."""
    project = Path(project).resolve()
    material_path, material = _read(project, material_bundle)
    generation_path, generation = _read(project, generation_receipt)
    qa_path, qa = _read(project, qa_receipt)
    contract_path, contract = _read(project, page_contract)
    style_path, _style = _read(project, style_execution)
    logo_path = _project_file(project, logo_svg)
    if not verify_page_material_bundle_seal(material):
        raise ValueError("material bundle seal is invalid")
    validate_v4_artifact("page_material_bundle_v4.schema.json", material)
    validate_v4_artifact("page_generation_v4.schema.json", generation)
    validate_v4_artifact("page_qa_v4.schema.json", qa)
    if qa["status"] != "pass":
        raise ValueError("reconstruction requires passing V4 QA")
    if not (
        material["page_number"] == generation["page_number"] == qa["page_number"] == contract.get("page_number")
        and generation["material_bundle_sha256"] == material["sealed_sha256"]
        and qa["material_bundle_sha256"] == material["sealed_sha256"]
        and qa["generation_receipt_sha256"] == _sha(generation_path)
    ):
        raise ValueError("reconstruction authorities do not bind to one page")
    body_path = _project_file(project, generation["body_image"]["path"])
    if _sha(body_path) != generation["body_image"]["sha256"]:
        raise ValueError("accepted body image was replaced")
    if _sha(style_path) != material["style_execution"]["sha256"]:
        raise ValueError("confirmed style artifact changed")
    if _sha(contract_path) != material["provenance"]["page_contract_sha256"]:
        raise ValueError("locked page contract changed")
    if _sha(logo_path) != material["provenance"]["logo_sha256"]:
        raise ValueError("locked original logo changed")
    content = material["authoritative_content"]
    work = {
        "artifact_version": WORK_ITEM_VERSION,
        "workflow_contract_version": "word-ppt-workflow-v4",
        "reconstruction_version": "editable-image-v3",
        "page_number": material["page_number"],
        "accepted_body_image": _artifact(project, body_path),
        "material_bundle": _artifact(project, material_path),
        "generation_receipt": _artifact(project, generation_path),
        "qa_receipt": _artifact(project, qa_path),
        "page_contract": _artifact(project, contract_path),
        "style_execution": _artifact(project, style_path),
        "logo_svg": _artifact(project, logo_path),
        "page_title": contract["page_title"],
        "authoritative_text": _paragraphs(content["body_text"]),
        "authoritative_tables": content["tables"],
        "page_images": material["page_images"],
        "attachment_evidence": material["attachment_evidence"],
        "search_evidence": material["search_evidence"],
        "required_presence_asset_ids": material["required_presence_asset_ids"],
        "geometry": {
            "version": CONTRACT_VERSION,
            "slide_cm": dict(SLIDE_SIZE_CM),
            "body_box_cm": dict(BODY_BOX_CM),
            "accepted_body_ratio": "17:8",
            "aspect_tolerance": 0.01,
        },
    }
    work["sealed_sha256"] = canonical_sha256(work)
    out = _project_file(project, output, must_exist=False)
    out.parent.mkdir(parents=True, exist_ok=True)
    out.write_text(json.dumps(work, ensure_ascii=False, sort_keys=True, separators=(",", ":")) + "\n", encoding="utf-8")
    return out


def validate_work_item(project: Path, path: Path) -> dict[str, Any]:
    work_path, work = _read(project, path)
    fields = {
        "artifact_version", "workflow_contract_version", "reconstruction_version", "page_number",
        "accepted_body_image", "material_bundle", "generation_receipt", "qa_receipt", "page_contract",
        "style_execution", "logo_svg", "page_title", "authoritative_text", "authoritative_tables",
        "page_images", "required_presence_asset_ids", "geometry", "sealed_sha256",
        "attachment_evidence", "search_evidence",
    }
    if set(work) != fields or work.get("artifact_version") != WORK_ITEM_VERSION:
        raise ValueError("reconstruction work item is not closed")
    seal = work.pop("sealed_sha256")
    valid = hmac.compare_digest(str(seal), canonical_sha256(work))
    work["sealed_sha256"] = seal
    if not valid:
        raise ValueError("reconstruction work item seal is invalid")
    for key in ("accepted_body_image", "material_bundle", "generation_receipt", "qa_receipt", "page_contract", "style_execution", "logo_svg"):
        record = work[key]
        file = _project_file(project, record["path"])
        if _sha(file) != record["sha256"]:
            raise ValueError(f"reconstruction authority changed: {key}")
    rebuilt = build_reconstruction_work_item(
        project,
        material_bundle=_project_file(project, work["material_bundle"]["path"]),
        generation_receipt=_project_file(project, work["generation_receipt"]["path"]),
        qa_receipt=_project_file(project, work["qa_receipt"]["path"]),
        page_contract=_project_file(project, work["page_contract"]["path"]),
        style_execution=_project_file(project, work["style_execution"]["path"]),
        logo_svg=_project_file(project, work["logo_svg"]["path"]),
        output=work_path.with_suffix(".revalidated.tmp.json"),
    )
    try:
        _, expected = _read(project, rebuilt)
    finally:
        rebuilt.unlink(missing_ok=True)
    if expected != work:
        raise ValueError("reconstruction work item is stale")
    return work


def _manifest_exact(manifest: Mapping[str, Any]) -> None:
    fields = {
        "artifact_version", "work_item_sha256", "workflow_contract_version", "reconstruction_contract_version",
        "slide", "content_box", "source", "text_boxes", "tables", "shapes", "images", "raster_components",
        "text_coverage", "table_coverage",
    }
    if set(manifest) != fields or manifest.get("artifact_version") != MANIFEST_VERSION:
        raise ValueError("reconstruction manifest is not closed")
    allowed = {
        "text_boxes": {"object_id", "name", "text", "box_px", "font_size", "font", "color", "bold", "italic", "align", "valign", "wrap", "fit_text", "z_index"},
        "tables": {"object_id", "name", "box_px", "rows", "font_size", "font_color", "cell_fill", "cell_margin_px", "z_index"},
        "shapes": {"object_id", "name", "type", "box_px", "fill", "stroke", "stroke_width", "z_index"},
        "images": {"object_id", "name", "path", "box_px", "alt", "z_index"},
        "raster_components": {"object_id", "sha256", "source_type", "source_id"},
        "text_coverage": {"source_id", "text", "object_name"},
        "table_coverage": {"table_id", "object_name"},
    }
    for section, keys in allowed.items():
        values = manifest.get(section)
        if not isinstance(values, list) or any(not isinstance(item, Mapping) or set(item) != keys for item in values):
            raise ValueError(f"reconstruction manifest {section} is not closed")
    identities: set[str] = set()
    names: set[str] = set()
    source = manifest.get("source")
    if not isinstance(source, Mapping) or set(source) != {"width_px", "height_px"}:
        raise ValueError("reconstruction manifest source is not closed")
    width, height = source["width_px"], source["height_px"]
    if type(width) is not int or type(height) is not int or width <= 0 or height <= 0:
        raise ValueError("reconstruction manifest source dimensions are invalid")
    for section in ("text_boxes", "tables", "shapes", "images"):
        for item in manifest[section]:
            if item["object_id"] in identities or item["name"] in names:
                raise ValueError("reconstruction manifest object identities are not unique")
            identities.add(item["object_id"]); names.add(item["name"])
            box = item["box_px"]
            if (not isinstance(box, list) or len(box) != 4 or any(type(value) is not int for value in box)
                    or box[0] < 0 or box[1] < 0 or box[2] <= 0 or box[3] <= 0
                    or box[0] + box[2] > width or box[1] + box[3] > height):
                raise ValueError("reconstruction manifest object box is invalid")
            for value in item.values():
                if isinstance(value, float) and not math.isfinite(value):
                    raise ValueError("reconstruction manifest numeric value is not finite")
    for section in ("text_boxes", "tables", "shapes", "images"):
        for item in manifest[section]:
            if str(item.get("name", "")).startswith("fixed-frame-"):
                raise ValueError("body manifest cannot create fixed layers")
    for table in manifest["tables"]:
        required = {"font_size", "font_color", "cell_fill", "cell_margin_px"}
        if not required <= set(table):
            raise ValueError("native table cells require explicit font, foreground, fill, and margins")
        if not isinstance(table["font_size"], (int, float)) or float(table["font_size"]) <= 0:
            raise ValueError("native table explicit font size is invalid")
        if not isinstance(table["cell_margin_px"], (int, float)) or float(table["cell_margin_px"]) < 0:
            raise ValueError("native table explicit cell margin is invalid")
        if any(not re.fullmatch(r"#[0-9A-Fa-f]{6}", str(table[field])) for field in ("font_color", "cell_fill")):
            raise ValueError("native table cell colors must be explicit RGB")


def _shape_by_name(slide) -> dict[str, Any]:
    result = {}
    for shape in slide.shapes:
        if shape.name in result:
            raise ValueError("reconstructed object names must be unique")
        result[shape.name] = shape
    return result


def _visible_body_shape(shape) -> bool:
    left, top = int(_BODY["x"] * _EMU_PER_CM), int(_BODY["y"] * _EMU_PER_CM)
    right = int((_BODY["x"] + _BODY["w"]) * _EMU_PER_CM)
    bottom = int((_BODY["y"] + _BODY["h"]) * _EMU_PER_CM)
    return shape.width > 0 and shape.height > 0 and shape.left >= left - 2 and shape.top >= top - 2 and shape.left + shape.width <= right + 2 and shape.top + shape.height <= bottom + 2


def _media_hashes(pptx: Path) -> set[str]:
    with ZipFile(pptx) as archive:
        return {hashlib.sha256(archive.read(name)).hexdigest() for name in archive.namelist() if name.startswith("ppt/media/")}


def _object_id(shape) -> str:
    nodes = shape._element.xpath(".//*[local-name()='cNvPr']")
    description = nodes[0].get("descr", "") if nodes else ""
    return description.removeprefix("object_id:") if description.startswith("object_id:") else ""


def _font_facts(shape) -> tuple[float, tuple[str, ...]]:
    sizes: list[float] = []
    colors: list[str] = []
    frames = []
    if shape.has_text_frame:
        frames = [shape.text_frame]
    elif shape.has_table:
        frames = [cell.text_frame for row in shape.table.rows for cell in row.cells]
    for frame in frames:
        for paragraph in frame.paragraphs:
            for run in paragraph.runs:
                if run.font.size is not None:
                    sizes.append(float(run.font.size.pt))
                try:
                    if run.font.color.rgb is not None:
                        colors.append(str(run.font.color.rgb).upper())
                except (AttributeError, TypeError):
                    pass
    return (max(sizes) if sizes else 12.0), tuple(colors)


def _theme_colors(pptx: Path) -> dict[str, str]:
    a = "http://schemas.openxmlformats.org/drawingml/2006/main"
    with ZipFile(pptx) as archive:
        theme_name = next((name for name in archive.namelist() if name.startswith("ppt/theme/theme") and name.endswith(".xml")), None)
        if theme_name is None:
            return {}
        root = ET.fromstring(archive.read(theme_name))
    scheme = root.find(f".//{{{a}}}clrScheme")
    colors: dict[str, str] = {}
    if scheme is not None:
        for item in list(scheme):
            color = next(iter(item), None)
            if color is None:
                continue
            value = color.get("val") or color.get("lastClr")
            if value and re.fullmatch(r"[0-9A-Fa-f]{6}", value):
                colors[item.tag.rsplit("}", 1)[-1]] = value.upper()
    aliases = {"tx1": "dk1", "tx2": "dk2", "bg1": "lt1", "bg2": "lt2"}
    for alias, source in aliases.items():
        if source in colors:
            colors[alias] = colors[source]
    return colors


def _solid_fill(element, theme: Mapping[str, str]) -> tuple[str, str | None]:
    a = "http://schemas.openxmlformats.org/drawingml/2006/main"
    sp_pr = element.find(f"./{{{a}}}spPr")
    if sp_pr is None:
        p = "http://schemas.openxmlformats.org/presentationml/2006/main"
        sp_pr = element.find(f"./{{{p}}}spPr")
    if sp_pr is None:
        return "unknown", None
    if sp_pr.find(f"./{{{a}}}noFill") is not None:
        return "none", None
    if sp_pr.find(f"./{{{a}}}gradFill") is not None:
        return "gradient", None
    if sp_pr.find(f"./{{{a}}}blipFill") is not None:
        return "image", None
    solid = sp_pr.find(f"./{{{a}}}solidFill")
    if solid is None:
        return "unknown", None
    color = next(iter(solid), None)
    if color is None:
        return "unknown", None
    local = color.tag.rsplit("}", 1)[-1]
    value = color.get("val")
    if local == "srgbClr" and value and re.fullmatch(r"[0-9A-Fa-f]{6}", value):
        return "solid", value.upper()
    if local == "sysClr":
        value = color.get("lastClr")
        if value and re.fullmatch(r"[0-9A-Fa-f]{6}", value):
            return "solid", value.upper()
    if local == "schemeClr" and value in theme:
        return "solid", theme[value]
    return "unknown", None


def _resolve_color_node(color: ET.Element | None, theme: Mapping[str, str]) -> str | None:
    if color is None:
        return None
    local = color.tag.rsplit("}", 1)[-1]
    value = color.get("val")
    if local == "srgbClr" and value and re.fullmatch(r"[0-9A-Fa-f]{6}", value):
        return value.upper()
    if local == "sysClr":
        value = color.get("lastClr")
        return value.upper() if value and re.fullmatch(r"[0-9A-Fa-f]{6}", value) else None
    if local == "schemeClr" and value in theme:
        return theme[value]
    return None


def _table_cell_facts(shape, theme: Mapping[str, str]) -> list[list[dict[str, Any]]]:
    a = "http://schemas.openxmlformats.org/drawingml/2006/main"
    facts = []
    for row_index, row in enumerate(shape.table.rows):
        fact_row = []
        for column_index, cell in enumerate(row.cells):
            tc_pr = cell._tc.tcPr
            solid = tc_pr.find(f"./{{{a}}}solidFill")
            fill = _resolve_color_node(next(iter(solid), None) if solid is not None else None, theme)
            props = [*cell._tc.findall(f".//{{{a}}}rPr"), *cell._tc.findall(f".//{{{a}}}endParaRPr")]
            sizes = [float(node.get("sz")) / 100 for node in props if node.get("sz")]
            colors = []
            for node in props:
                run_fill = node.find(f"./{{{a}}}solidFill")
                colors.append(_resolve_color_node(next(iter(run_fill), None) if run_fill is not None else None, theme))
            margins = {name: tc_pr.get(name) for name in ("marL", "marR", "marT", "marB")}
            fact_row.append({
                "text": _normalized(cell.text), "font_sizes": sizes, "font_colors": colors, "fill_color": fill,
                "margins": margins, "row_height": int(row.height), "column_width": int(shape.table.columns[column_index].width),
                "row": row_index, "column": column_index,
            })
        facts.append(fact_row)
    return facts


def _slide_background(pptx: Path, theme: Mapping[str, str]) -> str | None:
    p = "http://schemas.openxmlformats.org/presentationml/2006/main"
    a = "http://schemas.openxmlformats.org/drawingml/2006/main"
    with ZipFile(pptx) as archive:
        root = ET.fromstring(archive.read("ppt/slides/slide1.xml"))
    bg_pr = root.find(f"./{{{p}}}cSld/{{{p}}}bg/{{{p}}}bgPr")
    if bg_pr is None:
        return None
    solid = bg_pr.find(f"./{{{a}}}solidFill")
    if solid is None:
        return None
    color = next(iter(solid), None)
    if color is None:
        return None
    local = color.tag.rsplit("}", 1)[-1]
    value = color.get("val")
    if local == "srgbClr" and value and re.fullmatch(r"[0-9A-Fa-f]{6}", value):
        return value.upper()
    if local == "schemeClr" and value in theme:
        return theme[value]
    if local == "sysClr" and color.get("lastClr"):
        return str(color.get("lastClr")).upper()
    return None


def _semantic_inventory(pptx: Path, *, exclude_fixed: bool) -> list[dict[str, Any]]:
    deck = Presentation(pptx)
    if len(deck.slides) != 1:
        raise ValueError("editable page must contain exactly one slide")
    theme = _theme_colors(pptx)
    slide_background = _slide_background(pptx, theme)
    inventory = []
    for z_index, shape in enumerate(deck.slides[0].shapes):
        if exclude_fixed and shape.name in FRAME_NAMES:
            continue
        if not exclude_fixed and shape.name in FRAME_NAMES:
            raise ValueError("body PPTX contains fixed-layer objects")
        kind = "table" if shape.has_table else "image" if shape.shape_type == 13 else "text" if shape.has_text_frame and _normalized(shape.text) else "shape"
        font_size, colors = _font_facts(shape)
        xml = shape._element.xml
        fill_kind, fill_color = _solid_fill(shape._element, theme)
        record: dict[str, Any] = {
            "name": shape.name,
            "object_id": _object_id(shape),
            "kind": kind,
            "box": [int(shape.left), int(shape.top), int(shape.width), int(shape.height)],
            "z_index": z_index,
            "text": _normalized(shape.text) if shape.has_text_frame else "",
            "font_size": font_size,
            "font_colors": colors,
            "transparent": bool(re.search(r"<a:alpha[^>]+val=\"(?:0|[1-9][0-9]{0,3})\"", xml)),
            "fill_kind": fill_kind,
            "fill_color": fill_color,
            "slide_background": slide_background,
        }
        if shape.has_table:
            record["cells"] = [[_normalized(cell.text) for cell in row.cells] for row in shape.table.rows]
            record["cell_facts"] = _table_cell_facts(shape, theme)
        if kind == "image":
            record["media_sha256"] = hashlib.sha256(shape.image.blob).hexdigest()
        if kind == "shape":
            record["opaque"] = "<a:solidFill>" in xml and "<a:noFill" not in xml and not record["transparent"]
        else:
            record["opaque"] = kind == "image"
        inventory.append(record)
    return inventory


def _content_type_for(name: str, overrides: Mapping[str, str], defaults: Mapping[str, str]) -> str:
    return overrides.get("/" + name, defaults.get(Path(name).suffix.lstrip(".").casefold(), "application/octet-stream"))


def _canonical_node(element: ET.Element, relationships: Mapping[str, str]) -> Any:
    attributes = []
    for key, value in sorted(element.attrib.items()):
        if value in relationships:
            value = relationships[value]
        attributes.append((key, value))
    text = element.text or ""
    if not text.strip():
        text = ""
    return [element.tag, attributes, text, [_canonical_node(child, relationships) for child in list(element)]]


def _canonical_drawingml_body(pptx: Path, *, exclude_fixed: bool) -> dict[str, Any]:
    """Return complete ordered DrawingML with relationship IDs content-addressed."""
    p = "http://schemas.openxmlformats.org/presentationml/2006/main"
    rel_ns = "http://schemas.openxmlformats.org/package/2006/relationships"
    ct_ns = "http://schemas.openxmlformats.org/package/2006/content-types"
    with ZipFile(pptx) as archive:
        slide = ET.fromstring(archive.read("ppt/slides/slide1.xml"))
        rels = ET.fromstring(archive.read("ppt/slides/_rels/slide1.xml.rels"))
        content_types = ET.fromstring(archive.read("[Content_Types].xml"))
        overrides = {str(node.get("PartName")): str(node.get("ContentType")) for node in content_types.findall(f"{{{ct_ns}}}Override")}
        defaults = {str(node.get("Extension", "")).casefold(): str(node.get("ContentType")) for node in content_types.findall(f"{{{ct_ns}}}Default")}

        def part_digest(part_name: str, stack: frozenset[str] = frozenset()) -> str:
            if part_name in stack:
                raise ValueError("cyclic body relationship graph is forbidden")
            content_type = _content_type_for(part_name, overrides, defaults)
            data = archive.read(part_name)
            if not (part_name.endswith(".xml") or content_type.endswith("+xml") or content_type.endswith("/xml")):
                return hashlib.sha256(data).hexdigest()
            root = ET.fromstring(data)
            referenced = {value for node in root.iter() for value in node.attrib.values() if value.startswith("rId")}
            nested: dict[str, str] = {}
            directory, basename = posixpath.split(part_name)
            rel_name = posixpath.join(directory, "_rels", basename + ".rels")
            if rel_name in archive.namelist():
                nested_rels = ET.fromstring(archive.read(rel_name))
                for relation in nested_rels.findall(f"{{{rel_ns}}}Relationship"):
                    nested_id = str(relation.get("Id"))
                    if nested_id not in referenced:
                        continue
                    if relation.get("TargetMode") == "External":
                        raise ValueError("external body relationships are forbidden in editable reconstruction")
                    nested_target = posixpath.normpath(posixpath.join(directory, str(relation.get("Target"))))
                    if nested_target not in archive.namelist():
                        raise ValueError("body relationship target is missing")
                    nested[nested_id] = "relationship:" + canonical_sha256({
                        "relationship_type": str(relation.get("Type")),
                        "content_type": _content_type_for(nested_target, overrides, defaults),
                        "target_sha256": part_digest(nested_target, stack | {part_name}),
                    })
            return canonical_sha256(_canonical_node(root, nested))

        relationship_map: dict[str, str] = {}
        context_relationships = []
        for node in rels.findall(f"{{{rel_ns}}}Relationship"):
            rid = str(node.get("Id")); target = str(node.get("Target")); rel_type = str(node.get("Type"))
            if node.get("TargetMode") == "External":
                raise ValueError("external slide relationships are forbidden in editable reconstruction")
            target_name = posixpath.normpath(posixpath.join("ppt/slides", target))
            if target_name not in archive.namelist():
                raise ValueError("slide relationship target is missing")
            content_type = _content_type_for(target_name, overrides, defaults)
            descriptor = canonical_sha256({
                "relationship_type": rel_type,
                "content_type": content_type,
                "target_sha256": part_digest(target_name),
            })
            relationship_map[rid] = f"relationship:{descriptor}"
            if rel_type.endswith("/slideLayout"):
                context_relationships.append(relationship_map[rid])
        sp_tree = slide.find(f"./{{{p}}}cSld/{{{p}}}spTree")
        if sp_tree is None:
            raise ValueError("slide shape tree is missing")
        body = []
        fixed_seen = []
        root_nonvisual = None
        root_group_properties = None
        for child in list(sp_tree):
            local = child.tag.rsplit("}", 1)[-1]
            if local == "nvGrpSpPr":
                root_nonvisual = _canonical_node(child, relationship_map)
                continue
            if local == "grpSpPr":
                root_group_properties = _canonical_node(child, relationship_map)
                continue
            nv = child.find(f".//{{{p}}}cNvPr")
            name = nv.get("name") if nv is not None else None
            if name in FRAME_NAMES:
                fixed_seen.append(name)
                if exclude_fixed:
                    continue
                raise ValueError("body PPTX contains fixed-layer objects")
            body.append(_canonical_node(child, relationship_map))
        if exclude_fixed and sorted(fixed_seen) != sorted(FRAME_NAMES):
            raise ValueError("final page fixed object identity is incomplete")
        background = slide.find(f"./{{{p}}}cSld/{{{p}}}bg")
        color_map = slide.find(f"./{{{p}}}clrMapOvr")
        return {
            "root_nonvisual": root_nonvisual,
            "root_group_properties": root_group_properties,
            "body_sequence": body,
            "background": _canonical_node(background, relationship_map) if background is not None else None,
            "color_map": _canonical_node(color_map, relationship_map) if color_map is not None else None,
            "context_relationships": sorted(context_relationships),
        }


def _intersection_area(left: list[int], right: list[int]) -> int:
    x1 = max(left[0], right[0]); y1 = max(left[1], right[1])
    x2 = min(left[0] + left[2], right[0] + right[2]); y2 = min(left[1] + left[3], right[1] + right[3])
    return max(0, x2 - x1) * max(0, y2 - y1)


def _rect_union_area(rectangles: list[list[int]]) -> int:
    if not rectangles:
        return 0
    xs = sorted({value for rect in rectangles for value in (rect[0], rect[0] + rect[2])})
    total = 0
    for left, right in zip(xs, xs[1:]):
        intervals = sorted((rect[1], rect[1] + rect[3]) for rect in rectangles if rect[0] < right and rect[0] + rect[2] > left)
        covered = 0; start = end = None
        for top, bottom in intervals:
            if start is None:
                start, end = top, bottom
            elif top > end:
                covered += end - start; start, end = top, bottom
            else:
                end = max(end, bottom)
        if start is not None:
            covered += end - start
        total += (right - left) * covered
    return total


def _luminance(rgb: str) -> float:
    values = [int(rgb[offset:offset + 2], 16) / 255 for offset in (0, 2, 4)]
    values = [value / 12.92 if value <= 0.04045 else ((value + 0.055) / 1.055) ** 2.4 for value in values]
    return 0.2126 * values[0] + 0.7152 * values[1] + 0.0722 * values[2]


def _contrast_ratio(foreground: str, background: str) -> float:
    fore_lum, back_lum = _luminance(foreground), _luminance(background)
    return (max(fore_lum, back_lum) + 0.05) / (min(fore_lum, back_lum) + 0.05)


def _validate_table_cells(item: Mapping[str, Any]) -> None:
    facts = item.get("cell_facts")
    if not isinstance(facts, list) or not facts:
        raise ValueError(f"native table cell facts are missing: {item['name']}")
    for row in facts:
        for cell in row:
            label = f"{item['name']}[{cell['row']},{cell['column']}]"
            sizes = cell.get("font_sizes")
            colors = cell.get("font_colors")
            fill = cell.get("fill_color")
            margins = cell.get("margins")
            if not sizes or not colors or any(color is None for color in colors) or fill is None:
                raise ValueError(f"native table cell requires explicit resolvable font and fill RGB: {label}")
            if min(float(size) for size in sizes) < 8:
                raise ValueError(f"native table cell font is below the readable minimum: {label}")
            if not isinstance(margins, Mapping) or any(margins.get(name) is None for name in ("marL", "marR", "marT", "marB")):
                raise ValueError(f"native table cell margins must be explicit: {label}")
            margin_values = {name: int(margins[name]) for name in margins}
            width_pt = (int(cell["column_width"]) - margin_values["marL"] - margin_values["marR"]) / 914400 * 72
            height_pt = (int(cell["row_height"]) - margin_values["marT"] - margin_values["marB"]) / 914400 * 72
            font = max(float(size) for size in sizes)
            if width_pt < max(18, font * 1.5) or height_pt < font * 1.15:
                raise ValueError(f"native table cell has insufficient readable area: {label}")
            chars_per_line = max(1, int(width_pt / max(font * 0.55, 1)))
            line_count = max(1, int(height_pt / max(font * 1.2, 1)))
            if len(str(cell.get("text", ""))) > chars_per_line * line_count:
                raise ValueError(f"native table cell text exceeds actual capacity: {label}")
            if any(_contrast_ratio(str(color), str(fill)) < 4.5 for color in colors):
                raise ValueError(f"native table cell has insufficient contrast: {label}")


def _validate_readability(inventory: list[dict[str, Any]]) -> None:
    for index, item in enumerate(inventory):
        if item["kind"] not in {"text", "table"}:
            continue
        if item["kind"] == "table":
            _validate_table_cells(item)
        width_pt = item["box"][2] / 914400 * 72
        height_pt = item["box"][3] / 914400 * 72
        font = max(4.0, float(item["font_size"]))
        content = item["text"] if item["kind"] == "text" else "".join(cell for row in item.get("cells", []) for cell in row)
        if width_pt < 18 or height_pt < max(8, font * 0.7) or width_pt * height_pt < max(80, len(content) * font * font * 0.12):
            raise ValueError(f"visible editable object has insufficient capacity: {item['name']}")
        if item["transparent"]:
            raise ValueError(f"visible editable object has transparent text: {item['name']}")
        background = item.get("fill_color") if item.get("fill_kind") == "solid" else None
        if background is None and item.get("fill_kind") in {None, "none", "unknown"}:
            for lower in reversed(inventory[:index]):
                if lower.get("fill_kind") == "solid" and lower.get("fill_color") and _intersection_area(item["box"], lower["box"]) >= item["box"][2] * item["box"][3] * 0.9:
                    background = lower["fill_color"]
                    break
                if lower.get("kind") == "image" and _intersection_area(item["box"], lower["box"]) > 0:
                    break
            else:
                background = item.get("slide_background")
        if item["kind"] == "text" and background and item.get("font_colors"):
            if any(_contrast_ratio(color, str(background)) < 4.5 for color in item["font_colors"] if re.fullmatch(r"[0-9A-Fa-f]{6}", color)):
                raise ValueError(f"visible editable object has insufficient contrast: {item['name']}")
        area = item["box"][2] * item["box"][3]
        for cover in inventory[index + 1:]:
            if cover.get("opaque") and _intersection_area(item["box"], cover["box"]) >= area * 0.5:
                raise ValueError(f"visible editable object is occluded by a higher object: {item['name']}")


def _validate_fixed_logo_bytes(pptx: Path, logo_sha256: str) -> None:
    p = "http://schemas.openxmlformats.org/presentationml/2006/main"
    a = "http://schemas.openxmlformats.org/drawingml/2006/main"
    r = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    rel_ns = "http://schemas.openxmlformats.org/package/2006/relationships"
    with ZipFile(pptx) as archive:
        slide = ET.fromstring(archive.read("ppt/slides/slide1.xml"))
        rels = ET.fromstring(archive.read("ppt/slides/_rels/slide1.xml.rels"))
        targets = {node.get("Id"): node.get("Target") for node in rels.findall(f"{{{rel_ns}}}Relationship")}
        logos = []
        for picture in slide.findall(f".//{{{p}}}pic"):
            nv = picture.find(f"./{{{p}}}nvPicPr/{{{p}}}cNvPr")
            if nv is not None and nv.get("name") == "fixed-frame-logo":
                logos.append(picture)
        if len(logos) != 1:
            raise ValueError("fixed logo must exist exactly once in final OOXML")
        svg_nodes = logos[0].findall(f".//{{http://schemas.microsoft.com/office/drawing/2016/SVG/main}}svgBlip")
        svg_rel_ids = [node.get(f"{{{r}}}embed") for node in svg_nodes]
        svg_targets = [(rid, targets.get(rid)) for rid in svg_rel_ids if rid]
        all_svg_relationships = [(rid, target) for rid, target in targets.items() if isinstance(target, str) and target.lower().endswith(".svg")]
        if len(svg_targets) != 1 or len(all_svg_relationships) != 1 or svg_targets[0] != all_svg_relationships[0]:
            raise ValueError("fixed logo must have exactly one SVG media relationship")
        target = svg_targets[0][1]
        media_name = str(Path("ppt/slides") / target).replace("\\", "/")
        media_name = __import__("posixpath").normpath(media_name)
        if hashlib.sha256(archive.read(media_name)).hexdigest() != logo_sha256:
            raise ValueError("fixed logo embedded SVG bytes do not match the locked original")
        content_types = archive.read("[Content_Types].xml").decode("utf-8")
        if "image/svg+xml" not in content_types:
            raise ValueError("fixed logo SVG content type is missing")


def _derive_final_facts(work: Mapping[str, Any], manifest: Mapping[str, Any], inventory: list[dict[str, Any]]) -> dict[str, Any]:
    text_map, table_map = _strict_coverage_maps(work, manifest)
    by_name = {item["name"]: item for item in inventory}
    text_coverage = []
    for source in work["authoritative_text"]:
        mapping = text_map[source["source_id"]]
        actual = by_name.get(mapping.get("object_name")) if mapping else None
        if actual is None or actual["kind"] != "text" or actual["text"] != source["text"]:
            raise ValueError(f"final page Word text coverage changed: {source['source_id']}")
        text_coverage.append({"source_id": source["source_id"], "object_name": actual["name"]})
    table_coverage = []
    for source in work["authoritative_tables"]:
        mapping = table_map[source["table_id"]]
        actual = by_name.get(mapping.get("object_name")) if mapping else None
        expected_cells = [[_normalized(cell) for cell in row] for row in source["rows"]]
        if actual is None or actual["kind"] != "table" or actual.get("cells") != expected_cells:
            raise ValueError(f"final page Word table coverage changed: {source['table_id']}")
        table_coverage.append({"table_id": source["table_id"], "object_name": actual["name"], "cell_count": sum(map(len, expected_cells))})
    counts = {"text": 0, "table": 0, "image": 0, "shape": 0}
    for item in inventory:
        counts[item["kind"]] += 1
    return {"text_coverage": text_coverage, "table_coverage": table_coverage, "object_counts": counts}


def _strict_coverage_maps(work: Mapping[str, Any], manifest: Mapping[str, Any]) -> tuple[dict[str, Mapping[str, Any]], dict[str, Mapping[str, Any]]]:
    """Require a bijection from every Word source item to one matching native body object."""
    text_sources = work["authoritative_text"]
    table_sources = work["authoritative_tables"]
    expected_text = [item["source_id"] for item in text_sources]
    expected_tables = [item["table_id"] for item in table_sources]
    if len(set(expected_text)) != len(expected_text) or len(set(expected_tables)) != len(expected_tables):
        raise ValueError("authoritative Word source identities are not unique")

    text_entries = manifest["text_coverage"]
    table_entries = manifest["table_coverage"]
    text_ids = [item["source_id"] for item in text_entries]
    table_ids = [item["table_id"] for item in table_entries]
    if len(text_ids) != len(expected_text) or len(set(text_ids)) != len(text_ids) or set(text_ids) != set(expected_text):
        raise ValueError("text coverage must map every authoritative source exactly once")
    if len(table_ids) != len(expected_tables) or len(set(table_ids)) != len(table_ids) or set(table_ids) != set(expected_tables):
        raise ValueError("table coverage must map every authoritative table exactly once")
    if len({item["object_name"] for item in text_entries}) != len(text_entries):
        raise ValueError("text coverage cannot collapse multiple Word paragraphs into one object")
    if len({item["object_name"] for item in table_entries}) != len(table_entries):
        raise ValueError("table coverage cannot collapse multiple Word tables into one object")

    text_objects = {item["name"]: item for item in manifest["text_boxes"]}
    table_objects = {item["name"]: item for item in manifest["tables"]}
    covered_text_names = [item["object_name"] for item in text_entries]
    covered_table_names = [item["object_name"] for item in table_entries]
    if len(covered_text_names) != len(text_objects) or set(covered_text_names) != set(text_objects):
        raise ValueError("text coverage must reference every factual text object exactly once")
    if len(covered_table_names) != len(table_objects) or set(covered_table_names) != set(table_objects):
        raise ValueError("table coverage must reference every native table object exactly once")
    text_map = {item["source_id"]: item for item in text_entries}
    table_map = {item["table_id"]: item for item in table_entries}
    source_text = {item["source_id"]: item["text"] for item in text_sources}
    for source_id, mapping in text_map.items():
        actual = text_objects.get(mapping["object_name"])
        if actual is None or mapping["text"] != source_text[source_id] or actual["text"] != source_text[source_id]:
            raise ValueError(f"text coverage does not identify one matching text object: {source_id}")
    for table_id, mapping in table_map.items():
        if mapping["object_name"] not in table_objects:
            raise ValueError(f"table coverage does not identify one native table object: {table_id}")
    return text_map, table_map


def validate_reconstructed_page(project: Path, work_item: Path, manifest_path: Path, pptx: Path) -> dict[str, Any]:
    """Re-open the package and derive coverage/counts; worker claims are never trusted."""
    project = Path(project).resolve()
    work = validate_work_item(project, work_item)
    manifest_file, manifest = _read(project, manifest_path)
    _manifest_exact(manifest)
    if manifest["work_item_sha256"] != _sha(_project_file(project, work_item)):
        raise ValueError("reconstruction manifest is bound to a different work item")
    if manifest["workflow_contract_version"] != "fixed-canvas-cm-v2" or manifest["reconstruction_contract_version"] != "editable-image-v3":
        raise ValueError("reconstruction manifest version mismatch")
    with Image.open(_project_file(project, work["accepted_body_image"]["path"])) as accepted:
        expected_source = {"width_px": accepted.width, "height_px": accepted.height}
    if manifest["source"] != expected_source:
        raise ValueError("reconstruction source dimensions must match the accepted body image")
    text_map, table_map = _strict_coverage_maps(work, manifest)
    output = _project_file(project, pptx)
    inspect_editable_pptx(output)
    with ZipFile(output) as archive:
        slide_xml = archive.read("ppt/slides/slide1.xml").decode("utf-8")
    for section in ("text_boxes", "tables", "shapes", "images"):
        for item in manifest[section]:
            marker = f'descr="object_id:{html.escape(str(item["object_id"]), quote=True)}"'
            if marker not in slide_xml:
                raise ValueError(f"stable object identity is missing after package reopen: {item['object_id']}")
    deck = Presentation(output)
    if len(deck.slides) != 1:
        raise ValueError("reconstructed page must contain exactly one slide")
    slide = deck.slides[0]
    by_name = _shape_by_name(slide)
    if any(name.startswith("fixed-frame-") for name in by_name):
        raise ValueError("body reconstructor must not create fixed layers")
    for shape in slide.shapes:
        if not _visible_body_shape(shape):
            raise ValueError(f"body object is invisible or outside the body frame: {shape.name}")
    inventory = _semantic_inventory(output, exclude_fixed=False)
    _validate_readability(inventory)
    media_hashes = _media_hashes(output)
    if work["accepted_body_image"]["sha256"] in media_hashes:
        raise ValueError("accepted full-body image must not be embedded in the editable page")
    body_area = int(_BODY["w"] * _EMU_PER_CM) * int(_BODY["h"] * _EMU_PER_CM)
    raster_boxes = [item["box"] for item in inventory if item["kind"] == "image"]
    if _rect_union_area(raster_boxes) >= body_area * 0.80:
        raise ValueError("combined raster coverage reconstructs a flattened body")
    text_coverage = []
    for item in work["authoritative_text"]:
        mapping = text_map[item["source_id"]]
        if not mapping or mapping.get("text") != item["text"] or mapping.get("object_name") not in by_name:
            raise ValueError(f"authoritative text coverage is missing: {item['source_id']}")
        shape = by_name[mapping["object_name"]]
        if not shape.has_text_frame or _normalized(shape.text) != item["text"]:
            raise ValueError(f"authoritative text is not a visible editable object: {item['source_id']}")
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if run.font.size is not None and run.font.size.pt < 4:
                    raise ValueError("superficial tiny text cannot satisfy coverage")
        text_coverage.append({"source_id": item["source_id"], "object_name": shape.name})
    table_coverage = []
    for table in work["authoritative_tables"]:
        mapping = table_map[table["table_id"]]
        if not mapping or mapping.get("object_name") not in by_name:
            raise ValueError(f"authoritative table coverage is missing: {table['table_id']}")
        shape = by_name[mapping["object_name"]]
        if not shape.has_table:
            raise ValueError("authoritative table must be a native PowerPoint table")
        actual = [[_normalized(cell.text) for cell in row.cells] for row in shape.table.rows]
        expected = [[_normalized(cell) for cell in row] for row in table["rows"]]
        if actual != expected:
            raise ValueError(f"authoritative table cells changed: {table['table_id']}")
        table_coverage.append({"table_id": table["table_id"], "object_name": shape.name, "cell_count": sum(map(len, expected))})
    component_by_id = {entry["object_id"]: entry for entry in manifest["raster_components"]}
    page_images = {entry["asset_id"]: entry for entry in work["page_images"]}
    attachments = {
        key: entry
        for entry in work["attachment_evidence"]
        for key in (entry["evidence_id"], entry["asset_id"])
    }
    search_evidence = {entry["evidence_id"]: entry for entry in work["search_evidence"]}
    required_enterprise_ids = {
        entry["evidence_id"] for entry in work["search_evidence"]
        if entry.get("material_role") == "enterprise_logo"
    }
    present_enterprise_ids: set[str] = set()
    inventory_by_name = {item["name"]: item for item in inventory}
    present_required_ids = set()
    for image in manifest["images"]:
        component = component_by_id.get(image["object_id"])
        image_path = Path(image["path"])
        asset = _project_file(project, image_path if image_path.is_absolute() else manifest_file.parent / image_path)
        rendered = inventory_by_name.get(image["name"])
        if not component or not rendered or rendered.get("kind") != "image":
            raise ValueError("raster component has no matching PPTX picture relationship")
        actual_media_sha = rendered.get("media_sha256")
        if component.get("sha256") != _sha(asset) or component.get("sha256") != actual_media_sha:
            raise ValueError("raster component provenance is missing or stale")
        if component.get("source_type") not in {"page-image", "attachment", "search-evidence", "decorative-texture"}:
            raise ValueError("raster component source type is not allowed")
        if component.get("source_type") == "page-image":
            source = page_images.get(component["source_id"])
            if source is None or source["sha256"] != component["sha256"] or source["path"] != image_path.as_posix().removeprefix("../"):
                # Path spelling is supplemental; the relationship bytes/hash above
                # are authoritative. Accept manifest-relative paths that resolve to
                # the exact sealed source file.
                source_path = _project_file(project, source["path"]) if source is not None else None
                if source is None or source["sha256"] != component["sha256"] or source_path != asset:
                    raise ValueError("page-image raster provenance does not match the sealed material bundle")
            if source["presence_policy"] == "required_presence":
                present_required_ids.add(source["asset_id"])
        elif component.get("source_type") == "attachment":
            source = attachments.get(component["source_id"])
            if source is None or not str(source["media_type"]).startswith("image/") or source["sha256"] != actual_media_sha or _project_file(project, source["path"]) != asset:
                raise ValueError("attachment raster provenance does not match sealed evidence")
        elif component.get("source_type") == "search-evidence":
            # Search evidence does not carry a local path or MIME type in the
            # material contract. It may therefore appear as a raster only when
            # the sealed evidence digest is the digest of the exact embedded
            # media bytes. A screenshot/derivative needs a separately sealed
            # derived-asset receipt; the untrusted manifest cannot invent one.
            source = search_evidence.get(component["source_id"])
            if source is None or source["sha256"] != actual_media_sha:
                raise ValueError("search-evidence raster provenance does not match sealed evidence")
            if source.get("material_role") == "enterprise_logo":
                if _project_file(project, source["local_path"]) != asset:
                    raise ValueError("enterprise Logo raster path does not match sealed evidence")
                present_enterprise_ids.add(source["evidence_id"])
        else:
            if not str(component["source_id"]).startswith(f"decorative:{component['object_id']}") or component["source_id"] in page_images:
                raise ValueError("decorative raster cannot impersonate a sealed source")
    missing_required = sorted(set(work["required_presence_asset_ids"]) - present_required_ids)
    if missing_required:
        raise ValueError(f"required page image is absent from reconstruction: {missing_required[0]}")
    missing_enterprise = sorted(required_enterprise_ids - present_enterprise_ids)
    if missing_enterprise:
        raise ValueError(f"required enterprise Logo is absent from reconstruction: {missing_enterprise[0]}")
    counts = {"text": 0, "table": 0, "image": 0, "shape": 0}
    for shape in slide.shapes:
        if shape.has_table:
            counts["table"] += 1
        elif shape.shape_type == 13:
            counts["image"] += 1
        elif shape.has_text_frame and _normalized(shape.text):
            counts["text"] += 1
        else:
            counts["shape"] += 1
    return {
        "manifest": _artifact(project, manifest_file),
        "pptx": _artifact(project, output),
        "object_counts": counts,
        "text_coverage": text_coverage,
        "table_coverage": table_coverage,
        "raster_components": list(manifest["raster_components"]),
    }


def _secret(project: Path) -> bytes:
    path = _project_file(project, _SECRET_FILE, must_exist=False)
    path.parent.mkdir(parents=True, exist_ok=True)
    if path.is_symlink():
        raise ValueError("reconstruction integrity key cannot be a link")
    if not path.exists():
        path.write_bytes(secrets.token_bytes(32))
        try:
            os.chmod(path, 0o600)
        except OSError:
            pass
    payload = path.read_bytes()
    if len(payload) != 32:
        raise ValueError("reconstruction signing key is invalid")
    return payload


def build_and_sign_reconstruction(
    project: Path,
    *,
    work_item: Path,
    manifest: Path,
    gateway_invocation: Path,
    body_pptx: Path,
    final_pptx: Path,
    bundle_output: Path,
) -> Path:
    """Controlled local backend: validate manifest, build, reopen, apply fixed layers, attest."""
    project = Path(project).resolve()
    work = validate_work_item(project, work_item)
    from v4_reconstruction_gateway import verify_signed_bundle as verify_gateway_bundle
    gateway_path = _project_file(project, gateway_invocation)
    derived_manifest = verify_gateway_bundle(project, gateway_path, _project_file(project, work_item), consume=False)
    manifest_path, manifest_payload = _read(project, manifest)
    if derived_manifest != manifest_path:
        raise ValueError("manifest is not the signed reconstruction gateway result")
    _manifest_exact(manifest_payload)
    for image in manifest_payload["images"]:
        image_path = Path(image["path"])
        _project_file(project, image_path if image_path.is_absolute() else manifest_path.parent / image_path)
    body_path = _project_file(project, body_pptx, must_exist=False)
    write_pptx(manifest_payload, body_path, manifest_path)
    body_validation = validate_reconstructed_page(project, work_item, manifest_path, body_path)
    final_path = _project_file(project, final_pptx, must_exist=False)
    final_path.parent.mkdir(parents=True, exist_ok=True)
    final_path.write_bytes(body_path.read_bytes())
    _, style = _read(project, work["style_execution"]["path"])
    apply_fixed_frame(
        final_path, page_title=work["page_title"], page_number=work["page_number"],
        style_execution=style, logo_svg=_project_file(project, work["logo_svg"]["path"]),
    )
    fixed = inspect_fixed_frame(
        final_path, expected_title=work["page_title"], expected_page_number=work["page_number"],
        style_execution=style, logo_svg=_project_file(project, work["logo_svg"]["path"]),
    )
    if not fixed["passed"]:
        raise ValueError("fixed-layer validation failed: " + " | ".join(fixed["issues"]))
    inspection = inspect_editable_pptx(final_path)
    payload = {
        "artifact_version": SIGNED_BUNDLE_VERSION,
        "work_item_sha256": _sha(_project_file(project, work_item)),
        "manifest": body_validation["manifest"],
        "gateway_invocation": _artifact(project, gateway_path),
        "body_pptx": body_validation["pptx"],
        "editable_page": _artifact(project, final_path),
        "object_counts": body_validation["object_counts"],
        "text_coverage": body_validation["text_coverage"],
        "table_coverage": body_validation["table_coverage"],
        "raster_components": body_validation["raster_components"],
        "accepted_body_embedded": False,
        "fixed_layers": {"passed": True, "names": ["fixed-frame-title", "fixed-frame-logo", "fixed-frame-footer", "fixed-frame-page-number"]},
        "slide_fingerprint": inspection.slide_fingerprints[0],
        "nonce": secrets.token_hex(16),
    }
    payload["attestation"] = hmac.new(_secret(project), canonical_sha256(payload).encode("ascii"), hashlib.sha256).hexdigest()
    out = _project_file(project, bundle_output, must_exist=False)
    out.parent.mkdir(parents=True, exist_ok=True)
    out.write_text(json.dumps(payload, ensure_ascii=False, sort_keys=True, separators=(",", ":")) + "\n", encoding="utf-8")
    return out


def verify_signed_reconstruction(project: Path, work_item: Path, bundle: Path) -> dict[str, Any]:
    project = Path(project).resolve()
    work = validate_work_item(project, work_item)
    bundle_path, payload = _read(project, bundle)
    fields = {
        "artifact_version", "work_item_sha256", "manifest", "gateway_invocation", "body_pptx", "editable_page", "object_counts",
        "text_coverage", "table_coverage", "raster_components", "accepted_body_embedded", "fixed_layers",
        "slide_fingerprint", "nonce", "attestation",
    }
    if set(payload) != fields or payload.get("artifact_version") != SIGNED_BUNDLE_VERSION:
        raise ValueError("signed reconstruction bundle is not closed")
    signature = payload.pop("attestation")
    expected = hmac.new(_secret(project), canonical_sha256(payload).encode("ascii"), hashlib.sha256).hexdigest()
    payload["attestation"] = signature
    if not hmac.compare_digest(str(signature), expected):
        raise ValueError("reconstruction attestation is invalid")
    if payload["work_item_sha256"] != _sha(_project_file(project, work_item)) or payload["accepted_body_embedded"] is not False:
        raise ValueError("reconstruction bundle authority mismatch")
    from v4_reconstruction_gateway import verify_signed_bundle as verify_gateway_bundle
    gateway_path = _project_file(project, payload["gateway_invocation"]["path"])
    if _artifact(project, gateway_path) != payload["gateway_invocation"]:
        raise ValueError("reconstruction gateway invocation changed")
    derived_manifest = verify_gateway_bundle(project, gateway_path, _project_file(project, work_item), consume=False)
    if derived_manifest != _project_file(project, payload["manifest"]["path"]):
        raise ValueError("signed reconstruction is not bound to its gateway manifest")
    validation = validate_reconstructed_page(project, work_item, payload["manifest"]["path"], payload["body_pptx"]["path"])
    if payload["manifest"] != validation["manifest"] or payload["body_pptx"] != validation["pptx"]:
        raise ValueError("reconstruction bundle artifact identity changed")
    final_path = _project_file(project, payload["editable_page"]["path"])
    if _sha(final_path) != payload["editable_page"]["sha256"]:
        raise ValueError("editable page was replaced")
    _, style = _read(project, work["style_execution"]["path"])
    fixed = inspect_fixed_frame(final_path, expected_title=work["page_title"], expected_page_number=work["page_number"], style_execution=style, logo_svg=_project_file(project, work["logo_svg"]["path"]))
    if not fixed["passed"]:
        raise ValueError("fixed layers changed after reconstruction")
    _validate_fixed_logo_bytes(final_path, work["logo_svg"]["sha256"])
    body_path = _project_file(project, payload["body_pptx"]["path"])
    body_drawingml = _canonical_drawingml_body(body_path, exclude_fixed=False)
    final_body_drawingml = _canonical_drawingml_body(final_path, exclude_fixed=True)
    if canonical_sha256(final_body_drawingml) != canonical_sha256(body_drawingml):
        raise ValueError("final page body DrawingML is not equivalent to the verified body PPTX")
    body_inventory = _semantic_inventory(body_path, exclude_fixed=False)
    final_body_inventory = _semantic_inventory(final_path, exclude_fixed=True)
    for item in final_body_inventory:
        left = int(_BODY["x"] * _EMU_PER_CM); top = int(_BODY["y"] * _EMU_PER_CM)
        right = int((_BODY["x"] + _BODY["w"]) * _EMU_PER_CM); bottom = int((_BODY["y"] + _BODY["h"]) * _EMU_PER_CM)
        box = item["box"]
        if box[0] < left - 2 or box[1] < top - 2 or box[0] + box[2] > right + 2 or box[1] + box[3] > bottom + 2:
            raise ValueError(f"final page body object lies outside the body frame: {item['name']}")
    _validate_readability(final_body_inventory)
    _, manifest = _read(project, payload["manifest"]["path"])
    final_facts = _derive_final_facts(work, manifest, final_body_inventory)
    inspection = inspect_editable_pptx(final_path)
    final_artifact = _artifact(project, final_path)
    receipt = {
        "artifact_version": RECEIPT_VERSION,
        "workflow_contract_version": "word-ppt-workflow-v4",
        "reconstruction_version": "editable-image-v3",
        "page_number": work["page_number"],
        "work_item_sha256": payload["work_item_sha256"],
        "generation_sha256": work["generation_receipt"]["sha256"],
        "qa_sha256": work["qa_receipt"]["sha256"],
        "material_bundle_sha256": work["material_bundle"]["sha256"],
        "page_contract_sha256": work["page_contract"]["sha256"],
        "style_execution_sha256": work["style_execution"]["sha256"],
        "editable_page": final_artifact,
        "object_manifest": validation["manifest"],
        "object_counts": final_facts["object_counts"],
        "text_coverage": final_facts["text_coverage"],
        "table_coverage": final_facts["table_coverage"],
        "raster_components": validation["raster_components"],
        "flattened_body_image": False,
        "accepted_body_image_sha256": work["accepted_body_image"]["sha256"],
        "accepted_body_embedded": False,
        "fixed_layers_added": True,
        "fixed_layer_names": list(FRAME_NAMES),
        "slide_fingerprint": inspection.slide_fingerprints[0],
        "signed_bundle": _artifact(project, bundle_path),
        "gateway_invocation": payload["gateway_invocation"],
    }
    validate_v4_artifact("editable_receipt_v4.schema.json", receipt)
    return receipt


def write_editable_receipt(project: Path, work_item: Path, bundle: Path, output: Path) -> Path:
    receipt = verify_signed_reconstruction(project, work_item, bundle)
    out = _project_file(project, output, must_exist=False)
    out.parent.mkdir(parents=True, exist_ok=True)
    out.write_text(json.dumps(receipt, ensure_ascii=False, sort_keys=True, separators=(",", ":")) + "\n", encoding="utf-8")
    return out


def collect_reconstruction_closure(
    project: Path,
    *,
    seeds: list[Path],
    page_number: int,
    material_bundle_path: Path,
    material_bundle_file_sha256: str,
    material_bundle_sha256: str,
    authority_identity: str,
) -> dict[str, Path]:
    """Collect the complete project-local dependency graph plus V4 schemas."""
    project = Path(project).resolve()
    material = _load_authenticated_page_material_bundle(
        project,
        expected_path=material_bundle_path,
        expected_file_sha256=material_bundle_file_sha256,
        expected_sealed_sha256=material_bundle_sha256,
        expected_page_number=page_number,
        expected_authority_identity=authority_identity,
    )
    collected: dict[str, Path] = {}
    queue: list[Path] = []
    visited: set[Path] = set()
    authenticated_material_files = _authenticated_material_paths(project, material)
    logo = _project_file(project, "00_source/company_logo.svg")
    allowed_source_files = {
        path for path in authenticated_material_files
        if path.relative_to(project).as_posix().startswith("00_source/")
    }
    allowed_source_files.add(logo)

    def enqueue(path: Path) -> None:
        try:
            resolved = path.resolve(strict=True)
            relative = resolved.relative_to(project).as_posix()
        except (OSError, ValueError):
            return
        parts = Path(relative).parts
        if any(part in {".cache", ".private", ".workflow", "08_final"} for part in parts):
            return
        if any(token in relative.casefold() for token in ("secret", "token", "nonce", ".key")):
            return
        if relative in {
            "workflow_run.json", "01_page_contracts/source_lock.json",
            "confirm_ui/page_requirement_summary.json", "00_source/source.docx",
            "00_source/pages.json",
        }:
            return
        if relative.startswith("00_source/") and resolved not in allowed_source_files:
            return
        if resolved.is_file() and resolved not in visited:
            queue.append(resolved)

    for seed in seeds:
        enqueue(Path(seed))

    for source in {*authenticated_material_files, logo}:
        enqueue(source)

    while queue:
        candidate = Path(queue.pop()).resolve(strict=True)
        if candidate in visited:
            continue
        visited.add(candidate)
        try:
            relative = candidate.relative_to(project).as_posix()
        except ValueError:
            continue
        collected[relative] = candidate
    schemas = Path(__file__).resolve().parents[1] / "schemas"
    for schema in schemas.glob("*.schema.json"):
        collected[f"@schema/{schema.name}"] = schema.resolve()
    return collected


def _load_authenticated_page_material_bundle(
    project: Path,
    *,
    expected_path: str | Path,
    expected_file_sha256: str,
    expected_sealed_sha256: str,
    expected_page_number: int,
    expected_authority_identity: str,
) -> dict[str, Any]:
    """Load the one current project-owned bundle allowed to authorize source paths."""
    try:
        path, bundle = _read(project, expected_path)
        if (
            _sha(path) != expected_file_sha256
            or bundle.get("sealed_sha256") != expected_sealed_sha256
            or bundle.get("page_number") != expected_page_number
            or bundle.get("effective_page_authority", {}).get("sealed_sha256")
            != expected_authority_identity
        ):
            raise ValueError
        validate_v4_artifact("page_material_bundle_v4.schema.json", bundle)
        if not verify_page_material_bundle_seal(bundle, project):
            raise ValueError
    except (OSError, ValueError, KeyError, TypeError):
        raise ValueError("authenticated current page material bundle is invalid") from None
    return bundle


def _authenticated_material_paths(
    project: Path, bundle: Mapping[str, Any],
) -> set[Path]:
    paths: set[Path] = set()
    for collection_name, field in (
        ("page_images", "path"),
        ("attachment_evidence", "path"),
        ("search_evidence", "local_path"),
    ):
        for item in bundle[collection_name]:
            value = item.get(field)
            if not isinstance(value, str):
                raise ValueError("authenticated current page material bundle has an invalid path")
            path = _project_file(project, value)
            paths.add(path)
    return paths


def _expected_completed_logicals(
    project: Path,
    job: Mapping[str, Any],
    material: Mapping[str, Any],
    logical_files: Mapping[str, Any],
    hit_root: Path,
) -> set[str]:
    """Derive the current completed-page inventory without cache-manifest authority."""
    page_number = job.get("page_number")
    if type(page_number) is not int or page_number < 1:
        raise ValueError("completed page inventory has no valid page identity")
    logicals = {
        "00_source/company_logo.svg",
        str(job.get("contract_file")),
        str(job.get("material_bundle_file")),
        str(material["style_execution"]["path"]),
    }
    for path in _authenticated_material_paths(project, material):
        logicals.add(path.relative_to(project).as_posix())
    logicals.update(
        _completed_semantic_dependencies(
            project, job, logical_files=logical_files, hit_root=hit_root,
        )
    )
    for field in ("generation_receipt", "qa_receipt", "reconstruction_work_item", "editable_receipt"):
        record = job.get(field)
        if not isinstance(record, Mapping) or not isinstance(record.get("path"), str):
            raise ValueError(f"completed page inventory has no registered {field}")
        path = _current_or_cached_authority(
            project, record["path"], record.get("sha256"), logical_files, hit_root,
        )
        if not isinstance(record.get("sha256"), str):
            raise ValueError(f"completed page inventory {field} changed")
        logicals.add(_canonical_project_logical(project, record["path"]))
    coverage_logical, _coverage_path = _verified_coverage_dependency(
        project, job, logical_files=logical_files, hit_root=hit_root,
    )
    logicals.add(coverage_logical)
    package = job.get("page_package")
    if not isinstance(package, str):
        package = f"07_editable/page_{page_number:03d}/page-package.json"
    package_logical = _canonical_project_logical(project, package)
    logicals.add(package_logical)

    receipt_record = job["editable_receipt"]
    receipt_path = _current_or_cached_authority(
        project, receipt_record["path"], receipt_record["sha256"], logical_files, hit_root,
    )
    receipt = _read_json_path(receipt_path, "completed page receipt")
    validate_v4_artifact("editable_receipt_v4.schema.json", receipt)
    for field in ("editable_page", "object_manifest", "signed_bundle"):
        record = receipt.get(field)
        if not isinstance(record, Mapping) or not isinstance(record.get("path"), str):
            raise ValueError(f"completed page receipt has no closed {field}")
        logical = _canonical_project_logical(project, record["path"])
        if field != "editable_page":
            _current_or_cached_authority(
                project, logical, record.get("sha256"), logical_files, hit_root,
            )
        logicals.add(logical)
    bundle_logical = _canonical_project_logical(project, receipt["signed_bundle"]["path"])
    bundle_path = _current_or_cached_authority(
        project, bundle_logical, receipt["signed_bundle"]["sha256"], logical_files, hit_root,
    )
    signed_bundle = _read_json_path(bundle_path, "signed reconstruction bundle")
    body_record = signed_bundle.get("body_pptx")
    if not isinstance(body_record, Mapping) or not isinstance(body_record.get("path"), str):
        raise ValueError("completed page signed bundle has no body PPTX")
    body_logical = _canonical_project_logical(project, body_record["path"])
    logicals.add(body_logical)
    schemas = Path(__file__).resolve().parents[1] / "schemas"
    logicals.update(f"@schema/{path.name}" for path in schemas.glob("*.schema.json"))
    return logicals


def _read_json_path(path: Path, label: str) -> dict[str, Any]:
    try:
        value = json.loads(path.read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError) as exc:
        raise ValueError(f"{label} is invalid JSON") from exc
    if not isinstance(value, dict):
        raise ValueError(f"{label} must be an object")
    return value


def completed_semantic_dependency_paths(
    project: Path,
    job: Mapping[str, Any],
) -> set[Path]:
    """Return only paths authorized by current registered and signed page artifacts."""
    return set(_completed_semantic_dependencies(project, job).values())


def _verified_coverage_dependency(
    project: Path,
    job: Mapping[str, Any],
    *,
    logical_files: Mapping[str, Any] | None = None,
    hit_root: Path | None = None,
) -> tuple[str, Path]:
    """Resolve and semantically authenticate this job's coverage contract."""
    page_number = job.get("page_number")
    if type(page_number) is not int or page_number < 1:
        raise ValueError("coverage contract has no valid page identity")
    logical = _canonical_project_logical(project, job.get("coverage_contract_file"))
    expected_sha = job.get("coverage_sha256")
    if not isinstance(expected_sha, str):
        raise ValueError("coverage contract has no registered semantic digest")
    current = _project_file(project, logical, must_exist=False)
    if current.is_file():
        path = current
    else:
        if logical_files is None or hit_root is None or logical not in logical_files:
            raise ValueError("coverage contract dependency is missing")
        path = _cached_record_path(hit_root, logical_files[logical])
    coverage = _read_json_path(path, "coverage contract")
    verified = verify_coverage_contract(coverage, expected_page_number=page_number)
    if verified.get("sha256") != expected_sha:
        raise ValueError("coverage contract semantic digest changed")
    return logical, path


def _completed_semantic_dependencies(
    project: Path,
    job: Mapping[str, Any],
    *,
    logical_files: Mapping[str, Any] | None = None,
    hit_root: Path | None = None,
) -> dict[str, Path]:
    """Resolve the closed semantic role graph from live or exact cached records."""
    project = Path(project).resolve()
    dependencies: dict[str, Path] = {}

    coverage_logical, coverage_path = _verified_coverage_dependency(
        project, job, logical_files=logical_files, hit_root=hit_root,
    )
    dependencies[coverage_logical] = coverage_path

    def resolve(record: Mapping[str, Any], label: str) -> Path:
        if not isinstance(record, Mapping) or not isinstance(record.get("path"), str):
            raise ValueError(f"completed semantic inventory has no {label}")
        logical = _canonical_project_logical(project, record["path"])
        if logical_files is None:
            path = _project_file(project, logical)
            if _sha(path) != record.get("sha256"):
                raise ValueError(f"completed semantic inventory {label} changed")
        else:
            if hit_root is None:
                raise ValueError("completed semantic cache root is missing")
            path = _current_or_cached_authority(
                project, logical, record.get("sha256"), logical_files, hit_root,
            )
        prior = dependencies.get(logical)
        if prior is not None and _sha(prior) != _sha(path):
            raise ValueError("completed semantic roles disagree on one logical artifact")
        dependencies[logical] = path
        return path

    def registered(field: str) -> tuple[Path, dict[str, Any]]:
        record = job.get(field)
        path = resolve(record, field)
        return path, _read_json_path(path, field)

    _generation_path, generation = registered("generation_receipt")
    validate_v4_artifact("page_generation_v4.schema.json", generation)
    for field in ("body_image", "provider_trace"):
        resolve(generation[field], f"generation {field}")
    for record in generation.get("reference_images", []):
        if isinstance(record, Mapping) and isinstance(record.get("path"), str):
            resolve(record, "generation reference image")

    _qa_path, qa = registered("qa_receipt")
    validate_v4_artifact("page_qa_v4.schema.json", qa)
    for field in ("qa_work_item", "observation"):
        resolve(qa[field], f"QA {field}")
    observation = _read_json_path(resolve(qa["observation"], "QA observation"), "QA observation")
    invocation = observation.get("invocation")
    if isinstance(invocation, Mapping):
        for field in ("request", "raw_response", "signed_bundle"):
            record = invocation.get(field)
            if not isinstance(record, Mapping):
                raise ValueError("QA invocation closure is incomplete")
            resolve(record, f"QA invocation {field}")

    work_path, _work = registered("reconstruction_work_item")
    _receipt_path, receipt = registered("editable_receipt")
    validate_v4_artifact("editable_receipt_v4.schema.json", receipt)
    signed_record = receipt["signed_bundle"]
    signed_path = resolve(signed_record, "signed reconstruction bundle")
    signed = _read_json_path(signed_path, "signed reconstruction bundle")
    closed = {
        "artifact_version", "work_item_sha256", "manifest", "gateway_invocation",
        "body_pptx", "editable_page", "object_counts", "text_coverage",
        "table_coverage", "raster_components", "accepted_body_embedded",
        "fixed_layers", "slide_fingerprint", "nonce", "attestation",
    }
    if set(signed) != closed or signed.get("work_item_sha256") != _sha(work_path):
        raise ValueError("signed reconstruction bundle is not closed")
    signature = signed.pop("attestation")
    expected = hmac.new(
        _secret(project), canonical_sha256(signed).encode("ascii"), hashlib.sha256,
    ).hexdigest()
    signed["attestation"] = signature
    if not hmac.compare_digest(str(signature), expected):
        raise ValueError("reconstruction attestation is invalid")
    for field in ("manifest", "gateway_invocation", "body_pptx"):
        resolve(signed[field], f"signed reconstruction {field}")
    resolve(signed["editable_page"], "signed reconstruction editable page")
    object_record = receipt["object_manifest"]
    resolve(object_record, "reconstruction object manifest")
    gateway = _read_json_path(
        resolve(signed["gateway_invocation"], "signed reconstruction invocation"),
        "signed reconstruction invocation",
    )
    for field in ("request", "raw_response", "manifest"):
        record = gateway.get(field)
        if not isinstance(record, Mapping):
            raise ValueError("signed reconstruction invocation is incomplete")
        resolve(record, f"reconstruction invocation {field}")
    manifest_path = resolve(gateway["manifest"], "object manifest")
    manifest = _read_json_path(manifest_path, "object manifest")
    _manifest_exact(manifest)
    manifest_logical = _canonical_project_logical(project, gateway["manifest"]["path"])
    manifest_root_logical = PurePosixPath(manifest_logical).parent
    for image in manifest["images"]:
        value = PurePosixPath(image["path"])
        logical = _canonical_project_logical(
            project,
            value.as_posix() if value.is_absolute()
            else (manifest_root_logical / value).as_posix(),
        )
        current = _project_file(project, logical, must_exist=False)
        if current.is_file():
            dependencies[logical] = current
        elif logical_files is not None and logical in logical_files:
            dependencies[logical] = _cached_record_path(hit_root, logical_files[logical])
        else:
            raise ValueError("reconstruction manifest image dependency is missing")
    return dependencies


def _canonical_project_logical(project: Path, value: str | Path) -> str:
    if not isinstance(value, str) or not value or "\\" in value:
        raise ValueError("completed page authority path is non-canonical")
    pure = PurePosixPath(value)
    if pure.is_absolute() or pure.as_posix() != value or ".." in pure.parts:
        raise ValueError("completed page authority path is non-canonical")
    path = _project_file(project, Path(*pure.parts), must_exist=False)
    return path.relative_to(project).as_posix()


def _cached_record_path(hit_root: Path, record: Mapping[str, Any]) -> Path:
    if not isinstance(record, Mapping) or set(record) != {"path", "sha256"}:
        raise ValueError("editable-page cache logical record is invalid")
    record_path = record.get("path")
    if not isinstance(record_path, str) or "\\" in record_path:
        raise ValueError("editable-page cache payload path is invalid")
    pure = PurePosixPath(record_path)
    if pure.is_absolute() or pure.as_posix() != record_path or ".." in pure.parts:
        raise ValueError("editable-page cache payload path is non-canonical")
    path = (hit_root / Path(*pure.parts)).resolve(strict=True)
    if hit_root not in path.parents or not path.is_file():
        raise ValueError("editable-page cache payload path escapes the cache root")
    if _sha(path) != record.get("sha256"):
        raise ValueError("editable-page cache dependency was replaced")
    return path


def _current_or_cached_authority(
    project: Path,
    logical: str,
    expected_sha256: Any,
    logical_files: Mapping[str, Any],
    hit_root: Path,
) -> Path:
    canonical = _canonical_project_logical(project, logical)
    if not isinstance(expected_sha256, str):
        raise ValueError("completed page authority hash is invalid")
    current = _project_file(project, canonical, must_exist=False)
    if current.is_file():
        if _sha(current) != expected_sha256:
            raise ValueError(f"present project authority differs from completed cache: {canonical}")
        return current
    record = logical_files.get(canonical)
    cached = _cached_record_path(hit_root, record)
    if _sha(cached) != expected_sha256:
        raise ValueError(f"cached completed page authority differs: {canonical}")
    return cached


def _verify_completed_page_semantics(
    project: Path,
    job: Mapping[str, Any],
    material: Mapping[str, Any],
    receipt: Mapping[str, Any],
) -> dict[str, Any]:
    """Replay generation, QA, reconstruction, and page-package semantics."""
    from page_generation import validate_generation_receipt
    from page_pipeline import generation_request, load_contract, load_style
    from v4_qa import validate_qa_receipt
    from editppt.runtime.editable_page_cache import create_page_package

    run = _read_json_path(project / "workflow_run.json", "workflow state")
    attempt = job.get("attempt")
    if type(attempt) is not int or attempt < 1:
        raise ValueError("completed page generation attempt is invalid")
    request = generation_request(project, run, job, attempt).payload
    generation_record = job["generation_receipt"]
    generation_path = _project_file(project, generation_record["path"])
    generation = _read_json_path(generation_path, "generation receipt")
    generation_result = validate_generation_receipt(
        project,
        material,
        request,
        _project_file(project, generation["body_image"]["path"]),
        generation_path,
        cached_output=True,
    )
    style = load_style(project, run)
    qa_record = job["qa_receipt"]
    validate_qa_receipt(
        project,
        _project_file(project, qa_record["path"]),
        material_bundle=material,
        generation_receipt=generation_path,
        generation_receipt_sha256=generation_result["sha256"],
        style_execution=style["execution"],
        material_bundle_path=_project_file(project, job["material_bundle_file"]),
        page_contract=load_contract(project, job),
        logo_source=run["logo_source"],
    )
    expected_receipt = verify_signed_reconstruction(
        project,
        _project_file(project, job["reconstruction_work_item"]["path"]),
        _project_file(project, receipt["signed_bundle"]["path"]),
    )
    if receipt != expected_receipt:
        raise ValueError("completed page cache receipt is forged or stale")
    package_path = _project_file(project, job["page_package"])
    expected_package = package_path.with_name(".page-package.revalidated.tmp.json")
    try:
        create_page_package(
            project,
            page_number=int(job["page_number"]),
            cache_key=str(job["cache"]["key"]),
            pptx=_project_file(project, receipt["editable_page"]["path"]),
            output=expected_package,
        )
        if package_path.read_bytes() != expected_package.read_bytes():
            raise ValueError("completed page package differs from reconstructed semantics")
    finally:
        expected_package.unlink(missing_ok=True)
    return expected_receipt


def restore_and_validate_completed_cache(
    project: Path,
    job: Mapping[str, Any],
    hit,
    *,
    authority_identity: str,
) -> None:
    """Restore only missing copies from a v2 cache, then replay all semantics."""
    project = Path(project).resolve()
    if hit.manifest.get("artifact_version") != PAGE_CACHE_CONTRACT_VERSION:
        raise ValueError("legacy editable-page cache is not accepted by V4")
    logical_files = hit.manifest.get("logical_files")
    if not isinstance(logical_files, Mapping) or not logical_files:
        raise ValueError("editable-page cache dependency closure is missing")
    schema_root = Path(__file__).resolve().parents[1] / "schemas"
    material = _load_authenticated_page_material_bundle(
        project,
        expected_path=job.get("material_bundle_file"),
        expected_file_sha256=job.get("material_bundle_file_sha256"),
        expected_sealed_sha256=job.get("material_bundle_sha256"),
        expected_page_number=job.get("page_number"),
        expected_authority_identity=authority_identity,
    )
    hit_root = Path(hit.path).resolve(strict=True)
    if not hit_root.is_dir() or project not in hit_root.parents:
        raise ValueError("editable-page cache root is invalid")
    expected_logicals = _expected_completed_logicals(
        project, job, material, logical_files, hit_root,
    )
    if set(logical_files) != expected_logicals:
        raise ValueError("cache manifest does not match the exact completed-page inventory")
    allowed_source_logicals = {
        "00_source/company_logo.svg",
        *(
            path.relative_to(project).as_posix()
            for path in _authenticated_material_paths(project, material)
            if path.relative_to(project).as_posix().startswith("00_source/")
        ),
    }
    allowed_source_logicals_folded = {value.casefold() for value in allowed_source_logicals}
    forbidden_globals_folded = {
        "workflow_run.json", "01_page_contracts/source_lock.json",
        "confirm_ui/page_requirement_summary.json", "00_source/source.docx",
        "00_source/pages.json",
    }
    seen_logicals: set[str] = set()
    missing_copies: list[tuple[Path, Path]] = []
    for logical, record in logical_files.items():
        if not isinstance(logical, str):
            raise ValueError("editable-page cache logical path is invalid")
        logical_folded = logical.replace("\\", "/").removeprefix("./").casefold()
        if (
            logical_folded in forbidden_globals_folded
            or logical_folded.startswith("00_source/")
            and logical_folded not in allowed_source_logicals_folded
        ):
            raise ValueError(f"completed cache contains non-page-local global source: {logical}")
        normalized_logical = PurePosixPath(logical)
        if (
            "\\" in logical
            or normalized_logical.is_absolute()
            or normalized_logical.as_posix() != logical
            or ".." in normalized_logical.parts
            or logical.casefold() in seen_logicals
        ):
            raise ValueError("editable-page cache logical path is non-canonical")
        seen_logicals.add(logical.casefold())
        if not isinstance(record, Mapping) or set(record) != {"path", "sha256"}:
            raise ValueError("editable-page cache logical record is invalid")
        record_path = record.get("path")
        if not isinstance(record_path, str):
            raise ValueError("editable-page cache payload path is invalid")
        normalized_record_path = PurePosixPath(record_path)
        if (
            "\\" in record_path
            or normalized_record_path.is_absolute()
            or normalized_record_path.as_posix() != record_path
            or ".." in normalized_record_path.parts
        ):
            raise ValueError("editable-page cache payload path is non-canonical")
        cached_path = (hit_root / Path(*normalized_record_path.parts)).resolve(strict=True)
        if hit_root not in cached_path.parents or not cached_path.is_file():
            raise ValueError("editable-page cache payload path escapes the cache root")
        if _sha(cached_path) != record["sha256"]:
            raise ValueError("editable-page cache dependency was replaced")
        if logical.startswith("@schema/"):
            current = schema_root / logical.removeprefix("@schema/")
            if not current.is_file() or _sha(current) != record["sha256"]:
                raise ValueError("current V4 schema differs from the cached validation authority")
            continue
        destination = _project_file(project, logical, must_exist=False)
        if destination.exists():
            if not destination.is_file() or _sha(destination) != record["sha256"]:
                raise ValueError(f"present project authority differs from completed cache: {logical}")
        else:
            missing_copies.append((cached_path, destination))
    root = hit.path / "reconstruction"
    cached = {
        name: _project_file(project, root / name)
        for name in (
            "page.pptx", "editable-receipt.json", "object-manifest.json",
            "reconstruction-work-item.json", "signed-reconstruction.json", "body.pptx",
        )
    }
    receipt = json.loads(cached["editable-receipt.json"].read_text(encoding="utf-8"))
    bundle = json.loads(cached["signed-reconstruction.json"].read_text(encoding="utf-8"))
    work_record = job.get("reconstruction_work_item")
    receipt_record = job.get("editable_receipt")
    if not isinstance(work_record, Mapping) or not isinstance(receipt_record, Mapping):
        raise ValueError("completed page cache has no registered reconstruction authorities")
    destinations = {
        "page.pptx": receipt["editable_page"]["path"],
        "editable-receipt.json": receipt_record["path"],
        "object-manifest.json": receipt["object_manifest"]["path"],
        "reconstruction-work-item.json": work_record["path"],
        "signed-reconstruction.json": receipt["signed_bundle"]["path"],
        "body.pptx": bundle["body_pptx"]["path"],
    }
    reconstruction_copies: list[tuple[Path, Path]] = []
    for name, relative in destinations.items():
        destination = _project_file(project, relative, must_exist=False)
        if destination.exists():
            if not destination.is_file() or _sha(destination) != _sha(cached[name]):
                raise ValueError(f"present reconstruction output differs from completed cache: {name}")
        else:
            reconstruction_copies.append((cached[name], destination))
    copies_by_destination: dict[Path, Path] = {}
    for source, destination in [*missing_copies, *reconstruction_copies]:
        prior = copies_by_destination.get(destination)
        if prior is not None and _sha(prior) != _sha(source):
            raise ValueError("completed cache has conflicting payloads for one logical target")
        copies_by_destination[destination] = source
    published: list[tuple[Path, _OwnedIdentity]] = []
    created_directories: list[tuple[Path, _OwnedIdentity]] = []
    try:
        for destination, source in copies_by_destination.items():
            missing_parents: list[Path] = []
            parent = destination.parent
            while parent != project and not parent.exists():
                missing_parents.append(parent)
                parent = parent.parent
            for directory in reversed(missing_parents):
                directory.mkdir()
                identity = _owned_path_identity(directory)
                if identity is not None:
                    created_directories.append((directory, identity))
            if destination.exists():
                raise ValueError(f"restore destination appeared during validation: {destination}")
            with source.open("rb") as input_handle, destination.open("xb") as output_handle:
                identity = _owned_handle_identity(destination, output_handle)
                if identity is not None:
                    published.append((destination, identity))
                shutil.copyfileobj(input_handle, output_handle, length=1024 * 1024)
                output_handle.flush()
                os.fsync(output_handle.fileno())
        expected = _verify_completed_page_semantics(project, job, material, receipt)
        if (
            receipt != expected
            or _sha(_project_file(project, receipt_record["path"]))
            != receipt_record["sha256"]
        ):
            raise ValueError("completed page cache receipt is forged or stale")
    except BaseException:
        for destination, identity in reversed(published):
            if _still_owned(destination, identity):
                destination.unlink(missing_ok=True)
        for directory, identity in reversed(created_directories):
            try:
                if _still_owned(directory, identity):
                    directory.rmdir()
            except OSError:
                pass
        raise
