"""Assemble V5 final pages and run mandatory package/authentic-pixel checks."""

from __future__ import annotations

import hashlib
import json
import os
import posixpath
import uuid
import zipfile
from pathlib import Path
from typing import Any
from xml.etree import ElementTree

from pptx import Presentation

from workflow_v5_authentic_material import verify_authentic_placements
from workflow_v5_identity import ContentCatalog


EDITPPT_RUNTIME = Path(__file__).resolve().parents[2] / "reconstruct-editable-slide" / "cli" / "editppt" / "runtime"
if EDITPPT_RUNTIME.is_dir():
    if str(EDITPPT_RUNTIME) not in __import__("sys").path:
        __import__("sys").path.insert(0, str(EDITPPT_RUNTIME))
    from final_assembler import _copy_page_slide  # noqa: E402
else:
    # Installed workflow packages do not retain the repository sibling layout;
    # the editable runtime is an installed package in that environment.
    from editppt.runtime.final_assembler import _copy_page_slide  # noqa: E402


_REL_NS = "http://schemas.openxmlformats.org/package/2006/relationships"


def _sha(path: Path) -> str:
    return hashlib.sha256(path.read_bytes()).hexdigest()


def _placement_for_hash(pptx: Path, *, page: int, asset_id: str, artifact_id: str) -> dict[str, Any]:
    expected = artifact_id.removeprefix("sha256:")
    with zipfile.ZipFile(pptx) as archive:
        slide = f"ppt/slides/slide{page}.xml"
        rels = f"ppt/slides/_rels/slide{page}.xml.rels"
        root = ElementTree.fromstring(archive.read(rels))
        candidates = []
        for relation in root.findall(f"{{{_REL_NS}}}Relationship"):
            target = relation.get("Target")
            if relation.get("TargetMode") != "External" and isinstance(target, str):
                member = posixpath.normpath(posixpath.join(posixpath.dirname(slide), target))
                if member.startswith("ppt/media/"):
                    candidates.append(member)
        media = next((name for name in candidates if hashlib.sha256(archive.read(name)).hexdigest() == expected), None)
    if media is None:
        raise ValueError("authentic source bytes are absent from final deck")
    return {
        "page_number": page, "asset_id": asset_id,
        "source_artifact_id": artifact_id, "media_member": media,
        "custody": "immutable_image_object",
    }


def assemble_v5_deck(project: Path, *, page_numbers: list[int]) -> dict[str, Any]:
    root = Path(project).resolve()
    if page_numbers != sorted(set(page_numbers)) or not page_numbers:
        raise ValueError("V5 final page order must be unique and sorted")
    pages = [root / "04_v5" / "final-pages" / f"page_{page:03d}" / "page.pptx" for page in page_numbers]
    if any(not path.is_file() for path in pages):
        raise ValueError("a finalized V5 page package is missing")
    dag = json.loads((root / "04_v5" / "dag.json").read_text(encoding="utf-8"))
    by_id = {node["node_id"]: node for node in dag["nodes"]}
    if any(by_id[f"page:{page:03d}:visual_qa"]["status"] != "complete" for page in page_numbers):
        raise ValueError("every final page must pass the V5 visual QA gate before assembly")
    qa = json.loads((root / "09_reports" / "v5_final_qa.json").read_text(encoding="utf-8"))
    if qa.get("blocking_pages"):
        raise ValueError("V5 final QA still has blocking pages")
    output_dir = root / "08_final"
    output_dir.mkdir(parents=True, exist_ok=True)
    output = output_dir / "deck.pptx"
    temporary = output_dir / f".deck-v5-{uuid.uuid4().hex}.tmp"
    deck = Presentation(pages[0])
    layout = deck.slides[0].slide_layout
    for page, path in zip(page_numbers[1:], pages[1:]):
        _copy_page_slide(path, deck, layout, page)
    deck.save(temporary)
    reopened = Presentation(temporary)
    if len(reopened.slides) != len(page_numbers):
        raise ValueError("assembled V5 slide count is incorrect")
    if any(
        len([shape for shape in slide.shapes if shape.name.startswith("fixed-frame-")]) != 4
        for slide in reopened.slides
    ):
        raise ValueError("assembled V5 fixed-layer inventory is incorrect")
    os.replace(temporary, output)

    placements = []
    for compose in sorted((root / "04_v5" / "compose").glob("page_*.json")):
        payload = json.loads(compose.read_text(encoding="utf-8"))
        page = int(payload["page_number"])
        for item in payload.get("authentic_placements", []):
            placements.append(_placement_for_hash(
                output, page=page, asset_id=item["asset_id"],
                artifact_id=item["source_artifact_id"],
            ))
    authentic = verify_authentic_placements(output, placements)
    record = ContentCatalog(root).record_file("final-deck", output, boundary="before_delivery")
    report = {
        "artifact_version": "v5-final-assembly-v1",
        "workflow_contract_version": "word-ppt-workflow-v5",
        "status": "assembled_pending_office_validation",
        "page_count": len(page_numbers),
        "page_order": page_numbers,
        "output": output.relative_to(root).as_posix(),
        "output_sha256": _sha(output),
        "artifact_id": record["artifact_id"],
        "authentic_pixels": authentic,
        "fixed_layers_per_slide": 4,
    }
    (output_dir / "v5_assembly_report.json").write_text(
        json.dumps(report, ensure_ascii=False, indent=2) + "\n", encoding="utf-8",
    )
    return report
