"""Non-destructive, idempotent V4 to V5 project migration."""

from __future__ import annotations

import json
import hashlib
import os
import re
import uuid
import zipfile
from pathlib import Path
from typing import Any, Mapping

from PIL import Image
from pptx import Presentation

from workflow_v5_dag import DagStore, build_project_dag, validate_dag
from workflow_v5_identity import ContentCatalog
from workflow_v5_intent import compile_project_intents
from workflow_v5_ui import ConfirmationLifecycle
from v4_qa import validate_historical_qa_receipt


_V4 = "word-ppt-workflow-v4"
_MIGRATION = "v4-to-v5-migration-v1"
_MATERIAL_ID = re.compile(r"^[A-Za-z0-9._-]+$")


def _read_object(path: Path, label: str) -> dict[str, Any]:
    try:
        value = json.loads(path.read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError) as exc:
        raise ValueError(f"{label} is unavailable") from exc
    if not isinstance(value, dict):
        raise ValueError(f"{label} must be an object")
    return value


def _project_file(project: Path, value: Any, label: str) -> Path:
    if not isinstance(value, str) or not value:
        raise ValueError(f"{label} path is missing")
    path = Path(value)
    if path.is_absolute():
        raise ValueError(f"{label} path must be project-relative")
    actual = (project / path).resolve()
    try:
        actual.relative_to(project)
    except ValueError as exc:
        raise ValueError(f"{label} escapes the project") from exc
    if not actual.is_file():
        raise ValueError(f"{label} is missing")
    return actual


def _material_ids(job: Mapping[str, Any]) -> list[str]:
    values = job.get("material_ids", [])
    if not isinstance(values, list) or any(
        not isinstance(value, str) or not _MATERIAL_ID.fullmatch(value) for value in values
    ):
        raise ValueError("legacy page material_ids are invalid")
    return list(dict.fromkeys(values))


def _artifact_path(project: Path, job: Mapping[str, Any], kind: str) -> Path | None:
    if kind == "design":
        generation = job.get("generation")
        body = generation.get("body_image") if isinstance(generation, Mapping) else None
        value = body.get("path") if isinstance(body, Mapping) else None
        if not isinstance(value, str) and isinstance(generation, Mapping):
            value = generation.get("image")
        if not isinstance(value, str):
            value = job.get("expected_output")
    else:
        editable = job.get("editable_page")
        value = editable.get("path") if isinstance(editable, Mapping) else None
    if not isinstance(value, str):
        return None
    try:
        artifact = _project_file(project, value, f"legacy {kind} artifact")
    except ValueError:
        return None
    if kind == "design":
        try:
            with Image.open(artifact) as image:
                if image.format != "PNG" or image.size != (1904, 896):
                    return None
                image.verify()
        except (OSError, ValueError):
            return None
    if kind == "reconstruct":
        try:
            if not zipfile.is_zipfile(artifact):
                return None
            deck = Presentation(artifact)
            if len(deck.slides) != 1 or not any(
                shape.has_text_frame or shape.has_table for shape in deck.slides[0].shapes
            ):
                return None
        except Exception:
            return None
    return artifact


def _historical_passing_design(
    project: Path, job: Mapping[str, Any], design_path: Path | None,
) -> tuple[Path | None, str | None]:
    """Accept legacy Image2 only through its immutable passing QA closure."""
    if design_path is None:
        return None, "legacy_design_missing_or_failed_deterministic_preflight"
    if job.get("status") in {"technical_blocked", "repair", "content_blocked"}:
        return None, "historical_passing_qa_receipt_missing_or_invalid"
    qa_record = job.get("qa_receipt")
    if not isinstance(qa_record, Mapping):
        return None, "historical_passing_qa_receipt_missing_or_invalid"
    try:
        receipt_path = _project_file(
            project, qa_record.get("path"), "historical QA receipt",
        )
        validated = validate_historical_qa_receipt(project, receipt_path)
        artifact = validated.get("artifact")
        if not isinstance(artifact, Mapping):
            raise ValueError("historical QA receipt artifact is missing")
        if artifact.get("status") != "pass" or artifact.get("page_number") != job.get("page_number"):
            raise ValueError("historical QA receipt is not a passing decision for this page")
        recorded_sha = qa_record.get("sha256")
        if not isinstance(recorded_sha, str) or recorded_sha != validated.get("sha256"):
            raise ValueError("historical QA receipt state identity mismatch")
        work_record = artifact.get("qa_work_item")
        if not isinstance(work_record, Mapping):
            raise ValueError("historical QA work item reference is missing")
        work_path = _project_file(
            project, work_record.get("path"), "historical QA work item",
        )
        if work_record.get("sha256") != _sha256_file(work_path):
            raise ValueError("historical QA work item identity mismatch")
        work = _read_object(work_path, "historical QA work item")
        body_image = work.get("body_image")
        if not isinstance(body_image, Mapping):
            raise ValueError("historical QA body image reference is missing")
        closed_image = _project_file(
            project, body_image.get("path"), "historical QA body image",
        )
        if closed_image != design_path or body_image.get("sha256") != _sha256_file(design_path):
            raise ValueError("legacy design is not the Image2 closed by historical QA")
    except (KeyError, OSError, TypeError, ValueError):
        return None, "historical_passing_qa_receipt_missing_or_invalid"
    return design_path, None


def _sha256_file(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def _complete(node: dict[str, Any], result_key: str) -> None:
    node.update({
        "status": "complete", "attempts": max(1, node["attempts"]),
        "worker_id": None, "result_key": result_key, "failure": None,
    })


def _write_atomic(path: Path, value: Mapping[str, Any]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    data = json.dumps(
        value, ensure_ascii=False, sort_keys=True, separators=(",", ":"),
    ).encode("utf-8")
    if path.is_file() and path.read_bytes() == data:
        return
    temporary = path.with_name(f".{path.name}.{uuid.uuid4().hex}.tmp")
    temporary.write_bytes(data)
    os.replace(temporary, path)


def _publish_material_manifests(
    root: Path,
    *,
    compiled: Mapping[str, Any],
    catalog: ContentCatalog,
) -> dict[str, dict[str, Any]]:
    """Publish one immutable manifest per acquisition without network I/O."""
    grouped: dict[str, list[dict[str, Any]]] = {}
    for raw in compiled.get("available_assets", []):
        if not isinstance(raw, Mapping):
            continue
        for material_id in raw.get("material_ids", []):
            if isinstance(material_id, str):
                grouped.setdefault(material_id, []).append(dict(raw))

    requirements = {
        str(item["material_id"]): dict(item)
        for item in compiled.get("requirements", [])
        if isinstance(item, Mapping) and isinstance(item.get("material_id"), str)
    }
    published: dict[str, dict[str, Any]] = {}
    for material_id, raw_assets in sorted(grouped.items()):
        assets: list[dict[str, Any]] = []
        seen: set[tuple[str, str]] = set()
        for raw in sorted(
            raw_assets,
            key=lambda item: (str(item.get("asset_id", "")), str(item.get("artifact_id", ""))),
        ):
            relative = raw.get("relative_path")
            expected = raw.get("artifact_id")
            if not isinstance(relative, str) or not isinstance(expected, str):
                raise ValueError("compiled authentic asset identity is incomplete")
            source = _project_file(root, relative, "compiled authentic material")
            identity = (str(raw.get("asset_id", "")), expected)
            if identity in seen:
                continue
            seen.add(identity)
            record = catalog.record_file(
                f"material-{material_id}-asset-{expected.removeprefix('sha256:')[:16]}",
                source,
                boundary="ingestion",
            )
            if record["artifact_id"] != expected:
                raise ValueError("compiled authentic material hash does not match local evidence")
            assets.append({
                "asset_id": str(raw.get("asset_id") or expected),
                "evidence_id": str(raw.get("evidence_id") or ""),
                "relative_path": source.relative_to(root).as_posix(),
                "artifact_id": expected,
                "media_type": str(raw.get("media_type") or ""),
                "source_kind": str(raw.get("source_kind") or "search"),
                "entity": str(raw.get("entity") or ""),
                "material_role": str(raw.get("material_role") or "authentic_image"),
                "source": dict(raw.get("source") or {}),
            })
        if not assets:
            continue
        requirement = requirements.get(material_id, {})
        manifest = {
            "artifact_version": "v5-authentic-material-manifest-v2",
            "material_id": material_id,
            "page_numbers": list(requirement.get("page_numbers") or []),
            "required": bool(requirement.get("required", True)),
            "required_asset_count": len(assets),
            "discovery_reused": True,
            "new_search_performed": False,
            "assets": assets,
        }
        output = root / "04_v5" / "materials" / f"{material_id}.json"
        _write_atomic(output, manifest)
        manifest_record = catalog.record_file(
            f"material-{material_id}-manifest", output, boundary="ingestion",
        )
        published[material_id] = {
            **manifest,
            "path": output.relative_to(root).as_posix(),
            "artifact_id": manifest_record["artifact_id"],
        }
    return published


def migrate_v4_project(project: Path) -> dict[str, Any]:
    """Add V5 state in-place; never rename the project or alter V4 authorities."""
    root = Path(project).resolve()
    if not root.is_dir():
        raise ValueError("migration project directory is missing")
    state = _read_object(root / "workflow_run.json", "V4 workflow state")
    if state.get("workflow_contract_version") != _V4:
        raise ValueError("migration accepts only word-ppt-workflow-v4 projects")
    jobs = state.get("jobs")
    pagination = state.get("pagination")
    if not isinstance(jobs, list) or not jobs or not isinstance(pagination, Mapping):
        raise ValueError("V4 locked pages are invalid")
    pages = [job.get("page_number") for job in jobs if isinstance(job, Mapping)]
    if pages != pagination.get("locked_page_order"):
        raise ValueError("V4 page order is not locked consistently")

    catalog = ContentCatalog(root)
    source = state.get("word_source")
    logo = state.get("logo_source")
    gate = state.get("style_confirmation")
    if not isinstance(source, Mapping) or not isinstance(logo, Mapping) or not isinstance(gate, Mapping):
        raise ValueError("V4 project authorities are incomplete")
    source_record = catalog.record_file(
        "word-source", _project_file(root, source.get("path"), "Word source"), boundary="ingestion",
    )
    catalog.record_file(
        "fixed-logo", _project_file(root, logo.get("path"), "SVG logo"), boundary="ingestion",
    )
    style_record = None
    if gate.get("status") == "confirmed":
        style_record = catalog.record_file(
            "confirmed-style",
            _project_file(root, gate.get("execution_file"), "confirmed style"),
            boundary="ingestion",
        )
        ConfirmationLifecycle(root).confirm(style_record["artifact_id"])

    compiled = compile_project_intents(root, jobs)
    compiled_by_page = {item["page_number"]: item for item in compiled["pages"]}
    plans: list[dict[str, Any]] = []
    page_records: dict[int, dict[str, Any]] = {}
    for job in jobs:
        if not isinstance(job, Mapping) or type(job.get("page_number")) is not int:
            raise ValueError("V4 page job is invalid")
        page = int(job["page_number"])
        contract = catalog.record_file(
            f"page-{page:03d}-authority",
            _project_file(root, job.get("contract_file"), "page contract"),
            boundary="ingestion",
        )
        materials = list(dict.fromkeys([
            *_material_ids(job), *compiled_by_page[page]["material_ids"],
        ]))
        plans.append({
            "page_number": page,
            "authority_key": contract["artifact_id"],
            "material_ids": materials,
        })
        page_records[page] = {
            "job": job, "contract": contract, "material_ids": materials,
        }

    desired = build_project_dag(plans)
    by_id = {node["node_id"]: node for node in desired["nodes"]}
    _complete(by_id["project:source"], source_record["artifact_id"])
    if style_record is not None:
        _complete(by_id["project:style"], style_record["artifact_id"])
    material_manifests = _publish_material_manifests(
        root, compiled=compiled, catalog=catalog,
    )
    for material_id, manifest in material_manifests.items():
        _complete(by_id[f"material:{material_id}"], manifest["artifact_id"])
    for page, records in page_records.items():
        prefix = f"page:{page:03d}"
        _complete(by_id[f"{prefix}:intent"], records["contract"]["artifact_id"])
        design_path, design_invalidation_reason = _historical_passing_design(
            root, records["job"], _artifact_path(root, records["job"], "design"),
        )
        records["design_invalidation_reason"] = design_invalidation_reason
        design_record = None
        if design_path is not None:
            design_record = catalog.record_file(
                f"page-{page:03d}-design", design_path, boundary="after_external_output",
            )
            # Historical QA reviewed raw Image2 pixels, not the new exact-source
            # composed authority. Material pages must reuse the pixels through
            # generate_v5_design and rerun composed-body acceptance.
            if not records["material_ids"]:
                _complete(by_id[f"{prefix}:design"], design_record["artifact_id"])
    validate_dag(desired)

    dag_path = root / "04_v5" / "dag.json"
    if dag_path.is_file():
        dag = DagStore(root).reconcile_migration(desired)
    else:
        dag = desired
        DagStore(root).initialize(dag)

    by_id = {node["node_id"]: node for node in dag["nodes"]}
    page_report: list[dict[str, Any]] = []
    reused = 0
    for page, records in sorted(page_records.items()):
        prefix = f"page:{page:03d}"
        design_reused = by_id[f"{prefix}:design"]["status"] == "complete"
        reconstruct_reused = by_id[f"{prefix}:reconstruct"]["status"] == "complete"
        reused += int(design_reused) + int(reconstruct_reused)
        page_report.append({
            "page_number": page,
            "design_reused": design_reused,
            "reconstruction_reused": reconstruct_reused,
            "final_qa_required": by_id[f"{prefix}:visual_qa"]["status"] != "complete",
            "invalidation_reason": (
                records.get("design_invalidation_reason")
                if not design_reused else None
            ) or (
                "authentic_pixel_custody_requires_v5_compose"
                if records["material_ids"] and not reconstruct_reused else None
            ) or (
                "legacy_design_missing_or_failed_deterministic_preflight"
                if not design_reused else None
            ),
        })
    report = {
        "artifact_version": _MIGRATION,
        "source_contract": _V4,
        "target_contract": "word-ppt-workflow-v5",
        "project_directory_renamed": False,
        "successful_model_results_reused": reused,
        "compiled_material_requirements": compiled["requirements"],
        "compiled_available_assets": compiled["available_assets"],
        "published_material_manifests": [
            {
                "material_id": material_id,
                "artifact_id": manifest["artifact_id"],
                "asset_count": len(manifest["assets"]),
                "new_search_performed": False,
            }
            for material_id, manifest in sorted(material_manifests.items())
        ],
        "pages": page_report,
    }
    _write_atomic(root / "09_reports" / "v5_compatibility_report.json", report)
    return report
