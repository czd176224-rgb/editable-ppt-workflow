"""Mandatory final package, PowerPoint-open, and render validation for V5 delivery."""

from __future__ import annotations

import hashlib
import json
import shutil
import subprocess
from pathlib import Path
from typing import Any

from pptx import Presentation

from render_pptx import render_libreoffice, render_powerpoint
from workflow_v5_identity import ContentCatalog


def _sha(path: Path) -> str:
    return hashlib.sha256(path.read_bytes()).hexdigest()


def validate_v5_office(project: Path, *, page_numbers: list[int]) -> dict[str, Any]:
    root = Path(project).resolve()
    assembly_path = root / "08_final" / "v5_assembly_report.json"
    assembly = json.loads(assembly_path.read_text(encoding="utf-8"))
    if assembly.get("status") != "assembled_pending_office_validation":
        raise ValueError("V5 assembly is not awaiting Office validation")
    output = root / assembly["output"]
    before = _sha(output)
    if before != assembly.get("output_sha256"):
        raise ValueError("assembled V5 deck changed before Office validation")

    opened = Presentation(output)
    if len(opened.slides) != len(page_numbers):
        raise ValueError("PowerPoint package slide count is incorrect")
    editable_counts = [
        sum(1 for shape in slide.shapes if shape.has_text_frame or shape.has_table)
        for slide in opened.slides
    ]
    if any(count < 1 for count in editable_counts):
        raise ValueError("a final slide has no editable text or table objects")

    render_dir = root / "09_reports" / "v5_office_render" / before[:16]
    render_dir.mkdir(parents=True, exist_ok=True)
    try:
        rendered = render_powerpoint(output, render_dir)
        backend = "Microsoft PowerPoint COM"
    except Exception:
        rendered = render_libreoffice(output, render_dir)
        backend = "LibreOffice"
    proof_files = sorted(render_dir.glob("slide_*.png"))
    if rendered != len(page_numbers) or len(proof_files) != len(page_numbers):
        raise ValueError("Office renderer did not produce one proof image per slide")
    if _sha(output) != before:
        raise ValueError("Office validation changed the delivered deck")

    officecli = shutil.which("officecli")
    officecli_report: dict[str, Any] = {"available": bool(officecli), "passed": None}
    if officecli:
        checked = subprocess.run(
            [officecli, "validate", str(output), "--json"],
            capture_output=True, text=True, errors="replace", timeout=180, check=False,
        )
        officecli_report.update({
            "passed": checked.returncode == 0,
            "returncode": checked.returncode,
        })
        if checked.returncode != 0:
            try:
                details = json.loads(checked.stdout or "{}")
                warnings = details.get("warnings", []) if isinstance(details, dict) else []
                officecli_report["warning_count"] = len(warnings)
                officecli_report["first_warning"] = (
                    warnings[0].get("message") if warnings and isinstance(warnings[0], dict) else None
                )
            except json.JSONDecodeError:
                officecli_report["first_warning"] = (checked.stderr or checked.stdout)[:500]
            raise ValueError("officecli OpenXML validation failed")

    artifact = ContentCatalog(root).record_file(
        "final-deck-office-validated", output, boundary="before_delivery",
    )
    report = {
        "artifact_version": "v5-office-validation-v1",
        "workflow_contract_version": "word-ppt-workflow-v5",
        "passed": True,
        "output": output.relative_to(root).as_posix(),
        "output_sha256": before,
        "artifact_id": artifact["artifact_id"],
        "page_count": len(page_numbers),
        "page_order": page_numbers,
        "editable_object_counts": editable_counts,
        "open_backend": "python-pptx OpenXML",
        "render_backend": backend,
        "rendered_page_count": rendered,
        "render_proofs": [path.relative_to(root).as_posix() for path in proof_files],
        "officecli": officecli_report,
    }
    report_path = root / "09_reports" / "v5_office_validation.json"
    report_path.write_text(json.dumps(report, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    assembly["status"] = "complete"
    assembly["office_validation"] = report_path.relative_to(root).as_posix()
    assembly_path.write_text(json.dumps(assembly, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    return report
