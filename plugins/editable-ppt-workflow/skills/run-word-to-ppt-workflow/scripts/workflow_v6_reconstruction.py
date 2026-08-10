"""V6 reconstruction requests, fixed-layer finalization, and deck assembly."""

from __future__ import annotations

import hashlib
import json
import os
import shutil
import uuid
import zipfile
from pathlib import Path
from typing import Any, Mapping

from pptx import Presentation

from fixed_frame import apply_fixed_frame, inspect_fixed_frame
from workflow_v6_contract import geometry_contract, transition_page
from workflow_v6_state import load, update_page


EDITPPT_CLI = (
    Path(__file__).resolve().parents[2]
    / "reconstruct-editable-slide" / "cli"
)
if EDITPPT_CLI.is_dir():
    if str(EDITPPT_CLI) not in __import__("sys").path:
        __import__("sys").path.append(str(EDITPPT_CLI))
    from editppt.runtime.final_assembler import _copy_page_slide  # noqa: E402
else:
    from editppt.runtime.final_assembler import _copy_page_slide  # type: ignore # noqa: E402


def _sha256(path: Path) -> str:
    return hashlib.sha256(path.read_bytes()).hexdigest()


def _read_json(path: Path) -> dict[str, Any]:
    value = json.loads(path.read_text(encoding="utf-8"))
    if not isinstance(value, dict):
        raise ValueError(f"expected JSON object: {path}")
    return value


def _write_json(path: Path, value: Mapping[str, Any]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(value, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")


def build_reconstruction_request(project: Path, *, page_number: int) -> dict[str, Any]:
    root = Path(project).resolve()
    state = load(root)
    page = state["pages"][page_number - 1]
    if page["state"] not in {"accepted", "accepted_fallback_first"}:
        raise ValueError("V6 page must have a selected Image2 body before reconstruction")
    selected = page.get("selected_candidate")
    if not isinstance(selected, Mapping) or not isinstance(selected.get("path"), str):
        raise ValueError("V6 selected body is missing")
    image = root / selected["path"]
    if not image.is_file():
        raise ValueError("V6 selected body file is missing")
    effective = _read_json(
        root / "02_v6" / "effective_pages" / f"page_{page_number:03d}.json"
    )
    request = {
        "artifact_version": "reconstruction-request-v6",
        "workflow_contract_version": "word-ppt-workflow-v6",
        "operation": "reconstruct_editable_slide",
        "page_number": page_number,
        "page_title": page["title"],
        "source_body": {
            "path": image.relative_to(root).as_posix(),
            "sha256": _sha256(image),
            "pixels": {"width": 1904, "height": 896},
        },
        "effective_page": effective,
        "geometry": geometry_contract(),
        "requirements": {
            "object_level_editable": True,
            "body_only": True,
            "fixed_layers_added_after_reconstruction": True,
            "post_reconstruction_visual_qa": False,
            "exact_reference_material_custody": False,
        },
    }
    path = root / "05_v6" / "reconstruction_requests" / f"page_{page_number:03d}.json"
    _write_json(path, request)
    return request


def finalize_reconstructed_page(
    project: Path, *, page_number: int, reconstructed_body: Path
) -> dict[str, Any]:
    root = Path(project).resolve()
    reconstructed_body = Path(reconstructed_body).resolve()
    if not reconstructed_body.is_file() or reconstructed_body.suffix.lower() != ".pptx":
        raise ValueError("V6 reconstructed body must be an existing PPTX")
    opened = Presentation(reconstructed_body)
    if len(opened.slides) != 1:
        raise ValueError("V6 reconstructed body must contain exactly one slide")
    state = load(root)
    page_index = page_number - 1
    page = state["pages"][page_index]
    if page["state"] not in {"accepted", "accepted_fallback_first", "reconstructing"}:
        raise ValueError("V6 page is not ready for reconstruction finalization")
    if page["state"] != "reconstructing":
        page = transition_page(page, "reconstructing")
    output_dir = root / "06_v6" / "pages" / f"page_{page_number:03d}"
    output_dir.mkdir(parents=True, exist_ok=True)
    output = output_dir / "page.pptx"
    shutil.copy2(reconstructed_body, output)
    style = state["style_confirmation"]["contract"]
    if not isinstance(style, Mapping):
        raise ValueError("V6 confirmed style contract is missing")
    logo = root / state["logo_source"]["path"]
    apply_fixed_frame(
        output,
        page_title=page["title"],
        page_number=page_number,
        style_execution=style,
        logo_svg=logo,
    )
    fixed = inspect_fixed_frame(
        output,
        expected_title=page["title"],
        expected_page_number=page_number,
        style_execution=style,
        logo_svg=logo,
    )
    if fixed.get("passed") is not True:
        raise ValueError("V6 fixed-layer validation failed: " + "; ".join(fixed.get("issues", [])))
    page = transition_page(page, "page_complete")
    update_page(root, page_number, page)
    report = {
        "artifact_version": "final-page-v6",
        "page_number": page_number,
        "page_pptx": output.relative_to(root).as_posix(),
        "sha256": _sha256(output),
        "fixed_frame": fixed,
        "post_reconstruction_visual_qa": False,
    }
    _write_json(output_dir / "page.json", report)
    return report


def assemble_v6_deck(project: Path) -> dict[str, Any]:
    root = Path(project).resolve()
    state = load(root)
    if any(page["state"] != "page_complete" for page in state["pages"]):
        raise ValueError("every V6 page must be complete before assembly")
    pages = [
        root / "06_v6" / "pages" / f"page_{page['page_number']:03d}" / "page.pptx"
        for page in state["pages"]
    ]
    if any(not path.is_file() for path in pages):
        raise ValueError("a V6 finalized page package is missing")
    output_dir = root / "08_final"
    output_dir.mkdir(parents=True, exist_ok=True)
    output = output_dir / "deck.pptx"
    temporary = output_dir / f".deck-v6-{uuid.uuid4().hex[:8]}.tmp"
    deck = Presentation(pages[0])
    layout = deck.slides[0].slide_layout
    for page_number, path in enumerate(pages[1:], start=2):
        _copy_page_slide(path, deck, layout, page_number)
    deck.save(temporary)
    reopened = Presentation(temporary)
    if len(reopened.slides) != len(pages):
        raise ValueError("assembled V6 slide count is incorrect")
    if any(
        len([shape for shape in slide.shapes if shape.name.startswith("fixed-frame-")]) != 4
        for slide in reopened.slides
    ):
        raise ValueError("assembled V6 fixed-layer inventory is incorrect")
    if any(
        not any(shape.has_text_frame or shape.has_table for shape in slide.shapes)
        for slide in reopened.slides
    ):
        raise ValueError("assembled V6 slide has no editable text or table object")
    if not zipfile.is_zipfile(temporary):
        raise ValueError("assembled V6 output is not an OpenXML package")
    os.replace(temporary, output)
    report = {
        "artifact_version": "final-assembly-v6",
        "workflow_contract_version": "word-ppt-workflow-v6",
        "status": "complete",
        "page_count": len(pages),
        "page_order": [page["page_number"] for page in state["pages"]],
        "output": output.relative_to(root).as_posix(),
        "sha256": _sha256(output),
        "mechanical_validation": {
            "openxml_package": True,
            "slide_count": True,
            "fixed_layers": True,
            "editable_objects": True,
        },
        "office_render_required": False,
        "post_reconstruction_visual_qa": False,
    }
    _write_json(output_dir / "assembly.json", report)
    return report
