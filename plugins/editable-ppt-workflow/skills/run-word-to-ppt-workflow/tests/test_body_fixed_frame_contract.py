from __future__ import annotations

import hashlib
import json
import sys
import zipfile
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Cm
from docx import Document


ROOT = Path(__file__).resolve().parents[1]
SCRIPTS = ROOT / "scripts"
EDITPPT_CLI = ROOT.parent / "reconstruct-editable-slide" / "cli"
sys.path.insert(0, str(SCRIPTS))
sys.path.insert(0, str(EDITPPT_CLI))

from build_page_contracts import split_page_title_body  # noqa: E402
from fixed_frame import apply_fixed_frame, contained_logo_box, inspect_fixed_frame  # noqa: E402
from fixed_region_contract import BODY_BOX_CM, SLIDE_SIZE_CM, fixed_frame_execution  # noqa: E402
from page_generation import build_initial_request  # noqa: E402
from page_pipeline import generation_request  # noqa: E402
from prepare_run import prepare  # noqa: E402
from editppt.runtime.final_assembler import _copy_page_slide  # noqa: E402
from current_contract_fixture import install_current_page_artifacts  # noqa: E402


def _style() -> dict:
    execution = {
        "schema_version": "2.0",
        "canvas": "ppt169",
        "canvas_profile": {
            "aspect_ratio": "16:9",
            "slide_width_inches": 10,
            "slide_height_inches": 14.288 / 2.54,
            "fit": "reconstruct_to_body",
            "coordinate_space": "dynamic_source_normalized",
            "allow_crop": False,
        },
        "fixed_frame": {"title_color": "#0B1727", **fixed_frame_execution()},
        "hard_constraints": {
            "title_color": "#0B1727",
            "typography": {
                "heading": {"cjk": "Microsoft YaHei", "latin": "Arial"},
                "type_scale_pt": {"page_title": 28},
            },
        },
        "soft_preferences": {"visual_style": "formal-consulting", "information_density": "balanced"},
        "creative_freedom": {
            "layout": True,
            "composition": True,
            "visual_hierarchy": True,
            "content_visualization": True,
            "page_specific_emphasis": True,
        },
    }
    digest = hashlib.sha256(
        (json.dumps(execution, ensure_ascii=False, sort_keys=True, separators=(",", ":")) + "\n").encode("utf-8")
    ).hexdigest()
    return {"execution": execution, "sha256": digest}


def _body_pptx(path: Path) -> Path:
    deck = Presentation()
    deck.slide_width = Cm(SLIDE_SIZE_CM["w"])
    deck.slide_height = Cm(SLIDE_SIZE_CM["h"])
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Cm(BODY_BOX_CM["x"]),
        Cm(BODY_BOX_CM["y"]),
        Cm(5),
        Cm(2),
    )
    shape.name = "body-shape"
    deck.save(path)
    return path


def _logo(path: Path) -> Path:
    path.write_text('<svg xmlns="http://www.w3.org/2000/svg" width="120" height="40"><rect width="120" height="40"/></svg>', encoding="utf-8")
    return path


def test_logo_uses_maximum_box_as_contain_area_and_preserves_svg_ratio(tmp_path: Path) -> None:
    logo = tmp_path / "wide-logo.svg"
    logo.write_text('<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 607 133"/>', encoding="utf-8")
    fitted = contained_logo_box(logo)

    assert fitted["w"] == pytest.approx(2.667)
    assert fitted["h"] == pytest.approx(2.667 * 133 / 607)
    assert fitted["x"] + fitted["w"] == pytest.approx(21.844 + 2.667)
    assert fitted["y"] + fitted["h"] / 2 == pytest.approx(0.57152 + 1.0716 / 2)


def test_page_title_is_locked_separately_and_removed_from_body_generation_text() -> None:
    title, body = split_page_title_body("第2页\n软通动力收购事项取得关键进展\n交易结构与后续安排。", 2)

    assert title == "软通动力收购事项取得关键进展"
    assert body == "交易结构与后续安排。"


def test_physical_page_continuation_keeps_first_line_in_body_and_derives_short_title() -> None:
    title, body = split_page_title_body(
        "；其余仍有接近900万元尚未落实，正在持续推进。\n下一步将完成相关工作。",
        7,
        pagination_mode="physical_rendered_pages",
    )

    assert len(title) <= 28
    assert body.startswith("；其余仍有接近900万元")


def test_page_comment_can_override_title_without_becoming_body_content() -> None:
    title, body = split_page_title_body(
        "开头是跨页续文。\n本页展示后续工作。",
        3,
        pagination_mode="physical_rendered_pages",
        page_comments=[{"text": "PPT标题：项目后续工作安排"}],
    )

    assert title == "项目后续工作安排"
    assert body.startswith("开头是跨页续文。")


def test_physical_page_skips_table_header_and_uses_later_section_heading() -> None:
    title, body = split_page_title_body(
        "标准事项      统一方案                牵头方\n"
        "网络通信      无线方案\n"
        "（四）下一步计划及请求\n1、下一步重点工作",
        8,
        pagination_mode="physical_rendered_pages",
    )

    assert title == "下一步计划及请求"
    assert body.startswith("标准事项")


def test_image2_request_receives_sealed_body_and_fixed_layer_exclusions(tmp_path: Path) -> None:
    from test_v4_complete_body_generation import _write_generation_inputs

    project, bundle, style = _write_generation_inputs(tmp_path)
    request = build_initial_request(bundle, style, project / "page.png", project=project)

    prompt = request.payload["prompt"]
    assert bundle["authoritative_content"]["body_text"] in prompt
    assert "page_title is drawn by the fixed title layer" in prompt
    assert "original SVG logo" in prompt
    assert "footer" in prompt and "page_number" in prompt


def test_page_pipeline_rejects_a_page_without_its_sealed_material_bundle(tmp_path: Path) -> None:
    project = tmp_path / "project"
    (project / "01_page_contracts").mkdir(parents=True)
    (project / "02_style").mkdir()
    title = "本页结论标题"
    body = "第一项工作。\n第二项工作。"
    source = f"{title}\n{body}"
    contract = {
        "schema_version": "2.0",
        "page_number": 1,
        "page_title": title,
        "body_text": body,
        "body_hash": hashlib.sha256(body.encode("utf-8")).hexdigest(),
        "source_text": source,
        "source_hash": hashlib.sha256(source.encode("utf-8")).hexdigest(),
        "asset_bindings": [],
    }
    (project / "01_page_contracts/page_001.json").write_text(json.dumps(contract), encoding="utf-8")
    style = _style()
    (project / "02_style/style_execution.json").write_text(
        json.dumps(style["execution"], ensure_ascii=False, sort_keys=True, separators=(",", ":")) + "\n",
        encoding="utf-8",
    )
    run = {
        "style_confirmation": {
            "status": "confirmed",
            "execution_file": "02_style/style_execution.json",
            "execution_sha256": style["sha256"],
        }
    }
    job = {"page_number": 1, "status": "queued", "contract_file": "01_page_contracts/page_001.json"}
    install_current_page_artifacts(project, contract, job)

    with pytest.raises(ValueError, match="material bundle identity is incomplete"):
        generation_request(project, run, job, 1)


def test_fixed_frame_is_added_as_native_objects_with_exact_title_and_svg_logo(tmp_path: Path) -> None:
    pptx = _body_pptx(tmp_path / "page.pptx")
    logo = tmp_path / "logo.svg"
    logo.write_text(
        '<svg xmlns="http://www.w3.org/2000/svg" width="120" height="40">'
        '<rect width="120" height="40" fill="#0B1727"/><text x="12" y="27" fill="white">ACME</text></svg>',
        encoding="utf-8",
    )

    apply_fixed_frame(
        pptx,
        page_title="本页结论标题",
        page_number=3,
        style_execution=_style()["execution"],
        logo_svg=logo,
    )

    deck = Presentation(pptx)
    shapes = {shape.name: shape for shape in deck.slides[0].shapes}
    assert {"fixed-frame-title", "fixed-frame-logo", "fixed-frame-footer", "fixed-frame-page-number"} <= shapes.keys()
    assert shapes["fixed-frame-title"].text == "本页结论标题"
    assert shapes["fixed-frame-page-number"].text == "3"
    assert "body-shape" in shapes
    assert inspect_fixed_frame(
        pptx,
        expected_title="本页结论标题",
        expected_page_number=3,
        style_execution=_style()["execution"],
        logo_svg=logo,
    )["passed"] is True


def test_fixed_frame_geometry_is_identical_across_pages(tmp_path: Path) -> None:
    geometry = []
    for page_number, title in ((1, "短标题"), (2, "这是一个明显更长但仍然必须留在同一个固定标题框中的页面结论标题")):
        pptx = _body_pptx(tmp_path / f"page-{page_number}.pptx")
        apply_fixed_frame(
            pptx,
            page_title=title,
            page_number=page_number,
            style_execution=_style()["execution"],
            logo_svg=_logo(tmp_path / f"logo-{page_number}.svg"),
        )
        deck = Presentation(pptx)
        frame = {
            shape.name: (shape.left, shape.top, shape.width, shape.height)
            for shape in deck.slides[0].shapes
            if shape.name.startswith("fixed-frame-") and shape.name != "fixed-frame-logo"
        }
        geometry.append(frame)

    assert geometry[0] == geometry[1]


def test_fixed_frame_rejects_user_supplied_geometry_override(tmp_path: Path) -> None:
    pptx = _body_pptx(tmp_path / "confirmed-geometry.pptx")
    execution = _style()["execution"]
    execution["fixed_frame"]["title_bounds_cm"]["x"] = 3
    with pytest.raises(ValueError, match="title_bounds_cm"):
        apply_fixed_frame(
            pptx,
            page_title="确认区域",
            page_number=1,
            style_execution=execution,
            logo_svg=_logo(tmp_path / "logo.svg"),
        )


def test_body_objects_already_reconstructed_in_target_window_are_not_scaled_twice(tmp_path: Path) -> None:
    pptx = _body_pptx(tmp_path / "direct-target.pptx")
    before = next(shape for shape in Presentation(pptx).slides[0].shapes if shape.name == "body-shape")
    original = (before.left, before.top, before.width, before.height)

    apply_fixed_frame(
        pptx,
        page_title="直接目标重建",
        page_number=1,
        style_execution=_style()["execution"],
        logo_svg=_logo(tmp_path / "logo.svg"),
    )

    after = next(shape for shape in Presentation(pptx).slides[0].shapes if shape.name == "body-shape")
    assert (after.left, after.top, after.width, after.height) == original


def test_fixed_title_uses_the_exact_confirmed_heading_font_and_size(tmp_path: Path) -> None:
    pptx = _body_pptx(tmp_path / "typography.pptx")
    execution = _style()["execution"]
    execution["hard_constraints"] = {
        "title_color": "#123456",
        "typography": {
            "heading": {"cjk": "SimHei", "latin": "Arial"},
            "type_scale_pt": {"page_title": 31},
        },
    }

    apply_fixed_frame(
        pptx,
        page_title="固定标题字体",
        page_number=1,
        style_execution=execution,
        logo_svg=_logo(tmp_path / "logo.svg"),
    )

    title = next(shape for shape in Presentation(pptx).slides[0].shapes if shape.name == "fixed-frame-title")
    run = title.text_frame.paragraphs[0].runs[0]
    assert run.font.name == "SimHei"
    assert run.font.size.pt == 31


def test_prepare_copies_an_optional_svg_logo_without_sending_it_to_image2(tmp_path: Path) -> None:
    word = tmp_path / "source.docx"
    document = Document()
    document.add_paragraph("第1页")
    document.add_paragraph("本页结论")
    document.add_paragraph("本页正文。")
    document.save(word)
    logo = tmp_path / "company.svg"
    logo.write_text('<svg xmlns="http://www.w3.org/2000/svg" width="40" height="20"/>', encoding="utf-8")

    project = tmp_path / "project"
    prepare(word, project, logo=logo)

    state = json.loads((project / "workflow_run.json").read_text(encoding="utf-8"))
    assert state["logo_source"]["path"] == "00_source/company_logo.svg"
    assert len(state["logo_source"]["sha256"]) == 64
    assert (project / state["logo_source"]["path"]).read_bytes() == logo.read_bytes()


def test_final_assembler_preserves_the_native_svg_logo_relationship(tmp_path: Path) -> None:
    source = _body_pptx(tmp_path / "source-page.pptx")
    logo = tmp_path / "logo.svg"
    logo.write_text(
        '<svg xmlns="http://www.w3.org/2000/svg" width="120" height="40"><rect width="120" height="40"/></svg>',
        encoding="utf-8",
    )
    apply_fixed_frame(
        source,
        page_title="固定标题",
        page_number=1,
        style_execution=_style()["execution"],
        logo_svg=logo,
    )
    destination = Presentation()
    destination.slide_width = Cm(SLIDE_SIZE_CM["w"])
    destination.slide_height = Cm(SLIDE_SIZE_CM["h"])

    _copy_page_slide(source, destination, destination.slide_layouts[6], 1)
    output = tmp_path / "assembled.pptx"
    destination.save(output)

    assert len(Presentation(output).slides) == 1
    with zipfile.ZipFile(output) as archive:
        assert any(name.endswith(".svg") for name in archive.namelist())
