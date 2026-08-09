"""Deterministic multi-page acceptance coverage for V4 complete-body generation."""

from __future__ import annotations

import importlib.util
import hashlib
import json
import sys
import zipfile
from datetime import datetime, timezone
from pathlib import Path
from types import SimpleNamespace

from docx import Document
from lxml import etree
from PIL import Image
from pptx import Presentation


ROOT = Path(__file__).resolve().parents[1]
SCRIPTS = ROOT / "scripts"
SERVER_PATH = SCRIPTS / "confirm_ui" / "server.py"
sys.path.insert(0, str(SCRIPTS))

import batch_generation  # noqa: E402
from prepare_run import prepare  # noqa: E402
from production_runner import run_production  # noqa: E402
import workflow_state  # noqa: E402


W = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
O = "urn:schemas-microsoft-com:office:office"
PKG_REL = "http://schemas.openxmlformats.org/package/2006/relationships"
CONTENT_TYPES = "http://schemas.openxmlformats.org/package/2006/content-types"
PACKAGE_RELATIONSHIP = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/package"
FIXED_ZIP_TIME = (2026, 1, 1, 0, 0, 0)


def _fixed_zip(path: Path, replacements: dict[str, bytes] | None = None) -> None:
    replacements = replacements or {}
    with zipfile.ZipFile(path) as source:
        files = {name: source.read(name) for name in source.namelist()}
    files.update(replacements)
    temporary = path.with_suffix(path.suffix + ".fixed")
    with zipfile.ZipFile(temporary, "w", compression=zipfile.ZIP_DEFLATED, compresslevel=9) as output:
        for name in sorted(files):
            info = zipfile.ZipInfo(name, FIXED_ZIP_TIME)
            info.compress_type = zipfile.ZIP_DEFLATED
            info.external_attr = 0o600 << 16
            output.writestr(info, files[name])
    temporary.replace(path)


def _save_fixed(document: Document, path: Path) -> Path:
    fixed = datetime(2026, 1, 1, tzinfo=timezone.utc)
    document.core_properties.created = fixed
    document.core_properties.modified = fixed
    document.save(path)
    _fixed_zip(path)
    return path


def _attachment(path: Path) -> Path:
    document = Document()
    document.add_paragraph("投资额为80万元。")
    document.add_paragraph("计划日期为2027年12月31日。")
    document.add_paragraph("建议立即收购竞争对手，这是Word正文未涉及的新结论。")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "字段"
    table.cell(0, 1).text = "最新值"
    table.cell(1, 0).text = "投资额"
    table.cell(1, 1).text = "80万元"
    return _save_fixed(document, path)


def _inject_attachments(source: Path, attachment: Path) -> None:
    with zipfile.ZipFile(source) as archive:
        files = {name: archive.read(name) for name in archive.namelist()}

    document = etree.fromstring(files["word/document.xml"])
    relationships = etree.fromstring(files["word/_rels/document.xml.rels"])
    content_types = etree.fromstring(files["[Content_Types].xml"])
    anchors = {
        "INTERNAL_ATTACHMENT_PAGE_1": "rIdAttachmentA",
        "INTERNAL_ATTACHMENT_PAGE_3": "rIdAttachmentA",
        "EXTERNAL_ATTACHMENT_LINK": "rIdExternalAttachment",
    }
    found: set[str] = set()
    for paragraph in document.iter(f"{{{W}}}p"):
        text = "".join(node.text or "" for node in paragraph.iter(f"{{{W}}}t"))
        for anchor, relationship_id in anchors.items():
            if anchor not in text:
                continue
            found.add(anchor)
            for node in paragraph.iter(f"{{{W}}}t"):
                node.text = ""
            object_node = etree.SubElement(paragraph, f"{{{W}}}object")
            ole = etree.SubElement(object_node, f"{{{O}}}OLEObject")
            ole.set(f"{{{R}}}id", relationship_id)
    assert found == set(anchors)

    etree.SubElement(
        relationships,
        f"{{{PKG_REL}}}Relationship",
        Id="rIdAttachmentA",
        Type=PACKAGE_RELATIONSHIP,
        Target="embeddings/附件A.docx",
    )
    etree.SubElement(
        relationships,
        f"{{{PKG_REL}}}Relationship",
        Id="rIdExternalAttachment",
        Type=PACKAGE_RELATIONSHIP,
        Target="https://example.test/attachment-download",
        TargetMode="External",
    )
    etree.SubElement(
        content_types,
        f"{{{CONTENT_TYPES}}}Override",
        PartName="/word/embeddings/附件A.docx",
        ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    )
    files.update({
        "word/document.xml": etree.tostring(document, xml_declaration=True, encoding="UTF-8", standalone="yes"),
        "word/_rels/document.xml.rels": etree.tostring(relationships, xml_declaration=True, encoding="UTF-8", standalone="yes"),
        "[Content_Types].xml": etree.tostring(content_types, xml_declaration=True, encoding="UTF-8", standalone="yes"),
        "word/embeddings/附件A.docx": attachment.read_bytes(),
    })
    _fixed_zip(source, files)


def _mixed_word(path: Path, attachment: Path, inline_image: Path) -> Path:
    document = Document()
    document.add_paragraph("第1页")
    document.add_paragraph("园区升级路径")
    page_one = document.add_paragraph("园区升级应首先完成诊断，其次实施改造，最后验收；投资额为50万元。")
    document.add_comment(page_one.runs, text="本页仅使用时间轴，不改变Word主结论。", author="reviewer")
    document.add_paragraph("INTERNAL_ATTACHMENT_PAGE_1")
    document.add_paragraph("EXTERNAL_ATTACHMENT_LINK")

    document.add_paragraph("第2页")
    document.add_paragraph("运营数据底稿")
    dense = "；".join(f"指标{i}保持稳定并由Word正文直接给出" for i in range(1, 36)) + "。"
    page_two = document.add_paragraph(dense)
    document.add_comment(page_two.runs, text="文字表达图片化", author="reviewer")

    document.add_paragraph("第3页")
    document.add_paragraph("项目投资与现场证据")
    page_three = document.add_paragraph(
        "项目继续按既定方案推进，投资额为50万元，计划日期为2026年6月30日，现场设备已经到位。"
    )
    document.add_comment(
        page_three.runs,
        text="本页仅使用时间轴，不改变Word主结论。",
        author="reviewer",
    )
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "里程碑"
    table.cell(0, 1).text = "状态"
    table.cell(1, 0).text = "设备到位"
    table.cell(1, 1).text = "完成"
    document.add_picture(str(inline_image))
    document.add_paragraph("INTERNAL_ATTACHMENT_PAGE_3")

    _save_fixed(document, path)
    _inject_attachments(path, attachment)
    return path


def _logo(path: Path) -> Path:
    path.write_text(
        '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 120 48">'
        '<rect width="120" height="48" fill="#17365D"/></svg>',
        encoding="utf-8",
    )
    return path


def _load_server():
    spec = importlib.util.spec_from_file_location("mixed_e2e_confirm_ui", SERVER_PATH)
    assert spec and spec.loader
    module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(module)
    return module


def _confirm_and_freeze(project: Path) -> None:
    server = _load_server()
    client = server.create_app(project).test_client()
    recommendations = client.get("/api/recommendations").get_json()
    candidate = recommendations["design_directions"]["candidates"][recommendations["design_directions"]["selected"]]
    payload = {
        "stage": "final",
        "direction": recommendations["design_directions"]["selected"],
        "template_selection": candidate["template_selection"],
        "canvas": "ppt169",
        **{key: candidate[key] for key in (
            "visual_style", "color", "icons", "typography", "image_rendering", "style_axes",
            "layout_preferences", "information_density", "background_system", "image_role",
            "evidence_strength", "composition_tendency", "brand_device",
        )},
        "regional_style": {"enabled": False},
        "production_profile": "balanced",
        "additional_requirements": "保持Word主叙事和逐页证据边界",
    }
    response = client.post("/api/confirm", json=payload)
    assert response.status_code == 200, response.get_json()
    assert server._wait(project, "final", 1) == 0


def _fake_renderer(pptx: Path, output: Path) -> int:
    output.mkdir(parents=True, exist_ok=True)
    count = len(Presentation(pptx).slides)
    for number in range(1, count + 1):
        Image.new("RGB", (160, 90), "white").save(output / f"slide_{number:03d}.png")
    return count


def _all_text(slide) -> str:
    values: list[str] = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            values.append(shape.text)
        if getattr(shape, "has_table", False):
            values.extend(cell.text for row in shape.table.rows for cell in row.cells)
    return "\n".join(values)


def test_three_page_mixed_sample_generates_every_uncached_v4_body_then_stops_closed(
    tmp_path: Path, monkeypatch,
) -> None:
    attachment = _attachment(tmp_path / "附件A.docx")
    inline_image = tmp_path / "现场照片.png"
    Image.new("RGB", (320, 180), "#22577A").save(inline_image)
    word = _mixed_word(tmp_path / "mixed-source.docx", attachment, inline_image)
    project = tmp_path / "project"

    prepared = prepare(word, project, _logo(tmp_path / "logo.svg"))
    assert prepared["page_count"] == 3
    run = workflow_state.load(project)
    assert run["workflow_contract_version"] == "word-ppt-workflow-v4"
    assert all("route" not in job for job in run["jobs"])
    contracts = [
        json.loads((project / f"01_page_contracts/page_{number:03d}.json").read_text(encoding="utf-8"))
        for number in range(1, 4)
    ]
    assert [[comment["text"] for comment in contract["page_comments"]] for contract in contracts] == [
        ["本页仅使用时间轴，不改变Word主结论。"],
        ["文字表达图片化"],
        ["本页仅使用时间轴，不改变Word主结论。"],
    ]
    assert sum(binding["asset_role"] == "mandatory_inline_image" for binding in contracts[2]["asset_bindings"]) == 1
    assert not any(binding["asset_role"] == "mandatory_inline_image" for contract in contracts[:2] for binding in contract["asset_bindings"])

    _confirm_and_freeze(project)
    frozen = workflow_state.load(project)
    assert frozen["style_confirmation"]["status"] == "confirmed"
    assert (project / frozen["style_confirmation"]["execution_file"]).is_file()

    scheduled = workflow_state.next_action(project)
    assert scheduled["stage"] == "page_pipeline"
    assert sorted((item["page_number"], item["action"]) for item in scheduled["requests"]) == [
        (1, "generate"),
        (2, "generate"),
        (3, "generate"),
    ]
    assert all("route" not in item for item in scheduled["requests"])
    assert all(
        "complete editable-PPT body design" in item["generation_request"]["prompt"]
        for item in scheduled["requests"]
    )
    assert scheduled["requests"][2]["generation_request"]["image_roles"] == ["reference_only"]

    backend_calls: list[list[str]] = []

    def fake_backend(command, **_kwargs):
        backend_calls.append(list(command))
        output = Path(command[command.index("--out") + 1])
        trace = Path(command[command.index("--trace-out") + 1])
        size = command[command.index("--size") + 1]
        width, height = (int(value) for value in size.split("x"))
        output.parent.mkdir(parents=True, exist_ok=True)
        Image.new("RGB", (width, height), "#E8EEF4").save(output)
        images: list[str] = []
        roles: list[str] = []
        for index, value in enumerate(command):
            if value == "--image":
                images.append(command[index + 1])
            elif value == "--image-role":
                roles.append(command[index + 1])
        operation = command[2]
        trace.write_text(json.dumps({
            "operation": operation,
            "endpoint": "images/edits" if operation == "edit" else "images/generations",
            "model": command[command.index("--model") + 1],
            "auth": "codex_oauth",
            "input_images": [
                {
                    "role": role,
                    "path": str(Path(path).resolve()),
                    "sha256": hashlib.sha256(Path(path).read_bytes()).hexdigest(),
                }
                for path, role in zip(images, roles)
            ],
            "outputs": [{
                "path": str(output.resolve()),
                "sha256": hashlib.sha256(output.read_bytes()).hexdigest(),
            }],
        }), encoding="utf-8")
        return SimpleNamespace(returncode=0, stdout="", stderr="")

    monkeypatch.setattr(batch_generation.subprocess, "run", fake_backend)
    batch = batch_generation.run_batch(project, timeout=10)
    assert [item["status"] for item in batch["results"]] == ["qa", "qa", "qa"]
    completed = workflow_state.next_action(project)
    assert completed["stage"] == "qa_backend_pending"
    assert completed["pending_pages"] == [1, 2, 3]
    assert len(backend_calls) == 3
    assert all("token_expired" not in " ".join(call) for call in backend_calls)
    generated = workflow_state.load(project)
    assert [job["status"] for job in generated["jobs"]] == ["qa", "qa", "qa"]
    assert [job["generation_calls"] for job in generated["jobs"]] == [1, 1, 1]
    assert all(job["generation_receipt"]["artifact_version"] == "page-generation-v1" for job in generated["jobs"])
    receipts = [
        json.loads((project / job["generation_receipt"]["path"]).read_text(encoding="utf-8"))
        for job in generated["jobs"]
    ]
    assert [(item["body_image"]["width"], item["body_image"]["height"]) for item in receipts] == [
        (1904, 896),
        (1904, 896),
        (1904, 896),
    ]
    assert receipts[2]["reference_images"][0]["role"] == "reference_only"
    assert not (project / "08_final" / "deck.pptx").exists()
