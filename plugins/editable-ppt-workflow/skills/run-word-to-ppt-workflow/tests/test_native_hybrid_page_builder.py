from __future__ import annotations

import sys
import json
from pathlib import Path

from PIL import Image
import pytest
from pptx import Presentation
from pptx.util import Cm


SCRIPTS = Path(__file__).resolve().parents[1] / "scripts"
sys.path.insert(0, str(SCRIPTS))


def test_native_plan_uses_authoritative_word_text_and_editable_table(tmp_path):
    from native_page_plan import build_native_page_plan
    from native_page_builder import build_native_page

    contract = {
        "page_number": 1, "page_title": "经营数据", "body_text": "收入保持增长。",
        "source_tables": ["|指标|数值|\n|---|---|\n|收入|100亿元|"], "asset_bindings": [],
    }
    plan = build_native_page_plan(contract, {"mandatory_anchors": ["100亿元"]}, {"route": "native"})
    output = build_native_page(plan, {}, tmp_path / "page.pptx")
    deck = Presentation(output)
    slide = deck.slides[0]

    assert any(getattr(shape, "has_table", False) for shape in slide.shapes)
    assert "收入保持增长。" in "\n".join(getattr(shape, "text", "") for shape in slide.shapes)
    left, top = int(Cm(0.81)), int(Cm(2.3))
    right, bottom = left + int(Cm(23.78)), top + int(Cm(11.18))
    assert all(left <= shape.left and top <= shape.top and shape.left + shape.width <= right + 2 and shape.top + shape.height <= bottom + 2 for shape in slide.shapes)


def test_overlay_plan_keeps_required_images_page_local():
    from native_page_plan import build_native_page_plan

    contract = {
        "page_number": 2, "page_title": "现场进展", "body_text": "完成设备安装。", "source_tables": [],
        "asset_bindings": [{
            "asset_id": "word_asset_003", "asset_role": "mandatory_inline_image",
            "generation_input": {"relative_path": "00_source/word_assets/original/photo.png"},
        }],
    }
    plan = build_native_page_plan(contract, {"mandatory_anchors": []}, {"route": "hybrid"})

    assert plan["route"] == "hybrid"
    assert plan["required_images"] == [{"asset_id": "word_asset_003", "relative_path": "00_source/word_assets/original/photo.png"}]


def test_native_builder_renders_all_tables_and_all_required_images_with_complete_receipt(tmp_path):
    from native_page_plan import build_native_page_plan
    from native_page_builder import build_native_page, load_render_receipt
    from page_coverage import build_coverage_contract, validate_render_receipt

    for name, color in (("one.png", "red"), ("two.png", "blue")):
        Image.new("RGB", (320, 180), color).save(tmp_path / name)
    contract = {
        "page_number": 1, "page_title": "完整性", "body_text": "结论一。\n结论二。",
        "semantic_units": [
            {"unit_id": "unit_001", "kind": "sentence", "text": "结论一。"},
            {"unit_id": "unit_002", "kind": "sentence", "text": "结论二。"},
        ],
        "source_tables": [
            "|A|B|\n|---|---|\n|1|2|", "|C|D|\n|---|---|\n|3|4|",
        ],
        "asset_bindings": [
            {"asset_id": "img1", "asset_role": "mandatory_inline_image", "generation_input": {"relative_path": "one.png"}},
            {"asset_id": "img2", "asset_role": "mandatory_inline_image", "generation_input": {"relative_path": "two.png"}},
        ],
    }
    facts = {"mandatory_anchors": ["结论一"]}
    coverage = build_coverage_contract(contract, facts)
    plan = build_native_page_plan(contract, facts, {"route": "native"}, coverage)
    output = build_native_page(plan, {}, tmp_path / "complete.pptx", project_root=tmp_path)
    deck = Presentation(output)

    assert sum(1 for shape in deck.slides[0].shapes if getattr(shape, "has_table", False)) == 2
    assert sum(1 for shape in deck.slides[0].shapes if shape.shape_type == 13) == 2
    assert validate_render_receipt(coverage, load_render_receipt(output))["passed"] is True


def test_hybrid_builder_places_visual_background_before_authoritative_native_objects(tmp_path):
    from native_page_builder import build_overlay_page

    Image.new("RGB", (800, 400), "gray").save(tmp_path / "background.png")
    plan = {"page_number": 1, "route": "hybrid", "body_text": "权威正文", "tables": [], "required_images": [], "coverage_contract": {"required_items": []}}
    output = build_overlay_page(plan, {"background": tmp_path / "background.png"}, {}, tmp_path / "hybrid.pptx", project_root=tmp_path)
    shapes = list(Presentation(output).slides[0].shapes)
    names = [shape.name for shape in shapes]

    assert names[0] == "hybrid-visual-background"
    assert shapes[0].width / shapes[0].height == pytest.approx(23.78 / 11.18, rel=1e-4)
    assert shapes[0].left == pytest.approx(Cm(0.81), abs=1)
    assert "native-body-text" in names[1:]


def test_native_builder_executes_confirmed_palette_and_typography(tmp_path):
    from native_page_builder import build_native_page

    style = {
        "hard_constraints": {
            "palette": {
                "background": "#FFF8F0", "secondary_bg": "#F4E3D0", "primary": "#7A2E00",
                "body_text": "#352015", "table_header": "#7A2E00", "border": "#C9A98D",
            },
            "typography": {"body": {"cjk": "Arial"}, "type_scale_pt": {"body": 14, "caption": 10}},
        },
        "soft_preferences": {"information_density": "balanced"},
    }
    plan = {
        "page_number": 1, "route": "native", "body_text": "可编辑权威正文", "tables": [],
        "required_images": [], "semantic_units": [], "mandatory_anchors": [],
        "coverage_contract": {"required_items": []},
    }
    output = build_native_page(plan, style, tmp_path / "styled.pptx")
    deck = Presentation(output)
    text = next(shape for shape in deck.slides[0].shapes if shape.name == "native-body-text")
    run = text.text_frame.paragraphs[0].runs[0]

    assert run.font.name == "Arial"
    assert str(run.font.color.rgb) == "352015"
