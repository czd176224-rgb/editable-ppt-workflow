from __future__ import annotations

import json
import sys
from pathlib import Path

from PIL import Image
import pytest
from pptx import Presentation


ROOT = Path(__file__).resolve().parents[1]
SCRIPTS = ROOT / "scripts"
sys.path.insert(0, str(SCRIPTS))


def _source(filename: str, locator: str = "paragraph:1") -> dict:
    return {"file": filename, "locator": locator, "sha256": (filename.encode().hex() + "0" * 64)[:64]}


def _contract(body: str, comments: list[str], filenames: tuple[str, ...] = ("附件A.docx",)) -> dict:
    return {
        "page_number": 1,
        "page_title": "项目事实",
        "body_text": body,
        "source_text": body,
        "source_tables": [],
        "semantic_units": [{"unit_id": "unit_001", "kind": "sentence", "text": body}],
        "page_comments": [{"text": text} for text in comments],
        "asset_bindings": [
            {"asset_id": f"word_asset_{index:03d}", "original_filename": filename, "asset_role": "document_source"}
            for index, filename in enumerate(filenames, 1)
        ],
    }


def _selected(*chunks: tuple[str, str]) -> dict:
    return {
        "selected_chunks": [
            {"evidence_id": f"evidence-{index:03d}", "text": text, "source": _source(filename)}
            for index, (filename, text) in enumerate(chunks, 1)
        ],
    }


@pytest.mark.parametrize(("field", "word_value", "attachment_value"), [
    ("负责人", "待定", "https://example.com/owner"),
    ("建议", "继续观察", "收购竞争对手"),
    ("负责人", "待定", "立即收购竞争对手"),
])
def test_explicit_override_reuses_factual_non_normative_no_url_policy(
    field: str, word_value: str, attachment_value: str,
) -> None:
    from deterministic_qa import run_deterministic_qa
    from page_fact_plan import build_fact_plan

    contract = _contract(
        f"{field}：{word_value}", [f"{field}以附件A最新数据为准"],
    )
    facts = build_fact_plan(contract, _selected(("附件A.docx", f"{field}：{attachment_value}")))
    qa = run_deterministic_qa(None, contract, facts, {"route": "native"}, {})

    assert facts["field_overrides"] == []
    assert facts["attachment_supplements"] == []
    assert {issue["code"] for issue in qa["issues"]} == {"word_attachment_conflict"}


@pytest.mark.parametrize("reverse", [False, True])
def test_directive_sentences_bind_one_source_and_field_without_cross_expansion(reverse: bool) -> None:
    from page_fact_plan import build_fact_plan

    contract = _contract(
        "投资额：50万元；金额：10万元",
        ["以附件A最新数据为准，仅覆盖投资额字段；以附件B最新数据为准，仅覆盖金额字段"],
        ("附件A.docx", "附件B.docx"),
    )
    chunks = [
        ("附件A.docx", "投资额：80万元；金额：999万元"),
        ("附件B.docx", "投资额：777万元；金额：20万元"),
    ]
    if reverse:
        chunks.reverse()

    facts = build_fact_plan(contract, _selected(*chunks))

    assert {(item["field"], item["attachment_value"], item["source"]["file"]) for item in facts["field_overrides"]} == {
        ("投资额", "80万元", "附件A.docx"),
        ("金额", "20万元", "附件B.docx"),
    }


def test_one_field_authorized_to_two_attachments_is_ambiguous_and_word_wins() -> None:
    from deterministic_qa import run_deterministic_qa
    from page_fact_plan import build_fact_plan

    contract = _contract(
        "投资额：50万元",
        ["投资额以附件A最新数据为准；投资额以附件B最新数据为准"],
        ("附件A.docx", "附件B.docx"),
    )
    facts = build_fact_plan(
        contract,
        _selected(("附件A.docx", "投资额：80万元"), ("附件B.docx", "投资额：90万元")),
    )
    qa = run_deterministic_qa(None, contract, facts, {"route": "native"}, {})

    assert facts["field_overrides"] == []
    assert facts["attachment_supplements"] == []
    assert all(issue["code"] == "word_attachment_conflict" for issue in qa["issues"])
    assert len(qa["issues"]) == 2


@pytest.mark.parametrize(("field", "word_value", "attachment_value"), [
    ("负责人", "待定", "张三"),
    ("状态", "待审核", "已批准"),
    ("地点", "杭州", "上海"),
])
def test_unauthorized_text_field_conflicts_are_word_wins_not_supplements(
    field: str, word_value: str, attachment_value: str,
) -> None:
    from deterministic_qa import run_deterministic_qa
    from page_fact_plan import build_fact_plan

    contract = _contract(f"{field}：{word_value}", [])
    facts = build_fact_plan(contract, _selected(("附件A.docx", f"{field}：{attachment_value}")))
    qa = run_deterministic_qa(None, contract, facts, {"route": "native"}, {})

    assert facts["attachment_supplements"] == []
    assert facts["conflicts"][0]["resolution"] == "word_wins"
    assert qa["issues"][0]["code"] == "word_attachment_conflict"
    assert qa["issues"][0]["severity"] == "advisory"


@pytest.mark.parametrize("route_name", ["image", "hybrid"])
def test_background_text_detector_scans_once_and_blocks_any_text(route_name: str, tmp_path: Path) -> None:
    from qa_runtime import decide_page_qa

    image = tmp_path / f"{route_name}.png"
    Image.new("RGB", (320, 180), "white").save(image)
    calls = 0

    def detector(_image: Path) -> dict:
        nonlocal calls
        calls += 1
        return {"ocr_text": "A1 2026 100万元", "ocr_boxes": [{"text": "A1"}]}

    report = decide_page_qa(
        image,
        {"page_number": 1, "page_title": "标题", "asset_bindings": []},
        {"mandatory_anchors": []},
        {"route": route_name, "text_authority": "native_overlay"},
        ocr_provider=detector,
    )

    assert calls == 1
    assert report["semantic_calls"] == 0
    assert report["result"]["status"] == "repair"
    assert {item["code"] for item in report["result"]["issues"]} == {"background_text_detected"}


def test_background_text_detector_unavailable_is_blocking_not_advisory(tmp_path: Path) -> None:
    from qa_runtime import decide_page_qa

    image = tmp_path / "unavailable.png"
    Image.new("RGB", (320, 180), "white").save(image)

    def unavailable(_image: Path) -> dict:
        raise RuntimeError("local text detector unavailable")

    report = decide_page_qa(
        image,
        {"page_number": 1, "page_title": "标题", "asset_bindings": []},
        {"mandatory_anchors": []},
        {"route": "hybrid", "text_authority": "native_overlay"},
        ocr_provider=unavailable,
    )

    assert report["result"]["status"] == "repair"
    issue = report["result"]["issues"][0]
    assert issue["code"] == "background_text_detection_unavailable"
    assert issue["severity"] == "structural"


def test_second_visual_failure_is_blocked_after_configured_repair_budget(tmp_path: Path) -> None:
    from test_independent_page_workflow import _project
    import workflow_state
    from current_contract_fixture import write_valid_generation_receipt, write_valid_qa_observation

    project = _project(tmp_path, 1)
    state_path = project / "workflow_run.json"
    state = json.loads(state_path.read_text(encoding="utf-8"))
    state["runtime"]["automatic_repair_budget"] = 1
    state_path.write_text(json.dumps(state, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    first_request = workflow_state.next_action(project)["requests"][0]
    first = workflow_state.dispatch(project, 1, "image-worker", first_request["attempt"])
    first_image = project / "06_images/generated" / f"page_001_attempt_{first['attempt']:03d}.png"
    first_receipt = write_valid_generation_receipt(project, 1, first["attempt"], first_image)
    workflow_state.record_generation(
        project, 1, "image-worker", first["attempt"], first_image,
        generation_receipt=first_receipt,
    )
    assert workflow_state.record_qa(
        project, 1, "image-worker", first["attempt"],
        signed_invocation_bundle=write_valid_qa_observation(
            project, 1, failed_check="readable_no_overflow",
            failure_detail="Generated body content is clipped.",
        ),
    )["state"] == "repair"

    second_request = workflow_state.next_action(project)["requests"][0]
    second = workflow_state.dispatch(project, 1, "image-worker", second_request["attempt"])
    second_image = project / "06_images/generated" / f"page_001_attempt_{second['attempt']:03d}.png"
    second_receipt = write_valid_generation_receipt(project, 1, second["attempt"], second_image)
    workflow_state.record_generation(
        project, 1, "image-worker", second["attempt"], second_image,
        generation_receipt=second_receipt,
    )
    result = workflow_state.record_qa(
        project, 1, "image-worker", second["attempt"],
        signed_invocation_bundle=write_valid_qa_observation(
            project, 1, failed_check="readable_no_overflow",
            failure_detail="Generated body content remains clipped.",
        ),
    )
    job = workflow_state.load(project)["jobs"][0]

    assert result["state"] == "content_blocked"
    assert job["status"] == "content_blocked"
    assert job["qa_result"]["status"] == "blocked"
    assert job["page_failure"]["category"] == "qa_unresolved"
    assert "reconstruction_corrections" not in job


def test_builder_consumes_reconstruction_correction_and_restores_actual_overlay(tmp_path: Path) -> None:
    from native_page_builder import build_native_page
    from native_page_plan import build_native_page_plan
    from page_coverage import build_coverage_contract
    from page_qa import qa_issue

    contract = {
        "page_number": 1, "page_title": "标题", "body_text": "权威正文。",
        "semantic_units": [{"unit_id": "unit_001", "kind": "sentence", "text": "权威正文。"}],
        "source_tables": [], "asset_bindings": [],
    }
    facts = {"field_overrides": [], "attachment_supplements": [], "mandatory_anchors": []}
    coverage = build_coverage_contract(contract, facts)
    plan = build_native_page_plan(contract, facts, {"route": "native"}, coverage)
    output = build_native_page(plan, {}, tmp_path / "native.pptx")
    deck = Presentation(output)
    body = next(shape for shape in deck.slides[0].shapes if shape.name == "native-body-text")
    body._element.getparent().remove(body._element)
    deck.save(output)
    before = output.read_bytes()
    correction = qa_issue(
        "native_body_missing", "原生页缺少权威正文。", "structural",
        "native_pptx_scan", "native-body-text", "high",
    )

    build_native_page(plan, {}, output, corrections=[correction])

    assert output.read_bytes() != before
    assert any(shape.name == "native-body-text" for shape in Presentation(output).slides[0].shapes)
    receipt = json.loads(output.with_suffix(".coverage.json").read_text(encoding="utf-8"))
    assert receipt["corrections_consumed"] == ["native_body_missing"]


def test_doctor_reports_required_local_background_text_capability() -> None:
    import doctor

    status = doctor.background_text_detection_status()

    assert set(status) == {"available", "backend", "failure_behavior"}
    assert status["failure_behavior"] == "content_blocked"
