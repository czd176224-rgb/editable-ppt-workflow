from __future__ import annotations

import json
import sys
from pathlib import Path

from PIL import Image
import pytest
from jsonschema import Draft202012Validator
from pptx import Presentation


ROOT = Path(__file__).resolve().parents[1]
SCRIPTS = ROOT / "scripts"
sys.path.insert(0, str(SCRIPTS))


def _attachment_contract(
    *, body: str = "投资额：50万元", comments: list[str] | None = None,
    filename: str = "附件A.docx",
) -> dict:
    return {
        "page_number": 1,
        "page_title": "项目事实",
        "body_text": body,
        "source_text": body,
        "source_tables": [],
        "semantic_units": [{"unit_id": "unit_001", "kind": "sentence", "text": body}],
        "page_comments": [{"text": value} for value in (comments or [])],
        "asset_bindings": [{
            "asset_id": "word_asset_001",
            "original_filename": filename,
            "asset_role": "document_source",
        }],
        "detected_numbers": [],
        "detected_dates": [],
        "detected_amounts": [],
    }


def _evidence(text: str, *, filename: str = "附件A.docx") -> dict:
    return {"selected_chunks": [{
        "evidence_id": "evidence-001",
        "text": text,
        "source": {"file": filename, "locator": "paragraph:1", "sha256": "a" * 64},
    }]}


@pytest.mark.parametrize("comments", [
    ["附件A.docx 是本页数据来源", "投资额采用最新数据为准"],
    ["投资额以附件A-old最新数据为准"],
    ["投资额以Word为准，附件A.docx仅供参考"],
    ["不得采用附件A.docx的投资额"],
    ["附件A.docx的投资额仅供参考，不得覆盖"],
])
def test_override_authorization_requires_one_positive_comment_with_exact_bound_asset(comments: list[str]) -> None:
    from page_fact_plan import build_fact_plan

    plan = build_fact_plan(_attachment_contract(comments=comments), _evidence("投资额：80万元"))

    assert plan["field_overrides"] == []


def test_override_authorization_accepts_one_complete_exact_attachment_instruction() -> None:
    from page_fact_plan import build_fact_plan

    plan = build_fact_plan(
        _attachment_contract(comments=["以附件A.docx最新数据为准，仅覆盖投资额字段"]),
        _evidence("投资额：80万元"),
    )

    assert [(item["field"], item["word_value"], item["attachment_value"]) for item in plan["field_overrides"]] == [
        ("投资额", "50万元", "80万元")
    ]


@pytest.mark.parametrize(("field", "old", "new"), [
    ("日期", "2026-06-30", "2027-12-31"),
    ("金额", "1,234.50万元", "2,345.75万元"),
    ("比例", "12.5%", "18.75%"),
    ("负责人", "张三", "李四"),
])
def test_field_override_parses_and_replaces_the_complete_typed_value(field: str, old: str, new: str) -> None:
    from page_fact_plan import apply_field_overrides, build_fact_plan

    contract = _attachment_contract(
        body=f"{field}：{old}", comments=[f"{field}以附件A.docx最新值为准"],
    )
    plan = build_fact_plan(contract, _evidence(f"{field}：{new}"))

    assert [
        (item["field"], item["word_value"], item["attachment_value"])
        for item in plan["field_overrides"]
    ] == [(field, old, new)]
    rendered = apply_field_overrides(contract["body_text"], plan)
    assert rendered == f"{field}：{new}"
    assert old not in rendered


def _override_coverage_fixture() -> tuple[dict, dict]:
    contract = {
        "page_number": 1,
        "page_title": "投资计划",
        "body_text": "投资额：50万元。",
        "semantic_units": [{
            "unit_id": "unit_001", "kind": "sentence", "text": "投资额：50万元。",
            "source_block_index": 1,
        }],
        "source_tables": ["|字段|数值|\n|---|---|\n|投资额|50万元|"],
        "asset_bindings": [],
    }
    facts = {
        "mandatory_anchors": ["80万元"],
        "attachment_supplements": [],
        "field_overrides": [{
            "field": "投资额", "word_value": "50万元", "attachment_value": "80万元",
            "source": {"file": "附件A.docx", "locator": "paragraph:1", "sha256": "a" * 64},
            "evidence_id": "evidence-001",
        }],
    }
    return contract, facts


def test_coverage_required_sources_use_only_the_approved_override_value() -> None:
    from page_coverage import build_coverage_contract

    contract, facts = _override_coverage_fixture()
    coverage = build_coverage_contract(contract, facts)
    sources = json.dumps([item["source"] for item in coverage["required_items"]], ensure_ascii=False)

    assert "80万元" in sources
    assert "50万元" not in sources


def test_coverage_receipt_rejects_old_value_even_when_the_id_is_visible(tmp_path: Path) -> None:
    from native_page_builder import build_native_page, load_render_receipt
    from native_page_plan import build_native_page_plan
    from page_coverage import CoverageValidationError, build_coverage_contract, validate_render_receipt

    contract, facts = _override_coverage_fixture()
    coverage = build_coverage_contract(contract, facts)
    plan = build_native_page_plan(contract, facts, {"route": "native"}, coverage)
    output = build_native_page(plan, {}, tmp_path / "override.pptx")
    receipts = load_render_receipt(output)
    semantic = next(item for item in receipts if item["coverage_id"] == "semantic:unit_001")
    semantic["observed_text"] = "投资额：50万元。"

    with pytest.raises(CoverageValidationError):
        validate_render_receipt(coverage, receipts)


@pytest.mark.parametrize("route_name", ["native", "image", "hybrid"])
def test_approved_attachment_supplement_is_an_editable_provenanced_overlay_for_every_route(
    tmp_path: Path, route_name: str,
) -> None:
    from native_page_builder import build_native_page, build_overlay_page
    from native_page_plan import build_native_page_plan
    from page_coverage import build_coverage_contract

    source = {"file": "附件A.docx", "locator": "paragraph:2", "sha256": "b" * 64}
    contract = {
        "page_number": 1, "page_title": "客户情况", "body_text": "客户：甲公司。",
        "semantic_units": [{"unit_id": "unit_001", "kind": "sentence", "text": "客户：甲公司。"}],
        "source_tables": [], "asset_bindings": [],
    }
    facts = {
        "mandatory_anchors": [], "field_overrides": [],
        "attachment_supplements": [{
            "text": "客户：甲公司为核心客户", "field": "客户", "source": source,
            "evidence_id": "evidence-002", "authorization": "supplement_only",
        }],
    }
    coverage = build_coverage_contract(contract, facts)
    plan = build_native_page_plan(contract, facts, {"route": route_name}, coverage)

    assert plan.get("attachment_supplements") == [{
        "coverage_id": "supplement:evidence-002", "text": "客户：甲公司为核心客户",
        "field": "客户", "source": source, "evidence_id": "evidence-002",
    }]
    output = tmp_path / f"{route_name}.pptx"
    if route_name == "native":
        build_native_page(plan, {}, output)
    else:
        build_overlay_page(plan, {}, {}, output)
    supplements = [shape for shape in Presentation(output).slides[0].shapes if shape.name.startswith("native-attachment-supplement-")]
    assert [shape.text for shape in supplements] == ["客户：甲公司为核心客户"]
    assert "http" not in supplements[0].text.casefold()


@pytest.mark.parametrize(("removed_name", "expected_code"), [
    ("native-body-text", "native_body_missing"),
    ("native-data-table-001", "native_table_missing"),
    ("native-required-image-img1", "missing_inline_image"),
])
def test_native_qa_reads_the_completed_pptx_and_rejects_missing_body_table_or_image(
    tmp_path: Path, removed_name: str, expected_code: str,
) -> None:
    from native_page_builder import build_native_page
    from native_page_plan import build_native_page_plan
    from page_coverage import build_coverage_contract
    from qa_runtime import decide_page_qa

    Image.new("RGB", (160, 90), "blue").save(tmp_path / "image.png")
    contract = {
        "page_number": 1, "page_title": "固定标题不能冒充正文", "body_text": "权威正文必须存在。",
        "semantic_units": [{"unit_id": "unit_001", "kind": "sentence", "text": "权威正文必须存在。"}],
        "source_tables": ["|字段|数值|\n|---|---|\n|收入|100万元|"],
        "asset_bindings": [{
            "asset_id": "img1", "asset_role": "mandatory_inline_image",
            "generation_input": {"relative_path": "image.png"},
        }],
    }
    facts = {"mandatory_anchors": [], "attachment_supplements": [], "field_overrides": []}
    coverage = build_coverage_contract(contract, facts)
    plan = build_native_page_plan(contract, facts, {"route": "native"}, coverage)
    output = build_native_page(plan, {}, tmp_path / "native.pptx", project_root=tmp_path)
    deck = Presentation(output)
    shape = next(item for item in deck.slides[0].shapes if item.name == removed_name)
    shape._element.getparent().remove(shape._element)
    deck.save(output)
    report = decide_page_qa(
        output, contract, facts,
        {
            "route": "native", "text_authority": "native_overlay", "native_plan": plan,
            "coverage_contract": coverage, "coverage_receipt": str(output.with_suffix(".coverage.json")),
        },
    )

    assert report["result"]["status"] == "repair"
    assert expected_code in {item["code"] for item in report["result"]["issues"]}


def test_structured_qa_issue_is_the_only_repair_issue_contract() -> None:
    from deterministic_qa import run_deterministic_qa
    from page_qa import validate_qa_issue

    contract = {
        "page_number": 1, "page_title": "图片", "body_text": "必须有图",
        "asset_bindings": [{"asset_id": "img1", "asset_role": "mandatory_inline_image"}],
    }
    result = run_deterministic_qa(
        None, contract, {"mandatory_anchors": []}, {"route": "hybrid", "text_authority": "native_overlay"},
        {"inline_image_presence": {"img1": False}},
    )
    issue = result["issues"][0]

    assert set(issue) == {"code", "message", "severity", "trigger", "evidence", "confidence"}
    assert validate_qa_issue(issue) == issue
    with pytest.raises(ValueError, match="structured"):
        validate_qa_issue("missing image")


def test_page_qa_schema_accepts_structured_issues_and_rejects_strings() -> None:
    schema = json.loads((ROOT / "schemas" / "page_qa.schema.json").read_text(encoding="utf-8"))
    issue = {
        "code": "missing_inline_image", "message": "必需图片缺失。", "severity": "local",
        "trigger": "pptx_shape_scan", "evidence": "native-required-image-img1", "confidence": "high",
    }
    structured = {"status": "repair", "repair_scope": "local", "issues": [issue]}
    strings = {"status": "repair", "repair_scope": "local", "issues": ["必需图片缺失。"]}

    assert not list(Draft202012Validator(schema).iter_errors(structured))
    assert list(Draft202012Validator(schema).iter_errors(strings))
