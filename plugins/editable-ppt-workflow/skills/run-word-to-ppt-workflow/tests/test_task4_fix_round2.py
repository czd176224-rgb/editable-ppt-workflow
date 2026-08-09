from __future__ import annotations

import sys
from pathlib import Path

import pytest
from pptx import Presentation


ROOT = Path(__file__).resolve().parents[1]
SCRIPTS = ROOT / "scripts"
sys.path.insert(0, str(SCRIPTS))


def _selected_evidence(filename: str, text: str) -> dict:
    return {
        "selected_chunks": [{
            "evidence_id": "evidence-001",
            "text": text,
            "source": {
                "file": filename,
                "locator": "paragraph:1",
                "sha256": "a" * 64,
            },
        }],
    }


def _override_contract(filename: str, comment: str) -> dict:
    return {
        "page_number": 1,
        "page_title": "项目事实",
        "body_text": "投资额：50万元",
        "source_text": "投资额：50万元",
        "source_tables": [],
        "semantic_units": [{
            "unit_id": "unit_001", "kind": "sentence", "text": "投资额：50万元",
        }],
        "page_comments": [{"text": comment}],
        "asset_bindings": [{
            "asset_id": "word_asset_001",
            "original_filename": filename,
            "asset_role": "document_source",
        }],
    }


@pytest.mark.parametrize(("filename", "comment"), [
    ("附件A.docx", "投资额以附件A旧版最新数据为准"),
    ("附件A旧版.docx", "投资额以附件A最新数据为准"),
    ("附件A.docx", "投资额以附件A最新数据为准；以Word正文为准"),
    ("附件A.docx", "Word正文优先；投资额以附件A最新数据为准"),
    ("附件A.docx", "投资额以附件A最新数据为准；附件仅供参考"),
])
def test_attachment_override_directive_rejects_inexact_or_word_authority(
    filename: str, comment: str,
) -> None:
    from page_fact_plan import build_fact_plan

    plan = build_fact_plan(
        _override_contract(filename, comment),
        _selected_evidence(filename, "投资额：80万元"),
    )

    assert plan["field_overrides"] == []


def test_attachment_override_directive_accepts_exact_alias_with_explicit_latest_data_particle() -> None:
    from page_fact_plan import build_fact_plan

    filename = "附件A.docx"
    plan = build_fact_plan(
        _override_contract(filename, "投资额以附件A最新数据为准"),
        _selected_evidence(filename, "投资额：80万元"),
    )

    assert [(item["field"], item["attachment_value"]) for item in plan["field_overrides"]] == [
        ("投资额", "80万元"),
    ]


def _150_override_fixture() -> tuple[dict, dict]:
    contract = {
        "page_number": 1,
        "page_title": "投资情况",
        "body_text": "投资额：50万元。",
        "source_tables": [],
        "semantic_units": [{
            "unit_id": "unit_001", "kind": "sentence", "text": "投资额：50万元。",
        }],
        "asset_bindings": [],
    }
    facts = {
        "mandatory_anchors": ["150万元"],
        "attachment_supplements": [],
        "field_overrides": [{
            "field": "投资额", "word_value": "50万元", "attachment_value": "150万元",
            "source": {"file": "附件A.docx", "locator": "paragraph:1", "sha256": "b" * 64},
            "evidence_id": "evidence-001",
        }],
    }
    return contract, facts


def test_coverage_compares_complete_typed_field_values_not_old_value_substrings() -> None:
    from page_coverage import CoverageValidationError, build_coverage_contract, validate_render_receipt

    contract, facts = _150_override_fixture()
    coverage = build_coverage_contract(contract, facts)
    item = next(value for value in coverage["required_items"] if value["coverage_id"] == "semantic:unit_001")
    receipt = {
        "coverage_id": item["coverage_id"],
        "object_id": "native-body-text",
        "visible": True,
        "expected_sha256": item["expected"]["sha256"],
        "observed_text": "投资额：150万元。",
    }

    assert validate_render_receipt({**coverage, "required_items": [item]}, [receipt])["passed"] is True

    receipt["observed_text"] = "投资额：50万元。"
    with pytest.raises(CoverageValidationError):
        validate_render_receipt({**coverage, "required_items": [item]}, [receipt])


def _supplement_contract(field: str, word_value: str) -> dict:
    text = f"{field}：{word_value}"
    return {
        "page_number": 1,
        "page_title": "补充事实",
        "body_text": text,
        "source_text": text,
        "source_tables": [],
        "semantic_units": [{"unit_id": "unit_001", "kind": "sentence", "text": text}],
        "page_comments": [],
        "asset_bindings": [],
    }


@pytest.mark.parametrize(("field", "word_value", "attachment_value"), [
    ("结论", "维持现状", "收购竞争对手"),
    ("建议", "继续观察", "立即收购"),
    ("客户", "甲公司", "收购竞争对手"),
])
def test_supplement_only_rejects_conclusion_fields_and_decision_values(
    tmp_path: Path, field: str, word_value: str, attachment_value: str,
) -> None:
    from native_page_builder import build_native_page, build_overlay_page
    from native_page_plan import build_native_page_plan
    from page_coverage import build_coverage_contract
    from page_fact_plan import build_fact_plan

    contract = _supplement_contract(field, word_value)
    facts = build_fact_plan(contract, _selected_evidence("附件A.docx", f"{field}：{attachment_value}"))

    assert facts["attachment_supplements"] == []
    coverage = build_coverage_contract(contract, facts)
    for route_name in ("native", "image", "hybrid"):
        plan = build_native_page_plan(contract, facts, {"route": route_name}, coverage)
        output = tmp_path / f"{field}-{route_name}.pptx"
        if route_name == "native":
            build_native_page(plan, {}, output)
        else:
            build_overlay_page(plan, {}, {}, output)
        assert not any(
            shape.name.startswith("native-attachment-supplement-")
            for shape in Presentation(output).slides[0].shapes
        )


def test_factual_supplement_keeps_provenance_and_renders_for_all_routes(tmp_path: Path) -> None:
    from native_page_builder import build_native_page, build_overlay_page
    from native_page_plan import build_native_page_plan
    from page_coverage import build_coverage_contract
    from page_fact_plan import build_fact_plan

    contract = _supplement_contract("负责人", "")
    facts = build_fact_plan(contract, _selected_evidence("附件A.docx", "负责人：张三"))

    assert len(facts["attachment_supplements"]) == 1
    source = facts["attachment_supplements"][0]["source"]
    assert source == {"file": "附件A.docx", "locator": "paragraph:1", "sha256": "a" * 64}
    coverage = build_coverage_contract(contract, facts)
    for route_name in ("native", "image", "hybrid"):
        plan = build_native_page_plan(contract, facts, {"route": route_name}, coverage)
        output = tmp_path / f"factual-{route_name}.pptx"
        if route_name == "native":
            build_native_page(plan, {}, output)
        else:
            build_overlay_page(plan, {}, {}, output)
        supplements = [
            shape.text for shape in Presentation(output).slides[0].shapes
            if shape.name.startswith("native-attachment-supplement-")
        ]
        assert supplements == ["负责人：张三"]
