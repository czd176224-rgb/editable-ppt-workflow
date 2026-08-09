from __future__ import annotations

import copy
import hashlib
import hmac
import json
import os
import sys
from types import SimpleNamespace
from pathlib import Path
from zipfile import ZIP_DEFLATED, ZipFile
from xml.etree import ElementTree as ET

import pytest
from PIL import Image
from pptx import Presentation


ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT / "scripts"))
sys.path.insert(0, str(ROOT.parent / "reconstruct-editable-slide" / "cli"))

from editppt.runtime.fixed_region_runtime import CONTENT_BOX, SLIDE
from fixed_region_contract import BODY_BOX_CM, FOOTER_LINE, LOGO_BOX_CM, PAGE_NUMBER_BOX_CM, TITLE_BOX_CM
from v4_reconstruction import (
    _secret,
    _validate_readability,
    build_and_sign_reconstruction as _build_and_sign_reconstruction,
    build_reconstruction_work_item,
    collect_reconstruction_closure,
    verify_signed_reconstruction,
    write_editable_receipt,
    restore_and_validate_completed_cache,
)
import v4_reconstruction_gateway
import v4_reconstruction
from cache_key import canonical_sha256
from effective_page_authority import build_effective_page_authority
from page_material_bundle_v4 import _material_summary, _seal_digest
from style_contract import compile_style_execution
from test_style_contract import confirmed_result
from editppt.runtime.editable_page_cache import inspect_editable_pptx


def test_page_reconstruction_closure_excludes_global_word_and_unrelated_source_files(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch,
) -> None:
    project = tmp_path / "project"
    material = project / "04_v4/material/page_002.json"
    contract = project / "01_page_contracts/page_002.json"
    style = project / "02_style/style_execution.json"
    logo = project / "00_source/company_logo.svg"
    attachment = project / "00_source/word_assets/derived/page2.txt"
    word = project / "00_source/source.docx"
    unrelated = project / "00_source/word_assets/derived/page1.txt"
    for path, payload in (
        (contract, b"{}"), (style, b"{}"), (logo, b"<svg/>") ,
        (attachment, b"page two"), (word, b"whole word"), (unrelated, b"page one"),
    ):
        path.parent.mkdir(parents=True, exist_ok=True)
        path.write_bytes(payload)
    material_value = {
        "page_number": 2,
        "page_images": [],
        "attachment_evidence": [{"path": attachment.relative_to(project).as_posix()}],
        "search_evidence": [],
    }
    _write_json(material, material_value)
    monkeypatch.setattr(
        v4_reconstruction, "_load_authenticated_page_material_bundle",
        lambda *_args, **_kwargs: material_value,
    )

    closure = collect_reconstruction_closure(
        project,
        # A stale or defensive caller may still supply the whole Word file.
        # The closure boundary itself must reject that global input.
        seeds=[material, contract, style, logo, word],
        page_number=2,
        material_bundle_path=material,
        material_bundle_file_sha256=_sha(material),
        material_bundle_sha256="0" * 64,
        authority_identity="1" * 64,
    )

    assert attachment.relative_to(project).as_posix() in closure
    assert logo.relative_to(project).as_posix() in closure
    assert word.relative_to(project).as_posix() not in closure
    assert unrelated.relative_to(project).as_posix() not in closure


def test_page_reconstruction_closure_never_trusts_material_paths_from_other_seed_json(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch,
) -> None:
    project = tmp_path / "project"
    material_value = {
        "page_number": 2,
        "page_images": [],
        "attachment_evidence": [],
        "search_evidence": [],
    }
    material = _write_json(project / "04_v4/material/page_002.json", material_value)
    logo = project / "00_source/company_logo.svg"
    logo.parent.mkdir(parents=True, exist_ok=True)
    logo.write_text("<svg/>", encoding="utf-8")
    monkeypatch.setattr(
        v4_reconstruction, "_load_authenticated_page_material_bundle",
        lambda *_args, **_kwargs: material_value,
    )
    unrelated = project / "00_source/word_assets/derived/page1-only.txt"
    unrelated.parent.mkdir(parents=True, exist_ok=True)
    unrelated.write_text("page one only", encoding="utf-8")
    forged_seed = _write_json(project / "07_editable/page_002/forged.json", {
        "page_images": [{"path": unrelated.relative_to(project).as_posix()}],
    })

    closure = collect_reconstruction_closure(
        project,
        seeds=[material, forged_seed],
        page_number=2,
        material_bundle_path=material,
        material_bundle_file_sha256=_sha(material),
        material_bundle_sha256="0" * 64,
        authority_identity="1" * 64,
    )

    assert forged_seed.relative_to(project).as_posix() in closure
    assert unrelated.relative_to(project).as_posix() not in closure


def test_completed_cache_requires_authenticated_material_before_manifest_preflight(tmp_path: Path) -> None:
    project = tmp_path / "project"
    word = project / "00_source/source.docx"
    cached = project / ".cache/pages/forged/global-source.docx"
    word.parent.mkdir(parents=True, exist_ok=True)
    cached.parent.mkdir(parents=True, exist_ok=True)
    word.write_bytes(b"whole word")
    cached.write_bytes(word.read_bytes())
    hit = SimpleNamespace(
        path=project / ".cache/pages/forged",
        manifest={
            "artifact_version": "v4-editable-page-cache-v3",
            "logical_files": {
                "00_source/source.docx": {
                    "path": "global-source.docx",
                    "sha256": _sha(cached),
                }
            },
        },
    )

    with pytest.raises(ValueError, match="authenticated current page material bundle"):
        restore_and_validate_completed_cache(
            project, {}, hit, authority_identity="0" * 64,
        )


def test_completed_cache_manifest_cannot_self_authorize_source_from_forged_material(
    tmp_path: Path,
) -> None:
    project = tmp_path / "project"
    unrelated = project / "00_source/word_assets/derived/page1-only.txt"
    cached = project / ".cache/pages/forged/page1-only.txt"
    unrelated.parent.mkdir(parents=True, exist_ok=True)
    cached.parent.mkdir(parents=True, exist_ok=True)
    unrelated.write_text("page one only", encoding="utf-8")
    cached.write_bytes(unrelated.read_bytes())
    material = _write_json(project / "04_v4/material/page_002.json", {
        "page_number": 2,
        "page_images": [{"path": unrelated.relative_to(project).as_posix()}],
        "attachment_evidence": [],
        "search_evidence": [],
    })
    logical = unrelated.relative_to(project).as_posix()
    hit = SimpleNamespace(
        path=project / ".cache/pages/forged",
        manifest={
            "artifact_version": "v4-editable-page-cache-v3",
            "logical_files": {
                logical: {"path": "page1-only.txt", "sha256": _sha(cached)},
            },
        },
    )

    with pytest.raises(ValueError, match="authenticated current page material bundle"):
        restore_and_validate_completed_cache(
            project,
            {
                "page_number": 2,
                "material_bundle_file": material.relative_to(project).as_posix(),
                "material_bundle_file_sha256": _sha(material),
                "material_bundle_sha256": "0" * 64,
            },
            hit,
            authority_identity="0" * 64,
        )


def _write_json(path: Path, value: dict) -> Path:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(value, ensure_ascii=False, sort_keys=True, separators=(",", ":")) + "\n", encoding="utf-8")
    return path


def _sha(path: Path) -> str:
    return hashlib.sha256(path.read_bytes()).hexdigest()


def build_and_sign_reconstruction(project: Path, *, work_item: Path, manifest: Path, **kwargs):
    """Exercise Task5 only through the real signed reconstruction gateway boundary."""
    payload = json.loads(manifest.read_text(encoding="utf-8"))
    decision = {key: copy.deepcopy(payload[key]) for key in ("text_boxes", "tables", "shapes", "images", "text_coverage", "table_coverage")}
    for item in decision["text_boxes"]:
        item.update({key: item.get(key, value) for key, value in {
            "font": "Microsoft YaHei", "color": "#111111", "bold": False, "italic": False,
            "align": "left", "valign": "middle", "wrap": True, "fit_text": True, "z_index": 10,
        }.items()})
    for item in decision["tables"]: item.setdefault("z_index", 20)
    for item in decision["shapes"]:
        item.setdefault("stroke", item.get("fill", "#000000")); item.setdefault("stroke_width", 0)
    rasters = {item["object_id"]: item for item in payload["raster_components"]}
    for item in decision["images"]:
        raster = rasters[item["object_id"]]
        if raster["source_type"] == "page-image": item["source_id"] = f"page-image:{raster['source_id']}"
        elif raster["source_type"] == "attachment": item["source_id"] = f"attachment:{raster['source_id']}"
        elif raster["source_type"] == "decorative-texture": raise ValueError("decorative full-body raster is not a signed gateway source")
        else: raise ValueError(f"{raster['source_type']} raster provenance is not a signed gateway source")
        item.pop("path", None); item.setdefault("alt", ""); item.setdefault("z_index", 30)
    original_invoke = v4_reconstruction_gateway._invoke_structured
    v4_reconstruction_gateway._invoke_structured = lambda *_args, **_kwargs: SimpleNamespace(
        value=decision, turn_id="turn_task5_fixture", model="gpt-test",
    )
    try:
        gateway = v4_reconstruction_gateway.invoke_builtin_gateway(project, work_item, timeout=10)
        derived = v4_reconstruction_gateway.verify_signed_bundle(project, gateway, work_item)
        return _build_and_sign_reconstruction(project, work_item=work_item, manifest=derived, gateway_invocation=gateway, **kwargs)
    finally:
        v4_reconstruction_gateway._invoke_structured = original_invoke


def _fixture(
    project: Path, *, page_image_policy: str | None = None, search_raster: bool = False,
    body_size: tuple[int, int] = (1700, 800), declared_body_size: tuple[int, int] | None = None,
):
    (project / "authorities").mkdir(parents=True)
    style = compile_style_execution(confirmed_result())
    style["fixed_frame"] = {
        "geometry_version": "fixed-canvas-cm-v2",
        "body_bounds_cm": dict(BODY_BOX_CM),
        "title_bounds_cm": dict(TITLE_BOX_CM),
        "logo_bounds_cm": dict(LOGO_BOX_CM),
        "page_number_bounds_cm": dict(PAGE_NUMBER_BOX_CM),
        "footer_line": dict(FOOTER_LINE),
        "title_color": "#123456",
    }
    style_path = _write_json(project / "authorities/style.json", style)
    logo = project / "authorities/logo.svg"
    logo.write_text('<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 200 80"><rect width="200" height="80" fill="#123456"/></svg>', encoding="utf-8")
    contract = {"page_number": 1, "page_title": "项目结论", "body_text": "投资额：50万元", "source_tables": ["|项目|数值|\n|-|-|\n|投资额|50万元|"]}
    contract_path = _write_json(project / "authorities/page.json", contract)
    body = project / "generated/body.png"
    body.parent.mkdir()
    Image.new("RGB", body_size, "#ddeeff").save(body)
    page_images = []
    required_ids = []
    if page_image_policy is not None:
        source_image = project / "authorities/word-image.png"
        Image.new("RGB", (300, 180), "#cc8844").save(source_image)
        promotion = None
        if page_image_policy == "required_presence":
            promotion = {"source": "page_comment", "directive_type": "require_page_image", "asset_id": "word-image-1", "raw": "[require-page-image:word-image-1]"}
            required_ids = ["word-image-1"]
        page_images = [{"asset_id": "word-image-1", "path": source_image.relative_to(project).as_posix(), "sha256": _sha(source_image), "media_type": "image/png", "presence_policy": page_image_policy, "promotion": promotion}]
    search_evidence = []
    if search_raster:
        search_image = project / "authorities/search-result.png"
        Image.new("RGB", (320, 180), "#4466aa").save(search_image)
        search_evidence = [{
            "evidence_id": "search-1", "asset_id": "search-request-0123456789abcdef",
            "query": "market chart", "source_url": "https://example.test/chart.png",
            "excerpt": "Market chart", "retrieved_at": "2026-08-01T00:00:00Z", "sha256": _sha(search_image),
        }]
    source_text = "项目结论\n投资额：50万元\n|项目|数值|\n|-|-|\n|投资额|50万元|"
    authority = build_effective_page_authority(
        page_contract=contract, style_execution=style, directives=[], page_images=page_images,
        attachment_evidence=[], search_evidence=search_evidence,
    )
    material = {
        "artifact_version": "page-material-bundle-v4", "workflow_contract_version": "word-ppt-workflow-v4", "page_number": 1,
        "source_text": source_text, "source_hash": hashlib.sha256(source_text.encode("utf-8")).hexdigest(),
        "authoritative_content": {"body_text": "投资额：50万元", "tables": [{"table_id": "table-1", "rows": [["项目", "数值"], ["投资额", "50万元"]]}]},
        "style_execution": {"path": style_path.relative_to(project).as_posix(), "sha256": _sha(style_path)},
        "page_images": page_images, "required_presence_asset_ids": required_ids, "comment_intents": [],
        "resolved_directives": [], "effective_page_authority": authority,
        "required_directives": copy.deepcopy(authority["required_directives"]),
        "superseded_directives": copy.deepcopy(authority["superseded_directives"]),
        "generation_readiness": {"ready": True, "code": "ready", "directive_ids": [], "blocking_reasons": []},
        "attachment_evidence": [], "search_evidence": search_evidence,
        "material_summary": _material_summary(page_images, [], search_evidence),
        "provenance": {
            "project_id": "p1", "source_sha256": "1" * 64,
            "page_contract_sha256": _sha(contract_path), "logo_sha256": _sha(logo),
            "raw_page_comments": [], "resolution_receipts": [],
            "comment_resolution_artifact": {
                "path": "confirm_ui/page_requirement_summary.json",
                "page_entry_sha256": "b" * 64,
                "page_entry_signature": "c" * 64,
                "page_contract_sha256": _sha(contract_path),
                "page_lock_sha256": "d" * 64,
                "raw_comments_sha256": canonical_sha256([]),
            },
        },
        "bundle_attestation_signature": "5" * 64,
    }
    material["sealed_sha256"] = _seal_digest(material)
    material_path = _write_json(project / "materials/page.json", material)
    trace = _write_json(project / "generated/trace.json", {"ok": True})
    declared_body_size = declared_body_size or body_size
    generation = {
        "artifact_version": "page-generation-v1", "workflow_contract_version": "word-ppt-workflow-v4", "prompt_contract_version": "page-prompt-v8", "page_number": 1,
        "material_bundle_sha256": material["sealed_sha256"],
        "effective_authority_sha256": authority["sealed_sha256"], "required_directive_ids": [],
        "authority_reference_inputs_sha256": "6" * 64,
        "request": {"operation": "generate", "endpoint": "images/generations", "prompt_sha256": "2" * 64, "authority_prompt_sha256": "7" * 64, "model": "gpt-image-2", "size": f"{declared_body_size[0]}x{declared_body_size[1]}", "quality": "high"},
        "body_image": {"path": body.relative_to(project).as_posix(), "sha256": _sha(body), "width": declared_body_size[0], "height": declared_body_size[1]},
        "body_image_mapping": {"mode": "direct"}, "reference_images": [],
        "provider_trace": {"path": trace.relative_to(project).as_posix(), "sha256": _sha(trace)},
        "sealed_sha256": "8" * 64, "receipt_signature": "9" * 64,
    }
    generation_path = _write_json(project / "generated/receipt.json", generation)
    qa = {
        "artifact_version": "page-qa-v1", "workflow_contract_version": "word-ppt-workflow-v4", "qa_policy_version": "risk-qa-v5", "page_number": 1,
        "material_bundle_sha256": material["sealed_sha256"], "generation_receipt_sha256": _sha(generation_path), "qa_work_item_sha256": "3" * 64,
        "qa_work_item": {"path": "qa/work.json", "sha256": "3" * 64}, "observation": {"path": "qa/obs.json", "sha256": "4" * 64},
        "status": "pass", "issues": [], "observations": {"decodable": True, "width": body_size[0], "height": body_size[1], "aspect_error": 0, "aspect_within_tolerance": True, "visible_fraction": 1, "luminance_stddev": 1, "luminance_entropy": 1, "gross_content_present": True}, "repairs_used": 0,
    }
    qa_path = _write_json(project / "qa/receipt.json", qa)
    work = build_reconstruction_work_item(project, material_bundle=material_path, generation_receipt=generation_path, qa_receipt=qa_path, page_contract=contract_path, style_execution=style_path, logo_svg=logo, output=project / "reconstruction/work.json")
    return work, body


def _manifest(project: Path, work: Path, *, flattened: Path | None = None):
    payload = {
        "artifact_version": "editable-reconstruction-manifest-v1", "work_item_sha256": _sha(work),
        "workflow_contract_version": "fixed-canvas-cm-v2", "reconstruction_contract_version": "editable-image-v3",
        "slide": dict(SLIDE), "content_box": dict(CONTENT_BOX), "source": {"width_px": 1700, "height_px": 800},
        "text_boxes": [{"object_id": "word-p1", "name": "body-paragraph-1", "text": "投资额：50万元", "font_size": 18, "box_px": [100, 70, 700, 100]}],
        "tables": [{"object_id": "word-t1", "name": "body-table-1", "rows": [["项目", "数值"], ["投资额", "50万元"]], "box_px": [100, 250, 1000, 300],
                    "font_size": 12, "font_color": "#000000", "cell_fill": "#FFFFFF", "cell_margin_px": 8}],
        "shapes": [{"object_id": "panel-1", "name": "body-panel", "type": "rect", "box_px": [20, 20, 1600, 700], "fill": "#f5f5f5", "z_index": 0}],
        "images": [], "raster_components": [],
        "text_coverage": [{"source_id": "word-p1", "text": "投资额：50万元", "object_name": "body-paragraph-1"}],
        "table_coverage": [{"table_id": "table-1", "object_name": "body-table-1"}],
    }
    if flattened:
        payload["images"] = [{"object_id": "flat-1", "name": "flat-body", "path": "../generated/body.png", "box_px": [0, 0, 1700, 800]}]
        payload["raster_components"] = [{"object_id": "flat-1", "sha256": _sha(flattened), "source_type": "decorative-texture", "source_id": "flat"}]
    return _write_json(project / "reconstruction/manifest.json", payload)


def _provider_decision(project: Path, work: Path) -> dict:
    payload = json.loads(_manifest(project, work).read_text(encoding="utf-8"))
    decision = {key: copy.deepcopy(payload[key]) for key in (
        "text_boxes", "tables", "shapes", "images", "text_coverage", "table_coverage",
    )}
    for item in decision["text_boxes"]:
        item.update({
            "font": "Microsoft YaHei", "color": "#111111", "bold": False, "italic": False,
            "align": "left", "valign": "middle", "wrap": True, "fit_text": True, "z_index": 10,
        })
    for item in decision["tables"]:
        item["z_index"] = 20
    for item in decision["shapes"]:
        item.update({"stroke": item["fill"], "stroke_width": 0})
    return decision


def test_gateway_injects_source_canvas_from_sealed_work_item_with_one_model_call(tmp_path: Path, monkeypatch):
    """Dropping provider source must use authenticated image pixels, never a retry."""
    project = tmp_path.resolve()
    work, _body = _fixture(
        project, body_size=(1904, 896), declared_body_size=(1024, 1024),
    )
    decision = _provider_decision(project, work)
    calls = []

    def invoke(*args, **kwargs):
        calls.append((args, kwargs))
        return SimpleNamespace(value=decision, turn_id="turn-source-omitted", model="gpt-test")

    monkeypatch.setattr(v4_reconstruction_gateway, "_invoke_structured", invoke)
    bundle_path = v4_reconstruction_gateway.invoke_builtin_gateway(project, work, timeout=10)
    bundle = json.loads(bundle_path.read_text(encoding="utf-8"))
    manifest = json.loads((project / bundle["manifest"]["path"]).read_text(encoding="utf-8"))

    assert len(calls) == 1
    assert manifest["source"] == {"width_px": 1904, "height_px": 896}
    assert bundle["attestation"]["source_origin"] == "sealed_work_item"
    assert v4_reconstruction_gateway.verify_signed_bundle(project, bundle_path, work).is_file()


def test_gateway_accepts_exact_provider_source_from_sealed_image(tmp_path: Path, monkeypatch):
    """An exact provider echo remains compatible while sealed pixels stay authoritative."""
    project = tmp_path.resolve()
    work, _body = _fixture(project)
    decision = _provider_decision(project, work)
    decision["source"] = {"width_px": 1700, "height_px": 800}
    monkeypatch.setattr(v4_reconstruction_gateway, "_invoke_structured", lambda *_args, **_kwargs: SimpleNamespace(
        value=decision, turn_id="turn-source-exact", model="gpt-test",
    ))

    bundle = v4_reconstruction_gateway.invoke_builtin_gateway(project, work, timeout=10)

    assert v4_reconstruction_gateway.verify_signed_bundle(project, bundle, work).is_file()


def test_gateway_rejects_provider_source_that_disagrees_with_sealed_image(tmp_path: Path, monkeypatch):
    """A provider cannot redefine the canvas used for bounds validation."""
    project = tmp_path.resolve()
    work, _body = _fixture(project)
    decision = _provider_decision(project, work)
    decision["source"] = {"width_px": 1024, "height_px": 1024}
    monkeypatch.setattr(v4_reconstruction_gateway, "_invoke_structured", lambda *_args, **_kwargs: SimpleNamespace(
        value=decision, turn_id="turn-source-conflict", model="gpt-test",
    ))

    with pytest.raises(ValueError, match="provider source canvas conflicts"):
        v4_reconstruction_gateway.invoke_builtin_gateway(project, work, timeout=10)

    assert not (project / "07_editable/page_001/reconstruction-provider-request.json").exists()


def test_gateway_rejects_non_integer_provider_source_even_when_values_compare_equal(tmp_path: Path, monkeypatch):
    """Boolean/string dimensions must not bypass exact trusted-source comparison."""
    project = tmp_path.resolve()
    work, _body = _fixture(project)
    decision = _provider_decision(project, work)
    decision["source"] = {"width_px": "1700", "height_px": 800}
    monkeypatch.setattr(v4_reconstruction_gateway, "_invoke_structured", lambda *_args, **_kwargs: SimpleNamespace(
        value=decision, turn_id="turn-source-type-conflict", model="gpt-test",
    ))

    with pytest.raises(ValueError, match="provider source canvas conflicts"):
        v4_reconstruction_gateway.invoke_builtin_gateway(project, work, timeout=10)


def test_signed_gateway_verification_requires_sealed_source_origin(tmp_path: Path, monkeypatch):
    """Signed audit data cannot silently lose the local source provenance marker."""
    project = tmp_path.resolve()
    work, _body = _fixture(project)
    decision = _provider_decision(project, work)
    monkeypatch.setattr(v4_reconstruction_gateway, "_invoke_structured", lambda *_args, **_kwargs: SimpleNamespace(
        value=decision, turn_id="turn-source-audit", model="gpt-test",
    ))
    bundle_path = v4_reconstruction_gateway.invoke_builtin_gateway(project, work, timeout=10)
    payload = json.loads(bundle_path.read_text(encoding="utf-8"))
    payload["attestation"].pop("source_origin")
    unsigned = {key: value for key, value in payload.items() if key != "signature"}
    payload["signature"] = hmac.new(
        v4_reconstruction_gateway._key(project), canonical_sha256(unsigned).encode("ascii"), hashlib.sha256,
    ).hexdigest()
    bundle_path.write_bytes(v4_reconstruction_gateway._canonical(payload) + b"\n")
    registry = v4_reconstruction_gateway._registry(project)
    nonce = payload["attestation"]["nonce"]
    registry["nonces"][nonce]["bundle_sha256"] = _sha(bundle_path)
    v4_reconstruction_gateway._write_registry(project, registry)

    with pytest.raises(ValueError, match="source origin"):
        v4_reconstruction_gateway.verify_signed_bundle(project, bundle_path, work)


def test_controlled_backend_builds_native_text_table_and_exact_fixed_layers(tmp_path: Path):
    project = tmp_path.resolve()
    work, _body = _fixture(project)
    manifest = _manifest(project, work)
    bundle = build_and_sign_reconstruction(project, work_item=work, manifest=manifest, body_pptx=project / "reconstruction/body.pptx", final_pptx=project / "reconstruction/page.pptx", bundle_output=project / "reconstruction/bundle.json")
    receipt_path = write_editable_receipt(project, work, bundle, project / "reconstruction/receipt.json")
    receipt = json.loads(receipt_path.read_text(encoding="utf-8"))
    assert receipt["object_counts"]["table"] == 1
    assert receipt["flattened_body_image"] is False
    deck = Presentation(project / receipt["editable_page"]["path"])
    names = [shape.name for shape in deck.slides[0].shapes]
    assert names.count("fixed-frame-title") == names.count("fixed-frame-logo") == names.count("fixed-frame-footer") == names.count("fixed-frame-page-number") == 1


def test_rejects_accepted_full_body_even_with_editable_overlay(tmp_path: Path):
    project = tmp_path.resolve()
    work, body = _fixture(project)
    manifest = _manifest(project, work, flattened=body)
    with pytest.raises(ValueError, match="full-body"):
        build_and_sign_reconstruction(project, work_item=work, manifest=manifest, body_pptx=project / "reconstruction/body.pptx", final_pptx=project / "reconstruction/page.pptx", bundle_output=project / "reconstruction/bundle.json")


def test_signed_bundle_tamper_is_rejected(tmp_path: Path):
    project = tmp_path.resolve()
    work, _body = _fixture(project)
    manifest = _manifest(project, work)
    bundle = build_and_sign_reconstruction(project, work_item=work, manifest=manifest, body_pptx=project / "reconstruction/body.pptx", final_pptx=project / "reconstruction/page.pptx", bundle_output=project / "reconstruction/bundle.json")
    payload = json.loads(bundle.read_text(encoding="utf-8"))
    payload["object_counts"]["text"] += 1
    _write_json(bundle, payload)
    with pytest.raises(ValueError, match="attestation"):
        verify_signed_reconstruction(project, work, bundle)


def _resign_bundle_as_project_key_attacker(project: Path, bundle: Path) -> None:
    payload = json.loads(bundle.read_text(encoding="utf-8"))
    unsigned = dict(payload)
    unsigned.pop("attestation")
    payload["attestation"] = hmac.new(_secret(project), canonical_sha256(unsigned).encode("ascii"), hashlib.sha256).hexdigest()
    _write_json(bundle, payload)


def _rewrite_final_and_resign(project: Path, bundle: Path, rewrite) -> None:
    page = project / "reconstruction/page.pptx"
    replacement = page.with_suffix(".rewrite.pptx")
    with ZipFile(page) as source, ZipFile(replacement, "w", ZIP_DEFLATED) as target:
        files = {item.filename: source.read(item.filename) for item in source.infolist()}
        rewrite(files)
        for name, data in files.items():
            target.writestr(name, data)
    replacement.replace(page)
    payload = json.loads(bundle.read_text(encoding="utf-8"))
    payload["editable_page"]["sha256"] = _sha(page)
    payload["slide_fingerprint"] = inspect_editable_pptx(page).slide_fingerprints[0]
    _write_json(bundle, payload)
    _resign_bundle_as_project_key_attacker(project, bundle)


@pytest.mark.parametrize("attack", ["fill", "stroke", "geometry", "z-order", "text-props"])
def test_complete_drawingml_body_change_is_rejected_after_attacker_resigns(tmp_path: Path, attack: str):
    project = tmp_path.resolve()
    work, _body = _fixture(project)
    manifest = _manifest(project, work)
    bundle = build_and_sign_reconstruction(project, work_item=work, manifest=manifest, body_pptx=project / "reconstruction/body.pptx", final_pptx=project / "reconstruction/page.pptx", bundle_output=project / "reconstruction/bundle.json")

    def rewrite(files: dict[str, bytes]) -> None:
        p = "http://schemas.openxmlformats.org/presentationml/2006/main"
        a = "http://schemas.openxmlformats.org/drawingml/2006/main"
        root = ET.fromstring(files["ppt/slides/slide1.xml"])
        sp_tree = root.find(f".//{{{p}}}spTree")
        assert sp_tree is not None
        by_name = {}
        for child in list(sp_tree):
            nv = child.find(f".//{{{p}}}cNvPr")
            if nv is not None:
                by_name[nv.get("name")] = child
        if attack == "z-order":
            left, right = by_name["body-paragraph-1"], by_name["body-table-1"]
            children = list(sp_tree)
            li, ri = children.index(left), children.index(right)
            children[li], children[ri] = children[ri], children[li]
            sp_tree[:] = children
        elif attack == "fill":
            color = by_name["body-panel"].find(f".//{{{a}}}solidFill/{{{a}}}srgbClr")
            assert color is not None
            color.set("val", "000000")
        elif attack == "geometry":
            geometry = by_name["body-panel"].find(f".//{{{a}}}prstGeom")
            assert geometry is not None
            geometry.set("prst", "ellipse")
        elif attack == "stroke":
            sp_pr = by_name["body-panel"].find(f"./{{{p}}}spPr")
            assert sp_pr is not None
            line = ET.SubElement(sp_pr, f"{{{a}}}ln", {"w": "12700"})
            solid = ET.SubElement(line, f"{{{a}}}solidFill")
            ET.SubElement(solid, f"{{{a}}}srgbClr", {"val": "FF0000"})
        else:
            run_props = by_name["body-paragraph-1"].find(f".//{{{a}}}rPr")
            assert run_props is not None
            run_props.set("b", "1")
        files["ppt/slides/slide1.xml"] = ET.tostring(root, encoding="utf-8", xml_declaration=True)

    _rewrite_final_and_resign(project, bundle, rewrite)
    with pytest.raises(ValueError, match="body|equivalent|DrawingML"):
        verify_signed_reconstruction(project, work, bundle)


def test_body_relationship_media_change_is_rejected_after_attacker_resigns(tmp_path: Path):
    project = tmp_path.resolve()
    work, _body = _fixture(project, page_image_policy="required_presence")
    manifest = _manifest(project, work)
    payload = json.loads(manifest.read_text(encoding="utf-8"))
    image = project / "authorities/word-image.png"
    payload["images"] = [{"object_id": "photo-1", "name": "source-photo", "path": "../authorities/word-image.png", "box_px": [1200, 80, 380, 260]}]
    payload["raster_components"] = [{"object_id": "photo-1", "sha256": _sha(image), "source_type": "page-image", "source_id": "word-image-1"}]
    _write_json(manifest, payload)
    bundle = build_and_sign_reconstruction(project, work_item=work, manifest=manifest, body_pptx=project / "reconstruction/body.pptx", final_pptx=project / "reconstruction/page.pptx", bundle_output=project / "reconstruction/bundle.json")

    def rewrite(files: dict[str, bytes]) -> None:
        rels = ET.fromstring(files["ppt/slides/_rels/slide1.xml.rels"])
        rel_ns = "http://schemas.openxmlformats.org/package/2006/relationships"
        media_rel = next(node for node in rels.findall(f"{{{rel_ns}}}Relationship") if str(node.get("Target", "")).endswith(".png"))
        target = str(Path("ppt/slides") / str(media_rel.get("Target"))).replace("\\", "/")
        import posixpath
        target = posixpath.normpath(target)
        from io import BytesIO
        buffer = BytesIO()
        Image.new("RGB", (300, 180), "#112233").save(buffer, format="PNG")
        files[target] = buffer.getvalue()

    _rewrite_final_and_resign(project, bundle, rewrite)
    with pytest.raises(ValueError, match="body|equivalent|DrawingML"):
        verify_signed_reconstruction(project, work, bundle)


@pytest.mark.parametrize("transform", ["off", "ext", "chOff", "chExt"])
def test_root_group_transform_change_is_rejected_after_attacker_resigns(tmp_path: Path, transform: str):
    project = tmp_path.resolve()
    work, _body = _fixture(project)
    manifest = _manifest(project, work)
    bundle = build_and_sign_reconstruction(project, work_item=work, manifest=manifest, body_pptx=project / "reconstruction/body.pptx", final_pptx=project / "reconstruction/page.pptx", bundle_output=project / "reconstruction/bundle.json")

    def rewrite(files: dict[str, bytes]) -> None:
        p = "http://schemas.openxmlformats.org/presentationml/2006/main"
        a = "http://schemas.openxmlformats.org/drawingml/2006/main"
        root = ET.fromstring(files["ppt/slides/slide1.xml"])
        group = root.find(f".//{{{p}}}spTree/{{{p}}}grpSpPr")
        assert group is not None
        xfrm = ET.SubElement(group, f"{{{a}}}xfrm")
        for name in ("off", "ext", "chOff", "chExt"):
            attrs = {"x": "0", "y": "0"} if name in {"off", "chOff"} else {"cx": "1", "cy": "1"}
            if name == transform:
                attrs[next(iter(attrs))] = "12345"
            ET.SubElement(xfrm, f"{{{a}}}{name}", attrs)
        files["ppt/slides/slide1.xml"] = ET.tostring(root, encoding="utf-8", xml_declaration=True)

    _rewrite_final_and_resign(project, bundle, rewrite)
    with pytest.raises(ValueError, match="body|DrawingML|group"):
        verify_signed_reconstruction(project, work, bundle)


@pytest.mark.parametrize(("font_color", "cell_fill", "should_pass"), [
    ("#FFFFFF", "#333333", True),
    ("#000000", "#FFFFFF", True),
    ("#000000", "#000000", False),
    ("#FFFFFF", "#FFFFFF", False),
])
def test_native_table_cells_require_explicit_readable_contrast(tmp_path: Path, font_color: str, cell_fill: str, should_pass: bool):
    project = tmp_path.resolve()
    work, _body = _fixture(project)
    manifest = _manifest(project, work)
    payload = json.loads(manifest.read_text(encoding="utf-8"))
    payload["tables"][0].update({"font_color": font_color, "cell_fill": cell_fill})
    _write_json(manifest, payload)
    call = lambda: build_and_sign_reconstruction(project, work_item=work, manifest=manifest, body_pptx=project / "reconstruction/body.pptx", final_pptx=project / "reconstruction/page.pptx", bundle_output=project / "reconstruction/bundle.json")
    if should_pass:
        assert call().is_file()
    else:
        with pytest.raises(ValueError, match="table|cell|contrast"):
            call()


def test_native_table_without_explicit_cell_visuals_is_rejected(tmp_path: Path):
    project = tmp_path.resolve()
    work, _body = _fixture(project)
    manifest = _manifest(project, work)
    payload = json.loads(manifest.read_text(encoding="utf-8"))
    payload["shapes"][0]["fill"] = "#000000"
    for field in ("font_size", "font_color", "cell_fill", "cell_margin_px"):
        payload["tables"][0].pop(field)
    _write_json(manifest, payload)
    with pytest.raises(ValueError, match="table|cell|explicit"):
        build_and_sign_reconstruction(project, work_item=work, manifest=manifest, body_pptx=project / "reconstruction/body.pptx", final_pptx=project / "reconstruction/page.pptx", bundle_output=project / "reconstruction/bundle.json")


@pytest.mark.parametrize("mutation", [{"box_px": [100, 250, 1000, 30]}, {"cell_margin_px": 200}, {"font_size": 6}])
def test_native_table_cell_size_and_capacity_are_verified_from_pptx(tmp_path: Path, mutation: dict):
    project = tmp_path.resolve()
    work, _body = _fixture(project)
    manifest = _manifest(project, work)
    payload = json.loads(manifest.read_text(encoding="utf-8"))
    payload["tables"][0].update(mutation)
    _write_json(manifest, payload)
    with pytest.raises(ValueError, match="table|cell|capacity|font|margin|readable"):
        build_and_sign_reconstruction(project, work_item=work, manifest=manifest, body_pptx=project / "reconstruction/body.pptx", final_pptx=project / "reconstruction/page.pptx", bundle_output=project / "reconstruction/bundle.json")


@pytest.mark.parametrize(("foreground", "background", "should_pass"), [
    ("FFFFFF", "333333", True),
    ("FFFFFF", "FFFFFF", False),
    ("000000", "000000", False),
    ("FFFFFF", None, True),
])
def test_text_contrast_uses_actual_reliable_background(foreground: str, background: str | None, should_pass: bool):
    item = {"name": "text", "kind": "text", "box": [0, 0, 2000000, 500000], "z_index": 1, "text": "Readable", "font_size": 18.0,
            "font_colors": (foreground,), "transparent": False, "opaque": False, "fill_color": background, "fill_kind": "solid" if background else "unknown", "slide_background": None}
    call = lambda: _validate_readability([item])
    if should_pass:
        call()
    else:
        with pytest.raises(ValueError, match="contrast"):
            call()


def test_final_body_text_change_is_rejected_even_when_attacker_resigns_bundle(tmp_path: Path):
    project = tmp_path.resolve()
    work, _body = _fixture(project)
    manifest = _manifest(project, work)
    bundle = build_and_sign_reconstruction(project, work_item=work, manifest=manifest, body_pptx=project / "reconstruction/body.pptx", final_pptx=project / "reconstruction/page.pptx", bundle_output=project / "reconstruction/bundle.json")
    page = project / "reconstruction/page.pptx"
    deck = Presentation(page)
    next(shape for shape in deck.slides[0].shapes if shape.name == "body-paragraph-1").text = "篡改后的正文"
    deck.save(page)
    payload = json.loads(bundle.read_text(encoding="utf-8"))
    payload["editable_page"]["sha256"] = _sha(page)
    payload["slide_fingerprint"] = inspect_editable_pptx(page).slide_fingerprints[0]
    _write_json(bundle, payload)
    _resign_bundle_as_project_key_attacker(project, bundle)
    with pytest.raises(ValueError, match="final|body|coverage|equivalent"):
        verify_signed_reconstruction(project, work, bundle)


@pytest.mark.parametrize(("policy", "should_pass"), [("reference_only", True), ("required_presence", False)])
def test_page_image_direct_appearance_follows_sealed_policy(tmp_path: Path, policy: str, should_pass: bool):
    project = tmp_path.resolve()
    work, _body = _fixture(project, page_image_policy=policy)
    manifest = _manifest(project, work)
    call = lambda: build_and_sign_reconstruction(project, work_item=work, manifest=manifest, body_pptx=project / "reconstruction/body.pptx", final_pptx=project / "reconstruction/page.pptx", bundle_output=project / "reconstruction/bundle.json")
    if should_pass:
        assert call().is_file()
    else:
        with pytest.raises(ValueError, match="required page image"):
            call()


def test_required_page_image_can_ship_as_a_provenanced_local_raster_component(tmp_path: Path):
    project = tmp_path.resolve()
    work, _body = _fixture(project, page_image_policy="required_presence")
    manifest = _manifest(project, work)
    payload = json.loads(manifest.read_text(encoding="utf-8"))
    image = project / "authorities/word-image.png"
    payload["images"] = [{"object_id": "photo-1", "name": "source-photo", "path": "../authorities/word-image.png", "box_px": [1200, 80, 380, 260]}]
    payload["raster_components"] = [{"object_id": "photo-1", "sha256": _sha(image), "source_type": "page-image", "source_id": "word-image-1"}]
    _write_json(manifest, payload)
    bundle = build_and_sign_reconstruction(project, work_item=work, manifest=manifest, body_pptx=project / "reconstruction/body.pptx", final_pptx=project / "reconstruction/page.pptx", bundle_output=project / "reconstruction/bundle.json")
    assert verify_signed_reconstruction(project, work, bundle)["object_counts"]["image"] == 1


def test_search_raster_is_reference_material_not_a_direct_gateway_embed(tmp_path: Path):
    project = tmp_path.resolve()
    work, _body = _fixture(project, search_raster=True)
    manifest = _manifest(project, work)
    payload = json.loads(manifest.read_text(encoding="utf-8"))
    image = project / "authorities/search-result.png"
    payload["images"] = [{"object_id": "search-chart", "name": "search-chart", "path": "../authorities/search-result.png", "box_px": [1200, 80, 380, 260]}]
    payload["raster_components"] = [{"object_id": "search-chart", "sha256": _sha(image), "source_type": "search-evidence", "source_id": "search-1"}]
    _write_json(manifest, payload)
    with pytest.raises(ValueError, match="search-evidence raster provenance"):
        build_and_sign_reconstruction(project, work_item=work, manifest=manifest, body_pptx=project / "reconstruction/body-2.pptx", final_pptx=project / "reconstruction/page-2.pptx", bundle_output=project / "reconstruction/bundle-2.json")


def test_two_half_body_rasters_cannot_reassemble_a_flattened_body(tmp_path: Path):
    project = tmp_path.resolve()
    work, _body = _fixture(project)
    manifest = _manifest(project, work)
    left = project / "reconstruction/left.png"
    right = project / "reconstruction/right.png"
    Image.new("RGB", (850, 800), "#334455").save(left)
    Image.new("RGB", (850, 800), "#556677").save(right)
    payload = json.loads(manifest.read_text(encoding="utf-8"))
    payload["images"] = [
        {"object_id": "half-l", "name": "half-left", "path": "left.png", "box_px": [0, 0, 850, 800]},
        {"object_id": "half-r", "name": "half-right", "path": "right.png", "box_px": [850, 0, 850, 800]},
    ]
    payload["raster_components"] = [
        {"object_id": "half-l", "sha256": _sha(left), "source_type": "decorative-texture", "source_id": "decorative:half-l"},
        {"object_id": "half-r", "sha256": _sha(right), "source_type": "decorative-texture", "source_id": "decorative:half-r"},
    ]
    _write_json(manifest, payload)
    with pytest.raises(ValueError, match="raster|coverage|flatten"):
        build_and_sign_reconstruction(project, work_item=work, manifest=manifest, body_pptx=project / "reconstruction/body.pptx", final_pptx=project / "reconstruction/page.pptx", bundle_output=project / "reconstruction/bundle.json")


def test_decorative_raster_cannot_impersonate_required_page_image(tmp_path: Path):
    project = tmp_path.resolve()
    work, _body = _fixture(project, page_image_policy="required_presence")
    manifest = _manifest(project, work)
    payload = json.loads(manifest.read_text(encoding="utf-8"))
    image = project / "authorities/word-image.png"
    payload["images"] = [{"object_id": "fake-1", "name": "fake-required", "path": "../authorities/word-image.png", "box_px": [1200, 80, 380, 260]}]
    payload["raster_components"] = [{"object_id": "fake-1", "sha256": _sha(image), "source_type": "decorative-texture", "source_id": "word-image-1"}]
    _write_json(manifest, payload)
    with pytest.raises(ValueError, match="required|decorative"):
        build_and_sign_reconstruction(project, work_item=work, manifest=manifest, body_pptx=project / "reconstruction/body.pptx", final_pptx=project / "reconstruction/page.pptx", bundle_output=project / "reconstruction/bundle.json")


@pytest.mark.parametrize("attack", ["one-pixel", "occluded"])
def test_tiny_or_occluded_word_text_is_rejected(tmp_path: Path, attack: str):
    project = tmp_path.resolve()
    work, _body = _fixture(project)
    manifest = _manifest(project, work)
    payload = json.loads(manifest.read_text(encoding="utf-8"))
    if attack == "one-pixel":
        payload["text_boxes"][0]["box_px"] = [100, 70, 1, 1]
    else:
        payload["shapes"].append({"object_id": "cover-1", "name": "opaque-cover", "type": "rect", "box_px": [100, 70, 700, 100], "fill": "#000000", "z_index": 999})
    _write_json(manifest, payload)
    with pytest.raises(ValueError, match="visible|small|occlu|capacity"):
        build_and_sign_reconstruction(project, work_item=work, manifest=manifest, body_pptx=project / "reconstruction/body.pptx", final_pptx=project / "reconstruction/page.pptx", bundle_output=project / "reconstruction/bundle.json")


def test_replaced_embedded_svg_logo_is_rejected_even_when_bundle_is_resigned(tmp_path: Path):
    project = tmp_path.resolve()
    work, _body = _fixture(project)
    manifest = _manifest(project, work)
    bundle = build_and_sign_reconstruction(project, work_item=work, manifest=manifest, body_pptx=project / "reconstruction/body.pptx", final_pptx=project / "reconstruction/page.pptx", bundle_output=project / "reconstruction/bundle.json")
    page = project / "reconstruction/page.pptx"
    replacement = page.with_suffix(".tmp.pptx")
    import zipfile
    with zipfile.ZipFile(page) as source, zipfile.ZipFile(replacement, "w", zipfile.ZIP_DEFLATED) as target:
        for item in source.infolist():
            data = source.read(item.filename)
            if item.filename.endswith(".svg"):
                data = b'<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 1 1"><rect width="1" height="1"/></svg>'
            target.writestr(item, data)
    replacement.replace(page)
    payload = json.loads(bundle.read_text(encoding="utf-8"))
    payload["editable_page"]["sha256"] = _sha(page)
    payload["slide_fingerprint"] = inspect_editable_pptx(page).slide_fingerprints[0]
    _write_json(bundle, payload)
    _resign_bundle_as_project_key_attacker(project, bundle)
    with pytest.raises(ValueError, match="logo|SVG"):
        verify_signed_reconstruction(project, work, bundle)


@pytest.mark.parametrize(("mutation", "message"), [
    (lambda item: item.update({"box_px": [1800, 70, 100, 50]}), "outside"),
    (lambda item: item.update({"font_size": 1}), "tiny"),
])
def test_invisible_or_superficial_text_cannot_satisfy_word_coverage(tmp_path: Path, mutation, message: str):
    project = tmp_path.resolve()
    work, _body = _fixture(project)
    manifest = _manifest(project, work)
    payload = json.loads(manifest.read_text(encoding="utf-8"))
    mutation(payload["text_boxes"][0])
    _write_json(manifest, payload)
    with pytest.raises(ValueError, match=message):
        build_and_sign_reconstruction(project, work_item=work, manifest=manifest, body_pptx=project / "reconstruction/body.pptx", final_pptx=project / "reconstruction/page.pptx", bundle_output=project / "reconstruction/bundle.json")
