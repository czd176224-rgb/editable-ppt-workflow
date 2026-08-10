from __future__ import annotations

import hashlib
import importlib.util
import io
import json
import sys
import time
from pathlib import Path
from types import SimpleNamespace

from docx import Document
from PIL import Image
from pptx import Presentation
import pytest

pytestmark = pytest.mark.skip(reason="retired V4 one-command orchestration; production entry is V6-only")


ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT / "scripts"))

import workflow_state  # noqa: E402
import v4_qa_gateway  # noqa: E402
import production_runner  # noqa: E402
import page_material_bundle_v4  # noqa: E402
import run_workflow  # noqa: E402
from codex_subscription_runtime import CodexStructuredResult  # noqa: E402
from codex_web_material_gateway import (  # noqa: E402
    DownloadResponse,
    search_visual_material as production_search_visual_material,
    verify_search_material as production_verify_search_material,
)
from production_runner import run_production  # noqa: E402
from test_production_runner import (  # noqa: E402
    _confirm,
    _project,
    _raw_transport,
    _successful_backend,
)
from test_final_mechanical_assembly import _render_with_real_open  # noqa: E402
import v4_reconstruction_gateway  # noqa: E402


CONFIRM_UI_SERVER = ROOT / "scripts" / "confirm_ui" / "server.py"


def _four_page_word(path: Path) -> tuple[Path, dict[int, str], dict[int, list[str]]]:
    document = Document()
    titles = {
        1: "并购项目总体说明",
        2: "交易执行全流程",
        3: "浙江凤凰行动并购生态",
        4: "后续工作安排",
    }
    bodies = {
        1: "项目首先核验基础材料，其次形成交易方案，最后完成内部决策。",
        2: "本页是一段不拆分的完整长正文，" + "、".join(
            f"环节{i}必须保留Word原文事实并形成连续视觉叙事" for i in range(1, 31)
        ) + "，任何标题识别都不得删减或重写这段原文。",
        3: "浙江凤凰行动持续完善并购生态，新闻发布材料用于说明区域产业整合进展。",
        4: "下一阶段将完成资料复核、方案评审、交割准备与成果归档。",
    }
    comments = {1: [], 2: ["文字表达图片化"], 3: ["新闻稿图片"], 4: []}
    for page in range(1, 5):
        document.add_paragraph(f"第{page}页")
        document.add_paragraph(titles[page])
        body = document.add_paragraph(bodies[page])
        for comment in comments[page]:
            document.add_comment(body.runs, text=comment, author="sanitized-reviewer")
        if page < 4:
            document.add_page_break()
    document.save(path)
    source_text = {page: f"{titles[page]}\n\n{bodies[page]}" for page in range(1, 5)}
    return path, source_text, comments


def _four_page_logo(path: Path) -> Path:
    path.write_text(
        '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 120 48">'
        '<rect width="120" height="48" fill="#17365D"/></svg>',
        encoding="utf-8",
    )
    return path


def _load_confirm_ui():
    spec = importlib.util.spec_from_file_location("four_page_acceptance_confirm_ui", CONFIRM_UI_SERVER)
    assert spec and spec.loader
    module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(module)
    return module


def _submit_first_confirmation(project: Path) -> None:
    server = _load_confirm_ui()
    client = server.create_app(project).test_client()
    recommendations = client.get("/api/recommendations").get_json()
    selected = recommendations["design_directions"]["selected"]
    candidate = recommendations["design_directions"]["candidates"][selected]
    payload = {
        "stage": "final", "direction": selected,
        "template_selection": candidate["template_selection"], "canvas": "ppt169",
        **{key: candidate[key] for key in (
            "visual_style", "color", "icons", "typography", "image_rendering", "style_axes",
            "layout_preferences", "information_density", "background_system", "image_role",
            "evidence_strength", "composition_tendency", "brand_device",
        )},
        "regional_style": {"enabled": False}, "production_profile": "balanced",
        "additional_requirements": "保持Word原文、分页批注和逐页材料边界",
    }
    response = client.post("/api/confirm", json=payload)
    assert response.status_code == 200, response.get_json()


class _MockMaterialTransport:
    def __init__(self, payload: bytes) -> None:
        self.payload = payload

    def resolve(self, _hostname: str, _port: int):
        return ["93.184.216.34"]

    def get(self, _url: str, *, connect_ip: str, timeout: float, max_bytes: int):
        assert connect_ip == "93.184.216.34" and timeout > 0 and max_bytes >= len(self.payload)
        return DownloadResponse(200, {"content-type": "image/png"}, self.payload)


def test_sanitized_four_page_subscription_acceptance_is_one_command_and_closed(
    tmp_path: Path, monkeypatch,
) -> None:
    """The release fixture proves all page authorities without network or an API key."""
    word, expected_source_text, expected_comments = _four_page_word(tmp_path / "four-page-source.docx")
    logo = _four_page_logo(tmp_path / "logo.svg")
    project = tmp_path / "project"
    events: list[str] = []
    image_stream = io.BytesIO()
    Image.new("RGB", (1200, 800), "#2F6690").save(image_stream, format="PNG")
    image_payload = image_stream.getvalue()
    transport = _MockMaterialTransport(image_payload)
    produced_materials: list[dict] = []
    search_model_calls = 0

    def invoke_search(*_args, **_kwargs):
        nonlocal search_model_calls
        search_model_calls += 1
        usage = {"inputTokens": 10, "outputTokens": 20}
        return CodexStructuredResult(
            value={"candidates": [{
                "source_page_url": "https://news.example/zhejiang-release",
                "direct_image_url": "https://cdn.example/zhejiang-news.png",
                "title": "Zhejiang news release", "publisher": "Example Newsroom",
                "caption": "Sanitized Zhejiang news release image.",
                "matched_entities": ["浙江凤凰行动"], "retrieved_at": "2026-08-01T00:00:00Z",
            }]},
            thread_id="thread-search-fixture", turn_id="turn-search-fixture",
            model="gpt-test", model_provider="openai", auth_mode="chatgpt", plan_type="plus",
            usage=usage, safe_trace={
                "runtime": "codex-app-server", "role": "visual-material-search",
                "thread_id": "thread-search-fixture", "turn_id": "turn-search-fixture",
                "model": "gpt-test", "model_provider": "openai", "auth_mode": "chatgpt",
                "plan_type": "plus", "usage": usage, "image_count": 0, "web_search": "live",
            },
        )

    def searched(search_project, *, directives, page_context, timeout, budget, **_kwargs):
        events.append("search-page-3")
        outcomes = []
        for directive in directives:
            materials = production_search_visual_material(
                search_project,
                directive=directive,
                page_context=page_context,
                timeout=timeout,
                deadline=time.monotonic() + timeout,
                budget=budget,
                invoke=invoke_search,
                transport=transport,
            )
            if not produced_materials:
                produced_materials.extend(materials)
            else:
                assert materials == produced_materials
            outcomes.append(materials)
        return outcomes

    monkeypatch.setattr(page_material_bundle_v4, "search_visual_materials", searched)
    ui_actions: list[str] = []

    def confirm_ui_entry(command: list[str], **_kwargs):
        action = command[2]
        ui_actions.append(action)
        if action == "wait":
            _submit_first_confirmation(project)
        return SimpleNamespace(returncode=0, stdout="ok", stderr="")

    monkeypatch.setattr(run_workflow.subprocess, "run", confirm_ui_entry)
    first = run_workflow.run(
        word, logo, project, wait_ui=True, no_browser=True, execute=False,
    )
    assert first["stage"] != "awaiting_confirmation"
    assert ui_actions == ["start", "wait", "shutdown"]

    state_after_confirm = workflow_state.load(project)
    contracts = [
        json.loads((project / job["contract_file"]).read_text(encoding="utf-8"))
        for job in state_after_confirm["jobs"]
    ]
    assert [contract["source_text"] for contract in contracts] == [
        expected_source_text[page] for page in range(1, 5)
    ]
    assert [[comment["text"] for comment in contract["page_comments"]] for contract in contracts] == [
        expected_comments[page] for page in range(1, 5)
    ]
    contract_bytes = {
        job["contract_file"]: (project / job["contract_file"]).read_bytes()
        for job in state_after_confirm["jobs"]
    }
    source_lock = (project / "01_page_contracts/source_lock.json").read_bytes()
    style_confirmation = (project / "02_style/style_confirmation.json").read_bytes()
    style_execution = (project / "02_style/style_execution.json").read_bytes()
    run_workflow.run(word, logo, project, no_browser=True, execute=False)
    run_workflow.run(word, logo, project, no_browser=True, execute=False)
    assert ui_actions == ["start", "wait", "shutdown"]
    assert all((project / path).read_bytes() == payload for path, payload in contract_bytes.items())
    assert (project / "01_page_contracts/source_lock.json").read_bytes() == source_lock
    assert (project / "02_style/style_confirmation.json").read_bytes() == style_confirmation
    assert (project / "02_style/style_execution.json").read_bytes() == style_execution

    auth = tmp_path / "codex-auth.json"
    auth.write_text(json.dumps({
        "tokens": {"access_token": "fixture-access-token", "account_id": "fixture-account"},
        "last_refresh": "2026-08-01T00:00:00Z",
    }), encoding="utf-8")
    monkeypatch.setenv("CODEX_AUTH_FILE", str(auth))
    monkeypatch.delenv("OPENAI_API_KEY", raising=False)
    _authorize_raw_transport(monkeypatch)

    generated_commands: list[list[str]] = []
    backend = _successful_backend(generated_commands)

    def image_backend(command, **kwargs):
        output = Path(command[command.index("--out") + 1])
        events.append(f"image2-{output.stem}")
        return backend(command, **kwargs)

    monkeypatch.setattr("batch_generation.subprocess.run", image_backend)

    monkeypatch.setattr("final_mechanical_qa._default_renderer", _render_with_real_open)

    result = run_production(project, finalize=True)

    assert result["stage"] == "complete", json.dumps(result, ensure_ascii=False, indent=2)
    assert events.index("search-page-3") < next(
        index for index, item in enumerate(events) if item.startswith("image2-page_003_")
    )
    assert len(generated_commands) == 4
    for command in generated_commands:
        assert command[command.index("--size") + 1] == "1904x896"

    run = workflow_state.load(project)
    page2_job = run["jobs"][1]
    page2_bundle = json.loads((project / page2_job["material_bundle_file"]).read_text(encoding="utf-8"))
    assert page2_bundle["source_text"] == expected_source_text[2]
    assert page2_bundle["authoritative_content"]["body_text"] in expected_source_text[2]
    assert [item["text"] for item in page2_bundle["resolved_directives"]] == expected_comments[2]
    page3_bundle = json.loads((project / run["jobs"][2]["material_bundle_file"]).read_text(encoding="utf-8"))
    assert page3_bundle["source_text"] == expected_source_text[3]
    assert [item["text"] for item in page3_bundle["resolved_directives"]] == expected_comments[3]
    assert page3_bundle["search_evidence"][0]["source_url"] == "https://news.example/zhejiang-release"
    assert len(produced_materials) == 1
    assert search_model_calls == 1
    produced = produced_materials[0]
    verified = production_verify_search_material(
        project, produced,
        expected_material_id=produced["material_id"],
        expected_directive_id=produced["directive_id"],
        expected_query=produced["query"], deadline=time.monotonic() + 10,
    )
    assert verified["sha256"] == hashlib.sha256(image_payload).hexdigest()
    assert (project / produced["material_attestation_path"]).is_file()
    assert produced["material_attestation_sha256"] not in {"1" * 64, "2" * 64, "3" * 64}

    for job in run["jobs"]:
        generation = json.loads((project / job["generation_receipt"]["path"]).read_text(encoding="utf-8"))
        qa_work = json.loads((project / job["qa_work_item"]["path"]).read_text(encoding="utf-8"))
        assert generation["required_directive_ids"] == [
            item["directive_id"] for item in qa_work["required_directives"]
        ]
        assert generation["effective_authority_sha256"] == qa_work["effective_page_authority_sha256"]
        signed_invocation = json.loads(
            (project / "07_editable" / f"page_{job['page_number']:03d}" / "signed-reconstruction-invocation.json").read_text(encoding="utf-8")
        )
        assert signed_invocation["attestation"]["source_origin"] == "sealed_work_item"
        assert job["reconstruction_calls"] == 1
        assert (project / job["page_package"]).is_file()

    deck = Presentation(result["output"])
    assert len(deck.slides) == 4
    assert [job["page_number"] for job in run["jobs"]] == [1, 2, 3, 4]


def _authorize_raw_transport(monkeypatch, mode: str | None = None) -> None:
    fixture = Path(__file__).resolve().parent / "raw_transport"
    current = __import__("os").environ.get("PYTHONPATH", "")
    monkeypatch.setenv("PYTHONPATH", str(fixture) + (__import__("os").pathsep + current if current else ""))
    monkeypatch.setenv("EDITABLE_PPT_RAW_TRANSPORT_FIXTURE", "1")
    _raw_transport(monkeypatch, mode)


def _accepted_project(tmp_path: Path, monkeypatch, *, page_count: int = 1) -> Path:
    project = _project(tmp_path, page_count=page_count)
    _confirm(project)
    _authorize_raw_transport(monkeypatch)
    auth = tmp_path / "codex-auth.json"
    auth.write_text(json.dumps({"tokens": {"access_token": "fixture-access-token", "account_id": "fixture-account"}, "last_refresh": "2026-08-01T00:00:00Z"}), encoding="utf-8")
    monkeypatch.setenv("CODEX_AUTH_FILE", str(auth))
    _authorize_raw_transport(monkeypatch, "reconstruction_error")
    pending = run_production(project, finalize=False)
    assert pending["stage"] in {"reconstruction_backend_pending", "page_pipeline"}, pending.get("provider_error")
    assert pending.get("reconstruction_work_items"), pending
    for job in workflow_state.load(project)["jobs"]:
        generation_receipt = project / job["generation_receipt"]["path"]
        assert generation_receipt.is_file()
        generation = json.loads(generation_receipt.read_text(encoding="utf-8"))
        trace = project / generation["provider_trace"]["path"]
        assert trace.is_file()
        assert json.loads(trace.read_text(encoding="utf-8"))["auth"] == "codex_oauth"
        assert (project / job["qa_receipt"]["path"]).is_file()
        assert (project / job["qa_work_item"]["path"]).is_file()
    return project


def _provider_response(payload: dict) -> bytes:
    instructions = json.loads(payload["input"][0]["content"][0]["text"])
    text_boxes = []
    text_coverage = []
    for index, item in enumerate(instructions["authoritative_text"], start=1):
        name = f"body-text-{index:03d}"
        text_boxes.append({
            "object_id": item["source_id"], "name": name, "text": item["text"],
            "box_px": [100, 80 + 150 * (index - 1), 1500, 110], "font_size": 16,
            "font": "Microsoft YaHei", "color": "#111827", "bold": False, "italic": False,
            "align": "left", "valign": "middle", "wrap": True, "fit_text": True, "z_index": 10 + index,
        })
        text_coverage.append({"source_id": item["source_id"], "text": item["text"], "object_name": name})
    decision = {
        "text_boxes": text_boxes, "tables": [],
        "shapes": [{"object_id": "visual-panel", "name": "visual-panel", "type": "rect", "box_px": [20, 20, 1860, 850], "fill": "#E8EEF4", "stroke": "#E8EEF4", "stroke_width": 0, "z_index": 0}],
        "images": [], "text_coverage": text_coverage, "table_coverage": [],
    }
    return json.dumps({"id": "resp_reconstruct_fixture", "output": [{"content": [{"type": "output_text", "text": json.dumps(decision)}]}]}).encode()


def _authorize_gateway(monkeypatch, captured: list[dict] | None = None) -> None:
    del captured
    _authorize_raw_transport(monkeypatch)


def test_visual_gateway_request_contains_accepted_pixels_hash_and_layout_authority(tmp_path: Path, monkeypatch) -> None:
    project = _accepted_project(tmp_path, monkeypatch)
    action = workflow_state.next_action(project)
    work = project / action["reconstruction_work_items"][0]["path"]
    _authorize_gateway(monkeypatch)
    model_calls: list[dict] = []
    invoke = v4_reconstruction_gateway._invoke_structured

    def provider_without_source(*args, **kwargs):
        result = invoke(*args, **kwargs)
        assert "source" not in result.value
        model_calls.append(dict(result.value))
        return result

    monkeypatch.setattr(v4_reconstruction_gateway, "_invoke_structured", provider_without_source)

    bundle = v4_reconstruction_gateway.invoke_builtin_gateway(project, work, timeout=10)
    manifest = v4_reconstruction_gateway.verify_signed_bundle(project, bundle, work)

    signed = json.loads(bundle.read_text(encoding="utf-8"))
    request = json.loads((project / signed["request"]["path"]).read_text(encoding="utf-8"))
    content = request["input"][0]["content"]
    instructions = json.loads(content[0]["text"])
    assert instructions["body_image_sha256"] == json.loads(work.read_text(encoding="utf-8"))["accepted_body_image"]["sha256"]
    assert "determines composition" in " ".join(instructions["hard_rules"])
    assert content[2]["image_url"].startswith("data:image/png;base64,")
    final_manifest = json.loads(manifest.read_text(encoding="utf-8"))
    assert final_manifest["text_boxes"]
    assert final_manifest["source"] == {"width_px": 1904, "height_px": 896}
    assert signed["attestation"]["source_origin"] == "sealed_work_item"
    assert len(model_calls) == 1


def test_missing_reconstruction_codex_runtime_stays_pending_without_generic_deck(tmp_path: Path, monkeypatch) -> None:
    project = _accepted_project(tmp_path, monkeypatch)
    monkeypatch.setenv("EDITABLE_PPT_CODEX_EXECUTABLE", "missing-codex-test-executable")
    monkeypatch.setattr(production_runner, "invoke_reconstruction_gateway_worker", v4_reconstruction_gateway.invoke_gateway_worker)

    result = run_production(project, finalize=False)

    assert result["stage"] == "reconstruction_backend_pending"
    assert "Codex App Server" in result["provider_error"]
    assert not (project / "07_editable/page_001/page.pptx").exists()


def test_configured_run_reconstructs_assembles_and_resumes_without_repeating_pages(tmp_path: Path, monkeypatch) -> None:
    project = _accepted_project(tmp_path, monkeypatch, page_count=2)
    _authorize_gateway(monkeypatch)

    first = run_production(project, finalize=True)
    second = run_production(project, finalize=True)

    assert first["stage"] == second["stage"] == "complete"
    assert Path(first["output"]).is_file()
    deck = Presentation(first["output"])
    assert len(deck.slides) == 2
    assert [job["reconstruction_calls"] for job in workflow_state.load(project)["jobs"]] == [1, 1]
    for job in workflow_state.load(project)["jobs"]:
        page = int(job["page_number"])
        page_dir = project / "07_editable" / f"page_{page:03d}"
        assert (page_dir / "signed-reconstruction-invocation.json").is_file()
        assert (page_dir / "object-manifest.json").is_file()
        assert (page_dir / "signed-reconstruction.json").is_file()
        assert (page_dir / "editable-receipt.json").is_file()
        assert (project / job["page_package"]).is_file()
    for slide in deck.slides:
        names = [shape.name for shape in slide.shapes]
        assert names.count("fixed-frame-title") == 1
        assert names.count("fixed-frame-logo") == 1
        assert names.count("fixed-frame-footer") == 1
        assert names.count("fixed-frame-page-number") == 1
        assert not any(shape.shape_type == 13 and shape.width > 8_000_000 for shape in slide.shapes)
    gateway = project / "07_editable/page_001/signed-reconstruction-invocation.json"
    gateway.write_bytes(gateway.read_bytes() + b" ")
    assert run_production(project, finalize=True)["stage"] == "authority_changed"


def test_finalize_false_completes_pages_but_leaves_ordered_assembly_pending(tmp_path: Path, monkeypatch) -> None:
    project = _accepted_project(tmp_path, monkeypatch)
    _authorize_gateway(monkeypatch)

    result = run_production(project, finalize=False)

    assert result["stage"] == "assembly_pending"
    assert workflow_state.load(project)["jobs"][0]["status"] == "complete"
    assert workflow_state.load(project)["final_pptx"] is None


def test_default_renderer_produces_proof_pages_without_reconstructing_pages(tmp_path: Path, monkeypatch) -> None:
    project = _accepted_project(tmp_path, monkeypatch)
    _authorize_gateway(monkeypatch)
    monkeypatch.setattr("final_mechanical_qa._default_renderer", _render_with_real_open)

    result = run_production(project, finalize=True)
    report = json.loads((project / "08_final/final_mechanical_qa.json").read_text(encoding="utf-8"))
    assert result["stage"] == "complete"
    assert report["open_render_status"]["rendered"] is True
    assert report["open_render_status"]["rendered_page_count"] == 1
    assert workflow_state.load(project)["jobs"][0]["reconstruction_calls"] == 1
