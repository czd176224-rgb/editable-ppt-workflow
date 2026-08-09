from __future__ import annotations

import hashlib
import json
import sys
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches


ROOT = Path(__file__).resolve().parents[1]
SCRIPTS = ROOT / "scripts"
sys.path.insert(0, str(SCRIPTS))

import workflow_v5_migration as migration_module  # noqa: E402
from workflow_v5_migration import migrate_v4_project  # noqa: E402
from workflow_v5_report import build_runtime_report  # noqa: E402
from workflow_v5_request_ledger import RequestLedger  # noqa: E402
import workflow_state  # noqa: E402
from current_contract_fixture import write_valid_qa_observation  # noqa: E402
from test_v4_lightweight_qa import _generated_project, _work_item  # noqa: E402


def _write(path: Path, data: bytes) -> str:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_bytes(data)
    return hashlib.sha256(data).hexdigest()


def _write_png(path: Path, size: tuple[int, int] = (1904, 896)) -> str:
    path.parent.mkdir(parents=True, exist_ok=True)
    Image.new("RGB", size, "white").save(path, format="PNG")
    return hashlib.sha256(path.read_bytes()).hexdigest()


def _project(tmp_path: Path, *, authentic_material: bool) -> Path:
    project = tmp_path / "stable-project-name"
    word_sha = _write(project / "00_source/source.docx", b"word source")
    logo_sha = _write(project / "00_source/logo.svg", b"<svg/>")
    style_sha = _write(project / "02_style/style_execution.json", b'{"style":"formal"}')
    source_text = "第一页\n权威正文"
    contract = {
        "page_number": 1,
        "source_text": source_text,
        "source_hash": hashlib.sha256(source_text.encode()).hexdigest(),
    }
    contract_path = project / "01_page_contracts/page_001.json"
    contract_path.parent.mkdir(parents=True)
    contract_path.write_text(json.dumps(contract, ensure_ascii=False), encoding="utf-8")
    image = project / "06_images/generated/page_001.png"
    _write_png(image)
    editable = project / "07_editable_pages/page_001.pptx"
    editable.parent.mkdir(parents=True)
    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1)).text = "权威正文"
    deck.save(editable)
    job = {
        "page_number": 1,
        "status": "complete",
        "contract_file": "01_page_contracts/page_001.json",
        "generation": {"body_image": {"path": "06_images/generated/page_001.png"}},
        "editable_page": {"path": "07_editable_pages/page_001.pptx"},
        "material_ids": ["real-news-photo"] if authentic_material else [],
    }
    state = {
        "workflow_contract_version": "word-ppt-workflow-v4",
        "word_source": {"path": "00_source/source.docx", "sha256": word_sha},
        "logo_source": {"path": "00_source/logo.svg", "sha256": logo_sha},
        "style_confirmation": {
            "status": "confirmed", "execution_file": "02_style/style_execution.json",
            "execution_sha256": style_sha,
        },
        "pagination": {"page_count": 1, "locked_page_order": [1]},
        "jobs": [job],
    }
    (project / "workflow_run.json").write_text(json.dumps(state), encoding="utf-8")
    return project


def _node(dag: dict, node_id: str) -> dict:
    return next(item for item in dag["nodes"] if item["node_id"] == node_id)


def _attach_required_search_evidence(
    project: Path, *, acquisitions: list[tuple[str, int]], comment: str,
) -> None:
    contract_path = project / "01_page_contracts/page_001.json"
    contract = json.loads(contract_path.read_text(encoding="utf-8"))
    contract.update({
        "page_title": "材料迁移测试",
        "page_comments": [{"comment_id": "1", "text": comment}],
        "asset_bindings": [],
    })
    contract_path.write_text(json.dumps(contract, ensure_ascii=False), encoding="utf-8")
    directives = []
    evidence = []
    for material_index, (material_id, count) in enumerate(acquisitions, start=1):
        entity = f"企业{material_index}" if len(acquisitions) > 1 else ""
        directives.append({
            "directive_id": f"directive-{material_index}",
            "material_id": material_id,
            "action": "require",
            "entity": entity,
            "material_role": "enterprise_logo" if entity else "authentic_image",
        })
        for asset_index in range(1, count + 1):
            path = project / "03_evidence/page_001/search" / f"{material_id}-{asset_index}.png"
            path.parent.mkdir(parents=True, exist_ok=True)
            Image.new("RGB", (80, 50), (material_index * 20, asset_index * 40, 120)).save(path)
            digest = hashlib.sha256(path.read_bytes()).hexdigest()
            evidence.append({
                "asset_id": material_id,
                "evidence_id": f"evidence-{material_index}-{asset_index}",
                "local_path": path.relative_to(project).as_posix(),
                "sha256": digest,
                "media_type": "image/png",
                "entity": entity,
                "material_role": "enterprise_logo" if entity else "authentic_image",
                "query": f"{entity or '会议'} 官方图片",
                "source_url": f"https://example.test/{material_id}/{asset_index}",
                "publisher": "权威来源",
            })
    bundle = project / "04_v4/material/page_001.json"
    bundle.parent.mkdir(parents=True, exist_ok=True)
    bundle.write_text(json.dumps({
        "required_directives": directives,
        "search_evidence": evidence,
        "page_images": [],
        "attachment_evidence": [],
    }, ensure_ascii=False), encoding="utf-8")
    state_path = project / "workflow_run.json"
    state = json.loads(state_path.read_text(encoding="utf-8"))
    state["jobs"][0]["material_bundle_file"] = bundle.relative_to(project).as_posix()
    state_path.write_text(json.dumps(state), encoding="utf-8")


def _attach_passing_qa_stub(project: Path) -> Path:
    image = project / "06_images/generated/page_001.png"
    work = project / "04_v4/qa/page_001.work.json"
    work.parent.mkdir(parents=True, exist_ok=True)
    work.write_text(json.dumps({
        "body_image": {
            "path": image.relative_to(project).as_posix(),
            "sha256": hashlib.sha256(image.read_bytes()).hexdigest(),
        },
    }), encoding="utf-8")
    receipt = project / "04_v4/qa/page_001.qa.json"
    receipt.write_text("{}", encoding="utf-8")
    state_path = project / "workflow_run.json"
    state = json.loads(state_path.read_text(encoding="utf-8"))
    state["jobs"][0]["qa_receipt"] = {
        "path": receipt.relative_to(project).as_posix(),
        "sha256": hashlib.sha256(receipt.read_bytes()).hexdigest(),
    }
    state_path.write_text(json.dumps(state), encoding="utf-8")
    return work


def _stub_valid_passing_qa(monkeypatch, project: Path, work: Path) -> list[Path]:
    calls: list[Path] = []

    def validate(root: Path, receipt: Path) -> dict:
        calls.append(receipt)
        return {
            "artifact": {
                "status": "pass", "page_number": 1,
                "qa_work_item": {
                    "path": work.relative_to(project).as_posix(),
                    "sha256": hashlib.sha256(work.read_bytes()).hexdigest(),
                },
            },
            "path": receipt,
            "sha256": hashlib.sha256(receipt.read_bytes()).hexdigest(),
        }

    monkeypatch.setattr(migration_module, "validate_historical_qa_receipt", validate)
    return calls


def test_migration_keeps_unreviewed_legacy_image2_pending_even_when_png_is_valid(tmp_path: Path) -> None:
    project = _project(tmp_path, authentic_material=False)
    report = migrate_v4_project(project)
    dag = json.loads((project / "04_v5/dag.json").read_text(encoding="utf-8"))

    assert project.name == "stable-project-name"
    assert _node(dag, "page:001:design")["status"] == "pending"
    assert _node(dag, "page:001:compose")["status"] == "pending"
    assert _node(dag, "page:001:reconstruct")["status"] == "pending"
    assert _node(dag, "page:001:page_validate")["status"] == "pending"
    assert _node(dag, "page:001:visual_qa")["status"] == "pending"
    assert report["successful_model_results_reused"] == 0
    assert report["pages"][0]["invalidation_reason"] == "historical_passing_qa_receipt_missing_or_invalid"
    assert str(tmp_path) not in (project / "09_reports/v5_compatibility_report.json").read_text(encoding="utf-8")


def test_migration_reuses_only_image2_closed_by_valid_passing_historical_qa(
    tmp_path: Path, monkeypatch,
) -> None:
    project = _project(tmp_path, authentic_material=False)
    work = _attach_passing_qa_stub(project)
    calls = _stub_valid_passing_qa(monkeypatch, project, work)

    report = migrate_v4_project(project)
    dag = json.loads((project / "04_v5/dag.json").read_text(encoding="utf-8"))

    assert len(calls) == 1
    assert _node(dag, "page:001:design")["status"] == "complete"
    assert _node(dag, "page:001:compose")["status"] == "pending"
    assert _node(dag, "page:001:reconstruct")["status"] == "pending"
    assert report["successful_model_results_reused"] == 1


def test_migration_keeps_material_design_pending_even_with_historical_raw_qa(
    tmp_path: Path, monkeypatch,
) -> None:
    project = _project(tmp_path, authentic_material=True)
    work = _attach_passing_qa_stub(project)
    _stub_valid_passing_qa(monkeypatch, project, work)

    report = migrate_v4_project(project)
    dag = json.loads((project / "04_v5/dag.json").read_text(encoding="utf-8"))

    assert _node(dag, "page:001:design")["status"] == "pending"
    assert _node(dag, "page:001:compose")["status"] == "pending"
    assert report["successful_model_results_reused"] == 0


def test_migration_accepts_a_real_cryptographically_valid_passing_qa_closure(
    tmp_path: Path,
) -> None:
    project, attempt = _generated_project(tmp_path, size=(1904, 896))
    _work_item(project)
    invocation = write_valid_qa_observation(project, 1)
    workflow_state.record_qa(
        project, 1, "qa-worker", attempt, signed_invocation_bundle=invocation,
    )

    report = migrate_v4_project(project)
    dag = json.loads((project / "04_v5/dag.json").read_text(encoding="utf-8"))

    assert _node(dag, "page:001:design")["status"] == "complete"
    assert report["pages"][0]["design_reused"] is True


def test_migration_keeps_tampered_or_nonpassing_historical_qa_pending(
    tmp_path: Path, monkeypatch,
) -> None:
    for outcome in ("tampered", "repair", "technical_blocked"):
        project = _project(tmp_path / outcome, authentic_material=False)
        work = _attach_passing_qa_stub(project)
        if outcome == "tampered":
            monkeypatch.setattr(
                migration_module, "validate_historical_qa_receipt",
                lambda *_args, **_kwargs: (_ for _ in ()).throw(ValueError("tampered closure")),
            )
        else:
            _stub_valid_passing_qa(monkeypatch, project, work)
            state_path = project / "workflow_run.json"
            state = json.loads(state_path.read_text(encoding="utf-8"))
            if outcome == "repair":
                receipt = project / state["jobs"][0]["qa_receipt"]["path"]
                monkeypatch.setattr(
                    migration_module, "validate_historical_qa_receipt",
                    lambda _root, path: {
                        "artifact": {"status": "repair", "page_number": 1},
                        "path": path, "sha256": hashlib.sha256(path.read_bytes()).hexdigest(),
                    },
                )
            else:
                state["jobs"][0]["status"] = "technical_blocked"
                state_path.write_text(json.dumps(state), encoding="utf-8")

        report = migrate_v4_project(project)
        dag = json.loads((project / "04_v5/dag.json").read_text(encoding="utf-8"))
        assert _node(dag, "page:001:design")["status"] == "pending", outcome
        assert _node(dag, "page:001:reconstruct")["status"] == "pending", outcome
        assert report["successful_model_results_reused"] == 0, outcome


def test_migration_invalidates_only_changed_authentic_custody_descendants(tmp_path: Path) -> None:
    project = _project(tmp_path, authentic_material=True)
    # Authentic-material pages also require a passing historical QA closure
    # before their Image2 design may be retained.
    report = migrate_v4_project(project)
    dag = json.loads((project / "04_v5/dag.json").read_text(encoding="utf-8"))

    assert _node(dag, "page:001:design")["status"] == "pending"
    assert _node(dag, "page:001:compose")["status"] == "pending"
    assert _node(dag, "page:001:reconstruct")["status"] == "pending"
    assert report["pages"][0]["invalidation_reason"] == "historical_passing_qa_receipt_missing_or_invalid"


def test_migration_does_not_reuse_invalid_editable_package(tmp_path: Path) -> None:
    project = _project(tmp_path, authentic_material=False)
    (project / "07_editable_pages/page_001.pptx").write_bytes(b"not a pptx")

    report = migrate_v4_project(project)
    dag = json.loads((project / "04_v5/dag.json").read_text(encoding="utf-8"))

    assert _node(dag, "page:001:reconstruct")["status"] == "pending"
    assert _node(dag, "page:001:page_validate")["status"] == "pending"
    assert report["successful_model_results_reused"] == 0


def test_migration_rejects_wrong_size_legacy_design_before_dag_reuse(tmp_path: Path) -> None:
    project = _project(tmp_path, authentic_material=False)
    _write_png(project / "06_images/generated/page_001.png", (1536, 1024))

    report = migrate_v4_project(project)
    dag = json.loads((project / "04_v5/dag.json").read_text(encoding="utf-8"))

    assert _node(dag, "page:001:design")["status"] == "pending"
    assert _node(dag, "page:001:compose")["status"] == "pending"
    assert _node(dag, "page:001:reconstruct")["status"] == "pending"
    assert report["successful_model_results_reused"] == 0
    assert report["pages"][0]["invalidation_reason"] == (
        "legacy_design_missing_or_failed_deterministic_preflight"
    )


def test_migration_is_byte_idempotent_and_runtime_report_uses_ledger_facts(tmp_path: Path) -> None:
    project = _project(tmp_path, authentic_material=False)
    migrate_v4_project(project)
    tracked = [
        project / "04_v5/dag.json",
        project / "04_v5/content-catalog.json",
        project / "09_reports/v5_compatibility_report.json",
    ]
    before = {path: path.read_bytes() for path in tracked}
    migrate_v4_project(project)

    ledger = RequestLedger(project)
    claim = ledger.claim(
        "image2_design", {"page": 1, "authority": "page-one"}, worker_id="worker",
    )
    ledger.complete_success(
        claim["request_key"], worker_id="worker",
        result={"model": "gpt-image-2", "strength": "high", "cached": False},
    )
    runtime = build_runtime_report(project)

    assert {path: path.read_bytes() for path in tracked} == before
    assert runtime["external_calls"]["unique_requests"] == 1
    assert runtime["external_calls"]["successful"] == 1
    assert runtime["models"] == [{
        "purpose": "image2_design", "model": "gpt-image-2", "strength": "high",
        "auth_mode": None, "provider_backend_calls": 1,
    }]
    assert runtime["catalog"]["hash_operations"] > 0


def test_migration_reuses_one_search_need_with_all_three_local_assets(tmp_path: Path) -> None:
    project = _project(tmp_path, authentic_material=False)
    _attach_required_search_evidence(
        project, acquisitions=[("legacy-news-search", 3)], comment="新闻稿图片",
    )

    report = migrate_v4_project(project)
    dag = json.loads((project / "04_v5/dag.json").read_text(encoding="utf-8"))

    assert len(report["published_material_manifests"]) == 1
    published = report["published_material_manifests"][0]
    assert published["asset_count"] == 3
    assert published["new_search_performed"] is False
    manifest = json.loads(
        (project / "04_v5/materials" / f"{published['material_id']}.json")
        .read_text(encoding="utf-8")
    )
    assert len(manifest["assets"]) == 3
    assert manifest["new_search_performed"] is False
    assert _node(dag, f"material:{published['material_id']}")["status"] == "complete"


def test_migration_reuses_six_resolved_logo_needs_without_search_or_download(tmp_path: Path) -> None:
    project = _project(tmp_path, authentic_material=False)
    acquisitions = [(f"logo-{index}", 1) for index in range(1, 7)]
    _attach_required_search_evidence(
        project, acquisitions=acquisitions, comment="这页的企业Logo都要添加",
    )

    report = migrate_v4_project(project)
    dag = json.loads((project / "04_v5/dag.json").read_text(encoding="utf-8"))

    assert len(report["published_material_manifests"]) == 6
    assert sum(item["asset_count"] for item in report["published_material_manifests"]) == 6
    assert all(item["new_search_performed"] is False for item in report["published_material_manifests"])
    assert all(_node(dag, f"material:logo-{index}")["status"] == "complete" for index in range(1, 7))
