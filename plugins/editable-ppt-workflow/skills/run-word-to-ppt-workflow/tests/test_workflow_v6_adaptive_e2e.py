from __future__ import annotations

import hashlib
import importlib.util
import json
import sys
from collections import defaultdict
from pathlib import Path

from docx import Document
from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.util import Cm


ROOT = Path(__file__).resolve().parents[1]
SCRIPTS = ROOT / "scripts"
SERVER = SCRIPTS / "confirm_ui" / "server.py"
FIXTURE = Path(__file__).parent / "fixtures" / "v6_adaptive_project" / "fixture.json"
if str(SCRIPTS) not in sys.path:
    sys.path.insert(0, str(SCRIPTS))

import workflow_v6_source  # noqa: E402
from workflow_v6_contract import canonical_sha256  # noqa: E402
from workflow_v6_image import generate_page_body  # noqa: E402
from workflow_v6_reconstruction import (  # noqa: E402
    assemble_v6_deck,
    build_reconstruction_request,
    finalize_reconstructed_page,
)
from workflow_v6_source import fail_reference, import_reference, confirm_reference, initialize_v6_project  # noqa: E402
from workflow_v6_state import load, save  # noqa: E402


def _load_server():
    spec = importlib.util.spec_from_file_location("v6_adaptive_e2e_server", SERVER)
    assert spec and spec.loader
    module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(module)
    return module


def _png(path: Path, color: str, label: str, size: tuple[int, int] = (1200, 700)) -> Path:
    image = Image.new("RGB", size, color)
    ImageDraw.Draw(image).text((40, 40), label, fill="white")
    image.save(path)
    return path


def _source_fixture(tmp_path: Path) -> tuple[Path, Path, Path, Path]:
    licensed_meeting = FIXTURE.with_name("staff-meeting-public-domain.jpg")
    license_record = json.loads(FIXTURE.with_name("staff-meeting-public-domain.license.json").read_text(encoding="utf-8"))
    assert hashlib.sha256(licensed_meeting.read_bytes()).hexdigest() == license_record["sha256"]
    meeting = tmp_path / licensed_meeting.name
    meeting.write_bytes(licensed_meeting.read_bytes())
    company_logo = _png(tmp_path / "company-logo.png", "#B23A48", "REAL COMPANY LOGO", (900, 360))
    word = tmp_path / "four-pages.docx"
    document = Document()
    pages = [
        ("Strategy overview", "The approved strategy has three phases: diagnose, implement, review."),
        ("Meeting evidence", "The project team met on 12 August 2026 and approved the implementation plan."),
        ("Brand system", "The approved company identity must remain recognizable in the visual."),
        ("Revenue chart", "Revenue was 20 in 2025 and 30 in 2026."),
    ]
    for number, (title, body) in enumerate(pages, start=1):
        document.add_paragraph(f"第 {number} 页")
        document.add_paragraph(title)
        paragraph = document.add_paragraph(body)
        if number == 2:
            document.add_picture(str(meeting))
        if number == 3:
            document.add_comment(
                paragraph.runs,
                "[search-evidence:official company logo]",
                author="Reviewer",
                initials="RV",
            )
        if number == 4:
            document.add_comment(
                paragraph.runs,
                "[search-evidence:unavailable audit photograph]",
                author="Reviewer",
                initials="RV",
            )
    document.save(word)
    fixed_logo = tmp_path / "fixed-logo.svg"
    fixed_logo.write_text(
        '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 240 72">'
        '<rect width="240" height="72" rx="8" fill="#17365D"/>'
        '<text x="18" y="48" fill="white" font-size="32">FIXED SVG</text></svg>',
        encoding="utf-8",
    )
    return word, fixed_logo, meeting, company_logo


def _confirm_once(project: Path) -> dict:
    server = _load_server()
    client = server.create_app(project).test_client()
    recommendations = client.get("/api/recommendations").get_json()
    candidate = recommendations["design_directions"]["candidates"][0]
    pages = client.get("/api/pages").get_json()["pages"]
    editable = []
    for page in pages:
        editable.append({
            "page_number": page["page_number"],
            "effective_body": page["effective_body"],
            "attachment_extracts": page["attachment_extracts"],
            "chart_facts": page["chart_facts"],
            "image_requirements": page["image_requirements"],
            "degradations": page["degradations"],
            "reference_images": [
                {
                    "reference_id": item["reference_id"],
                    "purpose": item["purpose"],
                    "allow_crop": item["allow_crop"],
                    "allow_restyle": item["allow_restyle"],
                    "status": item["status"],
                    "decision": "keep",
                }
                for item in page["reference_images"]
            ],
            "reference_decisions": [],
        })
    payload = {
        "stage": "final",
        "revision": 0,
        "direction": 0,
        "template_selection": candidate["template_selection"],
        "canvas": "ppt169",
        "visual_style": candidate["visual_style"],
        "color": candidate["color"],
        "icons": candidate["icons"],
        "typography": candidate["typography"],
        "image_rendering": candidate["image_rendering"],
        "style_axes": candidate["style_axes"],
        "layout_preferences": ["auto", "editorial", "matrix"],
        "information_density": candidate["information_density"],
        "regional_style": {"enabled": False},
        "background_system": "light",
        "image_role": {"role": "evidence", "proportion": "medium-low"},
        "evidence_strength": "data-case",
        "image_usage_policy": "content-driven",
        "composition_tendency": "formal-consulting",
        "brand_device": "light",
        "production_profile": "balanced",
        "additional_requirements": "Use only confirmed facts and never generate the fixed title, fixed Logo, footer, or page number.",
        "confirmed_pages": editable,
    }
    response = client.post("/api/confirm", json=payload)
    assert response.status_code == 200, response.get_json()
    result = json.loads((project / "confirm_ui" / "result.json").read_text(encoding="utf-8"))
    state = load(project)
    state["style_confirmation"] = {"status": "confirmed", "contract": server.compile_style_execution(result)}
    state["confirmed_ui_revision"] = result["revision"]
    state["confirmed_ui_digest"] = canonical_sha256(result)
    state["page_materials_status"] = "confirmed"
    save(project, state)
    return result


def _trace(command: list[str], output: Path) -> None:
    trace = Path(command[command.index("--trace-out") + 1])
    images = [Path(command[index + 1]) for index, value in enumerate(command) if value == "--image"]
    roles = [command[index + 1] for index, value in enumerate(command) if value == "--image-role"]
    digests = [command[index + 1] for index, value in enumerate(command) if value == "--image-sha256"]
    trace.write_text(json.dumps({
        "operation": command[2],
        "model": "gpt-image-2",
        "quality": command[command.index("--quality") + 1],
        "size": command[command.index("--size") + 1],
        "input_images": [
            {"role": role, "path": str(path), "sha256": digest}
            for role, path, digest in zip(roles, images, digests)
        ],
        "outputs": [{
            "path": str(output.resolve()),
            "sha256": hashlib.sha256(output.read_bytes()).hexdigest(),
            "mime_type": "image/png",
        }],
    }), encoding="utf-8")


def _semantic_failure(score: int) -> dict:
    return {
        "accepted": False,
        "score": score,
        "checks": {
            "global_style_followed": {
                "result": "fail",
                "detail": "Contrast does not satisfy the frozen style.",
                "correction": {
                    "check": "global_style_followed",
                    "action": "increase",
                    "target": "contrast_relation",
                    "constraint": "contrast_relation",
                    "correction": "Increase contrast between confirmed body text and its background.",
                },
            }
        },
        "issues": [],
    }


def _editable_body(path: Path, page_number: int) -> None:
    deck = Presentation()
    deck.slide_width = Cm(25.4)
    deck.slide_height = Cm(14.288)
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    slide.shapes.add_textbox(Cm(1.5), Cm(3), Cm(16), Cm(2)).text = f"Editable approved body {page_number}"
    deck.save(path)


def test_four_page_adaptive_v6_runs_once_resumes_and_assembles_in_word_order(tmp_path: Path, monkeypatch) -> None:
    contract = json.loads(FIXTURE.read_text(encoding="utf-8"))
    license_record = json.loads(FIXTURE.with_name("staff-meeting-public-domain.license.json").read_text(encoding="utf-8"))
    assert license_record["license"].startswith("Public domain")
    word, fixed_logo, meeting, company_logo = _source_fixture(tmp_path)
    project = tmp_path / "brand-new-v6-project"
    real_extract = workflow_v6_source.extract_source_assets

    def extract_with_chart(*args, **kwargs):
        value = real_extract(*args, **kwargs)
        value.setdefault("chart_records", []).append({
            "page_numbers": [4],
            "title": "Revenue trend",
            "unit": "USD m",
            "series": [
                {"series": "Revenue", "time": "2025", "value": 20},
                {"series": "Revenue", "time": "2026", "value": 30},
            ],
            "image_path": "must-not-be-an-image-input.png",
        })
        return value

    monkeypatch.setattr(workflow_v6_source, "extract_source_assets", extract_with_chart)
    initialize_v6_project(word, fixed_logo, project)
    page3_receipt = json.loads((project / "02_v6/reference_materials/page_003.json").read_text(encoding="utf-8"))
    page4_receipt = json.loads((project / "02_v6/reference_materials/page_004.json").read_text(encoding="utf-8"))
    request3 = page3_receipt["reference_acquisitions"][0]["request_id"]
    request4 = page4_receipt["reference_acquisitions"][0]["request_id"]
    import_reference(project, page_number=3, request_id=request3, image=company_logo, source_url="https://example.test/logo")
    confirm_reference(project, page_number=3, request_id=request3)
    fail_reference(project, page_number=4, request_id=request4, reason="no correct source-backed material")
    frozen = _confirm_once(project)

    assert frozen["revision"] == 1
    assert [item["status"] for item in json.loads((project / "02_v6/reference_materials/page_003.json").read_text())["reference_acquisitions"]] == ["confirmed"]
    assert [item["status"] for item in json.loads((project / "02_v6/reference_materials/page_004.json").read_text())["reference_acquisitions"]] == ["failed_no_retry"]

    calls: list[list[str]] = []
    reviews: defaultdict[int, int] = defaultdict(int)

    def runner(command: list[str], _timeout: int) -> None:
        calls.append(list(command))
        output = Path(command[command.index("--out") + 1])
        page_number = int(output.name.split("_")[1].split(".")[0])
        attempt = int(output.name.split("candidate_")[1].split(".")[0])
        Image.new("RGB", (1904, 896), (35 * page_number, 25 * attempt, 80)).save(output)
        _trace(command, output)

    def reviewer(_project, *, image: Path, **_kwargs):
        page_number = int(image.name.split("_")[1].split(".")[0])
        reviews[page_number] += 1
        if page_number == 2 and reviews[page_number] == 1:
            return _semantic_failure(3)
        if page_number == 3:
            return _semantic_failure(3)
        return {"accepted": True, "score": 5, "checks": {}, "issues": []}

    receipts = [
        generate_page_body(project, page_number=number, runner=runner, reviewer=reviewer, retry_sleep=lambda _x: None)
        for number in range(1, 5)
    ]
    calls_before_resume = len(calls)
    receipt_bytes = [(project / f"04_v6/images/page_{number:03d}.json").read_bytes() for number in range(1, 5)]
    resumed = [
        generate_page_body(project, page_number=number, runner=runner, reviewer=reviewer, retry_sleep=lambda _x: None)
        for number in range(1, 5)
    ]
    assert len(calls) == calls_before_resume
    assert receipts == resumed
    assert receipt_bytes == [(project / f"04_v6/images/page_{number:03d}.json").read_bytes() for number in range(1, 5)]

    by_page = defaultdict(list)
    for call in calls:
        output = Path(call[call.index("--out") + 1])
        by_page[int(output.name.split("_")[1].split(".")[0])].append(call)
    expected = {item["page_number"]: item for item in contract["pages"]}
    for page_number, page_calls in by_page.items():
        assert all(call[2] == expected[page_number]["operation"] for call in page_calls)
        assert page_calls[0][page_calls[0].index("--quality") + 1] == expected[page_number]["quality"]
        assert len(page_calls) <= contract["candidate_limit"]
    assert len(by_page[1]) == 1
    assert len(by_page[2]) == 2
    assert len(by_page[3]) == 2
    assert len(by_page[4]) == 1
    assert "--image" not in by_page[1][0]
    assert "--image" not in by_page[4][0]
    assert [by_page[2][0][index + 1] for index, value in enumerate(by_page[2][0]) if value == "--image"] == [
        by_page[2][1][index + 1] for index, value in enumerate(by_page[2][1]) if value == "--image"
    ]
    assert all(str(meeting) != value for value in by_page[2][0])
    assert not any("candidate_1.png" in value for value in by_page[2][1])
    page2_inputs = [
        Path(by_page[2][0][index + 1]).resolve()
        for index, value in enumerate(by_page[2][0])
        if value == "--image"
    ]
    expected_page2_digest = frozen["confirmed_pages"][1]["reference_images"][0]["integrity"]["model_input_sha256"]
    assert [hashlib.sha256(path.read_bytes()).hexdigest() for path in page2_inputs] == [expected_page2_digest]
    assert all(path.is_relative_to(project) for path in page2_inputs)
    page3_roles = [
        by_page[3][0][index + 1]
        for index, value in enumerate(by_page[3][0])
        if value == "--image-role"
    ]
    assert any("logo" in role.lower() for role in page3_roles)

    prompt1 = Path(by_page[1][0][by_page[1][0].index("--prompt-file") + 1]).read_text(encoding="utf-8")
    prompt4 = Path(by_page[4][0][by_page[4][0].index("--prompt-file") + 1]).read_text(encoding="utf-8")
    for forbidden in ("Strategy overview", "FIXED SVG", "footer", "page number"):
        if forbidden == "Strategy overview":
            assert forbidden not in prompt1
    assert "Revenue trend" in prompt4 and '"value":20' in prompt4 and '"value":30' in prompt4
    assert "must-not-be-an-image-input" not in prompt4
    assert receipts[2]["selected"]["attempt"] == 1
    assert any("qa_no_effective_improvement" in reason for reason in receipts[2]["degraded_reasons"])
    for number, receipt in enumerate(receipts, start=1):
        selected = project / receipt["selected"]["path"]
        with Image.open(selected) as image:
            assert image.size == (1904, 896)
        assert receipt["request_identity"]
        assert receipt["request_prompt_sha256"]
        assert receipt["selected"]["output_sha256"]
        assert receipt["selected"]["trace_sha256"]

    for page_number in range(1, 5):
        request = build_reconstruction_request(project, page_number=page_number)
        assert request["page_title"] not in request["effective_page"].get("body_render_content", "")
        body = tmp_path / f"editable-body-{page_number}.pptx"
        _editable_body(body, page_number)
        finalized = finalize_reconstructed_page(project, page_number=page_number, reconstructed_body=body)
        assert finalized["fixed_frame"]["passed"] is True
    assembly = assemble_v6_deck(project)
    deck = Presentation(project / assembly["output"])
    assert assembly["page_order"] == [1, 2, 3, 4]
    assert len(deck.slides) == 4
    assert round(deck.slide_width / 360000, 3) == 25.4
    assert round(deck.slide_height / 360000, 3) == 14.288
    assert all(len([shape for shape in slide.shapes if shape.name == "fixed-frame-logo"]) == 1 for slide in deck.slides)
    assert all(
        next(
                shape._pic.blipFill.blip.rEmbed
            for shape in slide.shapes
            if shape.name == "fixed-frame-logo"
        ) in {
            relationship.rId
            for relationship in slide.part.rels.values()
            if relationship.reltype.endswith("/image")
            and relationship.target_part.content_type == "image/svg+xml"
        }
        for slide in deck.slides
    )
