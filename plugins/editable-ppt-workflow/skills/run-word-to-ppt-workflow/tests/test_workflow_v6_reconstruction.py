from __future__ import annotations

import json
import sys
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Cm


ROOT = Path(__file__).resolve().parents[1]
SCRIPTS = ROOT / "scripts"
if str(SCRIPTS) not in sys.path:
    sys.path.insert(0, str(SCRIPTS))

from workflow_v6_contract import new_page, new_project  # noqa: E402
from workflow_v6_reconstruction import (  # noqa: E402
    assemble_v6_deck,
    build_reconstruction_request,
    finalize_reconstructed_page,
)
from workflow_v6_state import create, load, save  # noqa: E402


def _style():
    return {
        "fixed_frame": {
            "geometry_version": "fixed-canvas-cm-v2",
            "body_bounds_cm": {"x": 0.81, "y": 2.3, "w": 23.78, "h": 11.18},
            "body_bounds": {"x": 0.81 / 25.4, "y": 2.3 / 14.288, "w": 23.78 / 25.4, "h": 11.18 / 14.288},
            "title_bounds_cm": {"x": 0.9, "y": 0.5, "w": 20.066, "h": 1.4288},
            "title_bounds": {"x": 0.9 / 25.4, "y": 0.5 / 14.288, "w": 20.066 / 25.4, "h": 1.4288 / 14.288},
            "logo_bounds_cm": {"x": 21.844, "y": 0.57152, "w": 2.667, "h": 1.0716},
            "logo_bounds": {"x": 21.844 / 25.4, "y": 0.57152 / 14.288, "w": 2.667 / 25.4, "h": 1.0716 / 14.288},
            "footer_line": {"x": 0.9, "y": 13.64504, "w": 23.6, "h": 0.028576, "color": "#B8C0CC"},
            "page_number_bounds_cm": {"x": 23.368, "y": 13.687904, "w": 1.143, "h": 0.3572},
            "page_number_bounds": {"x": 23.368 / 25.4, "y": 13.687904 / 14.288, "w": 1.143 / 25.4, "h": 0.3572 / 14.288},
            "page_number_style": {"font": "Microsoft YaHei", "size_pt": 9, "color": "#6B7280"},
            "title_color": "#0B1727",
        },
        "hard_constraints": {
            "title_color": "#0B1727",
            "typography": {
                "heading": {"cjk": "Microsoft YaHei"},
                "type_scale_pt": {"page_title": 28},
            },
        },
    }


def _body(path: Path, text: str):
    deck = Presentation()
    deck.slide_width = Cm(25.4)
    deck.slide_height = Cm(14.288)
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    box = slide.shapes.add_textbox(Cm(2), Cm(3), Cm(10), Cm(2))
    box.text = text
    deck.save(path)


def _project(tmp_path: Path, page_count: int = 2) -> Path:
    root = tmp_path / "project"
    (root / "00_source").mkdir(parents=True)
    (root / "00_source" / "logo.svg").write_text(
        '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 100 20"><rect width="100" height="20"/></svg>',
        encoding="utf-8",
    )
    pages = [new_page(number, title=f"标题{number}") for number in range(1, page_count + 1)]
    for page in pages:
        page["state"] = "accepted"
        image = root / "04_v6" / "images" / f"page_{page['page_number']:03d}.png"
        image.parent.mkdir(parents=True, exist_ok=True)
        Image.new("RGB", (1904, 896), "white").save(image)
        page["first_candidate"] = {"path": image.relative_to(root).as_posix(), "attempt": 1, "operation": "generate"}
        page["selected_candidate"] = dict(page["first_candidate"])
        effective = root / "02_v6" / "effective_pages" / f"page_{page['page_number']:03d}.json"
        effective.parent.mkdir(parents=True, exist_ok=True)
        effective.write_text(json.dumps({"page_number": page["page_number"], "word_original": "正文"}), encoding="utf-8")
    state = new_project(
        word_source={"path": "00_source/source.docx"},
        logo_source={"path": "00_source/logo.svg"},
        pages=pages,
    )
    state["style_confirmation"] = {"status": "confirmed", "contract": _style()}
    create(root, state)
    return root


def test_v6_reconstruction_request_has_no_exact_material_or_post_visual_qa(tmp_path: Path):
    project = _project(tmp_path, 1)
    request = build_reconstruction_request(project, page_number=1)
    assert request["workflow_contract_version"] == "word-ppt-workflow-v6"
    assert request["requirements"]["exact_reference_material_custody"] is False
    assert request["requirements"]["post_reconstruction_visual_qa"] is False


def test_finalize_and_assemble_add_fixed_layers_without_office_or_visual_qa(tmp_path: Path):
    project = _project(tmp_path, 2)
    for page in (1, 2):
        body = tmp_path / f"body-{page}.pptx"
        _body(body, f"可编辑正文{page}")
        report = finalize_reconstructed_page(project, page_number=page, reconstructed_body=body)
        assert report["post_reconstruction_visual_qa"] is False
        assert report["fixed_frame"]["passed"] is True

    report = assemble_v6_deck(project)
    output = project / report["output"]
    deck = Presentation(output)
    assert len(deck.slides) == 2
    assert report["office_render_required"] is False
    assert report["post_reconstruction_visual_qa"] is False
    assert all(page["state"] == "page_complete" for page in load(project)["pages"])
