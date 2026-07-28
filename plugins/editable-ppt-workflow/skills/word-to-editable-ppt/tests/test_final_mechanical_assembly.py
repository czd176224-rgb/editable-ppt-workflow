from __future__ import annotations

import hashlib
import json
import os
import subprocess
import sys
import zipfile
from contextlib import contextmanager
from dataclasses import replace
from pathlib import Path

import pytest
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
SCRIPTS = ROOT / "scripts"
EDITPPT_CLI = ROOT.parent / "image-to-editable-ppt" / "cli"
sys.path.insert(0, str(SCRIPTS))
sys.path.insert(0, str(EDITPPT_CLI))

import editppt.runtime.final_assembler as final_assembler  # noqa: E402
import editppt.runtime.editable_page_cache as editable_page_cache  # noqa: E402
import editppt.runtime.finalize_deck_run as finalize_deck_run  # noqa: E402
from cache_store import CacheStore  # noqa: E402
from page_pipeline import seal_completed_page  # noqa: E402
from page_qa import PageQAResult  # noqa: E402
from style_contract import canonical_json_bytes  # noqa: E402
from workflow_state import (  # noqa: E402
    dispatch,
    load,
    next_action,
    record_generation,
    record_qa,
    record_reconstruction,
    status,
)
from editppt.runtime.editable_page_cache import (  # noqa: E402
    PackageValidationError,
    create_page_package,
)
from editppt.runtime.final_assembler import assemble_deck  # noqa: E402
from editppt.runtime.finalize_deck_run import (  # noqa: E402
    build_current_assembly_plan,
    finalize_project,
)
from final_mechanical_qa import FINAL_QA_FIELDS, run_final_mechanical_qa  # noqa: E402


def _write_json(path: Path, value: dict) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(value, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")


def _project(
    tmp_path: Path,
    locked_order: tuple[int, ...] = (1, 2, 3),
    *,
    canvas: str = "ppt169",
) -> Path:
    project = tmp_path / "project"
    project.mkdir()
    jobs = []
    for page_number in locked_order:
        source = f"第{page_number}页唯一内容"
        contract = {
            "schema_version": "2.0",
            "page_number": page_number,
            "source_text": source,
            "source_hash": hashlib.sha256(source.encode("utf-8")).hexdigest(),
            "semantic_units": [],
            "source_tables": [],
            "explicit_relations": [],
            "asset_bindings": [],
            "detected_numbers": [],
            "detected_amounts": [],
        }
        contract_file = project / "01_page_contracts" / f"page_{page_number:03d}.json"
        _write_json(contract_file, contract)
        jobs.append({
            "slide_id": f"slide_{page_number:03d}",
            "page_number": page_number,
            "status": "pending_style_confirmation",
            "contract_file": contract_file.relative_to(project).as_posix(),
            "expected_output": f"06_images/generated/page_{page_number:03d}.png",
        })

    profile = {
        "aspect_ratio": "16:9" if canvas == "ppt169" else "4:3",
        "image_size": "1792x1008" if canvas == "ppt169" else "1536x1152",
        "slide_width_inches": 13.333333 if canvas == "ppt169" else 10.0,
        "slide_height_inches": 7.5,
        "fit": "contain",
        "allow_crop": False,
    }
    style = {
        "schema_version": "1.0",
        "direction": "editorial",
        "canvas": canvas,
        "canvas_profile": profile,
        "image_quality": "high",
        "generation_mode": "continuous",
        "max_concurrency": 2,
    }
    style_file = project / "02_style" / "style_execution.json"
    style_file.parent.mkdir()
    style_file.write_bytes(canonical_json_bytes(style))
    _write_json(project / "workflow_run.json", {
        "schema_version": "1.0",
        "workflow_contract_version": "word-only-v1",
        "project_name": "mechanical-assembly",
        "word_source": {},
        "pagination": {
            "page_count": len(locked_order),
            "locked_page_order": list(locked_order),
        },
        "style_confirmation": {
            "status": "confirmed",
            "confirmed_at": "2026-07-27T00:00:00Z",
            "execution_file": style_file.relative_to(project).as_posix(),
            "execution_sha256": hashlib.sha256(canonical_json_bytes(style)).hexdigest(),
        },
        "jobs": jobs,
        "final_pptx": None,
    })
    return project


def _request(action: dict, page_number: int) -> dict:
    return next(item for item in action["requests"] if item["page_number"] == page_number)


def _page_pptx(
    project: Path,
    page_number: int,
    *,
    suffix: str = "",
    canvas: str = "ppt169",
) -> Path:
    output = project / "07_editable" / f"page_{page_number:03d}{suffix}.pptx"
    output.parent.mkdir(parents=True, exist_ok=True)
    presentation = Presentation()
    presentation.slide_width = Inches(13.333333 if canvas == "ppt169" else 10.0)
    presentation.slide_height = Inches(7.5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(7), Inches(1))
    box.text_frame.paragraphs[0].text = f"Word page {page_number}"
    box.text_frame.paragraphs[0].font.size = Pt(24)
    presentation.save(output)
    return output


def _package_for_job(project: Path, job: dict) -> Path:
    pptx = _page_pptx(project, job["page_number"])
    output = project / "07_editable" / f"page_{job['page_number']:03d}.json"
    return create_page_package(
        project,
        page_number=job["page_number"],
        cache_key=job["cache"]["key"],
        pptx=pptx,
        output=output,
    )


def _complete_all(project: Path) -> None:
    next_action(project)
    run = load(project)
    for job in run["jobs"]:
        package = _package_for_job(project, job)
        hit = seal_completed_page(project, job, package)
        job["qa_result"] = PageQAResult("pass", "none").as_dict()
        job["status"] = "complete"
        job["assignment"] = None
        job["cache_hit"] = False
        job["reconstruction"] = {
            "artifact": package.relative_to(project).as_posix(),
            "attempt": 1,
            "cache_path": hit.path.relative_to(project).as_posix(),
        }
    _write_json(project / "workflow_run.json", run)


def _render_with_real_open(pptx: Path, output: Path) -> int:
    presentation = Presentation(pptx)
    output.mkdir(parents=True, exist_ok=True)
    for index, _slide in enumerate(presentation.slides, start=1):
        Image.new("RGB", (16, 9), "white").save(output / f"slide_{index:03d}.png")
    return len(presentation.slides)


def _slide_texts(path: Path) -> list[str]:
    presentation = Presentation(path)
    return [
        "\n".join(shape.text for shape in slide.shapes if hasattr(shape, "text_frame"))
        for slide in presentation.slides
    ]


def _rewrite_pptx(
    path: Path,
    *,
    replacements: dict[str, bytes] | None = None,
    additions: list[tuple[str, bytes, int]] | None = None,
) -> None:
    temporary = path.with_suffix(".rewrite.pptx")
    replacements = replacements or {}
    additions = additions or []
    with zipfile.ZipFile(path, "r") as source, zipfile.ZipFile(temporary, "w") as target:
        for item in source.infolist():
            target.writestr(item, replacements.get(item.filename, source.read(item.filename)))
        for name, payload, compression in additions:
            target.writestr(name, payload, compress_type=compression)
    os.replace(temporary, path)


def _symlink_or_skip(link: Path, target: Path, *, directory: bool = False) -> None:
    try:
        link.symlink_to(target, target_is_directory=directory)
    except (NotImplementedError, OSError) as exc:
        pytest.skip(f"symlinks are unavailable on this platform: {exc}")


def _assert_no_final_publication(project: Path) -> None:
    assert load(project)["final_pptx"] is None
    assert not (project / "08_final" / "deck.pptx").exists()
    assert not (project / "08_final" / "final_mechanical_qa.json").exists()
    assert not (project / "08_final" / "run_summary.json").exists()
    final_dir = project / "08_final"
    assert not final_dir.exists() or not list(final_dir.glob(".finalize-*"))


def test_passed_page_reconstructs_and_caches_while_another_page_remains_active(tmp_path: Path) -> None:
    project = _project(tmp_path, (1, 2))
    first = next_action(project)
    page_1_attempt = dispatch(project, 1, "generator-1", _request(first, 1)["attempt"])["attempt"]
    page_2_attempt = dispatch(project, 2, "generator-2", _request(first, 2)["attempt"])["attempt"]
    image = project / "06_images" / "generated" / "page_001.png"
    image.parent.mkdir(parents=True)
    image.write_bytes(b"page-one-image")
    record_generation(project, 1, "generator-1", page_1_attempt, image)
    record_qa(project, 1, "generator-1", page_1_attempt, PageQAResult("pass", "none"))

    reconstruction = _request(next_action(project), 1)
    attempt = dispatch(project, 1, "reconstructor-1", reconstruction["attempt"])["attempt"]
    package = create_page_package(
        project,
        page_number=1,
        cache_key=reconstruction["cache_key"],
        pptx=_page_pptx(project, 1),
        output=Path(reconstruction["output"]),
    )
    record_reconstruction(project, 1, "reconstructor-1", attempt, package)

    current = load(project)
    first_job = next(item for item in current["jobs"] if item["page_number"] == 1)
    assert status(project)["page_states"] == {"complete": [1], "generating": [2]}
    assert CacheStore(project).lookup("pages", first_job["cache"]["key"]) is not None
    assert page_2_attempt == 1


def test_editppt_record_cli_bridges_current_reconstruction_to_complete(tmp_path: Path) -> None:
    """The public editppt record command must seal and record a current Word page."""
    project = _project(tmp_path, (1,))
    first = next_action(project)
    generation_attempt = dispatch(project, 1, "page-worker", _request(first, 1)["attempt"])["attempt"]
    image = project / "06_images" / "generated" / "page_001_attempt_001.png"
    image.parent.mkdir(parents=True)
    image.write_bytes(b"generated-page")
    record_generation(project, 1, "page-worker", generation_attempt, image)
    record_qa(project, 1, "page-worker", generation_attempt, PageQAResult("pass", "none"))

    reconstruction = _request(next_action(project), 1)
    reconstruction_attempt = dispatch(
        project,
        1,
        "page-worker",
        reconstruction["attempt"],
    )["attempt"]
    pptx = _page_pptx(project, 1)
    descriptor = project / "07_editable" / "page_001.json"

    completed = subprocess.run(
        [
            sys.executable,
            "-m",
            "editppt.cli",
            "run",
            "record",
            str(project),
            "--page",
            "1",
            "--agent-id",
            "page-worker",
            "--attempt",
            str(reconstruction_attempt),
            "--pptx",
            str(pptx),
            "--artifact",
            str(descriptor),
        ],
        cwd=EDITPPT_CLI,
        env={
            **os.environ,
            "PYTHONPATH": os.pathsep.join((str(SCRIPTS), str(EDITPPT_CLI))),
        },
        check=False,
        capture_output=True,
        text=True,
    )

    assert completed.returncode == 0, completed.stderr or completed.stdout
    assert json.loads(completed.stdout)["state"] == "complete"
    assert descriptor.is_file()
    assert load(project)["jobs"][0]["status"] == "complete"


def test_ppt43_record_and_final_assembly_preserve_confirmed_slide_ratio(tmp_path: Path) -> None:
    project = _project(tmp_path, (1,), canvas="ppt43")
    first = next_action(project)
    generation_attempt = dispatch(project, 1, "page-worker", _request(first, 1)["attempt"])["attempt"]
    image = project / "06_images" / "generated" / "page_001.png"
    image.parent.mkdir(parents=True)
    image.write_bytes(b"image")
    record_generation(project, 1, "page-worker", generation_attempt, image)
    record_qa(project, 1, "page-worker", generation_attempt, PageQAResult("pass", "none"))
    reconstruction = _request(next_action(project), 1)
    attempt = dispatch(project, 1, "page-worker", reconstruction["attempt"])["attempt"]
    package = create_page_package(
        project,
        page_number=1,
        cache_key=reconstruction["cache_key"],
        pptx=_page_pptx(project, 1, canvas="ppt43"),
        output=Path(reconstruction["output"]),
    )
    record_reconstruction(project, 1, "page-worker", attempt, package)

    summary = finalize_project(project, renderer=_render_with_real_open)
    final = Presentation(summary["output"])

    assert final.slide_width / final.slide_height == pytest.approx(4 / 3, rel=1e-6)


def test_final_assembly_rejects_page_ratio_that_differs_from_confirmed_canvas(tmp_path: Path) -> None:
    project = _project(tmp_path, (1,), canvas="ppt43")
    next_action(project)
    run = load(project)
    job = run["jobs"][0]
    package = create_page_package(
        project,
        page_number=1,
        cache_key=job["cache"]["key"],
        pptx=_page_pptx(project, 1, canvas="ppt169"),
        output=project / "07_editable" / "page_001.json",
    )
    hit = seal_completed_page(project, job, package)
    job["qa_result"] = PageQAResult("pass", "none").as_dict()
    job["status"] = "complete"
    job["assignment"] = None
    job["cache_hit"] = False
    job["reconstruction"] = {
        "artifact": package.relative_to(project).as_posix(),
        "attempt": 1,
        "cache_path": hit.path.relative_to(project).as_posix(),
    }
    _write_json(project / "workflow_run.json", run)

    with pytest.raises(PackageValidationError, match="canvas|ratio|slide size"):
        build_current_assembly_plan(project)


def test_final_assembly_has_exactly_one_slide_per_locked_word_page_in_locked_order(tmp_path: Path) -> None:
    project = _project(tmp_path, (2, 1, 3))
    _complete_all(project)

    receipt = assemble_deck(build_current_assembly_plan(project), project / "08_final" / "deck.pptx")

    assert receipt.slide_order == (2, 1, 3)
    assert receipt.page_count == 3
    assert _slide_texts(receipt.output_path) == ["Word page 2", "Word page 1", "Word page 3"]
    assert len(Presentation(receipt.output_path).slides) == 3
    assert editable_page_cache.inspect_editable_pptx(receipt.output_path).slide_page_ids == (2, 1, 3)


def test_slide_fingerprint_ignores_serializer_only_xml_whitespace(tmp_path: Path) -> None:
    project = _project(tmp_path, (1,))
    original = _page_pptx(project, 1)
    rewritten = original.with_name("page_001-reserialized.pptx")
    with zipfile.ZipFile(original) as source, zipfile.ZipFile(rewritten, "w") as target:
        for item in source.infolist():
            payload = source.read(item.filename)
            if item.filename == "ppt/slides/slide1.xml":
                text = payload.decode("utf-8")
                text = text.replace("><", ">\n  <")
                payload = text.encode("utf-8")
            target.writestr(item, payload)

    first = editable_page_cache.inspect_editable_pptx(original)
    second = editable_page_cache.inspect_editable_pptx(rewritten)
    assert first.slide_fingerprints == second.slide_fingerprints


def test_final_mechanical_qa_contains_only_mechanical_gates(tmp_path: Path) -> None:
    project = _project(tmp_path, (1, 2))
    _complete_all(project)
    receipt = assemble_deck(build_current_assembly_plan(project), project / "08_final" / "deck.pptx")

    report = run_final_mechanical_qa(project, receipt, renderer=_render_with_real_open)

    assert set(report) == FINAL_QA_FIELDS
    assert report["passed"] is True
    assert report["page_count"] == {"expected": 2, "actual": 2, "passed": True}
    assert report["page_order"] == {"expected": [1, 2], "actual": [1, 2], "passed": True}
    assert report["artifact_existence"]["passed"] is True
    assert report["page_qa_status"]["passed"] is True
    assert report["editable_object_status"]["passed"] is True
    assert report["package_validity"]["passed"] is True
    assert report["open_render_status"]["passed"] is True
    serialized = json.dumps(report, ensure_ascii=False).lower()
    assert "style" not in serialized
    assert "similarity" not in serialized
    assert "visual_dna" not in serialized
    assert "global_visual" not in serialized


def test_final_render_proof_is_reused_only_for_the_same_strict_deck_hash(tmp_path: Path) -> None:
    project = _project(tmp_path, (1, 2))
    _complete_all(project)
    receipt = assemble_deck(build_current_assembly_plan(project), project / "08_final" / "deck.pptx")
    calls = 0

    def counted_renderer(pptx: Path, output: Path) -> int:
        nonlocal calls
        calls += 1
        return _render_with_real_open(pptx, output)

    first = run_final_mechanical_qa(project, receipt, renderer=counted_renderer)
    second = run_final_mechanical_qa(project, receipt, renderer=counted_renderer)

    assert first["open_render_status"]["cache_hit"] is False
    assert second["open_render_status"]["cache_hit"] is True
    assert calls == 1


def test_missing_optional_renderer_is_an_advisory_not_a_publication_blocker(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    project = _project(tmp_path, (1,))
    _complete_all(project)
    receipt = assemble_deck(build_current_assembly_plan(project), project / "08_final" / "deck.pptx")

    class MissingRenderer(RuntimeError):
        pass

    monkeypatch.setattr(
        "final_mechanical_qa._default_renderer",
        lambda _pptx, _output: (_ for _ in ()).throw(MissingRenderer("no local renderer installed")),
    )
    monkeypatch.setattr("final_mechanical_qa.NoRenderBackendError", MissingRenderer, raising=False)
    report = run_final_mechanical_qa(project, receipt)

    assert report["passed"] is True
    assert report["open_render_status"]["passed"] is True
    assert report["open_render_status"]["mode"] == "structural-only"
    assert "optional" in report["open_render_status"]["advisory"].lower()


def test_final_mechanical_qa_rejects_a_receipt_with_the_wrong_deck_hash(tmp_path: Path) -> None:
    project = _project(tmp_path, (1,))
    _complete_all(project)
    receipt = assemble_deck(build_current_assembly_plan(project), project / "08_final" / "deck.pptx")

    report = run_final_mechanical_qa(
        project,
        replace(receipt, output_sha256="0" * 64),
        renderer=_render_with_real_open,
    )

    assert report["passed"] is False
    assert report["package_validity"]["passed"] is False
    assert "changed" in (report["package_validity"]["error"] or "")


def test_final_mechanical_qa_standalone_cli_passes_a_valid_deck(tmp_path: Path) -> None:
    project = _project(tmp_path, (1,))
    _complete_all(project)
    receipt = assemble_deck(build_current_assembly_plan(project), project / "08_final" / "deck.pptx")
    report_path = project / "08_final" / "standalone-mechanical-qa.json"

    completed = subprocess.run(
        [
            sys.executable,
            str(SCRIPTS / "final_mechanical_qa.py"),
            "--project",
            str(project),
            "--input",
            str(receipt.output_path),
            "--report",
            str(report_path),
        ],
        check=False,
        capture_output=True,
        text=True,
        timeout=120,
    )

    assert completed.returncode == 0, completed.stderr or completed.stdout
    report = json.loads(report_path.read_text(encoding="utf-8"))
    assert report["passed"] is True
    assert report["output_sha256"] == receipt.output_sha256


def test_finalizer_writes_only_current_mechanical_outputs(tmp_path: Path) -> None:
    project = _project(tmp_path, (1, 2))
    _complete_all(project)

    summary = finalize_project(project, renderer=_render_with_real_open)

    assert summary["status"] == "complete"
    assert summary["page_count"] == 2
    assert summary["page_order"] == [1, 2]
    assert Path(summary["output"]).is_file()
    assert Path(summary["qa_report"]).is_file()
    assert load(project)["final_pptx"] == "08_final/deck.pptx"
    names = {path.name.lower() for path in (project / "08_final").iterdir()}
    assert not any("backup" in name or "recovery" in name or "legacy" in name for name in names)


def test_final_assembly_rejects_a_missing_sealed_page_package(tmp_path: Path) -> None:
    project = _project(tmp_path, (1, 2))
    _complete_all(project)
    run = load(project)
    first = run["jobs"][0]
    hit = CacheStore(project).lookup("pages", first["cache"]["key"])
    assert hit is not None
    (hit.path / hit.manifest["outputs"]["reconstruction"]).unlink()

    with pytest.raises(PackageValidationError, match="cache|package|artifact"):
        assemble_deck(build_current_assembly_plan(project), project / "08_final" / "deck.pptx")


def test_final_assembly_rejects_a_replaced_page_pptx(tmp_path: Path) -> None:
    project = _project(tmp_path, (1, 2))
    _complete_all(project)
    replaced = project / "07_editable" / "page_001.pptx"
    replacement = _page_pptx(project, 99, suffix="-replacement")
    replaced.write_bytes(replacement.read_bytes())

    with pytest.raises(PackageValidationError, match="replaced|SHA-256|identity"):
        assemble_deck(build_current_assembly_plan(project), project / "08_final" / "deck.pptx")


def test_final_assembly_waits_until_every_locked_page_is_complete(tmp_path: Path) -> None:
    project = _project(tmp_path, (1, 2))
    _complete_all(project)
    run = load(project)
    run["jobs"][1]["status"] = "accepted"
    _write_json(project / "workflow_run.json", run)

    with pytest.raises(PackageValidationError, match="complete"):
        build_current_assembly_plan(project)


def test_final_assembly_recomputes_the_current_six_input_cache_identity(tmp_path: Path) -> None:
    project = _project(tmp_path, (1, 2))
    _complete_all(project)
    style_path = project / "02_style" / "style_execution.json"
    style_path.write_text('{"schema_version":"1.0","direction":"changed"}\n', encoding="utf-8")

    with pytest.raises(PackageValidationError, match="current cache identity|style execution"):
        build_current_assembly_plan(project)


def test_failed_post_publish_authority_replay_leaves_no_final_output(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    project = _project(tmp_path, (1, 2))
    _complete_all(project)
    plan = build_current_assembly_plan(project)
    output = project / "08_final" / "deck.pptx"
    real_load = final_assembler._load_packages
    calls = 0

    def fail_final_replay(value):
        nonlocal calls
        calls += 1
        if calls == 3:
            raise PackageValidationError("page package disappeared during final replay")
        return real_load(value)

    monkeypatch.setattr(final_assembler, "_load_packages", fail_final_replay)
    with pytest.raises(PackageValidationError, match="disappeared"):
        assemble_deck(plan, output)

    assert calls == 3
    assert not output.exists()


def test_identical_word_pages_reuse_one_exact_six_input_cache_package(tmp_path: Path) -> None:
    project = _project(tmp_path, (1, 2))
    first_contract = json.loads(
        (project / "01_page_contracts" / "page_001.json").read_text(encoding="utf-8")
    )
    second_path = project / "01_page_contracts" / "page_002.json"
    second_contract = json.loads(second_path.read_text(encoding="utf-8"))
    second_contract["source_text"] = first_contract["source_text"]
    second_contract["source_hash"] = first_contract["source_hash"]
    _write_json(second_path, second_contract)
    _complete_all(project)
    run = load(project)
    assert run["jobs"][0]["cache"]["key"] == run["jobs"][1]["cache"]["key"]

    receipt = assemble_deck(
        build_current_assembly_plan(project),
        project / "08_final" / "deck.pptx",
    )

    assert receipt.slide_order == (1, 2)
    assert _slide_texts(receipt.output_path) == ["Word page 1", "Word page 1"]


def test_page_change_during_render_blocks_completion_and_cleans_staging(tmp_path: Path) -> None:
    project = _project(tmp_path, (1, 2))
    _complete_all(project)

    def mutate_page_then_render(pptx: Path, output: Path) -> int:
        current = load(project)
        current["jobs"][0]["status"] = "accepted"
        _write_json(project / "workflow_run.json", current)
        return _render_with_real_open(pptx, output)

    with pytest.raises(PackageValidationError, match="complete|authority|current"):
        finalize_project(project, renderer=mutate_page_then_render)

    assert load(project)["jobs"][0]["status"] == "accepted"
    _assert_no_final_publication(project)


def test_renderer_cannot_mutate_the_staged_deck_after_mechanical_inspection(tmp_path: Path) -> None:
    project = _project(tmp_path, (1, 2))
    _complete_all(project)

    def mutate_deck_then_render(pptx: Path, output: Path) -> int:
        count = _render_with_real_open(pptx, output)
        with pptx.open("ab") as handle:
            handle.write(b"renderer-mutated-the-inspected-deck")
        return count

    with pytest.raises(PackageValidationError, match="renderer|mechanical QA|changed"):
        finalize_project(project, renderer=mutate_deck_then_render)

    _assert_no_final_publication(project)


def test_final_authority_replay_compares_exact_descriptor_and_manifest_identity(tmp_path: Path) -> None:
    project = _project(tmp_path, (1,))
    _complete_all(project)
    plan = build_current_assembly_plan(project)
    receipt = assemble_deck(plan, project / "08_final" / "authority.pptx")
    before = receipt.page_packages[0]
    redirected = replace(before, cache_artifact=before.cache_artifact.with_name("equivalent.json"))

    assert finalize_deck_run._same_page_authority(before, redirected) is False


@pytest.mark.parametrize("failure", ["renderer", "qa"])
def test_renderer_or_qa_failure_publishes_nothing(tmp_path: Path, failure: str) -> None:
    project = _project(tmp_path, (1, 2))
    _complete_all(project)

    def renderer(pptx: Path, output: Path) -> int:
        if failure == "renderer":
            raise RuntimeError("injected renderer failure")
        output.mkdir(parents=True, exist_ok=True)
        Image.new("RGB", (16, 9), "white").save(output / "slide_001.png")
        return 1

    with pytest.raises(PackageValidationError, match="mechanical QA|renderer|render"):
        finalize_project(project, renderer=renderer)

    _assert_no_final_publication(project)


@pytest.mark.parametrize("failed_name", ["workflow_run.json", "run_summary.json"])
def test_state_or_summary_stage_write_failure_publishes_nothing(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
    failed_name: str,
) -> None:
    project = _project(tmp_path, (1, 2))
    _complete_all(project)
    real_write = finalize_deck_run._write_json_atomic

    def injected_write(path: Path, value: dict) -> None:
        if path.name == failed_name:
            raise OSError(f"injected {failed_name} write failure")
        real_write(path, value)

    monkeypatch.setattr(finalize_deck_run, "_write_json_atomic", injected_write)
    with pytest.raises((OSError, PackageValidationError), match="injected|write"):
        finalize_project(project, renderer=_render_with_real_open)

    _assert_no_final_publication(project)


def test_lock_release_failure_after_commit_does_not_rollback_outside_the_lock(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    project = _project(tmp_path, (1, 2))
    _complete_all(project)
    real_lock = finalize_deck_run.workflow_state.project_state_lock

    @contextmanager
    def failing_exit_lock(value: Path):
        with real_lock(value):
            yield
        raise OSError("injected state lock release failure")

    monkeypatch.setattr(finalize_deck_run.workflow_state, "project_state_lock", failing_exit_lock)
    summary = finalize_project(project, renderer=_render_with_real_open)

    assert summary["status"] == "complete"
    assert load(project)["final_pptx"] == "08_final/deck.pptx"
    assert Path(summary["output"]).is_file()


@pytest.mark.parametrize("failed_name", ["run_summary.json", "workflow_run.json"])
def test_publish_replace_failure_after_mutation_rolls_back_every_output(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
    failed_name: str,
) -> None:
    project = _project(tmp_path, (1, 2))
    _complete_all(project)
    real_replace = finalize_deck_run.os.replace
    injected = False

    def replace_then_fail(source: Path, destination: Path) -> None:
        nonlocal injected
        real_replace(source, destination)
        if not injected and Path(destination).name == failed_name:
            injected = True
            raise OSError(f"injected post-replace {failed_name} failure")

    monkeypatch.setattr(finalize_deck_run.os, "replace", replace_then_fail)
    with pytest.raises((OSError, PackageValidationError), match="injected|rollback|replace"):
        finalize_project(project, renderer=_render_with_real_open)

    _assert_no_final_publication(project)


def test_staging_cleanup_failure_rolls_back_completion(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    project = _project(tmp_path, (1, 2))
    _complete_all(project)
    real_cleanup = finalize_deck_run._cleanup_staging
    calls = 0

    def fail_once(value: Path, staging: Path) -> None:
        nonlocal calls
        calls += 1
        if calls == 1:
            raise OSError("injected staging cleanup failure")
        real_cleanup(value, staging)

    monkeypatch.setattr(finalize_deck_run, "_cleanup_staging", fail_once)
    with pytest.raises((OSError, PackageValidationError), match="injected|cleanup|rollback"):
        finalize_project(project, renderer=_render_with_real_open)

    _assert_no_final_publication(project)


def test_deck_mutation_after_qa_before_publish_is_rejected(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    project = _project(tmp_path, (1, 2))
    _complete_all(project)
    real_write = finalize_deck_run._write_json_atomic

    def mutate_after_summary(path: Path, value: dict) -> None:
        real_write(path, value)
        if path.name == "run_summary.json":
            with (path.parent / "deck.pptx").open("ab") as handle:
                handle.write(b"mutation-after-qa")

    monkeypatch.setattr(finalize_deck_run, "_write_json_atomic", mutate_after_summary)
    with pytest.raises(PackageValidationError, match="deck changed|hash|authority"):
        finalize_project(project, renderer=_render_with_real_open)

    _assert_no_final_publication(project)


def test_project_and_cache_artifact_symlinks_are_rejected_before_resolution(tmp_path: Path) -> None:
    project = _project(tmp_path, (1, 2))
    _complete_all(project)
    linked_project = tmp_path / "linked-project"
    _symlink_or_skip(linked_project, project, directory=True)

    with pytest.raises(PackageValidationError, match="link|reparse|project"):
        build_current_assembly_plan(linked_project)

    page_pptx = project / "07_editable" / "page_001.pptx"
    replacement = project / "07_editable" / "replacement.pptx"
    replacement.write_bytes(page_pptx.read_bytes())
    page_pptx.unlink()
    _symlink_or_skip(page_pptx, replacement)
    with pytest.raises(PackageValidationError, match="link|reparse|artifact"):
        build_current_assembly_plan(project)


def test_project_ancestor_reparse_is_rejected_before_resolve(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    ancestor = tmp_path / "redirected-ancestor"
    project = ancestor / "project"
    project.mkdir(parents=True)
    real_probe = editable_page_cache._is_link_or_reparse
    observed: list[Path] = []

    def flagged(path: Path) -> bool:
        observed.append(path)
        return path == ancestor or real_probe(path)

    monkeypatch.setattr(editable_page_cache, "_is_link_or_reparse", flagged)
    with pytest.raises(PackageValidationError, match="ancestor|link|reparse|project"):
        editable_page_cache.require_plain_project(project)

    assert ancestor in observed


def test_final_output_parent_symlink_cannot_redirect_publication(tmp_path: Path) -> None:
    project = _project(tmp_path, (1, 2))
    _complete_all(project)
    outside = tmp_path / "outside-final"
    outside.mkdir()
    _symlink_or_skip(project / "08_final", outside, directory=True)

    with pytest.raises(PackageValidationError, match="link|reparse|output"):
        finalize_project(project, renderer=_render_with_real_open)

    assert not (outside / "deck.pptx").exists()
    assert load(project)["final_pptx"] is None


def test_pptx_archive_entry_count_and_compression_ratio_are_bounded(tmp_path: Path) -> None:
    project = _project(tmp_path, (1,))
    pptx = _page_pptx(project, 1)
    _rewrite_pptx(
        pptx,
        additions=[
            (f"unused/entry-{index}.bin", b"", zipfile.ZIP_STORED)
            for index in range(4100)
        ],
    )
    with pytest.raises(PackageValidationError, match="entries|entry count"):
        editable_page_cache.inspect_editable_pptx(pptx)

    pptx = _page_pptx(project, 1, suffix="-ratio")
    _rewrite_pptx(
        pptx,
        additions=[("unused/compressed-bomb.bin", b"0" * (2 * 1024 * 1024), zipfile.ZIP_DEFLATED)],
    )
    with pytest.raises(PackageValidationError, match="compression ratio|archive"):
        editable_page_cache.inspect_editable_pptx(pptx)


def test_pptx_archive_per_entry_aggregate_and_xml_reads_are_bounded(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    project = _project(tmp_path, (1,))

    per_entry = _page_pptx(project, 1, suffix="-per-entry")
    monkeypatch.setattr(editable_page_cache, "MAX_ARCHIVE_ENTRY_BYTES", 64 * 1024)
    _rewrite_pptx(
        per_entry,
        additions=[("unused/oversized.bin", b"x" * (65 * 1024), zipfile.ZIP_STORED)],
    )
    with pytest.raises(PackageValidationError, match="entry is oversized"):
        editable_page_cache.inspect_editable_pptx(per_entry)

    aggregate = _page_pptx(project, 1, suffix="-aggregate")
    monkeypatch.setattr(editable_page_cache, "MAX_ARCHIVE_ENTRY_BYTES", 128 * 1024)
    monkeypatch.setattr(editable_page_cache, "MAX_ARCHIVE_UNCOMPRESSED_BYTES", 96 * 1024)
    _rewrite_pptx(
        aggregate,
        additions=[
            ("unused/aggregate-a.bin", b"a" * (48 * 1024), zipfile.ZIP_STORED),
            ("unused/aggregate-b.bin", b"b" * (48 * 1024), zipfile.ZIP_STORED),
        ],
    )
    with pytest.raises(PackageValidationError, match="aggregate uncompressed"):
        editable_page_cache.inspect_editable_pptx(aggregate)

    xml = _page_pptx(project, 1, suffix="-xml")
    monkeypatch.setattr(editable_page_cache, "MAX_ARCHIVE_ENTRY_BYTES", 128 * 1024 * 1024)
    monkeypatch.setattr(editable_page_cache, "MAX_ARCHIVE_UNCOMPRESSED_BYTES", 512 * 1024 * 1024)
    monkeypatch.setattr(editable_page_cache, "MAX_XML_PART_BYTES", 256 * 1024)
    _rewrite_pptx(
        xml,
        replacements={"ppt/presentation.xml": b"<presentation>" + b"x" * (300 * 1024)},
    )
    with pytest.raises(PackageValidationError, match="XML part is oversized|read bound"):
        editable_page_cache.inspect_editable_pptx(xml)


def test_extensionless_opc_part_declared_as_xml_uses_the_xml_bound(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    project = _project(tmp_path, (1,))
    pptx = _page_pptx(project, 1, suffix="-extensionless-xml")
    with zipfile.ZipFile(pptx) as archive:
        content_types = archive.read("[Content_Types].xml")
    content_types = content_types.replace(
        b"</Types>",
        (
            b'<Override PartName="/custom/extensionless" '
            b'ContentType="application/vnd.example.payload+xml"/></Types>'
        ),
    )
    monkeypatch.setattr(editable_page_cache, "MAX_XML_PART_BYTES", 64 * 1024)
    _rewrite_pptx(
        pptx,
        replacements={"[Content_Types].xml": content_types},
        additions=[("custom/extensionless", b"x" * (65 * 1024), zipfile.ZIP_STORED)],
    )

    with pytest.raises(PackageValidationError, match="XML part is oversized"):
        editable_page_cache.inspect_editable_pptx(pptx)


def test_external_slide_relationship_is_rejected_before_open_or_copy(tmp_path: Path) -> None:
    project = _project(tmp_path, (1,))
    pptx = _page_pptx(project, 1)
    rel_name = "ppt/slides/_rels/slide1.xml.rels"
    with zipfile.ZipFile(pptx) as archive:
        relationships = archive.read(rel_name)
    injected = relationships.replace(
        b"</Relationships>",
        (
            b'<Relationship Id="rIdExternal" '
            b'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink" '
            b'Target="https://example.invalid/" TargetMode="External"/>'
            b"</Relationships>"
        ),
    )
    _rewrite_pptx(pptx, replacements={rel_name: injected})

    with pytest.raises(PackageValidationError, match="external relationship"):
        editable_page_cache.inspect_editable_pptx(pptx)
